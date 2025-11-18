"""PunchITNow 9.0 Octopus - Cross-Platform Generator UI.
Reads a semicolon-separated CSV `Account generation.csv` from the same folder.
Provides sections for: Google Alerts, Gmail answer, Prompts, youtube comentary.

Each section shows a listbox with sampled items, Next & Copy buttons, and a shared Generate button.
Logs actions in a scrolled text widget.

🍎 macOS Port - Version 9.0 Octopus
🐙 New: Enhanced build system for Windows & macOS
🔧 Fixed: CSV profile search functionality, improved column detection
"""

# Виправлення кодування консолі для Windows
import sys
import io
if sys.platform == 'win32':
    # Перевіряємо чи stdout/stderr не None (може бути при компіляції в .exe)
    if sys.stdout is not None and hasattr(sys.stdout, 'buffer'):
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    if sys.stderr is not None and hasattr(sys.stderr, 'buffer'):
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# ⚡ ОСНОВНІ ІМПОРТИ - завантажуються відразу
import os
import json
import time
import threading
import platform
from functools import lru_cache

# 🔧 ПЛАТФОРМО-СПЕЦИФІЧНІ НАЛАШТУВАННЯ
PLATFORM = platform.system()  # 'Windows', 'Darwin' (macOS), 'Linux'
IS_MACOS = PLATFORM == 'Darwin'
IS_WINDOWS = PLATFORM == 'Windows'
IS_LINUX = PLATFORM == 'Linux'

# ⚡ БАЗОВІ МОДУЛІ - швидкі імпорти
import csv
import random
import shutil
import re
import datetime
import math
import subprocess

# ⚡ БАЗОВИЙ TKINTER - критично важливий
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import tkinter.scrolledtext as scrolledtext

# ⚡ ГЛОБАЛЬНІ НАЛАШТУВАННЯ
GLOBAL_SCROLL_SPEED = 1

# ⚡ LAZY LOADING СИСТЕМА
class LazyImports:
    """Система відкладеного завантаження важких бібліотек"""
    
    def __init__(self):
        self._imports = {}
        self._loading = set()
    
    def get_module(self, name, import_func):
        """Отримує модуль з кешу або завантажує його"""
        if name in self._imports:
            return self._imports[name]
        
        if name in self._loading:
            # Якщо модуль завантажується в іншому потоці, чекаємо
            while name in self._loading:
                time.sleep(0.01)
            return self._imports.get(name)
        
        self._loading.add(name)
        try:
            result = import_func()
            self._imports[name] = result
            return result
        except Exception as e:
            self._imports[name] = None
            print(f"⚠️ Не вдалося завантажити {name}: {e}")
            return None
        finally:
            self._loading.discard(name)

# Створюємо глобальний instance для lazy imports
_lazy = LazyImports()

# ⚡ LAZY IMPORT ФУНКЦІЇ
def get_customtkinter():
    """Lazy import для CustomTkinter"""
    def load_ctk():
        import customtkinter as ctk
        ctk.set_appearance_mode("dark")
        return ctk
    return _lazy.get_module('customtkinter', load_ctk)

def get_pygame():
    """Lazy import для pygame"""
    def load_pygame():
        try:
            import pygame
            pygame.mixer.init()
            return pygame
        except ImportError:
            try:
                # Fallback на pygame-ce якщо pygame не встановлено
                import pygame_ce as pygame
                pygame.mixer.init()
                return pygame
            except ImportError:
                # Якщо жоден варіант не працює, повертаємо None
                print("⚠️ pygame не встановлено - звуки вимкнено")
                return None
    return _lazy.get_module('pygame', load_pygame)

def get_selenium():
    """Lazy import для selenium (використовується рідко)"""
    def load_selenium():
        from selenium import webdriver
        from selenium.webdriver.common.by import By
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        from selenium.common.exceptions import TimeoutException
        from selenium.webdriver.chrome.options import Options
        from selenium.webdriver.chrome.service import Service
        from webdriver_manager.chrome import ChromeDriverManager
        return {
            'webdriver': webdriver,
            'By': By,
            'WebDriverWait': WebDriverWait,
            'EC': EC,
            'TimeoutException': TimeoutException,
            'Options': Options,
            'Service': Service,
            'ChromeDriverManager': ChromeDriverManager
        }
    return _lazy.get_module('selenium', load_selenium)

def get_requests():
    """Lazy import для requests"""
    def load_requests():
        import requests
        return requests
    return _lazy.get_module('requests', load_requests)

def get_pyperclip():
    """Lazy import для pyperclip"""
    def load_pyperclip():
        import pyperclip
        return pyperclip
    return _lazy.get_module('pyperclip', load_pyperclip)

def get_pyotp():
    """Lazy import для pyotp"""
    def load_pyotp():
        import pyotp
        return pyotp
    return _lazy.get_module('pyotp', load_pyotp)

def get_webbrowser():
    """Lazy import для webbrowser"""
    def load_webbrowser():
        import webbrowser
        return webbrowser
    return _lazy.get_module('webbrowser', load_webbrowser)

def get_keyboard():
    """Lazy import для keyboard (для global hotkeys)"""
    def load_keyboard():
        import keyboard
        return keyboard
    return _lazy.get_module('keyboard', load_keyboard)

def get_openai():
    """Lazy import для OpenAI"""
    def load_openai():
        import openai
        return openai
    return _lazy.get_module('openai', load_openai)

def get_file_generator():
    """Lazy import для file generator"""
    # Повертаємо вбудований клас FileGeneratorTab
    return FileGeneratorTab

# ⚡ ІНІЦІАЛІЗУЄМО КРИТИЧНО ВАЖLIVІ МОДУЛІ
ctk = get_customtkinter()
if not ctk:
    print("❌ CustomTkinter недоступний - програма не може працювати")
    sys.exit(1)

# 🔄 AUTO-UPDATE CHECK SYSTEM
def check_for_updates():
    """Перевіряє наявність оновлень на GitHub"""
    try:
        import urllib.request
        
        # GitHub Configuration
        GITHUB_USER = "qbyLive1"
        GITHUB_REPO = "PunchItNow"
        GITHUB_BRANCH = "main"
        VERSION_URL = f"https://raw.githubusercontent.com/{GITHUB_USER}/{GITHUB_REPO}/{GITHUB_BRANCH}/version.json"
        
        # Читаємо локальну версію
        local_version = "3.0.1"  # За замовчуванням
        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            version_file = os.path.join(script_dir, "version.json")
            if os.path.exists(version_file):
                with open(version_file, 'r', encoding='utf-8') as f:
                    local_data = json.load(f)
                    local_version = local_data.get("version", "3.0.1")
        except:
            pass
        
        # Перевіряємо віддалену версію
        with urllib.request.urlopen(VERSION_URL, timeout=3) as response:
            remote_data = json.loads(response.read().decode('utf-8'))
            remote_version = remote_data.get("version", "3.0.1")
            changelog = remote_data.get("changelog", "")
        
        # Порівнюємо версії
        local_parts = [int(x) for x in local_version.split('.')]
        remote_parts = [int(x) for x in remote_version.split('.')]
        
        if remote_parts > local_parts:
            print(f"\n{'='*60}")
            print(f"🎉 ДОСТУПНЕ ОНОВЛЕННЯ!")
            print(f"📦 Поточна версія: {local_version}")
            print(f"🆕 Нова версія: {remote_version}")
            print(f"📝 Зміни:\n{changelog}")
            print(f"💡 Закрийте програму та запустіть PunchIT Launcher.exe")
            print(f"{'='*60}\n")
            return True, remote_version, changelog
        else:
            print(f"✅ Версія актуальна: {local_version}")
            return False, None, None
            
    except Exception as e:
        # Тихо ігноруємо помилки перевірки (немає інтернету тощо)
        return False, None, None

def show_update_notification(version, changelog):
    """Показує повідомлення про доступне оновлення"""
    try:
        root = tk.Tk()
        root.withdraw()  # Ховаємо головне вікно
        
        msg = f"🎉 Доступна нова версія: {version}\n\n"
        msg += f"Зміни:\n{changelog}\n\n"
        msg += "Закрийте програму та запустіть PunchIT Launcher.exe\nщоб оновитися."
        
        messagebox.showinfo("Доступне оновлення", msg)
        root.destroy()
    except:
        pass

# ⚡ MEMORY MANAGEMENT СИСТЕМА
import gc
import weakref
from collections import deque

class MemoryManager:
    """Система управління пам'яттю для оптимізації роботи програми"""
    
    def __init__(self):
        self._widget_cache = {}
        self._cleanup_queue = deque(maxlen=100)
        self._gc_counter = 0
        
    def register_widget(self, widget_type, widget):
        """Реєструє віджет для відстеження пам'яті"""
        widget_id = id(widget)
        self._widget_cache[widget_id] = weakref.ref(widget, lambda ref: self._cleanup_queue.append(widget_id))
        
    def cleanup_widgets(self):
        """Очищує віджети що більше не використовуються"""
        while self._cleanup_queue:
            widget_id = self._cleanup_queue.popleft()
            self._widget_cache.pop(widget_id, None)
            
    def force_gc(self):
        """Примусове збирання сміття"""
        self._gc_counter += 1
        if self._gc_counter % 10 == 0:  # Кожні 10 викликів
            self.cleanup_widgets()
            collected = gc.collect()
            if collected > 0:
                print(f"🗑️ GC: очищено {collected} об'єктів")
                
    def get_memory_stats(self):
        """Повертає статистику використання пам'яті"""
        return {
            'widgets_tracked': len([w for w in self._widget_cache.values() if w() is not None]),
            'widgets_total': len(self._widget_cache),
            'gc_cycles': self._gc_counter
        }

# Глобальний менеджер пам'яті
_memory_manager = MemoryManager()

def optimized_widget_creation(widget_class, parent, **kwargs):
    """Оптимізоване створення віджетів з управлінням пам'яттю"""
    # Встановлюємо оптимальні настройки за замовчуванням
    default_kwargs = {
        'corner_radius': 8,
        'border_width': 1,
    }
    
    # Об'єднуємо з користувацькими параметрами
    final_kwargs = {**default_kwargs, **kwargs}
    
    # Створюємо віджет
    widget = widget_class(parent, **final_kwargs)
    
    # Реєструємо для відстеження пам'яті
    _memory_manager.register_widget(widget_class.__name__, widget)
    
    return widget

def cleanup_memory():
    """Публічна функція для очищення пам'яті"""
    _memory_manager.force_gc()

# ⚡ UI RENDERING OPTIMIZATIONS
class UIOptimizer:
    """Система оптимізації UI рендерингу"""
    
    def __init__(self):
        self._pending_updates = {}
        self._batch_delay = 50  # мс
        
    def debounced_update(self, key, update_func, *args, **kwargs):
        """Дебаунсинг для оновлень UI"""
        # Скасовуємо попередні оновлення
        if key in self._pending_updates:
            try:
                # Зупиняємо попередній after() якщо можливо
                pass
            except:
                pass
        
        # Плануємо нове оновлення
        def delayed_update():
            if key in self._pending_updates:
                del self._pending_updates[key]
                try:
                    update_func(*args, **kwargs)
                except Exception as e:
                    print(f"❌ UI update error: {e}")
        
        self._pending_updates[key] = delayed_update
        return self._batch_delay
    
    def batch_ui_updates(self, updates_dict):
        """Пакетне оновлення UI"""
        for key, (func, args, kwargs) in updates_dict.items():
            self.debounced_update(key, func, *args, **kwargs)

# Глобальний UI оптимізатор
_ui_optimizer = UIOptimizer()

@lru_cache(maxsize=32)
def get_optimized_font(family="Segoe UI", size=12, weight="normal"):
    """Кешований доступ до шрифтів"""
    return (family, size, weight)

def optimized_listbox_insert(listbox, items, clear=True):
    """Оптимізоване додавання елементів в listbox"""
    if clear:
        listbox.delete(0, tk.END)
    
    # Пакетне додавання для кращої продуктивності
    if len(items) > 100:
        # Для великих списків - показуємо тільки перші 100
        items = items[:100]
        items.append(f"... та ще {len(items) - 100} елементів")
    
    for item in items:
        listbox.insert(tk.END, item)

# ⚡ FILE OPERATIONS OPTIMIZATION
class FileCache:
    """Кеш для файлових операцій"""
    
    def __init__(self, max_size=50):
        self._cache = {}
        self._timestamps = {}
        self._max_size = max_size
        
    @lru_cache(maxsize=32)
    def _get_file_mtime(self, filepath):
        """Кешоване отримання часу модифікації файлу"""
        try:
            return os.path.getmtime(filepath)
        except:
            return 0
    
    def get_cached_data(self, filepath, loader_func):
        """Отримує дані з кешу або завантажує"""
        current_mtime = self._get_file_mtime(filepath)
        
        if (filepath in self._cache and 
            filepath in self._timestamps and
            self._timestamps[filepath] >= current_mtime):
            return self._cache[filepath]
        
        # Завантажуємо дані
        try:
            data = loader_func(filepath)
            
            # Управління розміром кешу
            if len(self._cache) >= self._max_size:
                # Видаляємо найстаріші записи
                oldest = min(self._timestamps.keys(), 
                           key=lambda k: self._timestamps[k])
                del self._cache[oldest]
                del self._timestamps[oldest]
            
            self._cache[filepath] = data
            self._timestamps[filepath] = current_mtime
            return data
        except Exception as e:
            print(f"❌ File loading error: {e}")
            return None
    
    def clear_cache(self):
        """Очищає кеш"""
        self._cache.clear()
        self._timestamps.clear()

# Глобальний файловий кеш
_file_cache = FileCache()

def optimized_csv_read(filepath):
    """Оптимізоване читання CSV"""
    def csv_loader(path):
        with open(path, 'r', encoding='utf-8') as f:
            return list(csv.reader(f, delimiter=';'))
    
    return _file_cache.get_cached_data(filepath, csv_loader)

def optimized_json_read(filepath):
    """Оптимізоване читання JSON"""
    def json_loader(path):
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    return _file_cache.get_cached_data(filepath, json_loader)

def clear_file_cache():
    """Очищає файловий кеш"""
    _file_cache.clear_cache()

# ⚡ OPTIMIZED ERROR HANDLING & LOGGING
class OptimizedLogger:
    """Оптимізований логер з мінімальним overhead"""
    
    def __init__(self, verbose=False):
        self._verbose = verbose
        self._error_count = 0
        self._last_errors = deque(maxlen=10)
        
    def log(self, level, message, *args):
        """Оптимізоване логування"""
        if level == "ERROR":
            self._error_count += 1
            self._last_errors.append(message)
        
        if self._verbose or level == "ERROR":
            if args:
                print(f"{level}: {message % args}")
            else:
                print(f"{level}: {message}")
                
    def error(self, msg, *args):
        self.log("ERROR", msg, *args)
        
    def info(self, msg, *args):
        if self._verbose:
            self.log("INFO", msg, *args)
            
    def get_stats(self):
        return {
            'error_count': self._error_count,
            'recent_errors': list(self._last_errors)
        }

# Глобальний логер
_logger = OptimizedLogger(verbose=False)

def safe_execute(func, *args, error_msg="Operation failed", **kwargs):
    """Безпечне виконання функцій з обробкою помилок"""
    try:
        return func(*args, **kwargs)
    except Exception as e:
        _logger.error(f"{error_msg}: {e}")
        return None

# ⚡ PERFORMANCE MONITORING
def get_performance_stats():
    """Повертає статистику продуктивності"""
    return {
        'memory': _memory_manager.get_memory_stats(),
        'logger': _logger.get_stats(),
        'lazy_imports': len(_lazy._imports),
        'file_cache_size': len(_file_cache._cache)
    }

def print_optimization_summary():
    """Виводить підсумок оптимізацій"""
    stats = get_performance_stats()
    print("\n🚀 OPTIMIZATION SUMMARY:")
    print(f"   📦 Lazy imports loaded: {stats['lazy_imports']}")
    print(f"   🗄️ File cache entries: {stats['file_cache_size']}")
    print(f"   🧠 Widgets tracked: {stats['memory']['widgets_tracked']}")
    print(f"   ♻️ GC cycles: {stats['memory']['gc_cycles']}")
    print(f"   ⚠️ Errors logged: {stats['logger']['error_count']}")
    print("   ✅ Optimizations active: Lazy loading, Memory management, File caching, UI debouncing")

# ⚡ HELPER ФУНКЦІЇ ДЛЯ ЗВОРОТНОЇ СУМІСНОСТІ
def copy_to_clipboard(text):
    """Оптимізована функція копіювання в буфер"""
    return safe_clipboard_operation("set", text)

def get_from_clipboard():
    """Оптимізована функція отримання з буфера"""
    return safe_clipboard_operation("get")

# Псевдоніми для зворотної сумісності (lazy)
pyperclip = type('LazyPyperclip', (), {
    'copy': lambda text: copy_to_clipboard(text),
    'paste': lambda: get_from_clipboard()
})()

pyotp = type('LazyPyotp', (), {
    'TOTP': lambda secret: get_pyotp().TOTP(secret) if get_pyotp() else None
})()

webbrowser = type('LazyWebbrowser', (), {
    'open': lambda url: get_webbrowser().open(url) if get_webbrowser() else None
})()

requests = type('LazyRequests', (), {
    'get': lambda *args, **kwargs: get_requests().get(*args, **kwargs) if get_requests() else None
})()

# Threading псевдонім
from threading import Thread

# Selenium псевдоніми (lazy loading)
def _get_selenium():
    selenium_dict = get_selenium()
    return selenium_dict if selenium_dict else {}

webdriver = type('LazyWebdriver', (), {
    'Chrome': lambda *args, **kwargs: _get_selenium().get('webdriver', type('NoWebDriver', (), {'Chrome': lambda *a, **k: None})).Chrome(*args, **kwargs) if _get_selenium() else None
})()

By = type('LazyBy', (), {
    'ID': lambda id_val: _get_selenium().get('By', type('NoBy', (), {'ID': lambda x: x})).ID if _get_selenium() else id_val,
    'TAG_NAME': lambda tag: _get_selenium().get('By', type('NoBy', (), {'TAG_NAME': lambda x: x})).TAG_NAME if _get_selenium() else tag
})()

WebDriverWait = lambda driver, timeout: _get_selenium().get('WebDriverWait', lambda d, t: None)(driver, timeout) if _get_selenium() else None
EC = _get_selenium().get('EC') if _get_selenium() else type('NoEC', (), {'presence_of_element_located': lambda x: None, 'visibility_of_element_located': lambda x: None})()
TimeoutException = _get_selenium().get('TimeoutException', Exception) if _get_selenium() else Exception
Options = _get_selenium().get('Options', type('NoOptions', (), {})) if _get_selenium() else type('NoOptions', (), {})
Service = _get_selenium().get('Service', type('NoService', (), {})) if _get_selenium() else type('NoService', (), {})
ChromeDriverManager = _get_selenium().get('ChromeDriverManager', type('NoChrome', (), {'install': lambda: ''})) if _get_selenium() else type('NoChrome', (), {'install': lambda: ''})

# OpenAI та FileGenerator псевдоніми
openai = type('LazyOpenAI', (), {
    'OpenAI': lambda *args, **kwargs: get_openai().OpenAI(*args, **kwargs) if get_openai() else None
})()

FileGeneratorTab = lambda *args, **kwargs: get_file_generator()(*args, **kwargs) if get_file_generator() else None

# Глобальні флаги
SOUND_AVAILABLE = True  # Перевіряємо динамічно через get_pygame()
HAS_CTK = True  # ctk вже ініціалізовано
HAS_OPENAI = lambda: get_openai() is not None
FILE_GENERATOR_AVAILABLE = lambda: get_file_generator() is not None

# ⚡ ОПТИМІЗОВАНІ CLIPBOARD ОПЕРАЦІЇ
@lru_cache(maxsize=1)
def _get_clipboard_backend():
    """Визначає оптимальний backend для clipboard операцій"""
    pyperclip = get_pyperclip()
    return pyperclip if pyperclip else 'tkinter'

def safe_clipboard_operation(operation, text=""):
    """Оптимізована робота з буфером обміну (cross-platform)"""
    
    # Швидка перевірка рекурсії без створення атрибутів щоразу
    if getattr(safe_clipboard_operation, '_lock', False):
        return "" if operation == "get" else False
    
    try:
        safe_clipboard_operation._lock = True
        
        # 🍎 macOS має нативні інструменти clipboard
        if IS_MACOS:
            import subprocess
            if operation == "get":
                try:
                    result = subprocess.run(['pbpaste'], capture_output=True, text=True, timeout=5)
                    return result.stdout if result.returncode == 0 else ""
                except:
                    pass
            elif operation == "set":
                try:
                    subprocess.run(['pbcopy'], input=text, text=True, timeout=5)
                    return True
                except:
                    pass
        
        # 🪟 Windows/Linux - використовуємо pyperclip або tkinter
        backend = _get_clipboard_backend()
        
        if operation == "get":
            if backend != 'tkinter':
                try:
                    return backend.paste()
                except:
                    pass
            
            # Tkinter fallback (створюємо root тільки якщо потрібно)
            try:
                root = tk.Tk()
                root.withdraw()
                result = root.clipboard_get()
                root.destroy()
                return result
            except:
                if 'root' in locals():
                    root.destroy()
                return ""
                
        elif operation == "set":
            if backend != 'tkinter':
                try:
                    backend.copy(text)
                    return True
                except:
                    pass
            
            # Tkinter fallback
            try:
                root = tk.Tk()
                root.withdraw()
                root.clipboard_clear()
                root.clipboard_append(text)
                root.update()
                root.destroy()
                return True
            except:
                if 'root' in locals():
                    root.destroy()
                return False
                
    except Exception:
        return "" if operation == "get" else False
    finally:
        safe_clipboard_operation._lock = False

def safe_text_input(widget, text):
    """Безпечне введення тексту в віджет (оптимізовано)"""
    try:
        # Швидке очищення та вставка
        if hasattr(widget, 'delete') and hasattr(widget, 'insert'):
            if isinstance(widget, tk.Entry):
                widget.delete(0, tk.END)
                widget.insert(0, text)
            else:
                widget.delete('1.0', tk.END)
                widget.insert('1.0', text)
        return True
    except Exception:
        return False

def setup_keyboard_locale_support(root):
    """Налаштовує підтримку розкладок (оптимізовано)"""
    try:
        import locale
        locale.setlocale(locale.LC_ALL, '')
        return True
    except Exception:
        return False

# ⚡ ОПТИМІЗОВАНІ SOUND ФУНКЦІЇ
def play_alert_sound():
    """Оптимізоване відтворення звуку з lazy loading"""
    pygame = get_pygame()
    if not pygame:
        return
    
    try:
        # Шукаємо звуковий файл
        sound_path = get_current_sound_file()
        
        # Fallback в Downloads
        if not sound_path:
            fallback_path = r"C:\Users\alexx\Downloads\Звуки - Внимание.mp3"
            if os.path.exists(fallback_path):
                try:
                    config_sound_path = get_config_path("sms_alert.mp3")
                    shutil.copy2(fallback_path, config_sound_path)
                    sound_path = config_sound_path
                except Exception:
                    sound_path = fallback_path
        
        if os.path.exists(sound_path):
            pygame.mixer.music.load(sound_path)
            pygame.mixer.music.play()
        else:
            # Системний звук fallback
            try:
                import winsound
                winsound.MessageBeep(winsound.MB_ICONEXCLAMATION)
            except:
                pass
    except Exception as e:
        print(f"Sound playback error: {e}")
ctk.set_default_color_theme("blue")

# Global variables for Gmail Hacks
sym = "abcdefghijklmnopqrstuvwxyz"

# 📁 Визначаємо робочу директорію (cross-platform)
if getattr(sys, 'frozen', False):
    # Пакований виконуваний файл (.exe/.app/.AppImage)
    if IS_MACOS:
        # macOS .app bundle - шукаємо поруч з .app
        # sys.executable = /path/to/App.app/Contents/MacOS/App
        # Нам потрібно /path/to/ (поруч з .app)
        if '.app/' in sys.executable or '.app\\' in sys.executable:
            # Знаходимо шлях до .app
            app_path = sys.executable
            while '.app' in os.path.basename(app_path):
                app_path = os.path.dirname(app_path)
            # Беремо папку, де знаходиться .app
            WORKDIR = os.path.dirname(app_path)
            print(f"🍎 macOS .app detected - WORKDIR: {WORKDIR}")
        else:
            WORKDIR = os.path.dirname(sys.executable)
    elif IS_WINDOWS:
        # Windows .exe файл
        WORKDIR = os.path.dirname(sys.executable)
    else:
        # Linux AppImage чи інший формат
        WORKDIR = os.path.dirname(sys.executable)
else:
    # Звичайний .py файл
    WORKDIR = os.path.dirname(os.path.abspath(__file__))

# Створюємо папку config для налаштувань
CONFIG_DIR = os.path.join(WORKDIR, 'config')
if not os.path.exists(CONFIG_DIR):
    os.makedirs(CONFIG_DIR)
    print(f"📁 Створено папку конфігурації: {CONFIG_DIR}")

CSV_PATH = os.path.join(WORKDIR, 'Account generation.csv')

# Функції для портативного збереження
def get_config_path(filename):
    """Повертає шлях до файлу конфігурації в папці config"""
    return os.path.join(CONFIG_DIR, filename)

def get_data_path(filename):
    """Повертає шлях до файлу даних в робочій директорії"""
    return os.path.join(WORKDIR, filename)

def get_global_scroll_speed():
    """Повертає поточну глобальну швидкість скролу"""
    global GLOBAL_SCROLL_SPEED
    return GLOBAL_SCROLL_SPEED

def set_global_scroll_speed(speed):
    """Встановлює глобальну швидкість скролу"""
    global GLOBAL_SCROLL_SPEED
    GLOBAL_SCROLL_SPEED = speed
    print(f"🔄 Глобальна швидкість скролу оновлена: {speed}x")

def load_global_scroll_speed():
    """Завантажує швидкість скролу з конфігурації"""
    global GLOBAL_SCROLL_SPEED
    try:
        with open(get_config_path("scroll_speed.json"), 'r') as f:
            config = json.load(f)
            GLOBAL_SCROLL_SPEED = config.get("scroll_speed", 1)
            print(f"📜 Завантажено глобальну швидкість скролу: {GLOBAL_SCROLL_SPEED}x")
    except (FileNotFoundError, json.JSONDecodeError):
        GLOBAL_SCROLL_SPEED = 1
        print("📜 Використовується швидкість скролу за замовчуванням: 1x")

def get_app_stuff_path(subfolder=""):
    """Повертає шлях до папки App Stuff для тимчасових файлів"""
    app_stuff_dir = os.path.join(WORKDIR, "App Stuff")
    if not os.path.exists(app_stuff_dir):
        os.makedirs(app_stuff_dir)
    
    if subfolder:
        full_path = os.path.join(app_stuff_dir, subfolder)
        if not os.path.exists(full_path):
            os.makedirs(full_path)
        return full_path
    
    return app_stuff_dir

def cleanup_temp_folders():
    """Переносить тимчасові папки в App Stuff для організації файлової структури"""
    try:
        import shutil
        import glob
        
        # Патерни тимчасових папок для перенесення
        temp_patterns = [
            "batch_*", "files_*", "auto_files", "generated_files", 
            "random_files", "test_files"
        ]
        
        moved_count = 0
        for pattern in temp_patterns:
            matching_folders = glob.glob(os.path.join(WORKDIR, pattern))
            for folder_path in matching_folders:
                if os.path.isdir(folder_path):
                    folder_name = os.path.basename(folder_path)
                    destination = get_app_stuff_path(folder_name)
                    
                    # Перенести лише якщо ще не існує в App Stuff
                    if not os.path.exists(destination):
                        shutil.move(folder_path, destination)
                        moved_count += 1
                        print(f"📁 Перенесено {folder_name} → App Stuff")
        
        if moved_count > 0:
            print(f"✅ Очищено {moved_count} тимчасових папок")
                        
    except Exception as e:
        print(f"⚠️  Помилка очищення папок: {e}")


# ═══════════════════════════════════════════════════════════════════
# 🎛️ FEATURES SETTINGS SYSTEM - Управління функціями програми
# ═══════════════════════════════════════════════════════════════════

# Дефолтні налаштування функцій
DEFAULT_FEATURES = {
    "generators": {
        "google_alerts": True,
        "google_sign": True,
        "email_subscription": True,
        "gmail_answer": True,
        "prompts": True,
        "youtube_commentary": True,
        "weirdo": True,
        "email_for_mail": True,
    },
    "daily_report": {
        "enabled": True,
        "farmer_name": True,
        "date": True,
        "copy_button": True,
    },
    "utilities": {
        "google_backup_codes": True,
        "2fa_generator": True,
        "password_generator": True,
    },
    "ui_sections": {
        "action_log": True,  # Права колонка з UI Generators
    },
    "windows_mode": {
        "two_windows": True,
        "four_windows": False,
        "six_windows": False,
    }
}

class FeaturesConfig:
    """Клас для управління конфігурацією функцій"""
    
    def __init__(self):
        self.config_file = get_config_path("features_config.json")
        self.features = self.load_config()
    
    def load_config(self):
        """Завантажує конфігурацію з файлу"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    loaded = json.load(f)
                    return self._merge_with_defaults(loaded, DEFAULT_FEATURES)
            else:
                return DEFAULT_FEATURES.copy()
        except Exception as e:
            print(f"❌ Помилка завантаження features config: {e}")
            return DEFAULT_FEATURES.copy()
    
    def _merge_with_defaults(self, loaded, defaults):
        """Мержить завантажену конфігурацію з дефолтними значеннями"""
        result = defaults.copy()
        for key, value in loaded.items():
            if key in result and isinstance(value, dict):
                result[key] = self._merge_with_defaults(value, result[key])
            else:
                result[key] = value
        return result
    
    def save_config(self):
        """Зберігає конфігурацію в файл"""
        try:
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(self.features, f, ensure_ascii=False, indent=2)
            return True
        except Exception as e:
            print(f"❌ Помилка збереження features config: {e}")
            return False
    
    def is_enabled(self, category, feature):
        """Перевіряє чи увімкнена функція"""
        return self.features.get(category, {}).get(feature, True)
    
    def set_enabled(self, category, feature, enabled):
        """Встановлює стан функції"""
        if category not in self.features:
            self.features[category] = {}
        self.features[category][feature] = enabled

# Глобальний instance для Features Config
_features_config = FeaturesConfig()

def is_feature_enabled(category, feature):
    """Глобальна функція для перевірки стану функції"""
    return _features_config.is_enabled(category, feature)


def load_csv_columns(path):
    """Load the first row of CSV headers and all non-empty values per column.
    CSV is semicolon separated.
    Returns dict: {header: [values...]}
    """
    if not os.path.exists(path):
        messagebox.showwarning('File Not Found', f'CSV file not found at {path}')
        return {}
    try:
        data = {}
        with open(path, 'r', encoding='cp1251', errors='ignore') as f:
            reader = csv.reader(f, delimiter=';')
            rows = list(reader)
        if not rows:
            messagebox.showwarning('Empty File', 'The CSV file is empty')
            return {}
    except Exception as e:
        messagebox.showerror('Error', f'Failed to read CSV file: {str(e)}')
        return {}
    headers = [h.strip() for h in rows[0]]
    for h in headers:
        data[h] = []
    for r in rows[1:]:
        for i, cell in enumerate(r):
            if i < len(headers):
                val = cell.strip()
                if val:
                    data[headers[i]].append(val)
    return data


class SectionFrame(ctk.CTkFrame):
    def __init__(
    self,
    master,
    title,
    min_items,
    max_items,
    column_values,
    logger,
    append_values=None,
    font=None,
    *args,
     **kwargs):
        """append_values: optional list of extra words to append to generated items (used by Email Subscription)."""
        super().__init__(master, *args, **kwargs)
        self.title = title
        self.min_items = min_items
        self.max_items = max_items
        self.column_values = column_values or []
        self.append_values = append_values or []
        self.logger = logger
        self.font = font
        self.items = []
        self.index = 0
        self._build()

    def _build(self):
        # Заголовок секції з підказкою
        title_frame = ctk.CTkFrame(self)
        title_frame.pack(fill='x', padx=2)
        ctk.CTkLabel(title_frame, text=self.title, font=self.font).pack(side='left')
        ctk.CTkLabel(title_frame, text="(клік=копій)", 
                    font=ctk.CTkFont(size=8), text_color="gray").pack(side='right')
        self.textbox = ctk.CTkTextbox(self, wrap="word", font=self.font, height=90)
        self.textbox.pack(fill='both', expand=True, padx=2, pady=2)
        
        # Додаємо click-to-copy функціонал до секції
        self._setup_click_to_copy()
        self.current_label = ctk.CTkLabel(
    self, text="Selected: None", font=self.font)
        self.current_label.pack(pady=2)
        button_frame = ctk.CTkFrame(self)
        button_frame.pack(fill='x', padx=2, pady=2)
        self.next_btn = ctk.CTkButton(
    button_frame,
    text='Next',
    command=self.next_item,
    width=50,
    height=20,
    corner_radius=4,
     font=self.font)
        self.next_btn.pack(side='left', padx=1)
        self.copy_btn = ctk.CTkButton(
    button_frame,
    text='Copy',
    command=self.copy_item,
    width=50,
    height=20,
    corner_radius=4,
     font=self.font)
        self.copy_btn.pack(side='left', padx=1)

    def _setup_click_to_copy(self):
        """Налаштування click-to-copy для секції"""
        try:
            # Отримуємо внутрішній Tkinter віджет
            if hasattr(self.textbox, '_textbox'):
                inner_text = self.textbox._textbox
                
                # Прив'язуємо події
                inner_text.bind("<Button-1>", self._on_click_copy)
                inner_text.bind("<Button-3>", self._show_section_menu)
                inner_text.bind('<Key>', lambda event: 'break')  # Блокуємо редагування
                
                # Створюємо контекстне меню
                self.context_menu = tk.Menu(inner_text, tearoff=0, bg='#2e2e2e', fg='white', 
                                          activebackground='#404040', activeforeground='white')
                self.context_menu.add_command(label="📋 Копіювати рядок", command=self._copy_current_line)
                self.context_menu.add_command(label="📋 Копіювати все", command=self._copy_all_items)
                self.context_menu.add_separator()
                self.context_menu.add_command(label="🔤 Вибрати все", command=self._select_all_items)
                
                print(f"Click-to-copy налаштовано для секції: {self.title}")
            
        except Exception as e:
            print(f"Помилка налаштування click-to-copy для {self.title}: {e}")

    def _on_click_copy(self, event):
        """Обробка кліку - копіює рядок під курсором"""
        try:
            # Отримуємо внутрішній текстовий віджет
            inner_text = self.textbox._textbox
            
            # Отримуємо позицію кліку
            click_index = inner_text.index(f"@{event.x},{event.y}")
            line_start = inner_text.index(f"{click_index} linestart")
            line_end = inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка
            line_text = inner_text.get(line_start, line_end).strip()
            
            if line_text:
                # Копіюємо рядок через оптимізовану функцію
                copy_to_clipboard(line_text)
                
                # Візуальний фідбек
                inner_text.tag_remove('copied_line', '1.0', 'end')
                inner_text.tag_add('copied_line', line_start, line_end)
                inner_text.tag_config('copied_line', background='#404040', foreground='#90EE90')
                
                # Прибираємо підсвітку через 500мс
                inner_text.after(500, lambda: inner_text.tag_remove('copied_line', '1.0', 'end'))
                
                # Логуємо копіювання
                print(f"📋 Скопійовано з {self.title}: {line_text}")
                self.logger.log(f"📋 Скопійовано: {line_text}")
                
        except Exception as e:
            print(f"Помилка копіювання в {self.title}: {e}")
            
    def _show_section_menu(self, event):
        """Показати контекстне меню секції"""
        try:
            self.context_menu.post(event.x_root, event.y_root)
        except:
            pass
            
    def _copy_current_line(self):
        """Копіювати поточний рядок"""
        try:
            inner_text = self.textbox._textbox
            current_pos = inner_text.index(tk.INSERT)
            line_start = inner_text.index(f"{current_pos} linestart")
            line_end = inner_text.index(f"{current_pos} lineend")
            line_text = inner_text.get(line_start, line_end).strip()
            
            if line_text:
                copy_to_clipboard(line_text)
                self.logger.log(f"📋 Скопійовано рядок: {line_text}")
        except:
            pass
            
    def _copy_all_items(self):
        """Копіювати всі елементи секції"""
        try:
            inner_text = self.textbox._textbox
            all_text = inner_text.get('1.0', 'end-1c')
            copy_to_clipboard(all_text)
            self.logger.log(f"📋 Скопійовано всю секцію {self.title}")
        except:
            pass
            
    def _select_all_items(self):
        """Вибрати всі елементи"""
        try:
            inner_text = self.textbox._textbox
            inner_text.tag_add(tk.SEL, '1.0', tk.END)
            inner_text.focus_set()
        except:
            pass

    def generate_items(self):
        # sample N items from column_values
        if not self.column_values:
            self.items = []
            self._refresh_listbox()
            self.logger.log(f"{self.title}: No data available to generate")
            return
            self._refresh_listbox()
            self.logger.log(f"{self.title}: No source values in CSV")
            return
        n = random.randint(self.min_items, self.max_items)
        vals = list(self.column_values)
        random.shuffle(vals)
        # If not enough distinct entries, allow repeats by sampling with replacement
        if len(vals) >= n:
            chosen = vals[:n]
        else:
            chosen = [random.choice(self.column_values) for _ in range(n)]
        # If append_values are provided, append one random append word to each
        # chosen site (крім Email Subscription - там потрібні лише лінки)
        if self.append_values and self.title != 'Email Subscription':
            combined = []
            for c in chosen:
                extra = random.choice(self.append_values)
                combined.append(f"{c} {extra}")
            self.items = combined
        else:
            self.items = chosen
            
        # Додаємо граматичні помилки для Prompts
        if self.title == 'Prompts':
            self.items = [self._add_typos_to_text_for_section(item) for item in self.items]
        self.index = 0
        self._refresh_listbox()
        self.logger.log(f"{self.title}: Generated {len(self.items)} items")

    def _add_typos_to_text_for_section(self, text):
        """Додає граматичні помилки в текст для унікальності"""
        if not text or len(text) < 10:
            return text
        
        words = text.split()
        if len(words) < 2:
            return text
        
        # Збільшуємо кількість помилок: 2-5 залежно від довжини тексту
        max_typos = min(5, max(2, len(words) // 3))
        num_typos = random.randint(2, max_typos)
        
        if num_typos == 0:
            return text
        
        # Вибираємо випадкові слова для помилок (не перше і не останнє)
        modifiable_indices = list(range(1, len(words) - 1)) if len(words) > 2 else list(range(len(words)))
        if not modifiable_indices:
            return text
            
        selected_indices = random.sample(modifiable_indices, min(num_typos, len(modifiable_indices)))
        
        for idx in selected_indices:
            word = words[idx]
            if len(word) > 3:  # Тільки довгі слова
                typo_type = random.choice([1, 2, 3, 4])  # 4 типи помилок: вставка, перестановка, заміна, подвоєння
                
                if typo_type == 1:  # Випадкова вставка символів a-z, 0-9
                    char_idx = random.randint(1, len(word) - 1)
                    chars = list(word)
                    
                    # Випадково обираємо тип символа для вставки
                    insert_type = random.choice(['letter', 'number'])
                    if insert_type == 'letter':
                        random_char = random.choice('abcdefghijklmnopqrstuvwxyz')
                    else:
                        random_char = random.choice('0123456789')
                    
                    chars.insert(char_idx, random_char)
                    words[idx] = ''.join(chars)
                
                elif typo_type == 2:  # Перестановка сусідніх символів
                    if len(word) > 3:
                        char_idx = random.randint(1, len(word) - 3)
                        chars = list(word)
                        chars[char_idx], chars[char_idx + 1] = chars[char_idx + 1], chars[char_idx]
                        words[idx] = ''.join(chars)
                
                elif typo_type == 3:  # Заміна символа на випадковий a-z, 0-9
                    char_idx = random.randint(1, len(word) - 2)
                    chars = list(word)
                    
                    # Випадково заміняємо символ на літеру або цифру
                    replace_type = random.choice(['letter', 'number'])
                    if replace_type == 'letter':
                        chars[char_idx] = random.choice('abcdefghijklmnopqrstuvwxyz')
                    else:
                        chars[char_idx] = random.choice('0123456789')
                    
                    words[idx] = ''.join(chars)
                
                elif typo_type == 4:  # Подвоєння символа
                    char_idx = random.randint(1, len(word) - 2)
                    chars = list(word)
                    chars.insert(char_idx, chars[char_idx])
                    words[idx] = ''.join(chars)
        
        return ' '.join(words)

    def _refresh_listbox(self):
        try:
            # Перевіряємо, чи існує textbox
            if not hasattr(self, 'textbox') or not self.textbox.winfo_exists():
                return
                
            self.textbox.configure(state='normal')
            self.textbox.delete('0.0', 'end')
            for it in self.items:
                self.textbox.insert('end', it + '\n')
            self.textbox.configure(state='disabled')
            
            # Перевіряємо, чи існує current_label
            if hasattr(self, 'current_label') and self.current_label.winfo_exists():
                if self.items:
                    self.current_label.configure(text=f"Selected: {self.items[self.index]}")
                else:
                    self.current_label.configure(text="Selected: None")
        except Exception as e:
            # Якщо помилка GUI, просто ігноруємо
            print(f"GUI refresh error in {self.title}: {e}")

    def next_item(self):
        try:
            if not self.items:
                return
            self.index = (self.index + 1) % len(self.items)
            if hasattr(self, 'current_label') and self.current_label.winfo_exists():
                self.current_label.configure(text=f"Selected: {self.items[self.index]}")
            self.logger.log(f"{self.title}: Next -> {self.items[self.index]}")
        except Exception as e:
            print(f"Error in next_item for {self.title}: {e}")

    def copy_item(self):
        sel = self.get_selected()
        if sel is None:
            return
        self._copy_to_clipboard(sel)
        self.logger.log(f"{self.title}: Copied -> {sel}")

    def get_selected(self):
        return self.items[self.index] if self.items else None

    def _copy_to_clipboard(self, text):
        safe_clipboard_operation("set", text)

    def save_selection(self):
        sel = self.get_selected()
        if sel is None:
            messagebox.showinfo('Save', 'No selection to save')
            return
        # save to a small JSON profile file
        profile = {
            'section': self.title,
            'value': sel
        }
        path = filedialog.asksaveasfilename(
    defaultextension='.json', filetypes=[
        ('JSON', '*.json')])
        if not path:
            return
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(profile, f, ensure_ascii=False, indent=2)
        self.logger.log(f"{self.title}: Saved selection to {os.path.basename(path)}")


class Logger:
    def __init__(self, text_widget):
        self.text = text_widget
        # Якщо це CustomTkinter textbox, отримуємо внутрішній Tkinter віджет
        if hasattr(text_widget, '_textbox'):
            self.inner_text = text_widget._textbox
        else:
            self.inner_text = text_widget
        self.setup_context_menu()

    def setup_context_menu(self):
        """Налаштування контекстного меню та single-click копіювання"""
        if self.inner_text:
            # Створюємо контекстне меню
            self.context_menu = tk.Menu(self.inner_text, tearoff=0, bg='#2e2e2e', fg='white', 
                                      activebackground='#404040', activeforeground='white')
            self.context_menu.add_command(label="📋 Копіювати рядок", command=self.copy_current_line)
            self.context_menu.add_command(label="📋 Копіювати вибране", command=self.copy_selected)
            self.context_menu.add_command(label="📋 Копіювати все", command=self.copy_all)
            self.context_menu.add_separator()
            self.context_menu.add_command(label="🔤 Вибрати все", command=self.select_all)
            self.context_menu.add_command(label="🗑 Очистити лог", command=self.clear_log)
            
    def show_context_menu(self, event):
        """Показати контекстне меню"""
        try:
            self.context_menu.post(event.x_root, event.y_root)
        except:
            pass
            
    def on_single_click(self, event):
        """Одинарний клік - миттєво копіює рядок під курсором"""
        try:
            # Отримуємо позицію кліку
            click_index = self.inner_text.index(f"@{event.x},{event.y}")
            line_start = self.inner_text.index(f"{click_index} linestart")
            line_end = self.inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка
            line_text = self.inner_text.get(line_start, line_end).strip()
            
            if line_text:
                # Копіюємо рядок
                copy_to_clipboard(line_text)
                
                # Візуальний фідбек - підсвітка рядка
                self.inner_text.tag_remove('copied_line', '1.0', 'end')
                self.inner_text.tag_add('copied_line', line_start, line_end)
                self.inner_text.tag_config('copied_line', background='#404040', foreground='#90EE90')
                
                # Прибираємо підсвітку через 500мс
                self.inner_text.after(500, lambda: self.inner_text.tag_remove('copied_line', '1.0', 'end'))
                
                # Показуємо повідомлення про копіювання
                self.show_copy_feedback(line_text)
                
        except Exception as e:
            print(f"Помилка копіювання рядка: {e}")
            
    def show_copy_feedback(self, text):
        """Показати фідбек про копіювання"""
        short_text = text[:30] + "..." if len(text) > 30 else text
        timestamp = datetime.datetime.now().strftime('%H:%M:%S')
        print(f"[{timestamp}] 📋 Скопійовано: {short_text}")
            
    def on_double_click(self, event):
        """Подвійний клік - виділяє слово (стандартна поведінка)"""
        # Дозволяємо стандартне виділення слова
        return
        
    def copy_current_line(self):
        """Копіювати рядок під курсором"""
        try:
            # Отримуємо поточну позицію курсору
            current_pos = self.inner_text.index(tk.INSERT)
            line_start = self.inner_text.index(f"{current_pos} linestart")
            line_end = self.inner_text.index(f"{current_pos} lineend")
            
            # Отримуємо текст рядка
            line_text = self.inner_text.get(line_start, line_end).strip()
            
            if line_text:
                copy_to_clipboard(line_text)
                print("📋 Скопійовано поточний рядок")
            else:
                print("⚠ Рядок порожній")
        except Exception as e:
            print("❌ Помилка копіювання рядка")
            
    def copy_selected(self):
        """Копіювати вибраний текст"""
        try:
            selected_text = self.inner_text.get(tk.SEL_FIRST, tk.SEL_LAST)
            copy_to_clipboard(selected_text)
            print("📋 Скопійовано вибраний текст")
        except tk.TclError:
            # Нічого не вибрано
            print("⚠ Нічого не вибрано для копіювання")
            
    def copy_all(self):
        """Копіювати весь текст логу"""
        try:
            all_text = self.inner_text.get('1.0', 'end-1c')
            copy_to_clipboard(all_text)
            print("📋 Скопійовано весь лог")
        except:
            print("❌ Помилка копіювання")
            
    def select_all(self):
        """Вибрати весь текст"""
        self.inner_text.tag_add(tk.SEL, '1.0', tk.END)
        self.inner_text.mark_set(tk.INSERT, '1.0')
        self.inner_text.see(tk.INSERT)
        self.inner_text.focus_set()
        
    def clear_log(self):
        """Очистити лог"""
        self.inner_text.delete('1.0', tk.END)
        self.log("🗑 Лог очищено")

    def log(self, msg):
        if self.text:
            timestamp = datetime.datetime.now().strftime('%H:%M:%S')
            # Текст завжди в normal стані, просто додаємо
            self.text.insert('end', f'[{timestamp}] {msg}\n')
            self.text.see('end')


# ================== Helper Functions ==================

def parse_accounts(content):
    """Parse accounts from text content - supports multiple formats:
    - email only: user@gmail.com
    - email:password: user@gmail.com:pass123
    - email:password:recovery: user@gmail.com:pass123:backup@mail.com
    """
    import re
    accounts_data = []
    lines = content.split('\n')
    
    # Email regex pattern
    email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Варіант 1: є двокрапка (email:password або email:password:recovery)
        if ':' in line:
            parts = line.split(':')
            if len(parts) >= 2:
                email = parts[0].strip()
                # Перевіряємо чи це валідний email
                if re.match(email_pattern, email):
                    account = {
                        'email': email,
                        'password': parts[1].strip() if len(parts) > 1 else '',
                        'recovery': parts[2].strip() if len(parts) > 2 else '',
                        'status': 'Не перевірено',
                        'has_2fa': False,
                        'phone': '',
                        'notes': ''
                    }
                    accounts_data.append(account)
        else:
            # Варіант 2: тільки email (без пароля)
            email = line.strip()
            if re.match(email_pattern, email):
                account = {
                    'email': email,
                    'password': '',
                    'recovery': '',
                    'status': 'Не перевірено',
                    'has_2fa': False,
                    'phone': '',
                    'notes': ''
                }
                accounts_data.append(account)
    
    return accounts_data


def check_emails_api(emails, progress_callback=None):
    """
    Check emails via gmailchecklive.com API (reverse engineered)
    API Endpoint: POST https://www.gmailchecklive.com/index.php
    """
    import requests
    import time
    import hashlib
    import random
    
    print(f"\n{'='*80}")
    print(f"🚀 check_emails_api STARTED")
    print(f"📧 Input emails count: {len(emails)}")
    print(f"{'='*80}\n")
    
    live = []
    die = []
    
    # 🔍 REVERSE ENGINEERED API ENDPOINT
    api_url = "https://www.gmailchecklive.com/index.php"
    
    # Headers для імітації браузера
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/141.0.0.0 Safari/537.36',
        'Accept': '*/*',
        'Accept-Language': 'en-US,en;q=0.9,uk;q=0.8',
        'Origin': 'https://www.gmailchecklive.com',
        'Referer': 'https://www.gmailchecklive.com/',
        'sec-ch-ua': '"Google Chrome";v="141", "Not?A_Brand";v="8", "Chromium";v="141"',
        'sec-ch-ua-mobile': '?0',
        'sec-ch-ua-platform': '"Windows"'
    }
    
    print(f"🔍 Checking {len(emails)} emails via gmailchecklive.com API...")
    
    try:
        # Використовуємо статичний токен зі сторінки (window._k9x)
        _t = 'bc56fbf0'
        
        # Об'єднуємо всі email-и
        emails_text = '\n'.join(emails)
        
        # Формуємо multipart/form-data (як на сайті)
        files = {
            'emails': (None, emails_text),
            'original_lines': (None, emails_text),
            '_t': (None, _t),
            'chunk_id': (None, 'chunk_1'),
            'chunk_total': (None, '1')
        }
        
        print(f"📡 Sending {len(emails)} emails to API...")
        
        # Відправляємо запит
        response = requests.post(
            api_url,
            files=files,
            headers=headers,
            timeout=30
        )
        
        print(f"📊 Response status: {response.status_code}")
        
        if response.status_code == 200:
            try:
                data = response.json()
                print(f"✅ Response data: {data}")
                
                # 🎯 ПРАВИЛЬНИЙ ФОРМАТ ВІДПОВІДІ:
                # {
                #   "success": true,
                #   "results": {
                #     "email@gmail.com": true,   // true = LIVE
                #     "email2@gmail.com": false  // false = DIE
                #   }
                # }
                
                if isinstance(data, dict) and 'results' in data:
                    results = data['results']
                    
                    # Парсимо кожен email
                    for email, is_live in results.items():
                        if is_live is True:
                            live.append(email)
                        elif is_live is False:
                            die.append(email)
                        else:
                            # Якщо незрозумілий статус - додаємо в die
                            die.append(email)
                
                else:
                    # Fallback: якщо формат інший
                    print("⚠️  Unexpected response format")
                    for email in emails:
                        die.append(email)
                
            except Exception as e:
                print(f"⚠️ Error parsing JSON response: {e}")
                print(f"Raw response: {response.text[:500]}")
                # Додаємо всі в die якщо не вдалося розпарсити
                for email in emails:
                    die.append(email)
        
        elif response.status_code == 401:
            print(f"❌ Invalid token! Need to update _t token")
            print(f"Response: {response.text}")
            # Додаємо всі в die
            for email in emails:
                die.append(email)
        
        else:
            print(f"❌ API returned status {response.status_code}")
            print(f"Response: {response.text[:200]}")
            # Додаємо всі в die
            for email in emails:
                die.append(email)
        
        # Оновлюємо прогрес (progress, current, total)
        if progress_callback:
            try:
                progress_callback(100, len(emails), len(emails))
            except Exception as prog_err:
                print(f"⚠️ Progress callback error (ignored): {prog_err}")
    
    except Exception as e:
        print(f"❌ Network Error: {e}")
        import traceback
        traceback.print_exc()
        # При помилці мережі додаємо всі в die ТІЛЬКИ якщо списки порожні
        if not live and not die:
            for email in emails:
                die.append(email)
    
    print(f"\n{'='*80}")
    print(f"📊 FINAL RESULTS:")
    print(f"   LIVE: {len(live)} emails")
    print(f"   DIE: {len(die)} emails")
    print(f"   LIVE list: {live[:5]}..." if len(live) > 5 else f"   LIVE list: {live}")
    print(f"   DIE list: {die[:5]}..." if len(die) > 5 else f"   DIE list: {die}")
    print(f"{'='*80}\n")
    
    return live, die


def get_csv_path():
    """Get path for CSV file"""
    import os
    return os.path.join(os.path.expanduser("~"), "Desktop", "gmail_parser_results.csv")


def save_to_csv(csv_path, live, die):
    """Save results to CSV file"""
    import csv
    
    with open(csv_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(['Email', 'Status'])
        
        for email in live:
            writer.writerow([email, 'Live'])
        
        for email in die:
            writer.writerow([email, 'Die'])


class ScrollableFrame(ttk.Frame):
    def __init__(self, container, *args, **kwargs):
        super().__init__(container, *args, **kwargs)
        self.canvas = tk.Canvas(self)
        self._window_id = None
        self._is_destroyed = False
        self.scrollbar = ttk.Scrollbar(
    self, orient="vertical", command=self.canvas.yview)
        self.scrollable_frame = ttk.Frame(self.canvas)
        self.scrollable_frame.bind(
            "<Configure>",
            self._on_frame_configure
        )
        self._window_id = self.canvas.create_window(
    (0, 0), window=self.scrollable_frame, anchor="nw")
        self.canvas.configure(yscrollcommand=self.scrollbar.set)
        self.canvas.pack(side="left", fill="both", expand=True)
        self.scrollbar.pack(side="right", fill="y")
        self.scrollable_frame.bind('<Enter>', self._bind_to_mousewheel)
        self.scrollable_frame.bind('<Leave>', self._unbind_from_mousewheel)
    
    def _on_frame_configure(self, event):
        """Безпечне оновлення scrollregion"""
        try:
            if not self._is_destroyed and self.canvas.winfo_exists():
                self.canvas.configure(scrollregion=self.canvas.bbox("all"))
        except tk.TclError:
            pass

    def _bind_to_mousewheel(self, event):
        try:
            if not self._is_destroyed and self.canvas.winfo_exists():
                self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)
        except tk.TclError:
            pass

    def _unbind_from_mousewheel(self, event):
        try:
            if not self._is_destroyed and self.canvas.winfo_exists():
                self.canvas.unbind_all("<MouseWheel>")
        except tk.TclError:
            pass

    def _on_mousewheel(self, event):
        try:
            if not self._is_destroyed and self.canvas.winfo_exists():
                # Використовуємо глобальну швидкість скролу
                scroll_speed = get_global_scroll_speed()
                scroll_amount = int(-1 * (event.delta / 120)) * scroll_speed
                self.canvas.yview_scroll(scroll_amount, "units")
                print(f"🔄 ScrollableFrame скрол: {scroll_amount} (швидкість: {scroll_speed}x)")
        except tk.TclError:
            pass
    
    def destroy(self):
        """Безпечне видалення з очищенням прив'язок"""
        self._is_destroyed = True
        try:
            self._unbind_from_mousewheel(None)
        except:
            pass
        try:
            if self.canvas.winfo_exists():
                self.canvas.unbind_all("<MouseWheel>")
        except:
            pass
        super().destroy()


# ================== Gmail Hacks ==================
# (Removed duplicate class - using main version below)


class GmailHacksTab(ctk.CTkFrame):
    AUTOSAVE_FILE = get_config_path("account_info_autosave.json")
    LOCAL_CSV_CONFIG = get_config_path("sheet_local_data.csv")
    LOCAL_PROFILES_CSV = get_config_path("profiles_data.csv")  # Локальний CSV з профілями
    LOG_FILE = get_config_path("account_manager_log.txt")
    DATABASE_URL = "https://happeening.com/db/db-cloudflare.php"

    def __init__(self, master, font=None):
        super().__init__(master)
        self.font = font
        self.pack(fill="both", expand=True)
        self.octo_indicator = None  # Індикатор для Octo Browser hotkey
        self.octo_toggle_button = None  # Кнопка для toggle hotkey
        self.hotkey_enabled = False  # Стан hotkey (вимкнено за замовчуванням)
        self.conversion_monitor_active = False  # Стан автомоніторингу конверсії
        self.conversion_monitor_thread = None  # Потік для моніторингу
        self.conversion_entry = None  # Посилання на поле конверсії
        self.conversion_auto_btn = None  # Кнопка автопошуку
        self.last_clipboard_content = ""  # Останній вміст буфера
        self.card_autofill_enabled = False  # Стан автозаповнення картки
        self.card_entry = None  # Посилання на поле кредитної картки
        self.card_toggle_btn = None  # Кнопка toggle автозаповнення
        self.card_sequence = []  # Послідовність даних для вставки [номер, дата, cvv]
        self.card_sequence_index = 0  # Поточний індекс в послідовності
        self.original_clipboard = ""  # Оригінальний вміст буфера
        self.card_paste_in_progress = False  # Флаг блокування подвійної обробки
        self.last_paste_time = 0  # Час останньої вставки для debounce
        self.generator_visible = True  # Стан видимості Profile Generator
        self.generator_content_frame = None  # Фрейм контенту генератора
        self.generator_toggle_btn = None  # Кнопка згортання генератора
        self.main_container = None  # Головний контейнер з grid
        self.left_frame = None  # Ліва колонка
        self.right_frame = None  # Права колонка
        self.multi_org_mode = False  # Режим відображення 4 організацій
        self.multi_org_toggle_btn = None  # Кнопка toggle для 4 орг
        self.org_displays = []  # Список з 4 текстових полів для орг
        self.multi_org_start_index = 0  # Початковий індекс для відображення 4 орг
        self.multi_org_nav_frame = None  # Фрейм для стрілочок навігації
        self.saved_single_org_state = None  # Збережений стан для Undo
        self.last_moved_folders = []  # Останні переміщені папки для Undo
        self.setup_ui()
        self.init_csv_config()  # Ініціалізуємо CSV конфіг

    def setup_octo_hotkey(self):
        """Налаштовує глобальний hotkey для Octo Browser (з підтримкою custom hotkeys)"""
        try:
            # Завантажуємо custom hotkey з конфігу
            config_file = get_config_path("hotkeys_config.json")
            octo_key = "f"  # За замовчуванням
            try:
                if os.path.exists(config_file):
                    with open(config_file, 'r', encoding='utf-8') as f:
                        hotkeys_config = json.load(f)
                        octo_key = hotkeys_config.get("octo_browser", "f")
            except:
                pass
            
            # Зберігаємо поточний hotkey
            self.current_octo_hotkey = octo_key
            
            keyboard = get_keyboard()
            if keyboard and IS_WINDOWS:
                try:
                    keyboard.unhook_key(octo_key)
                except:
                    pass
                keyboard.add_hotkey(octo_key, self.toggle_octo_browser_silent, suppress=False)
                print(f"✓ Глобальний hotkey '{octo_key.upper()}' для Octo Browser активовано")
                if self.octo_indicator:
                    self.octo_indicator.configure(text_color="#4CAF50")
                self.hotkey_enabled = True
                if self.octo_toggle_button:
                    self.octo_toggle_button.configure(text=f"🐙 Hotkey ON [{octo_key.upper()}]", fg_color="#4CAF50", hover_color="#45a049")
        except Exception as e:
            print(f"⚠️ Не вдалося налаштувати hotkey для Octo Browser: {e}")
            if self.octo_indicator:
                self.octo_indicator.configure(text_color="#f44336")
            self.hotkey_enabled = False
            if self.octo_toggle_button:
                octo_key = getattr(self, 'current_octo_hotkey', 'f')
                self.octo_toggle_button.configure(text=f"🐙 Hotkey OFF [{octo_key.upper()}]", fg_color="#f44336", hover_color="#d32f2f")

    def disable_octo_hotkey(self):
        """Вимикає глобальний hotkey для Octo Browser"""
        try:
            octo_key = getattr(self, 'current_octo_hotkey', 'f')
            keyboard = get_keyboard()
            if keyboard and IS_WINDOWS:
                try:
                    keyboard.remove_hotkey(octo_key)
                    print(f"✓ Глобальний hotkey '{octo_key.upper()}' для Octo Browser вимкнено")
                except:
                    try:
                        keyboard.unhook_key(octo_key)
                        print(f"✓ Глобальний hotkey '{octo_key.upper()}' для Octo Browser вимкнено (unhook)")
                    except Exception as e:
                        print(f"⚠️ Помилка вимкнення hotkey: {e}")
                
                if self.octo_indicator:
                    self.octo_indicator.configure(text_color="gray50")
                self.hotkey_enabled = False
                if self.octo_toggle_button:
                    self.octo_toggle_button.configure(text=f"🐙 Hotkey OFF [{octo_key.upper()}]", fg_color="#f44336", hover_color="#d32f2f")
        except Exception as e:
            print(f"⚠️ Не вдалося вимкнути hotkey для Octo Browser: {e}")

    def toggle_hotkey_state(self):
        """Перемикає стан hotkey (увімкнено/вимкнено)"""
        if self.hotkey_enabled:
            self.disable_octo_hotkey()
        else:
            self.setup_octo_hotkey()

    def start_conversion_monitor(self):
        """Запускає автоматичний моніторинг буфера для конверсії"""
        if self.conversion_monitor_active:
            return
        
        self.conversion_monitor_active = True
        if self.conversion_auto_btn:
            self.conversion_auto_btn.configure(text="⚡", fg_color="#4CAF50", hover_color="#45a049")
        
        def monitor_clipboard():
            """Моніторить буфер обміну та автоматично обробляє конверсію"""
            import re
            while self.conversion_monitor_active:
                try:
                    current_clipboard = get_from_clipboard()
                    
                    # Перевіряємо чи змінився буфер і чи містить він паттерн конверсії
                    if current_clipboard != self.last_clipboard_content:
                        self.last_clipboard_content = current_clipboard
                        
                        # Шукаємо паттерн конверсії в буфері
                        pattern = r"'send_to':\s*'AW-([^']+)'"
                        match = re.search(pattern, current_clipboard)
                        
                        if match:
                            conversion_id = match.group(1)
                            print(f"✅ Автопошук: Знайдено конверсію: {conversion_id}")
                            
                            # Оновлюємо поле в UI потоці
                            if self.conversion_entry:
                                self.after(0, lambda: self.update_conversion_field(conversion_id))
                    
                    time.sleep(0.5)  # Перевірка кожні 0.5 секунди
                except Exception as e:
                    print(f"⚠️ Помилка моніторингу буфера: {e}")
                    time.sleep(1)
        
        self.conversion_monitor_thread = threading.Thread(target=monitor_clipboard, daemon=True)
        self.conversion_monitor_thread.start()
        print("✓ Автопошук конверсії активовано")

    def stop_conversion_monitor(self):
        """Зупиняє автоматичний моніторинг буфера"""
        self.conversion_monitor_active = False
        if self.conversion_auto_btn:
            self.conversion_auto_btn.configure(text="⚡", fg_color="gray50", hover_color="gray40")
        print("✓ Автопошук конверсії вимкнено")

    def toggle_conversion_monitor(self):
        """Перемикає стан автоматичного моніторингу"""
        if self.conversion_monitor_active:
            self.stop_conversion_monitor()
        else:
            self.start_conversion_monitor()

    def update_conversion_field(self, conversion_id):
        """Оновлює поле конверсії та копіює в буфер"""
        try:
            if self.conversion_entry:
                self.conversion_entry.delete(0, tk.END)
                self.conversion_entry.insert(0, conversion_id)
                safe_clipboard_operation("set", conversion_id)
                print(f"📋 Конверсія оброблена та скопійована: {conversion_id}")
        except Exception as e:
            print(f"⚠️ Помилка оновлення поля: {e}")

    def toggle_card_autofill(self):
        """Перемикає стан автозаповнення кредитної картки"""
        self.card_autofill_enabled = not self.card_autofill_enabled
        
        if self.card_autofill_enabled:
            if self.card_toggle_btn:
                self.card_toggle_btn.configure(text="💳", fg_color="#4CAF50", hover_color="#45a049")
            self.setup_card_autofill_hotkey()
            print("=" * 50)
            print("✅ ГЛОБАЛЬНЕ АВТОЗАПОВНЕННЯ КАРТКИ УВІМКНЕНО")
            print("📌 Натискайте клавішу V (без Ctrl) - картка заповниться автоматично!")
            print("📌 Формат: 4262890017681197 11/28 232")
            print("📌 V раз #1 → номер, V раз #2 → дата, V раз #3 → CVV")
            print("=" * 50)
        else:
            if self.card_toggle_btn:
                self.card_toggle_btn.configure(text="💳", fg_color="gray50", hover_color="gray40")
            self.disable_card_autofill_hotkey()
            print("=" * 50)
            print("❌ ГЛОБАЛЬНЕ АВТОЗАПОВНЕННЯ КАРТКИ ВИМКНЕНО")
            print("📌 Клавіша V працює як звичайна")
            print("=" * 50)

    def setup_card_autofill_hotkey(self):
        """Встановлює глобальний перехоплювач V для автозаповнення карток"""
        try:
            keyboard = get_keyboard()
            if keyboard and IS_WINDOWS:
                # Видаляємо старий hotkey якщо є
                try:
                    keyboard.remove_hotkey('v')
                except:
                    pass
                
                # Встановлюємо hotkey з suppress=True щоб перехопити лише V
                keyboard.add_hotkey('v', self.handle_card_paste_sequence, suppress=True)
                print("✓ Глобальний перехоплювач V активовано (без Ctrl)")
        except Exception as e:
            print(f"⚠️ Не вдалося встановити hotkey для карток: {e}")

    def disable_card_autofill_hotkey(self):
        """Вимикає глобальний перехоплювач V"""
        try:
            keyboard = get_keyboard()
            if keyboard and IS_WINDOWS:
                try:
                    keyboard.remove_hotkey('v')
                    print("✓ Глобальний перехоплювач V вимкнено")
                except:
                    pass
                
                # Очищаємо послідовність
                self.card_sequence = []
                self.card_sequence_index = 0
        except Exception as e:
            print(f"⚠️ Помилка вимкнення hotkey: {e}")

    def handle_card_paste_sequence(self):
        """Обробляє послідовне натискання V для вставки частин картки"""
        try:
            import time
            
            # DEBOUNCE: ігноруємо повторні натискання швидше ніж 300ms (0.3 сек)
            current_time = time.time()
            time_since_last = current_time - self.last_paste_time
            
            if time_since_last < 0.3:
                print(f"⚠️ Debounce: ігноруємо швидке повторне натискання (пройшло {time_since_last*1000:.0f}ms, потрібно 300ms)")
                return
            
            print(f"⏱️ Часова перевірка OK: пройшло {time_since_last*1000:.0f}ms з попереднього натискання")
            self.last_paste_time = current_time
            
            # LOCK: якщо вже обробляємо - ігноруємо
            if self.card_paste_in_progress:
                print("⚠️ Lock: вставка вже виконується, ігноруємо подвійний виклик")
                return
            
            self.card_paste_in_progress = True
            
            keyboard = get_keyboard()
            
            # Якщо послідовність порожня - перевіряємо буфер на наявність картки
            if not self.card_sequence:
                clipboard_text = get_from_clipboard().strip()
                
                # Перевіряємо формат картки
                import re
                pattern1 = r'^(\d{16})\s+(\d{2}/\d{2})\s+(\d{3,4})$'
                pattern2 = r'^(\d{4}\s+\d{4}\s+\d{4}\s+\d{4})\s+(\d{2}/\d{2})\s+(\d{3,4})$'
                
                match = re.match(pattern1, clipboard_text)
                if not match:
                    match = re.match(pattern2, clipboard_text)
                
                if match:
                    card_number = match.group(1).replace(' ', '')
                    expiry_date = match.group(2).replace('/', '')  # Видаляємо /
                    cvv = match.group(3)
                    
                    # Зберігаємо послідовність: [номер, дата без /, CVV]
                    self.card_sequence = [card_number, expiry_date, cvv]
                    self.card_sequence_index = 0
                    self.original_clipboard = clipboard_text
                    
                    print("\n" + "=" * 60)
                    print("💳 ДЕТЕКТОВАНО КАРТКУ В БУФЕРІ!")
                    print(f"   Послідовність: {card_number[:4]}**** → {expiry_date} → ***")
                    print(f"   V раз #1 → {card_number[:4]}****{card_number[-4:]}")
                    print(f"   V раз #2 → {expiry_date}")
                    print(f"   V раз #3 → ***")
                    print("=" * 60 + "\n")
                    
                    # Вставляємо перший елемент (номер картки) через typewrite
                    current_data = self.card_sequence[0]
                    print(f"✓ V раз #1: {current_data[:4]}****{current_data[-4:]}")
                    keyboard.write(current_data)
                    self.card_sequence_index = 1
                    self.card_paste_in_progress = False  # Розблокування
                    return
                else:
                    # Не картка - вставляємо оригінальний текст через typewrite
                    keyboard.write(clipboard_text)
                    self.card_paste_in_progress = False  # Розблокування
                    return
            
            # Якщо є послідовність - вставляємо наступний елемент
            if self.card_sequence and self.card_sequence_index < len(self.card_sequence):
                current_data = self.card_sequence[self.card_sequence_index]
                
                print(f"✓ V раз #{self.card_sequence_index + 1}: {current_data if self.card_sequence_index != 2 else '***'} (індекс: {self.card_sequence_index}/{len(self.card_sequence)})")
                
                # Вставляємо через typewrite
                keyboard.write(current_data)
                
                self.card_sequence_index += 1
                
                # Якщо це була остання вставка - очищаємо
                if self.card_sequence_index >= len(self.card_sequence):
                    print("✅ Послідовність завершено\n")
                    self.card_sequence = []
                    self.card_sequence_index = 0
                
                self.card_paste_in_progress = False  # Розблокування після кожної вставки
            else:
                # Якась помилка в послідовності - вставляємо те що в буфері
                self.card_sequence = []
                self.card_sequence_index = 0
                clipboard_text = get_from_clipboard().strip()
                keyboard.write(clipboard_text)
                self.card_paste_in_progress = False  # Розблокування
                
        except Exception as e:
            print(f"⚠️ Помилка обробки послідовності: {e}")
            # При помилці - вставляємо те що в буфері
            self.card_sequence = []
            self.card_sequence_index = 0
            self.card_paste_in_progress = False  # Розблокування при помилці
            try:
                keyboard = get_keyboard()
                clipboard_text = get_from_clipboard().strip()
                keyboard.write(clipboard_text)
            except:
                pass

    def paste_card_with_autofill(self, entry):
        """Обробляє вставку для поля кредитної картки з автоматичним Tab"""
        print("\n" + "="*60)
        print("🔍 КНОПКА PASTE ДЛЯ CREDIT CARD НАТИСНУТА")
        print(f"📊 Стан автозаповнення: {self.card_autofill_enabled}")
        
        try:
            # Отримуємо дані з буфера
            clipboard_text = get_from_clipboard().strip()
            print(f"📋 Буфер обміну: '{clipboard_text[:60]}...'")
            
            if not self.card_autofill_enabled:
                # Звичайна вставка
                print("❌ Автозаповнення ВИМКНЕНО - звичайна вставка")
                print("="*60 + "\n")
                safe_text_input(entry, clipboard_text)
                return
            
            print("✅ Автозаповнення УВІМКНЕНО - обробляємо дані")
            
            # Перевіряємо чи це схоже на дані картки (номер пробіл дата пробіл CVV)
            import re
            
            # Паттерн 1: 16 цифр без пробілів, пробіл, дата MM/YY, пробіл, 3-4 цифри CVV
            pattern1 = r'^(\d{16})\s+(\d{2}/\d{2})\s+(\d{3,4})$'
            # Паттерн 2: 16 цифр з пробілами (4 4 4 4), пробіл, дата, пробіл, CVV
            pattern2 = r'^(\d{4}\s+\d{4}\s+\d{4}\s+\d{4})\s+(\d{2}/\d{2})\s+(\d{3,4})$'
            
            match = re.match(pattern1, clipboard_text)
            if not match:
                match = re.match(pattern2, clipboard_text)
            
            if match:
                card_number = match.group(1).replace(' ', '')  # Видаляємо пробіли з номера
                expiry_date = match.group(2)
                cvv = match.group(3)
                
                print(f"✅ КАРТКУ РОЗПІЗНАНО!")
                print(f"   💳 Номер: {card_number[:4]}****{card_number[-4:]}")
                print(f"   📅 Дата: {expiry_date}")
                print(f"   🔒 CVV: ***")
                print("🚀 ЗАПУСКАЄМО АВТОЗАПОВНЕННЯ...")
                print("="*60 + "\n")
                
                # Запускаємо автозаповнення
                self.autofill_card_fields(card_number, expiry_date, cvv, entry)
            else:
                print(f"❌ ФОРМАТ НЕ ПІДХОДИТЬ")
                print(f"   Очікується: 4262890017681197 11/28 232")
                print(f"   Отримано: {clipboard_text}")
                print("   Використовується звичайна вставка")
                print("="*60 + "\n")
                # Звичайна вставка
                safe_text_input(entry, clipboard_text)
        except Exception as e:
            print(f"❌ ПОМИЛКА: {e}")
            print("="*60 + "\n")
            messagebox.showerror("Помилка", f"Не вдалося вставити дані: {e}")

    def handle_card_paste(self, event):
        """Обробляє Ctrl+V для поля кредитної картки з автоматичним Tab"""
        print("\n" + "="*60)
        print("🔍 ОБРОБНИК Ctrl+V ВИКЛИКАНО")
        print(f"📊 Стан автозаповнення: {self.card_autofill_enabled}")
        
        if not self.card_autofill_enabled:
            print("❌ Автозаповнення ВИМКНЕНО - використовується звичайна вставка")
            print("="*60 + "\n")
            return  # Дозволяємо стандартну поведінку
        
        print("✅ Автозаповнення УВІМКНЕНО - обробляємо дані")
        
        try:
            # Отримуємо дані з буфера
            clipboard_text = get_from_clipboard().strip()
            print(f"📋 Буфер обміну: '{clipboard_text[:60]}...'")
            
            # Перевіряємо чи це схоже на дані картки (номер пробіл дата пробіл CVV)
            import re
            
            # Паттерн 1: 16 цифр без пробілів, пробіл, дата MM/YY, пробіл, 3-4 цифри CVV
            pattern1 = r'^(\d{16})\s+(\d{2}/\d{2})\s+(\d{3,4})$'
            # Паттерн 2: 16 цифр з пробілами (4 4 4 4), пробіл, дата, пробіл, CVV
            pattern2 = r'^(\d{4}\s+\d{4}\s+\d{4}\s+\d{4})\s+(\d{2}/\d{2})\s+(\d{3,4})$'
            
            match = re.match(pattern1, clipboard_text)
            if not match:
                match = re.match(pattern2, clipboard_text)
            
            if match:
                card_number = match.group(1).replace(' ', '')  # Видаляємо пробіли з номера
                expiry_date = match.group(2)
                cvv = match.group(3)
                
                print(f"✅ КАРТКУ РОЗПІЗНАНО!")
                print(f"   💳 Номер: {card_number[:4]}****{card_number[-4:]}")
                print(f"   📅 Дата: {expiry_date}")
                print(f"   🔒 CVV: ***")
                print("🚀 ЗАПУСКАЄМО АВТОЗАПОВНЕННЯ...")
                print("="*60 + "\n")
                
                # Блокуємо стандартну поведінку paste
                event.widget.after(10, lambda: self.autofill_card_fields(card_number, expiry_date, cvv, event.widget))
                return "break"  # Зупиняємо стандартну обробку
            else:
                print(f"❌ ФОРМАТ НЕ ПІДХОДИТЬ")
                print(f"   Очікується: 4262890017681197 11/28 232")
                print(f"   Отримано: {clipboard_text}")
                print("   Використовується звичайна вставка")
                print("="*60 + "\n")
        except Exception as e:
            print(f"❌ ПОМИЛКА: {e}")
            print("="*60 + "\n")
        
        return None  # Дозволяємо стандартну поведінку

    def global_autofill_card(self, card_number, expiry_date, cvv):
        """Глобальне автозаповнення картки - працює будь-де"""
        try:
            import pyautogui
            
            print("✓ Крок 1: Номер картки вже в буфері, перехід до дати...")
            
            # Номер картки вже вставлений стандартним Ctrl+V
            # Просто чекаємо і переходимо далі
            self.after(1500, lambda: self.continue_card_fill_step2(card_number, expiry_date, cvv))
            
        except ImportError:
            print("⚠️ pyautogui не встановлено. Встановіть: pip install pyautogui")
        except Exception as e:
            print(f"⚠️ Помилка глобального автозаповнення: {e}")

    def autofill_card_fields(self, card_number, expiry_date, cvv, start_widget):
        """Автоматично заповнює поля картки з Tab між ними (для кнопки)"""
        try:
            import pyautogui
            
            # Вставляємо номер картки в текстове поле
            start_widget.delete(0, tk.END)
            start_widget.insert(0, card_number)
            start_widget.update()  # Примусово оновлюємо UI
            print(f"✓ Вставлено номер картки в поле: {card_number}")
            
            # Копіюємо номер картки в буфер для вставки в браузер
            safe_clipboard_operation("set", card_number)
            print(f"📋 Номер картки скопійовано в буфер: {card_number[:4]}****{card_number[-4:]}")
            
            # Чекаємо та переходимо до наступного кроку
            self.after(375, lambda: self.continue_card_fill_step2(card_number, expiry_date, cvv))
            
        except ImportError:
            print("⚠️ pyautogui не встановлено. Встановіть: pip install pyautogui")
            messagebox.showerror("Помилка", "Встановіть pyautogui: pip install pyautogui")
        except Exception as e:
            print(f"⚠️ Помилка автозаповнення картки: {e}")

    def continue_card_fill_step2(self, card_number, expiry_date, cvv):
        """Крок 2: Вставляємо номер картки через typewrite (імітація вводу)"""
        try:
            import pyautogui
            
            # Використовуємо typewrite для повільного надійного вводу
            pyautogui.typewrite(card_number, interval=0.25)  # 250ms між символами (~4/сек)
            print(f"✓ typewrite -> введено номер картки: {card_number[:4]}****{card_number[-4:]}")
            
            # Затримка перед Tab
            self.after(561, lambda: self.continue_card_fill_step3(expiry_date, cvv))
            
        except Exception as e:
            print(f"⚠️ Помилка вставки номера: {e}")

    def continue_card_fill_step3(self, expiry_date, cvv):
        """Крок 3: Tab та дата закінчення"""
        try:
            import pyautogui
            
            # Натискаємо Tab для переходу до поля дати
            pyautogui.press('tab')
            print("✓ Tab -> поле дати")
            
            # Затримка перед вставкою дати
            self.after(375, lambda: self.continue_card_fill_step4(expiry_date, cvv))
            
        except Exception as e:
            print(f"⚠️ Помилка Tab до дати: {e}")

    def continue_card_fill_step4(self, expiry_date, cvv):
        """Крок 4: Вставка дати через typewrite"""
        try:
            import pyautogui
            
            # Вводимо дату посимвольно повільніше
            pyautogui.typewrite(expiry_date.replace('/', ''), interval=0.25)  # Вводимо без слешу (~4/сек)
            print(f"✓ typewrite -> введено дату: {expiry_date}")
            
            # Затримка перед Tab до CVV
            self.after(561, lambda: self.continue_card_fill_step5(cvv))
            
        except Exception as e:
            print(f"⚠️ Помилка вставки дати: {e}")

    def continue_card_fill_step5(self, cvv):
        """Крок 5: Tab та CVV"""
        try:
            import pyautogui
            
            # Натискаємо Tab для переходу до CVV
            pyautogui.press('tab')
            print("✓ Tab -> CVV поле")
            
            # Затримка перед вставкою CVV
            self.after(375, lambda: self.continue_card_fill_step6(cvv))
            
        except Exception as e:
            print(f"⚠️ Помилка Tab до CVV: {e}")

    def continue_card_fill_step6(self, cvv):
        """Крок 6: Вставка CVV через typewrite"""
        try:
            import pyautogui
            
            # Вводимо CVV посимвольно повільніше
            pyautogui.typewrite(cvv, interval=0.25)  # 250ms між цифрами (~4/сек)
            print(f"✓ typewrite -> введено CVV: ***")
            print("✅ Автозаповнення картки завершено!")
            
        except Exception as e:
            print(f"⚠️ Помилка вставки CVV: {e}")

    def toggle_octo_browser_silent(self):
        """Toggle Octo Browser через клік по taskbar іконці - найстабільніший метод"""
        try:
            if IS_WINDOWS:
                import win32gui
                import win32con
                import win32process
                import psutil
                from ctypes import windll
                
                def callback(hwnd, windows):
                    if win32gui.IsWindowVisible(hwnd):
                        try:
                            _, pid = win32process.GetWindowThreadProcessId(hwnd)
                            proc = psutil.Process(pid)
                            if "Octo Browser.exe" == proc.name():
                                windows.append(hwnd)
                        except:
                            pass
                    return True
                
                windows = []
                win32gui.EnumWindows(callback, windows)
                
                if windows:
                    hwnd = windows[0]
                    
                    # Отримуємо інформацію про поточний стан вікна
                    is_minimized = win32gui.IsIconic(hwnd)
                    current_fg = win32gui.GetForegroundWindow()
                    is_foreground = (hwnd == current_fg)
                    
                    # Якщо вікно активне і не згорнуте - згортаємо
                    if is_foreground and not is_minimized:
                        win32gui.ShowWindow(hwnd, win32con.SW_MINIMIZE)
                        print("🐙 Octo Browser згорнуто")
                    else:
                        # Інакше розгортаємо і активуємо (симуляція кліку по taskbar)
                        # Метод 1: Restore якщо згорнуте
                        if is_minimized:
                            win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                        
                        # Метод 2: Використовуємо Shell API для активації через taskbar
                        # Це найточніше імітує клік по іконці
                        windll.user32.SwitchToThisWindow(hwnd, True)
                        time.sleep(0.1)
                        
                        # Додатково форсуємо активацію
                        win32gui.SetForegroundWindow(hwnd)
                        
                        print("🐙 Octo Browser розгорнуто")
                else:
                    print("⚠️ Octo Browser не знайдено")
        except Exception as e:
            print(f"❌ Помилка toggle Octo Browser: {e}")

    def toggle_octo_browser(self):
        """Toggle Octo Browser з повідомленнями (для кнопки)"""
        try:
            if IS_WINDOWS:
                import win32gui
                import win32con
                import win32process
                import psutil
                
                def callback(hwnd, windows):
                    if win32gui.IsWindowVisible(hwnd):
                        try:
                            _, pid = win32process.GetWindowThreadProcessId(hwnd)
                            proc = psutil.Process(pid)
                            if "Octo Browser.exe" == proc.name():
                                windows.append(hwnd)
                        except:
                            pass
                    return True
                
                windows = []
                win32gui.EnumWindows(callback, windows)
                
                if windows:
                    hwnd = windows[0]
                    if win32gui.IsIconic(hwnd):
                        win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                        win32gui.SetForegroundWindow(hwnd)
                    else:
                        win32gui.ShowWindow(hwnd, win32con.SW_MINIMIZE)
                    messagebox.showinfo("🐙 Octo Browser", "Octo Browser toggle виконано!")
                else:
                    messagebox.showwarning("⚠️ Увага", "Octo Browser не знайдено!\nПереконайтеся що Octo Browser.exe запущено.")
            else:
                messagebox.showinfo("ℹ️ Info", "Ця функція доступна тільки на Windows")
        except ImportError:
            messagebox.showerror("❌ Помилка", "Не вдалося імпортувати pywin32.\nВстановіть: pip install pywin32 psutil")
        except Exception as e:
            messagebox.showerror("❌ Помилка", f"Помилка toggle Octo Browser:\n{e}")

    def toggle_generator_visibility(self):
        """Згортає/розгортає Octo Profile Generator"""
        if self.generator_visible:
            # Згортаємо - ховаємо праву колонку повністю
            self.right_frame.grid_forget()
            # Ліва колонка займає всю ширину
            self.left_frame.grid(row=0, column=0, columnspan=2, sticky="nsew", padx=5, pady=5)
            self.generator_toggle_btn.configure(text="▶ Generator")
            self.generator_visible = False
            print("◀ Octo Profile Generator згорнуто - ліва колонка розтягнута")
        else:
            # Розгортаємо - повертаємо дві колонки
            self.left_frame.grid(row=0, column=0, sticky="nsew", padx=(0, 5), pady=5)
            self.right_frame.grid(row=0, column=1, sticky="nsew", padx=(5, 0), pady=5)
            self.generator_toggle_btn.configure(text="◀ Generator")
            self.generator_visible = True
            print("▶ Octo Profile Generator розгорнуто - дві колонки")

    def toggle_multi_org_mode(self):
        """Перемикає між звичайним режимом (1 орга) та режимом 4 орг"""
        if not self.multi_org_mode:
            # Увімкнути режим 4 орг
            self.enable_multi_org_mode()
        else:
            # Вимкнути режим 4 орг
            self.disable_multi_org_mode()

    def enable_multi_org_mode(self):
        """Показує 4 вікна з першими 4 організаціями"""
        try:
            if not self.organisation_folders or len(self.organisation_folders) == 0:
                print("⚠️ Немає організацій для відображення. Оберіть робочу папку.")
                return
            
            # Зберігаємо попередній стан для можливості повернення
            if hasattr(self, 'current_selected_object'):
                self.saved_single_org_state = {
                    'object': self.current_selected_object,
                    'folder': self.current_folder_path if hasattr(self, 'current_folder_path') else None
                }
            
            # Скидаємо індекс на початок
            self.multi_org_start_index = 0
            self.display_multi_org_group()
            
            # Оновлюємо стан
            self.multi_org_mode = True
            self.multi_org_toggle_btn.configure(text="x1", fg_color="#4CAF50", hover_color="#45a049")
            print(f"✅ Режим 4 орг увімкнено")
            
        except Exception as e:
            print(f"❌ Помилка enable_multi_org_mode: {e}")
            import traceback
            traceback.print_exc()

    def display_multi_org_group(self):
        """Відображає поточну групу з 4 організацій"""
        try:
            # Перевіряємо що org_content існує
            if not hasattr(self, 'org_content') or not self.org_content:
                print("⚠️ org_content не ініціалізовано")
                return
                
            # Безпечно очищаємо org_content
            try:
                if self.org_content.winfo_exists():
                    for widget in list(self.org_content.winfo_children()):
                        try:
                            if widget.winfo_exists():
                                widget.destroy()
                        except:
                            pass
            except:
                pass
            
            # Очищаємо список дисплеїв
            self.org_displays = []
            
            # Створюємо фрейм для навігації (стрілочки та кнопки)
            nav_frame = ctk.CTkFrame(self.org_content)
            nav_frame.pack(fill="x", pady=5)
            self.multi_org_nav_frame = nav_frame
            
            # Лейбл з інформацією про поточну групу
            total_orgs = len(self.organisation_folders)
            end_index = min(self.multi_org_start_index + 4, total_orgs)
            info_label = ctk.CTkLabel(nav_frame, 
                                     text=f"Орги {self.multi_org_start_index + 1}-{end_index} з {total_orgs}",
                                     font=ctk.CTkFont(size=10))
            info_label.pack(side="left", padx=10)
            
            # Кнопка відкриття 4 папок (у зворотньому порядку)
            open_4_btn = ctk.CTkButton(nav_frame, text="📂 Open x4", width=80, height=25,
                                      command=self.open_4_folders_reverse,
                                      font=ctk.CTkFont(size=10, weight="bold"),
                                      fg_color="#FF6B35", hover_color="#CC5529")
            open_4_btn.pack(side="left", padx=5)
            
            # Кнопка Undo для відновлення переміщених папок
            undo_btn = ctk.CTkButton(nav_frame, text="↩️ Undo Move", width=90, height=25,
                                    command=self.undo_last_move,
                                    font=ctk.CTkFont(size=10, weight="bold"),
                                    fg_color="#FFA500", hover_color="#FF8C00")
            undo_btn.pack(side="left", padx=5)
            
            # Кнопка для повернення до режиму x1
            exit_x4_btn = ctk.CTkButton(nav_frame, text="❌ Exit x4", width=70, height=25,
                                    command=self.disable_multi_org_mode,
                                    font=ctk.CTkFont(size=10, weight="bold"),
                                    fg_color="#666666", hover_color="#555555")
            exit_x4_btn.pack(side="left", padx=5)
            
            # Стрілочки
            arrows_frame = ctk.CTkFrame(nav_frame)
            arrows_frame.pack(side="right", padx=10)
            
            up_btn = ctk.CTkButton(arrows_frame, text="◀", width=30, height=25, corner_radius=3,
                                  command=self.multi_org_prev, font=ctk.CTkFont(size=12))
            up_btn.pack(side="left", padx=2)
            
            down_btn = ctk.CTkButton(arrows_frame, text="▶", width=30, height=25, corner_radius=3,
                                    command=self.multi_org_next, font=ctk.CTkFont(size=12))
            down_btn.pack(side="left", padx=2)
            
            # Беремо поточні 4 орги
            orgs_to_display = self.organisation_folders[self.multi_org_start_index:self.multi_org_start_index + 4]
            self.org_displays = []  # Очищаємо список
            
            # Контейнер для grid
            grid_container = ctk.CTkFrame(self.org_content)
            grid_container.pack(fill="both", expand=True, pady=5)
            
            # Створюємо 4 вікна (2x2 grid)
            for i, org_info in enumerate(orgs_to_display):
                row = i // 2
                col = i % 2
                
                # Фрейм для кожної орги
                org_display_frame = ctk.CTkFrame(grid_container)
                org_display_frame.grid(row=row, column=col, padx=5, pady=5, sticky="nsew")
                
                # Заголовок з номером та назвою
                header = ctk.CTkFrame(org_display_frame)
                header.pack(fill="x", padx=5, pady=5)
                
                actual_number = self.multi_org_start_index + i + 1
                ctk.CTkLabel(header, text=f"#{actual_number} {org_info['name'][:20]}...", 
                           font=ctk.CTkFont(size=10, weight="bold"),
                           anchor="w").pack(side="left", padx=5)
                
                # Текстове поле з company.txt
                textbox = ctk.CTkTextbox(org_display_frame, height=100, 
                                        font=ctk.CTkFont(size=9),
                                        wrap="word")
                textbox.pack(fill="both", expand=True, padx=5, pady=5)
                
                # Завантажуємо вміст company.txt
                company_txt_path = os.path.join(org_info['path'], "company.txt")
                # Додаємо шлях до папки організації внизу
                folder_path_line = f"\n\n{org_info['path']}"
                
                if org_info['has_company']:
                    try:
                        with open(company_txt_path, 'r', encoding='utf-8') as f:
                            content = f.read().strip()
                            if content:
                                parsed_content, _ = self.parse_postal_code(content)
                                textbox.insert("0.0", parsed_content + folder_path_line)
                            else:
                                textbox.insert("0.0", "📄 Порожній файл" + folder_path_line)
                    except Exception as e:
                        textbox.insert("0.0", f"❌ Помилка: {e}" + folder_path_line)
                else:
                    textbox.insert("0.0", "❌ company.txt відсутній" + folder_path_line)
                
                # Налаштовуємо click-to-copy для цього textbox
                # Прив'язуємо до внутрішнього _textbox віджета з безпечною перевіркою
                try:
                    if hasattr(textbox, '_textbox') and textbox.winfo_exists():
                        inner_widget = textbox._textbox
                        if hasattr(inner_widget, 'winfo_exists') and inner_widget.winfo_exists():
                            inner_widget.bind("<Button-1>", lambda event, tb=textbox: self.copy_line_on_click(event, tb))
                except:
                    pass
                
                self.org_displays.append({'frame': org_display_frame, 'textbox': textbox, 'org': org_info})
            
            # Налаштовуємо grid
            grid_container.grid_rowconfigure(0, weight=1)
            grid_container.grid_rowconfigure(1, weight=1)
            grid_container.grid_columnconfigure(0, weight=1)
            grid_container.grid_columnconfigure(1, weight=1)
            
            print(f"✅ Показано орги {self.multi_org_start_index + 1}-{end_index} з {total_orgs}")
            
        except Exception as e:
            print(f"❌ Помилка display_multi_org_group: {e}")
            import traceback
            traceback.print_exc()

    def multi_org_prev(self):
        """Показує попередні 4 орги"""
        if self.multi_org_start_index >= 4:
            self.multi_org_start_index -= 4
            self.display_multi_org_group()
        else:
            print("⚠️ Це перша група організацій")

    def multi_org_next(self):
        """Показує наступні 4 орги"""
        total_orgs = len(self.organisation_folders)
        if self.multi_org_start_index + 4 < total_orgs:
            self.multi_org_start_index += 4
            self.display_multi_org_group()
        else:
            print("⚠️ Це остання група організацій")

    def open_4_folders_reverse(self):
        """Відкриває поточні 4 папки у зворотньому порядку (4-3-2-1)"""
        try:
            if not self.org_displays:
                print("⚠️ Немає організацій для відкриття")
                return
            
            # Відкриваємо у зворотньому порядку
            for display in reversed(self.org_displays):
                org_info = display['org']
                folder_path = org_info['path']
                folder_name = org_info['name']
                
                if os.path.exists(folder_path):
                    # Відкриваємо папку залежно від ОС
                    if os.name == 'nt':  # Windows
                        os.startfile(folder_path)
                    elif sys.platform == 'darwin':  # macOS
                        subprocess.call(['open', folder_path])
                    else:  # Linux
                        subprocess.call(['xdg-open', folder_path])
                    
                    print(f"📂 Відкрито папку: {folder_name}")
                    time.sleep(0.3)  # Невелика затримка між відкриттям
                else:
                    print(f"⚠️ Папка не існує: {folder_name}")
            
            # Логування
            if hasattr(self, 'log_display'):
                self.log_display.configure(state='normal')
                self.log_display.insert('end', f"📂 Відкрито {len(self.org_displays)} папок (у зворотньому порядку)\n")
                self.log_display.configure(state='disabled')
                self.log_display.see('end')
            
        except Exception as e:
            print(f"❌ Помилка open_4_folders_reverse: {e}")
            import traceback
            traceback.print_exc()

    def disable_multi_org_mode(self):
        """Повертає звичайний режим з 1 організацією"""
        try:
            # Перевіряємо що org_content існує
            if not hasattr(self, 'org_content') or not self.org_content:
                print("⚠️ org_content не ініціалізовано")
                return
                
            # Безпечно очищаємо прив'язки перед знищенням віджетів
            if hasattr(self, 'org_displays'):
                for display_info in self.org_displays:
                    try:
                        textbox = display_info.get('textbox')
                        if textbox and textbox.winfo_exists():
                            # Відв'язуємо події
                            if hasattr(textbox, '_textbox'):
                                inner_widget = textbox._textbox
                                try:
                                    if inner_widget.winfo_exists():
                                        inner_widget.unbind("<Button-1>")
                                except:
                                    pass
                    except:
                        pass
            
            # Безпечно очищаємо org_content
            try:
                if self.org_content.winfo_exists():
                    for widget in list(self.org_content.winfo_children()):
                        try:
                            if widget.winfo_exists():
                                widget.destroy()
                        except:
                            pass
            except:
                pass
            
            # Очищаємо список дисплеїв
            self.org_displays = []
            
            # Відновлюємо стандартний UI
            # Вибір папки
            folder_frame = ctk.CTkFrame(self.org_content)
            folder_frame.pack(fill="x", pady=5)
            
            ctk.CTkLabel(folder_frame, text="Робоча папка:", font=self.font).pack(anchor="w", padx=5, pady=2)
            
            folder_path_frame = ctk.CTkFrame(folder_frame)
            folder_path_frame.pack(fill="x", padx=5, pady=2)
            
            # Відновлюємо org_folder_path у новий контейнер
            self.org_folder_path = ctk.CTkEntry(folder_path_frame, placeholder_text="Оберіть папку...", font=self.font)
            self.org_folder_path.pack(side="left", fill="x", expand=True, padx=(0,5))
            if hasattr(self, 'current_folder_path') and self.current_folder_path:
                self.org_folder_path.delete(0, tk.END)
                self.org_folder_path.insert(0, self.current_folder_path)
            
            ctk.CTkButton(folder_path_frame, text="📁", width=30, height=28,
                         command=self.select_organisation_folder).pack(side="right")
            
            # Вибір об'єкта
            selection_frame = ctk.CTkFrame(self.org_content)
            selection_frame.pack(fill="x", pady=5)
            
            ctk.CTkLabel(selection_frame, text="Обрати об'єкт:", font=self.font).pack(anchor="w", padx=5, pady=2)
            
            object_container = ctk.CTkFrame(selection_frame)
            object_container.pack(fill="x", padx=5, pady=2)
            
            # Пересоздаємо dropdown
            self.object_dropdown = ctk.CTkComboBox(object_container, values=self.object_names,
                                                  width=200, font=self.font, state="readonly",
                                                  command=self.on_object_selected)
            self.object_dropdown.pack(side="left", padx=(0, 5))
            self.object_dropdown.set("Не обрано")
            
            # Пересоздаємо стрілочки
            arrows_frame = ctk.CTkFrame(object_container)
            arrows_frame.pack(side="left", padx=5)
            
            self.obj_up_btn = ctk.CTkButton(arrows_frame, text="▲", width=25, height=20, corner_radius=3,
                                           command=self.object_up, font=ctk.CTkFont(size=10))
            self.obj_up_btn.pack(side="top", pady=1)
            
            self.obj_down_btn = ctk.CTkButton(arrows_frame, text="▼", width=25, height=20, corner_radius=3,
                                             command=self.object_down, font=ctk.CTkFont(size=10))
            self.obj_down_btn.pack(side="top", pady=1)
            
            # Статус
            self.object_status_label = ctk.CTkLabel(selection_frame, text="", font=ctk.CTkFont(size=10), text_color="gray")
            self.object_status_label.pack(anchor="w", padx=5, pady=2)
            
            # Область перегляду
            preview_frame = ctk.CTkFrame(self.org_content)
            preview_frame.pack(fill="both", expand=True, pady=5)
            
            preview_header = ctk.CTkFrame(preview_frame)
            preview_header.pack(fill="x", padx=5, pady=2)
            
            ctk.CTkLabel(preview_header, text="company.txt:", font=self.font).pack(side="left", padx=5)
            ctk.CTkButton(preview_header, text="📋 Copy All", width=80, height=25,
                         command=self.copy_org_content, font=self.font).pack(side="right", padx=5)
            
            # Пересоздаємо текстове поле
            self.org_content_display = ctk.CTkTextbox(preview_frame, height=150, font=self.font)
            self.org_content_display.pack(fill="both", expand=True, padx=5, pady=2)
            
            # Відновлюємо click-to-copy
            self._setup_org_click_to_copy()
            
            # Оновлюємо стан
            self.multi_org_mode = False
            self.multi_org_toggle_btn.configure(text="x4", fg_color="gray50", hover_color="gray40")
            print("✅ Режим 1 орг увімкнено - звичайний вигляд")
            
            # Відновлюємо попередньо вибраний об'єкт (якщо є збережений стан)
            if hasattr(self, 'saved_single_org_state') and self.saved_single_org_state:
                saved_obj = self.saved_single_org_state.get('object')
                if saved_obj and saved_obj in self.object_names:
                    self.object_dropdown.set(saved_obj)
                    self.on_object_selected(saved_obj)
                    print(f"✅ Відновлено об'єкт: {saved_obj}")
            elif hasattr(self, 'current_selected_object') and self.current_selected_object:
                self.display_selected_object()
            else:
                self.org_content_display.insert("0.0", "Оберіть об'єкт для перегляду...")
            
        except Exception as e:
            print(f"❌ Помилка disable_multi_org_mode: {e}")
            import traceback
            traceback.print_exc()

    def show_instruction(self):
        instruction_text = """
📧 GMAIL HACKS - ІНСТРУКЦІЯ З ВИКОРИСТАННЯ

🔵 ACCOUNT MANAGER:
• Завантаження з Google Sheets - введіть URL таблиці
• Profile Search - пошук по назві профілю в CSV даних  
• Ручний пошук - знаходить схожі назви профілів
• Автозаповнення - дані розподіляються по полям автоматично
• Click-to-copy - клік по полю копіює дані в буфер

📊 УВАГА: CSV файли тепер управляються через вкладку Generators!
• Для роботи з CSV файлами перейдіть на вкладку Generators
• Там є повноцінний CSV File Manager з усіма функціями

🏢 ORGANISATION CONTROL:
• Обрати папку - вибираємо директорію з організацією
• Dropdown об'єктів - список всіх підпапок
• Click-to-copy в company.txt області
• Done кнопка - переносить об'єкт в папку used/
• Open Folder - відкриває папку об'єкта

🎯 OCTO PROFILE GENERATOR:
• Profile Template - шаблон назви профілю
• Proxy Template - шаблон проксі з портом
• Автогенерація Excel/CSV файлів для OctoBrowser
• Папка Octo Import/ створюється автоматично

💡 ПРИНЦИП РОБОТИ: Завантажуємо → Парсимо → Копіюємо → Використовуємо
        """
        messagebox.showinfo("📧 Gmail Hacks - Інструкція", instruction_text)

    def setup_ui(self):
        # --- Основний layout ---
        # Заголовок і інструкція
        header_frame = ctk.CTkFrame(self)
        header_frame.pack(fill=tk.X, padx=5, pady=5)
        
        ctk.CTkLabel(header_frame, text="Gmail Hacks - Account Manager", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(side=tk.LEFT, padx=10)
        
        # Кнопка Octo Browser Hotkey Toggle з індикатором
        octo_frame = ctk.CTkFrame(header_frame, fg_color="transparent")
        octo_frame.pack(side=tk.RIGHT, padx=5)
        
        self.octo_indicator = ctk.CTkLabel(octo_frame, text="●", 
                                          font=ctk.CTkFont(size=16), 
                                          text_color="gray50")
        self.octo_indicator.pack(side=tk.LEFT, padx=(0, 3))
        
        self.octo_toggle_button = ctk.CTkButton(octo_frame, text="🐙 Hotkey OFF [F]", 
                     command=self.toggle_hotkey_state, 
                     width=150, height=28, corner_radius=6,
                     fg_color="#f44336", hover_color="#d32f2f",
                     font=self.font)
        self.octo_toggle_button.pack(side=tk.LEFT)
        
        # Кнопка згортання/розгортання Profile Generator
        self.generator_toggle_btn = ctk.CTkButton(header_frame, text="◀ Generator", 
                     command=self.toggle_generator_visibility, 
                     width=120, height=28, corner_radius=6,
                     font=self.font)
        self.generator_toggle_btn.pack(side=tk.RIGHT, padx=5)
        
        ctk.CTkButton(header_frame, text="Інструкція", 
                     command=self.show_instruction, width=100, height=28, 
                     corner_radius=6, font=self.font).pack(side=tk.RIGHT, padx=10)
        
        # Основна скролювальна область
        main_scrollable = ctk.CTkScrollableFrame(self)
        main_scrollable.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # Контейнер для двох колонок всередині скролювальної області
        main_container = ctk.CTkFrame(main_scrollable)
        main_container.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        self.main_container = main_container  # Зберігаємо посилання
        
        # Налаштування grid для контейнера
        main_container.grid_columnconfigure(0, weight=3)  # Ліва колонка більша
        main_container.grid_columnconfigure(1, weight=2)  # Права колонка менша
        main_container.grid_rowconfigure(0, weight=1)
        
        # Ліва колонка - Account Manager та Organisation Control
        left_frame = ctk.CTkFrame(main_container)
        left_frame.grid(row=0, column=0, sticky="nsew", padx=(0, 5), pady=5)
        self.left_frame = left_frame  # Зберігаємо посилання
        
        # Права колонка - Profile Generator  
        right_frame = ctk.CTkFrame(main_container)
        right_frame.grid(row=0, column=1, sticky="nsew", padx=(5, 0), pady=5)
        self.right_frame = right_frame  # Зберігаємо посилання

        # --- AUTOSAVE: load data if exists ---
        def load_autosave():
            try:
                if os.path.exists(self.AUTOSAVE_FILE):
                    with open(self.AUTOSAVE_FILE, "r", encoding="utf-8") as f:
                        return json.load(f)
            except Exception:
                pass
            return {}

        def save_autosave():
            try:
                data = {label: entry.get()
                                         for label, entry in zip([f[0] for f in fields], entries)}
                with open(self.AUTOSAVE_FILE, "w", encoding="utf-8") as f:
                    json.dump(data, f, ensure_ascii=False)
            except Exception:
                pass

        def on_entry_change(*args):
            save_autosave()

        def set_entry(entry, value):
            entry.delete(0, tk.END)
            entry.insert(0, value)

        # Налаштування pack для лівої колонки
        left_frame.grid_rowconfigure(0, weight=1)
        left_frame.grid_columnconfigure(0, weight=1)
        left_frame.grid_columnconfigure(1, weight=1)
        
        # Створюємо внутрішній контейнер для правильного layout
        left_content = ctk.CTkFrame(left_frame)
        left_content.pack(fill="both", expand=True, padx=5, pady=5)
        
        # Налаштування для двох колонок в лівому фреймі
        left_content.grid_columnconfigure(0, weight=1)
        left_content.grid_columnconfigure(1, weight=1)
        left_content.grid_rowconfigure(0, weight=1)
        
        # --- Organisation Control Section (Ліва частина) ---
        org_frame = ctk.CTkFrame(left_content)
        org_frame.grid(row=0, column=0, padx=(0, 5), pady=5, sticky="nsew")
        self.org_frame = org_frame  # Зберігаємо посилання
        
        # Header з заголовком та кнопкою toggle
        org_header = ctk.CTkFrame(org_frame)
        org_header.pack(fill="x", pady=(10,5))
        
        ctk.CTkLabel(org_header, text="📁 Organisation Control", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(side="left", padx=10)
        
        # Кнопка toggle для 4 організацій
        self.multi_org_toggle_btn = ctk.CTkButton(org_header, text="x4", 
                     command=self.toggle_multi_org_mode, 
                     width=35, height=25, corner_radius=4,
                     font=ctk.CTkFont(size=10, weight="bold"),
                     fg_color="gray50", hover_color="gray40")
        self.multi_org_toggle_btn.pack(side="right", padx=10)
        
        # Контейнер для основного контенту Organisation Control
        org_content = ctk.CTkFrame(org_frame)
        org_content.pack(fill="both", expand=True, padx=10, pady=10)
        self.org_content = org_content  # Зберігаємо посилання
        
        # Вибір папки
        folder_frame = ctk.CTkFrame(org_content)
        folder_frame.pack(fill="x", pady=5)
        
        ctk.CTkLabel(folder_frame, text="Робоча папка:", font=self.font).pack(anchor="w", padx=5, pady=2)
        
        folder_path_frame = ctk.CTkFrame(folder_frame)
        folder_path_frame.pack(fill="x", padx=5, pady=2)
        
        self.org_folder_path = ctk.CTkEntry(folder_path_frame, placeholder_text="Оберіть папку...", font=self.font)
        self.org_folder_path.pack(side="left", fill="x", expand=True, padx=(0,5))
        
        ctk.CTkButton(folder_path_frame, text="📁", width=30, height=28,
                     command=self.select_organisation_folder).pack(side="right")
        
        # Вибір об'єкта (папки)
        selection_frame = ctk.CTkFrame(org_content)
        selection_frame.pack(fill="x", pady=5)
        
        ctk.CTkLabel(selection_frame, text="Обрати об'єкт:", font=self.font).pack(anchor="w", padx=5, pady=2)
        
        # Контейнер для dropdown та стрілочок (як у Profile Search)
        object_container = ctk.CTkFrame(selection_frame)
        object_container.pack(fill="x", padx=5, pady=2)
        
        # Dropdown список об'єктів (папок)
        self.object_names = ["Не обрано"]  # Початкове значення
        self.object_dropdown = ctk.CTkComboBox(object_container, values=self.object_names,
                                              width=200, font=self.font, state="readonly",
                                              command=self.on_object_selected)
        self.object_dropdown.pack(side="left", padx=(0, 5))
        
        # Стрілочки для навігації (як у Profile Search)
        arrows_frame = ctk.CTkFrame(object_container)
        arrows_frame.pack(side="left", padx=5)
        
        self.obj_up_btn = ctk.CTkButton(arrows_frame, text="▲", width=25, height=20, corner_radius=3,
                                       command=self.object_up, font=ctk.CTkFont(size=10))
        self.obj_up_btn.pack(side="top", pady=1)
        
        self.obj_down_btn = ctk.CTkButton(arrows_frame, text="▼", width=25, height=20, corner_radius=3,
                                         command=self.object_down, font=ctk.CTkFont(size=10))
        self.obj_down_btn.pack(side="top", pady=1)
        
        # Статус обраного об'єкта
        self.object_status_label = ctk.CTkLabel(selection_frame, text="", font=ctk.CTkFont(size=10), text_color="gray")
        self.object_status_label.pack(anchor="w", padx=5, pady=2)
        
        # Область перегляду .txt файлів
        preview_frame = ctk.CTkFrame(org_content)
        preview_frame.pack(fill="both", expand=True, pady=5)
        
        preview_header = ctk.CTkFrame(preview_frame)
        preview_header.pack(fill="x", padx=5, pady=2)
        
        ctk.CTkLabel(preview_header, text="company.txt:", font=self.font).pack(side="left", padx=5)
        ctk.CTkButton(preview_header, text="📋 Copy All", width=80, height=25,
                     command=self.copy_org_content, font=self.font).pack(side="right", padx=5)
        
        self.org_content_display = ctk.CTkTextbox(preview_frame, height=150, font=self.font)
        self.org_content_display.pack(fill="both", expand=True, padx=5, pady=2)
        
        # --- Account Manager Fields (Права частина, компактні) ---
        fields_frame = ctk.CTkFrame(left_content)
        fields_frame.grid(row=0, column=1, padx=(5, 0), pady=5, sticky="nsew")
        
        ctk.CTkLabel(fields_frame, text="👤 Account Manager", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(10,5))
        
        # Контейнер для полів
        account_content = ctk.CTkFrame(fields_frame)
        account_content.pack(fill="both", expand=True, padx=10, pady=10)
        
        fields = [
            ("Email/Login", "Дані для логіна"),
            ("Password", "Дані для пароля"),
            ("2FA", "Дані для 2FA"),
            ("Backup Codes", "Дані для резервних кодів"),
            ("Reserve Mail", "Дані для резервної пошти"),
            ("Credit Card", "Дані для кредитної картки"),
            ("Конверсія", "Дані для конверсії"),
            ("API Cloudflare", "API ключ Cloudflare"),
            ("Cloudflare ID", "ID аккаунту Cloudflare"),
            ("Cloudflare Password", "Пароль Cloudflare (з @ за умовчуванням)"),
        ]


        
        entries = []
        # Оголошуємо функції зараніє
        def paste_to_entry(entry):
            try:
                clipboard_text = safe_clipboard_operation("get")
                safe_text_input(entry, clipboard_text)
            except Exception as e:
                messagebox.showerror("Помилка", f"Не вдалося вставити дані: {e}")

        def copy_from_entry(entry):
            try:
                text = entry.get()
                safe_clipboard_operation("set", text)
                # messagebox.showinfo("Успіх", "Дані скопійовано в буфер обміну")
            except Exception as e:
                messagebox.showerror("Помилка", f"Не вдалося скопіювати дані: {e}")
        
        def record_data_to_database(field_name, entry):
            """Записує дані з поля в базу даних"""
            try:
                data = entry.get().strip()
                if not data:
                    messagebox.showwarning("Попередження", f"Поле '{field_name}' порожнє")
                    return
                
                # Отримуємо поточну назву профілю
                profile_name = getattr(self, 'current_profile_name', 'Unknown Profile')
                
                # Записуємо в CSV файл
                csv_file = get_config_path("account_manager_records.csv")
                
                # Створюємо заголовки якщо файл новий
                if not os.path.exists(csv_file):
                    with open(csv_file, 'w', newline='', encoding='utf-8') as f:
                        writer = csv.writer(f)
                        writer.writerow(['Timestamp', 'Profile', 'Field', 'Data'])
                
                # Записуємо дані
                with open(csv_file, 'a', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                    writer.writerow([timestamp, profile_name, field_name, data])
                
                # Якщо це конверсія - додаємо в внутрішню пам'ять
                if field_name == "Конверсія":
                    self.add_conversion_to_memory(profile_name, data)
                
                # Показуємо підтвердження (тимчасово відключено)
                # messagebox.showinfo("Успіх", f"Дані '{field_name}' записано в базу")
                print(f"Recorded {field_name}: {data[:20]}... for profile: {profile_name}")
                
            except Exception as e:
                messagebox.showerror("Помилка", f"Не вдалося записати дані: {e}")
        
        def extract_conversion_id(text):
            """Витягує conversion ID з HTML скрипта Google Ads (regex-based, як у Lite версії)"""
            import re
            
            # Regex для пошуку 'send_to': 'AW-XXXXXXXXX'
            # Паттерн: 'send_to':\s*'AW-([^']+)'
            pattern = r"'send_to':\s*'AW-([^']+)'"
            
            match = re.search(pattern, text)
            
            if match:
                conversion_id = match.group(1)  # Витягуємо все після AW-
                print(f"✅ Parsed conversion ID: {conversion_id}")
                return conversion_id
            else:
                # Якщо не знайдено - повертаємо текст як є
                print(f"⚠️ Pattern not found, returning original text")
                return text.strip()
        
        def paste_from_buffer_to_all():
            """Вставляє дані з буфера у всі поля (parse_and_fill)"""
            parse_and_fill()
        
        # Компактне створення полів
        for i, (label, default) in enumerate(fields):
            # Контейнер для кожного поля
            field_container = ctk.CTkFrame(account_content)
            field_container.pack(fill="x", pady=2)
            
            # Спеціальна обробка для поля 2FA - додаємо кнопку генерації зліва
            if label == "2FA":
                # Кнопка генерації 2FA (зліва від лейбла) - генерує код з секретного ключа
                self.gen_2fa_btn = ctk.CTkButton(
                    field_container, text="🔑",
                    command=lambda: self.generate_2fa_code(),
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8),
                    hover_color="darkgreen", fg_color="green")
                self.gen_2fa_btn.pack(side="left", padx=(2, 1))
                
                # Створюємо простий tooltip ефект при наведенні
                def on_enter(event):
                    self.gen_2fa_btn.configure(fg_color="darkgreen")
                def on_leave(event):
                    self.gen_2fa_btn.configure(fg_color="green")
                
                self.gen_2fa_btn.bind("<Enter>", on_enter)
                self.gen_2fa_btn.bind("<Leave>", on_leave)
            
            # Спеціальна обробка для поля Credit Card - додаємо кнопку toggle автозаповнення
            if label == "Credit Card":
                # Кнопка toggle автозаповнення картки (зліва від лейбла)
                self.card_toggle_btn = ctk.CTkButton(
                    field_container, text="💳",
                    command=self.toggle_card_autofill,
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8),
                    hover_color="gray40", fg_color="gray50")
                self.card_toggle_btn.pack(side="left", padx=(2, 1))
            
            # Спеціальна обробка для поля Конверсії - додаємо кнопку парсингу HTML
            if label == "Конверсія":
                # Кнопка автопошуку конверсії в буфері (зліва від парсингу)
                self.conversion_auto_btn = ctk.CTkButton(
                    field_container, text="⚡",
                    command=self.toggle_conversion_monitor,
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8),
                    hover_color="gray40", fg_color="gray50")
                self.conversion_auto_btn.pack(side="left", padx=(2, 1))
                
                # Кнопка парсингу HTML конверсії (зліва від лейбла)
                def parse_conversion_from_buffer():
                    try:
                        clipboard_text = get_from_clipboard()
                        conversion_id = extract_conversion_id(clipboard_text)
                        entry.delete(0, tk.END)
                        entry.insert(0, conversion_id)
                        
                        # 📋 Автоматично копіюємо готовий ключ в буфер
                        safe_clipboard_operation("set", conversion_id)
                        print(f"🔗 Parsed conversion: {conversion_id}")
                        print(f"📋 Conversion key copied to clipboard: {conversion_id}")
                    except Exception as e:
                        messagebox.showerror("Помилка", f"Не вдалося розпарсити конверсію: {e}")
                
                self.parse_conv_btn = ctk.CTkButton(
                    field_container, text="🔗",
                    command=parse_conversion_from_buffer,
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8),
                    hover_color="darkorange", fg_color="orange")
                self.parse_conv_btn.pack(side="left", padx=(2, 1))
            
            # Спеціальна обробка для поля Cloudflare ID - додаємо кнопку парсингу ID з URL
            if label == "Cloudflare ID":
                # Зберігаємо індекс цього entry для використання в функції
                cf_id_entry_ref = entry
                current_index = len(entries)  # Індекс поточного entry (буде додано після циклу)
                
                def parse_cloudflare_id_from_buffer():
                    try:
                        clipboard_text = get_from_clipboard()
                        # Витягуємо Cloudflare ID з URL типу: https://dash.cloudflare.com/549658e1f7428eb8d1c7609f1a8e07f9/welcome
                        import re
                        match = re.search(r'dash\.cloudflare\.com/([a-f0-9]{32})', clipboard_text)
                        if match:
                            cf_id = match.group(1)
                        else:
                            # Якщо це просто ID - використовуємо як є
                            cf_id = clipboard_text.strip()
                        
                        # Вставляємо ID в поле Cloudflare ID
                        cf_id_entry_ref.delete(0, tk.END)
                        cf_id_entry_ref.insert(0, cf_id)
                        
                        # Автоматично заповнюємо Cloudflare Password = Password + "@"
                        # Password це поле [1] в масиві entries
                        if len(entries) > 1:  # Password entry
                            password_value = entries[1].get().strip()  # entries[1] = password_entry
                            if password_value:
                                # Cloudflare Password буде entries[9] (останнє поле)
                                # Але воно ще не додане, тому використаємо entries напряму після створення всіх полів
                                # Поки що просто виводимо в консоль
                                cloudflare_password = f"{password_value}@"
                                print(f"🔐 Cloudflare Password буде: {cloudflare_password}")
                        
                        # 📋 Автоматично копіюємо ID в буфер
                        safe_clipboard_operation("set", cf_id)
                        print(f"☁️ Parsed Cloudflare ID: {cf_id}")
                        print(f"📋 Cloudflare ID copied to clipboard: {cf_id}")
                    except Exception as e:
                        messagebox.showerror("Помилка", f"Не вдалося розпарсити Cloudflare ID: {e}")
                
                # Кнопка парсингу Cloudflare ID (зліва від лейбла)
                self.parse_cf_id_btn = ctk.CTkButton(
                    field_container, text="☁️",
                    command=parse_cloudflare_id_from_buffer,
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8),
                    hover_color="darkblue", fg_color="steelblue")
                self.parse_cf_id_btn.pack(side="left", padx=(2, 1))
            
            # Label (менший розмір)
            ctk.CTkLabel(field_container, text=label, font=ctk.CTkFont(size=10), width=80).pack(side="left", padx=2)
                
            # Entry field (компактний)
            entry = ctk.CTkEntry(field_container, width=180, height=24, font=ctk.CTkFont(size=10), placeholder_text=default)
            entry.pack(side="left", padx=2, fill="x", expand=True)
            entries.append(entry)
            
            # Зберігаємо посилання на entry поле конверсії для автоматичного оновлення
            if label == "Конверсія":
                self.conversion_entry = entry
            
            # Зберігаємо посилання на entry поле Credit Card
            if label == "Credit Card":
                self.card_entry = entry
                print("✓ Credit Card поле створено - використовуйте кнопку 📥 для вставки")
            
            # Кнопки для полів: record data, paste from buffer (центр), copy
            record_btn = ctk.CTkButton(
                field_container, text="💾",
                command=lambda l=label, e=entry: record_data_to_database(l, e),
                width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8),
                fg_color="darkgreen", hover_color="green")
            
            # Спеціальна кнопка paste для Credit Card
            if label == "Credit Card":
                paste_btn = ctk.CTkButton(
                    field_container, text="📥",
                    command=lambda e=entry: self.paste_card_with_autofill(e),
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8))
            else:
                paste_btn = ctk.CTkButton(
                    field_container, text="📥",
                    command=lambda e=entry: paste_to_entry(e),
                    width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8))
            
            copy_btn = ctk.CTkButton(
                field_container, text="📋",
                command=lambda e=entry: copy_from_entry(e),
                width=25, height=24, corner_radius=4, font=ctk.CTkFont(size=8))
            
            # Розташування кнопок: record data, paste from buffer (центр), copy
            record_btn.pack(side='left', padx=1)
            paste_btn.pack(side='left', padx=1)  # Paste from buffer в центрі
            copy_btn.pack(side='left', padx=2)

        # --- AUTOSAVE: load values into entries ---
        autosave_data = load_autosave()
        for (label, _), entry in zip(fields, entries):
            if label in autosave_data:
                set_entry(entry, autosave_data[label])

        # --- AUTOSAVE: trace changes ---
        for entry in entries:
            entry_var = tk.StringVar(value=entry.get())
            entry.configure(textvariable=entry_var)
            entry_var.trace_add(
    "write",
    lambda *args,
     e=entry: on_entry_change())

        email_entry, password_entry, fa_entry, codes_entry, backup_entry, card_entry, conversion_entry, api_cf_entry, cf_id_entry, cf_pass_entry = entries

        # Зберігаємо посилання на password і cloudflare password для функції парсингу
        self._password_entry_ref = password_entry
        self._cf_pass_entry_ref = cf_pass_entry
        self._cf_id_entry_ref = cf_id_entry

        # Оновлюємо функцію парсингу Cloudflare ID щоб вона автоматично заповнювала Cloudflare Password
        def parse_cloudflare_id_from_buffer_final():
            try:
                clipboard_text = get_from_clipboard()
                # Витягуємо Cloudflare ID з URL типу: https://dash.cloudflare.com/549658e1f7428eb8d1c7609f1a8e07f9/welcome
                import re
                match = re.search(r'dash\.cloudflare\.com/([a-f0-9]{32})', clipboard_text)
                if match:
                    cf_id = match.group(1)
                else:
                    # Якщо це просто ID - використовуємо як є
                    cf_id = clipboard_text.strip()
                
                # Вставляємо ID в поле Cloudflare ID
                self._cf_id_entry_ref.delete(0, tk.END)
                self._cf_id_entry_ref.insert(0, cf_id)
                
                # Автоматично заповнюємо Cloudflare Password = Password + "@"
                password_value = self._password_entry_ref.get().strip()
                if password_value:
                    cloudflare_password = f"{password_value}@"
                    self._cf_pass_entry_ref.delete(0, tk.END)
                    self._cf_pass_entry_ref.insert(0, cloudflare_password)
                    print(f"🔐 Auto-filled Cloudflare Password: {cloudflare_password}")
                
                # 📋 Автоматично копіюємо ID в буфер
                safe_clipboard_operation("set", cf_id)
                print(f"☁️ Parsed Cloudflare ID: {cf_id}")
                print(f"📋 Cloudflare ID copied to clipboard: {cf_id}")
            except Exception as e:
                messagebox.showerror("Помилка", f"Не вдалося розпарсити Cloudflare ID: {e}")
        
        # Оновлюємо команду кнопки ☁️
        if hasattr(self, 'parse_cf_id_btn'):
            self.parse_cf_id_btn.configure(command=parse_cloudflare_id_from_buffer_final)

        # Кнопка Parse from Buffer для всіх полів одразу
        parse_all_frame = ctk.CTkFrame(account_content)
        parse_all_frame.pack(fill="x", pady=(10, 5))
        
        ctk.CTkButton(parse_all_frame, text="📋 Paste from Buffer (All Fields)", 
                     command=paste_from_buffer_to_all, 
                     width=300, height=32, corner_radius=6, 
                     font=ctk.CTkFont(size=12, weight="bold"),
                     fg_color="blue", hover_color="darkblue").pack(pady=5)

        # --- Функції для роботи з даними ---

        def parse_and_fill():
            try:
                clipboard_text = get_from_clipboard()
                
                # Перевіряємо чи це HTML скрипт з конверсією
                if 'gtag(' in clipboard_text and 'send_to' in clipboard_text:
                    # Це HTML скрипт - парсимо тільки конверсію
                    conversion_id = extract_conversion_id(clipboard_text)
                    conversion_entry.delete(0, tk.END)
                    conversion_entry.insert(0, conversion_id)
                    
                    # 📋 Автоматично копіюємо готовий ключ в буфер
                    safe_clipboard_operation("set", conversion_id)
                    print(f"🔗 Parsed conversion from HTML: {conversion_id}")
                    print(f"📋 Conversion key auto-copied: {conversion_id}")
                else:
                    # Це табульовані дані - парсимо як зазвичай
                    parts = clipboard_text.split('\t')
                    for entry in entries:
                        entry.delete(0, tk.END)
                    if len(parts) > 0: email_entry.insert(0, parts[0].strip())
                    if len(parts) > 1: password_entry.insert(0, parts[1].strip())
                    if len(parts) > 2: fa_entry.insert(0, parts[2].strip())
                    if len(parts) > 3: codes_entry.insert(0, parts[3].strip())
                    if len(parts) > 4: backup_entry.insert(0, parts[4].strip())
                    if len(parts) > 7: card_entry.insert(0, parts[7].strip())
                    if len(parts) > 8: 
                        # Обробляємо конверсію через парсер
                        raw_conversion = parts[8].strip()
                        parsed_conversion = extract_conversion_id(raw_conversion)
                        conversion_entry.insert(0, parsed_conversion)
                        
                        # 📋 Автоматично копіюємо готовий ключ в буфер якщо це не порожній рядок
                        if parsed_conversion and parsed_conversion.strip():
                            safe_clipboard_operation("set", parsed_conversion)
                            print(f"📋 Parsed conversion from tabbed data auto-copied: {parsed_conversion}")
                
                update_log()
                    
            except Exception as e:
                messagebox.showerror("Помилка", f"Не вдалося розпарсити дані: {e}")

        def update_log():
            if hasattr(self, 'log_display'):
                self.log_display.configure(state='normal')
                self.log_display.delete('0.0', 'end')
                for entry in entries:
                    self.log_display.insert('end', f"{entry.get()}\n")
                self.log_display.configure(state='disabled')

        def copy_to_clipboard():
            all_data = "\n".join([entry.get() for entry in entries])
            safe_clipboard_operation("set", all_data)
            # messagebox.showinfo("Успіх", "Всі дані скопійовано в буфер обміну")

        # === Нижня секція для всіх додаткових елементів ===
        bottom_section = ctk.CTkFrame(left_content)
        bottom_section.grid(row=1, column=0, columnspan=2, sticky="ew", padx=5, pady=5)
        
        # Profile Search секція
        search_frame = ctk.CTkFrame(bottom_section)
        search_frame.pack(fill="x", pady=10, padx=10)
        
        ctk.CTkLabel(search_frame, text="Profile Search (CSV Name Column):", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        
        # Контейнер для dropdown та стрілочок
        profile_container = ctk.CTkFrame(search_frame)
        profile_container.pack(padx=10, pady=5, fill="x")
        
        # Dropdown список профілів
        self.profile_names = ["Nothing"]  # Початкове значення
        self.profile_dropdown = ctk.CTkComboBox(profile_container, values=self.profile_names,
                                               width=250, font=self.font, state="readonly",
                                               command=self.on_profile_selected, height=32,
                                               dropdown_font=self.font)
        self.profile_dropdown.pack(side="left", padx=(0, 5))
        
        # Кнопка оновлення dropdown
        ctk.CTkButton(profile_container, text="🔄", 
                     command=self.refresh_profile_dropdown, width=30, height=32, 
                     corner_radius=6, font=self.font).pack(side="left", padx=2)
        
        # Кнопка показу всіх профілів
        ctk.CTkButton(profile_container, text="📋", 
                     command=self.show_all_profiles, width=30, height=32, 
                     corner_radius=6, font=self.font).pack(side="left", padx=2)
        
        # Стрілочки для навігації
        arrows_frame = ctk.CTkFrame(profile_container)
        arrows_frame.pack(side="left", padx=5)
        
        self.up_btn = ctk.CTkButton(arrows_frame, text="▲", width=25, height=20, corner_radius=3,
                                   command=self.profile_up, font=ctk.CTkFont(size=10))
        self.up_btn.pack(side="top", pady=1)
        
        self.down_btn = ctk.CTkButton(arrows_frame, text="▼", width=25, height=20, corner_radius=3,
                                     command=self.profile_down, font=ctk.CTkFont(size=10))
        self.down_btn.pack(side="top", pady=1)
        
        # Кнопки швидкого вибору аккаунтів 1,2,3,4
        quick_select_frame = ctk.CTkFrame(profile_container)
        quick_select_frame.pack(side="left", padx=5)
        
        # Створюємо 4 кнопки в 2 ряди
        quick_top_frame = ctk.CTkFrame(quick_select_frame)
        quick_top_frame.pack(side="top", pady=1)
        quick_bottom_frame = ctk.CTkFrame(quick_select_frame)
        quick_bottom_frame.pack(side="top", pady=1)
        
        # Кнопки 1 і 2 (верхній ряд)
        self.quick_btn_1 = ctk.CTkButton(quick_top_frame, text="1", width=20, height=20, corner_radius=3,
                                        command=lambda: self.select_account_by_offset_from_current(0), 
                                        font=ctk.CTkFont(size=10, weight="bold"),
                                        fg_color="#4CAF50", hover_color="#45a049")
        self.quick_btn_1.pack(side="left", padx=1)
        
        self.quick_btn_2 = ctk.CTkButton(quick_top_frame, text="2", width=20, height=20, corner_radius=3,
                                        command=lambda: self.select_account_by_offset_from_current(1), 
                                        font=ctk.CTkFont(size=10, weight="bold"),
                                        fg_color="#2196F3", hover_color="#1976D2")
        self.quick_btn_2.pack(side="left", padx=1)
        
        # Кнопки 3 і 4 (нижній ряд)  
        self.quick_btn_3 = ctk.CTkButton(quick_bottom_frame, text="3", width=20, height=20, corner_radius=3,
                                        command=lambda: self.select_account_by_offset_from_current(2), 
                                        font=ctk.CTkFont(size=10, weight="bold"),
                                        fg_color="#FF9800", hover_color="#F57C00")
        self.quick_btn_3.pack(side="left", padx=1)
        
        self.quick_btn_4 = ctk.CTkButton(quick_bottom_frame, text="4", width=20, height=20, corner_radius=3,
                                        command=lambda: self.select_account_by_offset_from_current(3), 
                                        font=ctk.CTkFont(size=10, weight="bold"),
                                        fg_color="#9C27B0", hover_color="#7B1FA2")
        self.quick_btn_4.pack(side="left", padx=1)
        
        # Поле для ручного вводу (показується тільки коли обрано "Nothing")
        self.profile_search_entry = ctk.CTkEntry(search_frame, width=300, font=self.font, 
                                               placeholder_text="Введіть назву профілю для пошуку...")
        self.profile_search_entry.pack(padx=10, pady=5, fill="x")
        
        # Bind Enter key для автоматичного пошуку
        self.profile_search_entry.bind('<Return>', lambda event: self.parse_csv_data())
        
        # Bind колесика миші до dropdown
        self.profile_dropdown.bind("<MouseWheel>", self.on_mouse_wheel)
        
        # Дублікат ручного пошуку видалено - використовуйте основний пошук зверху
        
        # Кнопки управління
        btn_control_frame = ctk.CTkFrame(bottom_section)
        btn_control_frame.pack(fill="x", pady=10, padx=10)
        
        ctk.CTkButton(btn_control_frame, text="🔍 Parse CSV Data", 
                     command=self.parse_csv_data, width=150, height=32, corner_radius=6, font=self.font).pack(side='left', padx=5, pady=5)
        ctk.CTkButton(btn_control_frame, text="📋 Copy Formatted", 
                     command=self.copy_formatted_data, width=140, height=32, corner_radius=6, font=self.font).pack(side='left', padx=5, pady=5)
        
        # Залишаємо тільки Google Sheets функціональність, CSV Manager є в Generators
        
        # Google Sheets управління
        sheets_frame = ctk.CTkFrame(bottom_section)
        sheets_frame.pack(fill="x", pady=10, padx=10)
        
        ctk.CTkLabel(sheets_frame, text="Google Sheets Manager", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(5,0))
        
        sheets_btn_frame = ctk.CTkFrame(sheets_frame)
        sheets_btn_frame.pack(pady=5, fill="x")
        
        ctk.CTkButton(sheets_btn_frame, text="📄 Update from Sheets", 
                     command=self.open_google_sheets_dialog, width=150, height=32, 
                     corner_radius=6, font=self.font).pack(side='left', padx=5, pady=5)
        ctk.CTkButton(sheets_btn_frame, text="⚙️ Change URL", 
                     command=self.force_change_url_dialog, width=120, height=32, 
                     corner_radius=6, font=self.font).pack(side='left', padx=5, pady=5)
        
        # Результати CSV
        self.sheets_result_frame = ctk.CTkFrame(sheets_frame)
        self.sheets_result_frame.pack(fill="both", expand=True, padx=5, pady=5)
        
        ctk.CTkLabel(self.sheets_result_frame, text="CSV Data:", 
                    font=self.font).pack(anchor="w", padx=10, pady=(10,5))
        
        self.sheets_textarea = ctk.CTkTextbox(self.sheets_result_frame, width=600, height=150, font=self.font)
        self.sheets_textarea.pack(padx=10, pady=5, fill="both", expand=True)
        
        # Кнопки для роботи з CSV даними
        sheets_csv_btn_frame = ctk.CTkFrame(self.sheets_result_frame)
        sheets_csv_btn_frame.pack(pady=5, fill="x", padx=10)
        
        ctk.CTkButton(sheets_csv_btn_frame, text="Copy CSV Data", 
                     command=self.copy_csv_data, width=120, height=28, 
                     corner_radius=6, font=self.font).pack(side='left', padx=5)
        ctk.CTkButton(sheets_csv_btn_frame, text="Clear CSV", 
                     command=self.clear_csv_data, width=100, height=28, 
                     corner_radius=6, font=self.font).pack(side='left', padx=5)
        
        # CSV File Management винесено в Generators вкладку
        
        # Ініціалізуємо змінні для Google Sheets
        self.save_url_enabled = False
        self.saved_sheets_url = ""
        self.saved_sheet_id = ""
        self.cached_csv_data = ""  # Кеш для CSV даних
        self.sheets_config_file = get_config_path("gmail_hacks_config.json")
        self.octo_profile_config_file = get_config_path("octo_profile_last.json")
        
        # Зберігаємо посилання на entries для використання в методах
        self.entries_dict = {
            "Email/Login": email_entry,
            "Password": password_entry,
            "2FA": fa_entry,
            "Backup Codes": codes_entry,
            "Reserve Mail": backup_entry,
            "Credit Card": card_entry,
            "Конверсія": conversion_entry,
            "API Cloudflare": api_cf_entry,
            "Cloudflare ID": cf_id_entry,
            "Cloudflare Password": cf_pass_entry
        }
        
        # Внутрішня пам'ять для конверсій (профіль -> conversion_id)
        self.conversion_memory = {}
        self.load_conversion_memory()
        
        # Завантажуємо збережені налаштування
        self.load_sheets_config()
        # CSV конфіг перенесено в Generators
        
        # Ініціалізуємо стан кнопки Choose Last +1
        self.update_choose_last_button_state()
        
        # Ініціалізуємо dropdown з початковим значенням
        self.profile_dropdown.set("Nothing")
        self.on_profile_selected("Nothing")  # Показати поле для ручного вводу
        
        # Ініціалізуємо Organisation Control
        self.__init_org_control__()
        
        # Ініціалізуємо dropdown для об'єктів з початковим значенням
        self.object_dropdown.set("Не обрано")
        self.on_object_selected("Не обрано")
        
        # Додатково налаштовуємо click-to-copy для Organisation Control
        self._setup_org_click_to_copy()
        # Account Manager Log (в fields_frame)
        ctk.CTkLabel(account_content, text="Log:", font=ctk.CTkFont(size=10, weight="bold")).pack(anchor="w", pady=(10,2))
                    
        self.log_display = ctk.CTkTextbox(account_content, height=80, font=ctk.CTkFont(size=9))
        self.log_display.pack(fill="both", expand=True, pady=2)
        
        # Done кнопка під логом
        buttons_frame = ctk.CTkFrame(account_content)
        buttons_frame.pack(pady=5)
        
        done_btn = ctk.CTkButton(buttons_frame, text="✅ Done", 
                               command=self.mark_folder_as_done,
                               width=80, height=30, corner_radius=6, 
                               font=ctk.CTkFont(size=11, weight="bold"))
        done_btn.pack(side="left", padx=2)
        
        # Done x4 кнопка - переносить перші 4 орги
        done_x4_btn = ctk.CTkButton(buttons_frame, text="✅ x4", 
                               command=self.mark_first_4_as_done,
                               width=60, height=30, corner_radius=6, 
                               font=ctk.CTkFont(size=11, weight="bold"),
                               fg_color="#2196F3", hover_color="#1976D2")
        done_x4_btn.pack(side="left", padx=2)
        
        # Open Folder кнопка
        open_folder_btn = ctk.CTkButton(buttons_frame, text="📂 Open Folder", 
                                      command=self.open_selected_organisation_folder,
                                      width=120, height=30, corner_radius=6, 
                                      font=ctk.CTkFont(size=11, weight="bold"))
        open_folder_btn.pack(side="left", padx=2)
        
        # ===== Profile Generator (Права колонка) =====
        generator_frame = ctk.CTkFrame(right_frame)
        generator_frame.pack(fill="both", expand=True, padx=10, pady=10)
        
        ctk.CTkLabel(generator_frame, text="Octo Profile Generator", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(pady=(10,15))
        
        # Profile Name секція
        profile_section = ctk.CTkFrame(generator_frame)
        profile_section.pack(fill="x", padx=10, pady=8)
        
        ctk.CTkLabel(profile_section, text="Profile Name Template:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        self.octo_profile_entry = ctk.CTkEntry(profile_section, placeholder_text="Alex298(USA)", 
                                             font=self.font, height=32)
        self.octo_profile_entry.pack(fill="x", padx=10, pady=(0,10))
        
        # Proxy секція  
        proxy_section = ctk.CTkFrame(generator_frame)
        proxy_section.pack(fill="x", padx=10, pady=8)
        
        ctk.CTkLabel(proxy_section, text="Proxy Template:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        self.octo_proxy_entry = ctk.CTkEntry(proxy_section, 
                                           placeholder_text="socks5://user:pass@proxy.host.io:9298", 
                                           font=self.font, height=32)
        self.octo_proxy_entry.pack(fill="x", padx=10, pady=(0,10))
        
        # Кількість профілів секція
        count_section = ctk.CTkFrame(generator_frame)
        count_section.pack(fill="x", padx=10, pady=8)
        
        ctk.CTkLabel(count_section, text="Кількість профілів:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        self.octo_count_entry = ctk.CTkEntry(count_section, width=150, placeholder_text="20", 
                                           font=self.font, height=32)
        self.octo_count_entry.pack(anchor="w", padx=10, pady=(0,10))
        self.octo_count_entry.insert(0, "20")
        
        # Кнопки управління
        buttons_section = ctk.CTkFrame(generator_frame)
        buttons_section.pack(fill="x", padx=10, pady=15)
        
        self.choose_last_btn = ctk.CTkButton(buttons_section, text="Choose Last +1", 
                     command=self.choose_last_plus_one, width=180, height=40, 
                     corner_radius=8, font=ctk.CTkFont(size=13, weight="bold"), 
                     fg_color="orange", hover_color="darkorange")
        self.choose_last_btn.pack(pady=4)
        ctk.CTkButton(buttons_section, text="Генерувати Профілі", 
                     command=self.generate_octo_profile, width=180, height=40, 
                     corner_radius=8, font=ctk.CTkFont(size=13, weight="bold")).pack(pady=4)
        ctk.CTkButton(buttons_section, text="Відкрити папку", 
                     command=self.open_octo_folder, width=180, height=40, 
                     corner_radius=8, font=ctk.CTkFont(size=13, weight="bold")).pack(pady=4)
        
        # ===== Database Section (під Profile Generator) =====
        db_section = ctk.CTkFrame(generator_frame)
        db_section.pack(fill="x", padx=10, pady=15)
        
        ctk.CTkLabel(db_section, text="Database Upload", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(10,10))
        
        db_buttons = ctk.CTkFrame(db_section)
        db_buttons.pack(fill="x", padx=10, pady=5)
        
        ctk.CTkButton(db_buttons, text="➕ Add to Database", 
                     command=self.upload_to_database, width=170, height=40, 
                     corner_radius=8, font=ctk.CTkFont(size=12, weight="bold"),
                     fg_color="#7B2CBF", hover_color="#5A189A").pack(pady=4)
        
        ctk.CTkButton(db_buttons, text="☁️ Upload Cloudflare Accounts", 
                     command=self.upload_cloudflare_accounts, width=170, height=36, 
                     corner_radius=8, font=ctk.CTkFont(size=11, weight="bold"),
                     fg_color="#FF6B35", hover_color="#CC5529").pack(pady=2)
        
        ctk.CTkButton(db_buttons, text="💾 Save to CSV", 
                     command=self.save_account_to_csv, width=170, height=32, 
                     corner_radius=6, font=ctk.CTkFont(size=11, weight="bold"),
                     fg_color="#06A77D", hover_color="#048860").pack(pady=2)
        
        ctk.CTkButton(db_buttons, text="📂 Load from CSV", 
                     command=self.load_account_from_csv, width=170, height=32, 
                     corner_radius=6, font=ctk.CTkFont(size=11, weight="bold"),
                     fg_color="#0077B6", hover_color="#00538A").pack(pady=2)
        
        ctk.CTkButton(db_buttons, text="📋 View Log", 
                     command=self.open_log_file, width=170, height=32, 
                     corner_radius=6, font=ctk.CTkFont(size=11, weight="bold"),
                     fg_color="#444", hover_color="#555").pack(pady=2)
        
        ctk.CTkButton(db_buttons, text="📄 View Config", 
                     command=self.open_csv_config, width=170, height=32, 
                     corner_radius=6, font=ctk.CTkFont(size=11, weight="bold"),
                     fg_color="#666", hover_color="#777").pack(pady=2)
        
        # Лог панель для Database Upload
        db_log_frame = ctk.CTkFrame(db_section)
        db_log_frame.pack(fill="both", expand=True, padx=10, pady=(5, 10))
        
        ctk.CTkLabel(db_log_frame, text="Database Upload Log:", 
                    font=ctk.CTkFont(size=11, weight="bold"), 
                    anchor="w").pack(padx=5, pady=(5, 2), fill="x")
        
        # Текстове поле для логів
        self.db_log_display = ctk.CTkTextbox(db_log_frame, height=120, 
                                             font=ctk.CTkFont(size=10),
                                             fg_color="#1a1a1a")
        self.db_log_display.pack(padx=5, pady=5, fill="both", expand=True)
        self.db_log_display.insert("0.0", "💡 Лог операцій з базою даних буде відображатися тут...\n")
        self.db_log_display.configure(state="disabled")

    # ===== БЕЗПЕЧНІ ДОПОМІЖНІ МЕТОДИ =====
    def _safe_log_to_display(self, message):
        """Безпечне логування в log_display з перевіркою існування віджета"""
        try:
            if hasattr(self, 'log_display') and self.log_display.winfo_exists():
                self.log_display.configure(state='normal')
                self.log_display.insert('end', message)
                self.log_display.configure(state='disabled')
                self.log_display.see('end')
        except (tk.TclError, AttributeError):
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"⚠️ Помилка логування: {e}")
    
    def _safe_widget_update(self, widget, method_name, *args, **kwargs):
        """Безпечне оновлення віджета з перевіркою існування"""
        try:
            if widget and widget.winfo_exists():
                method = getattr(widget, method_name)
                method(*args, **kwargs)
        except (tk.TclError, AttributeError):
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"⚠️ Помилка оновлення віджета: {e}")

    def copy_csv_data(self):
        """Копіює дані CSV в буфер обміну"""
        try:
            csv_data = self.sheets_textarea.get("0.0", "end-1c")
            if csv_data.strip():
                safe_clipboard_operation("set", csv_data)
                # messagebox.showinfo("Успіх", "CSV дані скопійовано в буфер обміну")
            else:
                messagebox.showwarning("Попередження", "Немає даних для копіювання")
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка копіювання: {str(e)}")

    def clear_csv_data(self):
        """Очищає область CSV даних"""
        self.sheets_textarea.delete("0.0", "end")
        self.cached_csv_data = ""
        self.save_sheets_config()

    def open_google_sheets_dialog(self):
        """Відкриває діалог для роботи з Google Sheets"""
        # Якщо є збережений URL і галочка активна, використовуємо його
        if hasattr(self, 'save_url_enabled') and self.save_url_enabled and hasattr(self, 'saved_sheets_url') and self.saved_sheets_url:
            self.quick_update_with_saved_url()
            return
            
        # Інакше відкриваємо діалог
        self.show_sheets_dialog()
        
    def force_change_url_dialog(self):
        """Примусово відкриває діалог для зміни URL"""
        self.show_sheets_dialog()
        
    def show_sheets_dialog(self):
        """Показує діалог для введення/зміни URL Google Sheets"""
            
        # Створюємо діалогове вікно
        popup = ctk.CTkToplevel(self)
        popup.title("Google Sheets Update")
        popup.geometry("500x250")
        popup.resizable(False, False)
        popup.attributes('-topmost', True)
        
        # Центруємо вікно
        popup.transient(self)
        popup.grab_set()
        
        # Заголовок
        title_label = ctk.CTkLabel(popup, text="Завантаження даних з Google Таблиці", 
                                  font=ctk.CTkFont(size=16, weight="bold"))
        title_label.pack(pady=10)
        
        # Поле для введення URL
        url_frame = ctk.CTkFrame(popup)
        url_frame.pack(fill="x", padx=20, pady=10)
        
        ctk.CTkLabel(url_frame, text="Посилання на Google Таблицю:", 
                    font=self.font).pack(anchor="w", padx=10, pady=(10,5))
        
        self.sheets_url_entry = ctk.CTkEntry(url_frame, width=420, font=self.font,
                                           placeholder_text="https://docs.google.com/spreadsheets/d/...")
        self.sheets_url_entry.pack(padx=10, pady=(0,10))
        
        # Завантажити збережену URL якщо є
        if hasattr(self, 'saved_sheets_url') and self.saved_sheets_url:
            self.sheets_url_entry.insert(0, self.saved_sheets_url)
        
        # Галочка для збереження URL
        checkbox_frame = ctk.CTkFrame(popup)
        checkbox_frame.pack(fill="x", padx=20, pady=5)
        
        # Змінна для галочки
        if not hasattr(self, 'save_url_var'):
            self.save_url_var = tk.BooleanVar()
        self.save_url_var.set(self.save_url_enabled)
        
        save_checkbox = ctk.CTkCheckBox(checkbox_frame, text="Зберегти URL (наступного разу відкриється автоматично)", 
                                       variable=self.save_url_var, font=self.font,
                                       command=self.toggle_save_url)
        save_checkbox.pack(padx=10, pady=5)
        
        # Кнопки
        btn_frame = ctk.CTkFrame(popup)
        btn_frame.pack(pady=10)
        
        update_btn = ctk.CTkButton(btn_frame, text="Завантажити дані", 
                                  command=lambda: self.update_from_google_sheets(popup),
                                  width=140, height=32, corner_radius=6, font=self.font)
        update_btn.pack(side="left", padx=5)
        
        cancel_btn = ctk.CTkButton(btn_frame, text="Скасувати", 
                                  command=popup.destroy,
                                  width=120, height=32, corner_radius=6, font=self.font)
        cancel_btn.pack(side="left", padx=5)

    def toggle_save_url(self):
        """Перемикач збереження URL"""
        self.save_url_enabled = self.save_url_var.get()
        self.save_sheets_config()  # Автоматично зберігаємо
        
    def quick_update_with_saved_url(self):
        """Швидке оновлення з збереженим URL без діалогу"""
        try:
            # Показуємо повідомлення про завантаження
            loading_popup = self.show_loading_popup()
            
            # Формуємо URL для завантаження CSV
            csv_url = f"https://docs.google.com/spreadsheets/d/{self.saved_sheet_id}/export?format=csv&gid=0"
            
            # Завантажуємо файл напряму
            self.after(100, lambda: self.download_csv_directly(csv_url, loading_popup))
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка швидкого оновлення: {str(e)}")
    
    def show_loading_popup(self):
        """Показує popup з інформацією про завантаження"""
        loading_popup = ctk.CTkToplevel(self)
        loading_popup.title("Завантаження...")
        loading_popup.geometry("350x120")
        loading_popup.resizable(False, False)
        loading_popup.attributes('-topmost', True)
        
        # Центруємо вікно
        loading_popup.transient(self)
        loading_popup.grab_set()
        
        # Іконка та текст
        main_frame = ctk.CTkFrame(loading_popup)
        main_frame.pack(fill="both", expand=True, padx=20, pady=20)
        
        loading_label = ctk.CTkLabel(main_frame, text="⬇️ Завантажуємо дані з Google Таблиці...", 
                                   font=ctk.CTkFont(size=14, weight="bold"))
        loading_label.pack(pady=10)
        
        info_label = ctk.CTkLabel(main_frame, text="Зачекайте, йде завантаження...", 
                                font=self.font, text_color="gray")
        info_label.pack(pady=5)
        
        return loading_popup

    def update_from_google_sheets(self, popup_window=None):
        """Оновлення даних з Google Sheets"""
        url = self.sheets_url_entry.get().strip()
        if not url:
            messagebox.showerror("Помилка", "Введіть посилання на Google Таблицю")
            return
        
        try:
            # Витягуємо ID таблиці з різних форматів URL
            sheet_id = None
            patterns = [
                r'/spreadsheets/d/([a-zA-Z0-9-_]+)',  # Стандартний формат
                r'id=([a-zA-Z0-9-_]+)',               # Формат з параметрами
                r'^([a-zA-Z0-9-_]{44})$'              # Прямий ID
            ]
            
            for pattern in patterns:
                match = re.search(pattern, url)
                if match:
                    sheet_id = match.group(1)
                    break
            
            if not sheet_id:
                messagebox.showerror("Помилка", "Не вдається знайти ID таблиці в посиланні")
                return
            
            # Зберігаємо ID та URL якщо галочка активна
            self.saved_sheet_id = sheet_id
            if hasattr(self, 'save_url_var') and self.save_url_var.get():
                self.saved_sheets_url = url
                self.save_url_enabled = True
                self.save_sheets_config()  # Автоматично зберігаємо
            
            # Показуємо повідомлення про завантаження
            loading_popup = self.show_loading_popup()
            
            # Формуємо URL для завантаження CSV
            csv_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid=0"
            
            # Закриваємо основний діалог
            if popup_window:
                popup_window.destroy()
            
            # Завантажуємо файл напряму
            self.after(100, lambda: self.download_csv_directly(csv_url, loading_popup))
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка завантаження: {str(e)}")
    
    def download_csv_directly(self, csv_url, loading_popup=None):
        """Завантажує CSV файл напряму з Google Sheets"""
        try:
            import requests
            
            # Завантажуємо CSV дані
            print(f"🌐 Завантаження CSV з: {csv_url}")
            response = requests.get(csv_url, timeout=30)
            response.raise_for_status()
            
            # Отримуємо CSV контент і спробуємо виправити кодування
            csv_content = response.text
            print(f"✅ Завантажено {len(csv_content)} символів")
            
            # Спробуємо виправити кодування якщо є кракозябри
            if 'Ð' in csv_content or 'Ñ' in csv_content:
                print("🔄 Виявлено проблеми з кодуванням, спроба виправлення...")
                try:
                    # Спроба перекодувати з latin1 в UTF-8
                    csv_bytes = response.content
                    for encoding in ['utf-8', 'cp1251', 'windows-1251']:
                        try:
                            fixed_content = csv_bytes.decode(encoding)
                            if 'Название' in fixed_content or 'Name' in fixed_content:
                                csv_content = fixed_content
                                print(f"✅ Виправлено кодування з {encoding}")
                                break
                        except:
                            continue
                except Exception as e:
                    print(f"⚠️  Не вдалося виправити кодування: {e}")
                    print("💡 Продовжуємо з оригінальним контентом")
            
            # Перевіряємо, чи це дійсно CSV дані
            if not csv_content or len(csv_content) < 10:
                raise ValueError("Отримано порожній або некоректний CSV файл")
            
            # Діагностика структури CSV
            lines = csv_content.strip().split('\n')
            if lines:
                headers = lines[0].split(',') if ',' in lines[0] else lines[0].split(';')
                print(f"📋 Знайдено колонки: {headers}")
                print(f"📊 Рядків даних: {len(lines)-1}")
                
                # Перевіряємо наявність колонки "Название" для профілів
                if "Название" in headers:
                    print("✅ Знайдено колонку 'Название' для профілів")
                else:
                    print("⚠️  Колонка 'Название' не знайдена. Доступні колонки:")
                    for i, header in enumerate(headers):
                        print(f"   {i+1}. '{header.strip()}'")
                    print("💡 Для роботи з профілями потрібна колонка 'Название'")
            
            # Відображаємо в текстовому полі
            self.sheets_textarea.delete("0.0", "end")
            self.sheets_textarea.insert("0.0", csv_content)
            
            # Зберігаємо в локальний файл
            with open(self.LOCAL_PROFILES_CSV, 'w', encoding='utf-8', newline='') as f:
                f.write(csv_content)
            print(f"💾 Дані збережені в локальний файл: {self.LOCAL_PROFILES_CSV}")
            
            # Кешуємо CSV дані
            self.cached_csv_data = csv_content
            self.save_sheets_config()
            
            # Оновлюємо список профілів
            self.update_profile_list()
            
            # Закриваємо loading popup
            if loading_popup:
                loading_popup.destroy()
            
            messagebox.showinfo("Успіх", "Дані успішно завантажені з Google Sheets!")
            print("🎉 CSV дані успішно завантажені та оброблені")
            
        except requests.RequestException as e:
            print(f"❌ Помилка запиту: {str(e)}")
            if loading_popup:
                loading_popup.destroy()
            messagebox.showerror("Помилка мережі", f"Не вдалося завантажити дані:\n{str(e)}\n\nПеревірте з'єднання та права доступу до таблиці.")
        except Exception as e:
            print(f"❌ Помилка завантаження: {str(e)}")
            if loading_popup:
                loading_popup.destroy()
            messagebox.showerror("Помилка", f"Помилка обробки CSV: {str(e)}")

    def auto_load_downloaded_csv(self):
        """Автоматично завантажує останній CSV файл з папки Downloads"""
        try:
            # Знаходимо папку Downloads
            downloads_folder = os.path.join(os.path.expanduser("~"), "Downloads")
            if not os.path.exists(downloads_folder):
                messagebox.showerror("Помилка", "Папка Downloads не знайдена")
                return
            
            # Шукаємо CSV файли
            csv_files = []
            for file in os.listdir(downloads_folder):
                if file.lower().endswith('.csv'):
                    file_path = os.path.join(downloads_folder, file)
                    csv_files.append((file_path, os.path.getctime(file_path)))
            
            if not csv_files:
                messagebox.showwarning("Попередження", "CSV файли не знайдені в папці Downloads")
                return
            
            # Сортуємо за часом створення (найновіший перший)
            csv_files.sort(key=lambda x: x[1], reverse=True)
            latest_csv = csv_files[0][0]
            
            # Читаємо CSV файл з правильним кодуванням
            try:
                with open(latest_csv, 'r', encoding='utf-8-sig') as file:
                    csv_content = file.read()
            except UnicodeDecodeError:
                # Fallback до cp1251 якщо utf-8-sig не працює
                with open(latest_csv, 'r', encoding='cp1251') as file:
                    csv_content = file.read()
            
            # Відображаємо в текстовому полі
            self.sheets_textarea.delete("0.0", "end")
            self.sheets_textarea.insert("0.0", csv_content)
            
            # Зберігаємо в локальний файл
            with open(self.LOCAL_PROFILES_CSV, 'w', encoding='utf-8', newline='') as f:
                f.write(csv_content)
            print(f"💾 Дані збережені в локальний файл: {self.LOCAL_PROFILES_CSV}")
            
            # Кешуємо CSV дані
            self.cached_csv_data = csv_content
            self.save_sheets_config()
            
            # Оновлюємо список профілів в dropdown
            self.update_profile_list()
            
            # messagebox.showinfo("Успіх", f"Дані успішно завантажені з файлу:\n{os.path.basename(latest_csv)}")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка автозавантаження CSV: {str(e)}")
    
    def load_sheets_config(self):
        """Завантажує збережені налаштування Google Sheets"""
        try:
            if os.path.exists(self.sheets_config_file):
                with open(self.sheets_config_file, "r", encoding="utf-8") as f:
                    config = json.load(f)
                    self.save_url_enabled = config.get("save_url_enabled", False)
                    self.saved_sheets_url = config.get("saved_sheets_url", "")
                    self.saved_sheet_id = config.get("saved_sheet_id", "")
                    self.cached_csv_data = config.get("cached_csv_data", "")
                    
                    # Якщо є кешовані дані, відображаємо їх
                    if self.cached_csv_data:
                        self.sheets_textarea.delete("0.0", "end")
                        self.sheets_textarea.insert("0.0", self.cached_csv_data)
                        # Оновлюємо список профілів
                        self.update_profile_list()
        except Exception as e:
            print(f"Помилка завантаження конфігурації: {str(e)}")
    
    def save_sheets_config(self):
        """Зберігає налаштування Google Sheets"""
        try:
            config = {
                "save_url_enabled": self.save_url_enabled,
                "saved_sheets_url": self.saved_sheets_url,
                "saved_sheet_id": self.saved_sheet_id,
                "cached_csv_data": self.cached_csv_data
            }
            with open(self.sheets_config_file, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Помилка збереження конфігурації: {str(e)}")
    
    def load_conversion_memory(self):
        """Завантажує внутрішню пам'ять конверсій"""
        try:
            memory_file = get_config_path("conversion_memory.json")
            if os.path.exists(memory_file):
                with open(memory_file, "r", encoding="utf-8") as f:
                    self.conversion_memory = json.load(f)
            else:
                self.conversion_memory = {}
            print(f"Завантажено {len(self.conversion_memory)} конверсій з пам'яті")
        except Exception as e:
            print(f"Помилка завантаження пам'яті конверсій: {str(e)}")
            self.conversion_memory = {}
    
    def save_conversion_memory(self):
        """Зберігає внутрішню пам'ять конверсій"""
        try:
            memory_file = get_config_path("conversion_memory.json")
            with open(memory_file, "w", encoding="utf-8") as f:
                json.dump(self.conversion_memory, f, ensure_ascii=False, indent=2)
            print(f"Збережено {len(self.conversion_memory)} конверсій в пам'ять")
        except Exception as e:
            print(f"Помилка збереження пам'яті конверсій: {str(e)}")
    
    def add_conversion_to_memory(self, profile_name, conversion_id):
        """Додає конверсію в пам'ять"""
        if profile_name and conversion_id:
            self.conversion_memory[profile_name] = conversion_id
            self.save_conversion_memory()
            print(f"Додано в пам'ять: {profile_name} -> {conversion_id}")
    
    def get_conversion_from_memory(self, profile_name):
        """Отримує конверсію з пам'яті"""
        return self.conversion_memory.get(profile_name, "")
    
    # ================== CSV File Management Methods ==================
    # ===== CSV File Management методи видалені - використовуйте Generators вкладку =====
    

    def parse_csv_data(self):
        """Парсить CSV дані на основі введеного Profile Name"""
        try:
            profile_name = self.profile_search_entry.get().strip()
            if not profile_name:
                messagebox.showwarning("Попередження", "Введіть назву профілю для пошуку")
                return
            
            # Отримуємо CSV дані з текстового поля
            csv_content = self.sheets_textarea.get("0.0", "end-1c").strip()
            if not csv_content:
                messagebox.showwarning("Попередження", "Спочатку завантажте CSV дані з Google Sheets")
                return
            
            # Парсимо CSV з автоматичним визначенням роздільника
            import csv
            import io
            
            # Автоматичне визначення роздільника (як в update_profile_list)
            delimiter = ';'  # За замовчуванням
            first_line = csv_content.split('\n')[0] if '\n' in csv_content else csv_content
            
            if first_line.count(',') > first_line.count(';'):
                delimiter = ','
                print("📋 parse_csv_data: Використовуємо кому як роздільник")
            else:
                print("📋 parse_csv_data: Використовуємо крапку з комою як роздільник")
            
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            fieldnames = csv_reader.fieldnames
            print(f"📊 parse_csv_data: Знайдені колонки: {fieldnames}")
            
            # Визначаємо колонку для профілів (як в update_profile_list)
            profile_column = None
            possible_profile_columns = ["Название", "Name", "Profile", "Account", "Username", "Email", "Login", "User"]
            
            # Додаємо варіанти з проблемами кодування
            encoded_variants = ["ÐÐ°Ð·Ð²Ð°Ð½Ð¸Ðµ", "Ð\\x9dÐ°Ð·Ð²Ð°Ð½Ð¸Ðµ", "ÐÐ°Ð²ÐµÐ´ÐµÐ½Ñ"]
            possible_profile_columns.extend(encoded_variants)
            
            for col in possible_profile_columns:
                if col in fieldnames:
                    profile_column = col
                    print(f"✅ parse_csv_data: Знайдено колонку профілів: '{col}'")
                    break
            
            if not profile_column and fieldnames:
                profile_column = fieldnames[0]
                print(f"💡 parse_csv_data: Використовуємо першу колонку: '{profile_column}'")
            
            # Завантажуємо користувацьку конфігурацію парсингу
            app = self.winfo_toplevel()
            parsing_config = app.load_sheets_parsing_config()
            
            # Створюємо динамічний мапінг на основі користувацької конфігурації
            column_mapping = {}
            field_mapping = {
                "profile_name": "profile_name",
                "email": "Email/Login",
                "password": "Password",
                "2fa": "2FA",
                "backup_codes": "Backup Codes",
                "reserve_mail": "Reserve Mail",
                "credit_card": "Credit Card",
                "conversion": "Конверсія",
                "api_cf": "API Cloudflare",
                "cf_id": "Cloudflare ID",
                "cf_password": "Cloudflare Password"
            }
            
            # Будуємо мапінг: назва_колонки_з_таблиці -> назва_поля_у_формі
            for config_key, form_field in field_mapping.items():
                column_name = parsing_config.get(config_key)
                if column_name:
                    column_mapping[column_name] = form_field
            
            print(f"📊 parse_csv_data: Використовуємо динамічний мапінг: {column_mapping}")
            
            # Перевіряємо, чи знайшли колонку профілів
            if not profile_column:
                messagebox.showerror("Помилка", "Не вдалося визначити колонку з профілями в CSV")
                return
            
            # Шукаємо рядок з відповідним Profile Name (покращена логіка пошуку)
            found_row = None
            profile_search_name = profile_name.lower().strip()
            
            # Прибираємо дужки і їх вміст для пошуку (наприклад "Alex17(E5)" -> "Alex17")
            import re
            clean_profile_name = re.sub(r'\([^)]*\)', '', profile_search_name).strip()
            
            print(f"🔍 Шукаємо профіль '{profile_name}' у колонці '{profile_column}'")
            
            # Спочатку шукаємо точний збіг
            for row in csv_reader:
                csv_profile_name = row.get(profile_column, "").strip().lower()
                
                # Точний збіг - найвищий пріоритет
                if csv_profile_name == profile_search_name:
                    found_row = row
                    print(f"✅ Знайдено точний збіг: '{csv_profile_name}'")
                    break
            
            # Якщо точного збігу немає, шукаємо збіг без дужок
            if not found_row:
                csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)  # Перезавантажуємо reader з правильним роздільником
                for row in csv_reader:
                    csv_profile_name = row.get(profile_column, "").strip().lower()
                    if csv_profile_name == clean_profile_name:
                        found_row = row
                        print(f"✅ Знайдено збіг без дужок: '{csv_profile_name}'")
                        break
            
            # Тільки в крайньому випадку шукаємо часткове входження (але більш обережно)
            if not found_row:
                csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)  # Перезавантажуємо reader з правильним роздільником
                for row in csv_reader:
                    csv_profile_name = row.get(profile_column, "").strip().lower()
                    # Часткове входження тільки якщо довжина збігу > 50% від найкоротшого рядка
                    min_length = min(len(clean_profile_name), len(csv_profile_name))
                    if min_length > 5:  # Тільки для достатньо довгих назв
                        if (clean_profile_name in csv_profile_name or csv_profile_name in clean_profile_name):
                            # Додаткова перевірка - чи не закінчується один рядок цифрою, а інший - іншою
                            if not (clean_profile_name[-1].isdigit() and csv_profile_name[-1].isdigit() and 
                                   clean_profile_name[-1] != csv_profile_name[-1]):
                                found_row = row
                                break
            
            if not found_row:
                messagebox.showwarning("Не знайдено", f"Профіль з назвою '{profile_name}' не знайдено в CSV")
                return
            
            # Заповнюємо поля даними
            for csv_column, form_field in column_mapping.items():
                if form_field != "profile_name" and csv_column in found_row:
                    value = found_row[csv_column].strip() if found_row[csv_column] else ""
                    if form_field in self.entries_dict:
                        entry = self.entries_dict[form_field]
                        entry.delete(0, tk.END)
                        entry.insert(0, value)
            
            # Додаткова обробка специфічних полів з нового CSV
            additional_info = []
            if found_row.get("ID"):
                additional_info.append(f"ID: {found_row['ID']}")
            if found_row.get("ФИО"):
                additional_info.append(f"ФИО: {found_row['ФИО']}")
            if found_row.get("Прокси"):
                additional_info.append(f"Прокси: {found_row['Прокси']}")
            if found_row.get("Организация"):
                additional_info.append(f"Орг: {found_row['Организация']}")
            if found_row.get("Статус"):
                additional_info.append(f"Статус: {found_row['Статус']}")
            if found_row.get("День фарма"):
                additional_info.append(f"День: {found_row['День фарма']}")
            
            # Логуємо додаткову інформацію
            if additional_info:
                if hasattr(self, 'log_display'):
                    self.log_display.delete("0.0", "end")
                    self.log_display.insert("0.0", "\n".join(additional_info))
            
            # Зберігаємо знайдену назву профілю для форматованого копіювання (використовуємо реальну назву з CSV)
            self.current_profile_name = found_row.get(profile_column, profile_name)
            
            # НОВЕ: Auto-select в dropdown якщо профіль знайдено (гнучкий пошук)
            if hasattr(self, 'profile_names') and self.profile_names:
                # Шукаємо відповідний профіль в dropdown
                matching_dropdown_profile = None
                search_name_lower = profile_name.lower()
                
                for dropdown_profile in self.profile_names:
                    dropdown_lower = dropdown_profile.lower()
                    # Точний збіг
                    if dropdown_lower == search_name_lower:
                        matching_dropdown_profile = dropdown_profile
                        break
                    # Часткове входження (Alex16 знайде Alex16(ES))
                    elif search_name_lower in dropdown_lower or dropdown_lower in search_name_lower:
                        matching_dropdown_profile = dropdown_profile
                        break
                
                if matching_dropdown_profile:
                    print(f"🎯 Auto-select: '{profile_name}' -> '{matching_dropdown_profile}'")
                    self.profile_dropdown.set(matching_dropdown_profile)
                    # Залишаємо поле пошуку видимим для пошуку наступних профілів
                    # self.profile_search_entry.pack_forget()
            
            # messagebox.showinfo("Успіх", f"Дані для профілю '{profile_name}' успішно завантажені!")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка парсингу CSV: {str(e)}")

    def copy_formatted_data(self):
        """Копіює дані у красивому форматі з урахуванням нового CSV"""
        try:
            # Отримуємо значення з полів
            profile_name = getattr(self, 'current_profile_name', 'Unknown Profile')
            email = self.entries_dict["Email/Login"].get().strip()
            password = self.entries_dict["Password"].get().strip()
            fa_code = self.entries_dict["2FA"].get().strip()
            backup_codes = self.entries_dict["Backup Codes"].get().strip()
            reserve_mail = self.entries_dict["Reserve Mail"].get().strip()
            credit_card = self.entries_dict["Credit Card"].get().strip()
            conversion = self.entries_dict["Конверсія"].get().strip()
            
            # Отримуємо додаткову інформацію з логу
            log_content = ""
            if hasattr(self, 'log_display'):
                log_content = self.log_display.get("0.0", "end-1c").strip()
            
            # Форматуємо дані у стовпчик (розширений формат)
            formatted_data = f"""📧 === {profile_name} ===
Почта: {email}
Пароль: {password}
2фа: {fa_code}
Бэкап: {backup_codes}
Резервка: {reserve_mail}"""
            
            # Додаємо картку і конверсію якщо є
            if credit_card:
                formatted_data += f"\nКарта: {credit_card}"
            if conversion:
                formatted_data += f"\nКонверсія: {conversion}"
            
            # Додаємо додаткову інформацію з логу
            if log_content:
                formatted_data += f"\n\n📋 Додатково:\n{log_content}"
            
            # Копіюємо в буфер
            safe_clipboard_operation("set", formatted_data)
            
            # Короткий фідбек в логах
            timestamp = datetime.datetime.now().strftime('%H:%M:%S')
            print(f"[{timestamp}] 📋 Скопійовано дані профілю: {profile_name}")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка копіювання: {str(e)}")
    
    def update_profile_list(self):
        """Оновлює список профілів з CSV даних (покращена версія)"""
        try:
            # Спочатку пробуємо завантажити з локального файлу
            csv_content = ""
            if os.path.exists(self.LOCAL_PROFILES_CSV):
                with open(self.LOCAL_PROFILES_CSV, 'r', encoding='utf-8') as f:
                    csv_content = f.read().strip()
                    # Оновлюємо textarea з локального файлу
                    self.sheets_textarea.delete("0.0", "end")
                    self.sheets_textarea.insert("0.0", csv_content)
                    print(f"📂 Завантажено з локального файлу: {self.LOCAL_PROFILES_CSV}")
            else:
                # Якщо файлу немає - беремо з textarea
                csv_content = self.sheets_textarea.get("0.0", "end-1c").strip()
            
            if not csv_content:
                return
            
            import csv
            import io
            
            # Автоматичне визначення роздільника
            delimiter = ';'  # За замовчуванням
            first_line = csv_content.split('\n')[0] if '\n' in csv_content else csv_content
            
            # Перевіряємо, який роздільник використовується
            if first_line.count(',') > first_line.count(';'):
                delimiter = ','
                print("📋 Використовуємо кому як роздільник")
            else:
                print("📋 Використовуємо крапку з комою як роздільник")
            
            print(f"🔍 Перша лінія CSV: {first_line[:100]}...")
            
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            profile_names = ["Nothing"]  # Початкове значення
            
            # Логуємо доступні колонки
            fieldnames = csv_reader.fieldnames
            print(f"📊 Знайдені колонки: {fieldnames}")
            
            # Визначаємо колонку для профілів
            profile_column = None
            possible_profile_columns = ["Название", "Name", "Profile", "Account", "Username", "Email", "Login", "User"]
            
            # Додаємо варіанти з проблемами кодування
            encoded_variants = [
                "ÐÐ°Ð·Ð²Ð°Ð½Ð¸Ðµ",  # Одна з можливих варіацій кодування
                "Ð\x9dÐ°Ð·Ð²Ð°Ð½Ð¸Ðµ",  # Інша варіація
                "ÐÐ°Ð²ÐµÐ´ÐµÐ½Ñ",  # З тестового файлу
            ]
            possible_profile_columns.extend(encoded_variants)
            
            for col in possible_profile_columns:
                if col in fieldnames:
                    profile_column = col
                    print(f"✅ Знайдено колонку профілів: '{col}'")
                    break
            
            if not profile_column:
                print("❌ Колонка профілів не знайдена!")
                print(f"💡 Доступні колонки: {fieldnames}")
                print("💡 Створюємо тестові профілі з першої колонки...")
                profile_column = fieldnames[0] if fieldnames else None
            
            if not profile_column:
                print("❌ Немає жодної колонки для обробки")
                return
            
            profile_count = 0
            for row in csv_reader:
                name = row.get(profile_column, "").strip()
                if name and name not in profile_names:
                    profile_names.append(name)
                    profile_count += 1
                    # Детальне логування для діагностики
                    if profile_count <= 5:  # Логуємо перші 5 профілів
                        print(f"✅ Додано профіль {profile_count}: {name}")
                    elif profile_count == 6:
                        print("... (інші профілі додано без логування)")
            
            # Показуємо загальну кількість та останні профілі для перевірки
            print(f"Загальна кількість профілів: {len(profile_names)-1} (без 'Nothing')")
            if len(profile_names) > 10:
                print("Останні 5 профілів:")
                for i, name in enumerate(profile_names[-5:]):
                    if name != "Nothing":
                        print(f"  {len(profile_names)-5+i}: {name}")
            
            # Перевіряємо, чи профіль "Alex_FarmPunch_USA_1_081025_11" є в списку
            target_profile = "Alex_FarmPunch_USA_1_081025_11"
            if target_profile in profile_names:
                index = profile_names.index(target_profile)
                print(f"🎯 Профіль '{target_profile}' знайдено на позиції {index}")
            else:
                print(f"❌ Профіль '{target_profile}' НЕ знайдено в списку!")
            
            self.profile_names = profile_names
            self.profile_dropdown.configure(values=profile_names)
            
            # Логуємо результат
            if profile_count > 0:
                print(f"✅ Завантажено {profile_count} профілів з CSV")
                # Оновлюємо статус в лозі
                if hasattr(self, 'log_display'):
                    self.log_display.delete("0.0", "end")
                    self.log_display.insert("0.0", f"📊 Завантажено {profile_count} профілів\n🔄 Оберіть профіль зі списку або введіть вручну")
            
        except Exception as e:
            print(f"Помилка оновлення списку профілів: {str(e)}")
    
    def refresh_profile_dropdown(self):
        """Оновлює список профілів у dropdown"""
        try:
            # Викликаємо метод оновлення списку профілів
            self.update_profile_list()
            messagebox.showinfo("Успіх", "Список профілів оновлено!")
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка оновлення: {str(e)}")
    
    def show_all_profiles(self):
        """Показує всі профілі у окремому вікні зі скролом"""
        try:
            if not hasattr(self, 'profile_names') or len(self.profile_names) <= 1:
                messagebox.showinfo("Інформація", "Спочатку завантажте CSV дані")
                return
            
            # Створюємо нове вікно
            profile_window = ctk.CTkToplevel(self)
            profile_window.title("Всі профілі")
            profile_window.geometry("400x500")
            profile_window.transient(self)
            profile_window.grab_set()
            
            # Заголовок
            ctk.CTkLabel(profile_window, text=f"Всього профілів: {len(self.profile_names)-1}", 
                        font=ctk.CTkFont(size=14, weight="bold")).pack(pady=10)
            
            # Поле пошуку
            search_frame = ctk.CTkFrame(profile_window)
            search_frame.pack(fill="x", padx=10, pady=5)
            
            ctk.CTkLabel(search_frame, text="Пошук:").pack(side="left", padx=5)
            search_entry = ctk.CTkEntry(search_frame, placeholder_text="Введіть назву профілю...")
            search_entry.pack(side="left", fill="x", expand=True, padx=5)
            
            # Scrollable frame для списку профілів
            scrollable_frame = ctk.CTkScrollableFrame(profile_window)
            scrollable_frame.pack(fill="both", expand=True, padx=10, pady=5)
            
            # Змінна для кнопок
            profile_buttons = []
            
            def update_profile_list(filter_text=""):
                # Очищуємо попередні кнопки
                for btn in profile_buttons:
                    btn.destroy()
                profile_buttons.clear()
                
                # Додаємо кнопки для профілів (крім "Nothing")
                filtered_profiles = [p for p in self.profile_names[1:] 
                                   if filter_text.lower() in p.lower()]
                
                for profile in filtered_profiles:
                    btn = ctk.CTkButton(scrollable_frame, text=profile, 
                                      command=lambda p=profile: self.select_profile_from_list(p, profile_window),
                                      anchor="w", height=30)
                    btn.pack(fill="x", pady=1, padx=5)
                    profile_buttons.append(btn)
            
            # Bind пошуку
            def on_search_change(event=None):
                update_profile_list(search_entry.get())
            
            search_entry.bind('<KeyRelease>', on_search_change)
            
            # Початкове заповнення
            update_profile_list()
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка показу профілів: {str(e)}")
    
    def select_profile_from_list(self, profile_name, window):
        """Вибирає профіль зі списку та закриває вікно"""
        try:
            # Встановлюємо профіль у dropdown
            self.profile_dropdown.set(profile_name)
            
            # Викликаємо обробку вибору профілю
            self.on_profile_selected(profile_name)
            
            # Закриваємо вікно
            window.destroy()
            
            print(f"✅ Обрано профіль: {profile_name}")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка вибору профілю: {str(e)}")
    
    def on_profile_selected(self, selected_profile):
        """Обробляє вибір профілю з dropdown (покращена версія)"""
        # 💾 Зберігаємо поточні дані перед переключенням на інший профіль
        try:
            email = self.entries_dict["Email/Login"].get().strip()
            if email:  # Зберігаємо тільки якщо є email
                self.save_account_to_csv()
                print(f"💾 Поточний акаунт автоматично збережено перед переключенням")
        except Exception as save_error:
            print(f"⚠️ Помилка автозбереження перед переключенням: {save_error}")
        
        if selected_profile == "Nothing":
            # Показуємо поле для ручного вводу
            self.profile_search_entry.pack(padx=10, pady=5, fill="x")
            self.profile_search_entry.focus()
            # Очищуємо поля
            for entry in self.entries_dict.values():
                entry.delete(0, tk.END)
            if hasattr(self, 'log_display'):
                self.log_display.delete("0.0", "end")
                self.log_display.insert("0.0", "💡 Введіть назву профілю для пошуку")
        else:
            # Ховаємо поле для ручного вводу та автоматично парсимо
            self.profile_search_entry.pack_forget()
            # Автоматично заповнюємо поле пошуку вибраним профілем та парсимо
            self.profile_search_entry.delete(0, tk.END)
            self.profile_search_entry.insert(0, selected_profile)
            # Зберігаємо як новий базовий аккаунт тільки при ручному виборі
            if not getattr(self, '_programmatic_selection', False):
                self.base_selected_account = selected_profile
            # Автоматично парсимо дані для обраного профілю
            self.auto_parse_selected_profile(selected_profile)
        
        # Оновлюємо кнопки швидкого вибору відносно нового вибраного аккаунта
        self.update_quick_select_buttons(selected_profile)
    
    def auto_parse_selected_profile(self, profile_name):
        """Автоматично парсить дані для обраного профілю"""
        try:
            # print(f"🔍 Початок парсингу профілю: '{profile_name}'")
            
            # Отримуємо CSV дані
            csv_content = self.sheets_textarea.get("0.0", "end-1c").strip()
            if not csv_content:
                messagebox.showwarning("Попередження", "Спочатку завантажте CSV дані з Google Sheets")
                return
            
            # Парсимо CSV (використовуємо той же код що в parse_csv_data)
            import csv
            import io
            
            # Шукаємо рядок з відповідним Profile Name (покращена логіка пошуку)
            found_row = None
            profile_search_name = profile_name.lower().strip()
            
            # Прибираємо дужки і їх вміст для пошуку (наприклад "Alex17(E5)" -> "Alex17")
            import re
            clean_profile_name = re.sub(r'\([^)]*\)', '', profile_search_name).strip()
            
            row_count = 0
            # Автоматичне визначення роздільника (як в update_profile_list)
            lines = csv_content.split('\n')
            delimiter = ','
            if lines:
                first_line = lines[0]
                if first_line.count(';') > first_line.count(','):
                    delimiter = ';'
            
            # Визначаємо колонку профілю (як в update_profile_list)
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            headers = csv_reader.fieldnames or []
            profile_column = None
            
            for header in headers:
                if any(name.lower() in header.lower() for name in ["название", "назва", "profile", "профіль", "name", "ім'я"]):
                    profile_column = header
                    break
                if any(encoded.lower() in header.lower() for encoded in ["РќР°Р·РІР°РЅРёРµ", "РќР°Р·РІР°"]):
                    profile_column = header
                    break
            
            if not profile_column and headers:
                profile_column = headers[0]
            
            print(f"🔍 Використовуємо колонку профілю: '{profile_column}', роздільник: '{delimiter}'")
            
            # Завантажуємо користувацьку конфігурацію парсингу
            app = self.winfo_toplevel()
            parsing_config = app.load_sheets_parsing_config()
            
            # Створюємо динамічний мапінг на основі користувацької конфігурації
            column_mapping = {profile_column: "profile_name"}  # Колонка профілю завжди перша
            
            field_mapping = {
                "email": "Email/Login",
                "password": "Password",
                "2fa": "2FA",
                "backup_codes": "Backup Codes",
                "reserve_mail": "Reserve Mail",
                "credit_card": "Credit Card",
                "conversion": "Конверсія",
                "api_cf": "API Cloudflare",
                "cf_id": "Cloudflare ID",
                "cf_password": "Cloudflare Password"
            }
            
            # Будуємо мапінг: назва_колонки_з_таблиці -> назва_поля_у_формі
            for config_key, form_field in field_mapping.items():
                column_name = parsing_config.get(config_key)
                if column_name:
                    column_mapping[column_name] = form_field
            
            print(f"📊 load_profile_from_csv: Використовуємо динамічний мапінг: {column_mapping}")
            
            # Спочатку шукаємо точний збіг
            for row in csv_reader:
                row_count += 1
                csv_profile_name = row.get(profile_column, "").strip().lower()
                
                # Логування для діагностики (вимкнуто для релізу)
                # if row_count <= 15:
                #     print(f"Рядок {row_count}: '{csv_profile_name}' vs '{profile_search_name}'")
                
                # Точний збіг - найвищий пріоритет
                if csv_profile_name == profile_search_name:
                    found_row = row
                    print(f"✅ Знайдено точний збіг: '{profile_name}'")
                    break
            
            # Якщо точного збігу немає, шукаємо збіг без дужок
            if not found_row:
                csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)  # Перезавантажуємо reader з правильним роздільником
                for row in csv_reader:
                    csv_profile_name = row.get(profile_column, "").strip().lower()
                    if csv_profile_name == clean_profile_name:
                        found_row = row
                        print(f"✅ Знайдено збіг без дужок: '{csv_profile_name}'")
                        break
            
            # Тільки в крайньому випадку шукаємо часткове входження (але більш обережно)
            if not found_row:
                csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)  # Перезавантажуємо reader з правильним роздільником
                for row in csv_reader:
                    csv_profile_name = row.get(profile_column, "").strip().lower()
                    # Часткове входження тільки якщо довжина збігу > 50% від найкоротшого рядка
                    min_length = min(len(clean_profile_name), len(csv_profile_name))
                    if min_length > 5:  # Тільки для достатньо довгих назв
                        if (clean_profile_name in csv_profile_name or csv_profile_name in clean_profile_name):
                            # Додаткова перевірка - чи не закінчується один рядок цифрою, а інший - іншою
                            if not (clean_profile_name[-1].isdigit() and csv_profile_name[-1].isdigit() and 
                                   clean_profile_name[-1] != csv_profile_name[-1]):
                                found_row = row
                                print(f"✅ Знайдено безпечне часткове входження: '{csv_profile_name}'")
                                break
            
            # print(f"Загальна кількість рядків у CSV: {row_count}")
            
            if not found_row:
                print(f"❌ Профіль '{profile_name}' не знайдено в CSV")
                messagebox.showwarning("Не знайдено", f"Профіль з назвою '{profile_name}' не знайдено в CSV")
                return
            
            # print(f"📋 Заповнюємо поля для профілю '{profile_name}'")
            
            # Спочатку очищаємо ВСІ поля перед завантаженням нових даних
            for entry in self.entries_dict.values():
                entry.delete(0, tk.END)
            
            # Заповнюємо поля даними
            for csv_column, form_field in column_mapping.items():
                if form_field != "profile_name" and csv_column in found_row:
                    value = found_row[csv_column].strip() if found_row[csv_column] else ""
                    if form_field in self.entries_dict:
                        entry = self.entries_dict[form_field]
                        entry.delete(0, tk.END)
                        entry.insert(0, value)
                        # if value:
                        #     print(f"  {form_field}: {value}")
            
            # Автоматична генерація Cloudflare Password (password@)
            try:
                if hasattr(self, '_password_entry_ref') and hasattr(self, '_cf_pass_entry_ref'):
                    password_value = self._password_entry_ref.get().strip()
                    if password_value:
                        cf_password = f"{password_value}@"
                        self._cf_pass_entry_ref.delete(0, tk.END)
                        self._cf_pass_entry_ref.insert(0, cf_password)
                        print(f"✅ Cloudflare Password автоматично згенеровано: {cf_password}")
            except Exception as cf_error:
                print(f"⚠️ Помилка генерації Cloudflare Password: {cf_error}")
            
            # Додаткова обробка специфічних полів з нового CSV
            additional_info = []
            if found_row.get("ID"):
                additional_info.append(f"ID: {found_row['ID']}")
            if found_row.get("ФИО"):
                additional_info.append(f"ФИО: {found_row['ФИО']}")
            if found_row.get("Прокси"):
                additional_info.append(f"Прокси: {found_row['Прокси']}")
            if found_row.get("Организация"):
                additional_info.append(f"Орг: {found_row['Организация']}")
            if found_row.get("Статус"):
                additional_info.append(f"Статус: {found_row['Статус']}")
            if found_row.get("День фарма"):
                additional_info.append(f"День: {found_row['День фарма']}")
            
            # Логуємо додаткову інформацію
            if additional_info:
                if hasattr(self, 'log_display'):
                    self.log_display.delete("0.0", "end")
                    self.log_display.insert("0.0", "\n".join(additional_info))
            
            # Зберігаємо знайдену назву профілю для форматованого копіювання (використовуємо реальну назву з CSV)
            self.current_profile_name = found_row.get(profile_column, profile_name)
            
            print(f"✅ Профіль '{profile_name}' завантажено успішно")
            
        except Exception as e:
            print(f"❌ Помилка автоматичного парсингу: {str(e)}")
            messagebox.showerror("Помилка", f"Помилка автоматичного парсингу: {str(e)}")
    
    def manual_search_profile(self):
        """Ручний пошук профілю по назві"""
        try:
            profile_name = self.manual_search_entry.get().strip()
            if not profile_name:
                messagebox.showwarning("Попередження", "Введіть назву профілю для ручного пошуку")
                return
            
            # Отримуємо CSV дані
            csv_content = self.sheets_textarea.get("0.0", "end-1c").strip()
            if not csv_content:
                messagebox.showwarning("Попередження", "Спочатку завантажте CSV дані з Google Sheets")
                return
            
            # Парсимо CSV
            import csv
            import io
            
            csv_reader = csv.DictReader(io.StringIO(csv_content))
            
            # Завантажуємо користувацьку конфігурацію парсингу
            app = self.winfo_toplevel()
            parsing_config = app.load_sheets_parsing_config()
            
            # Створюємо динамічний мапінг
            column_mapping = {parsing_config.get("profile_name", "Название"): "profile_name"}
            
            field_mapping = {
                "email": "Email/Login",
                "password": "Password",
                "2fa": "2FA",
                "backup_codes": "Backup Codes",
                "reserve_mail": "Reserve Mail",
                "credit_card": "Credit Card",
                "conversion": "Конверсія",
                "api_cf": "API Cloudflare",
                "cf_id": "Cloudflare ID",
                "cf_password": "Cloudflare Password"
            }
            
            for config_key, form_field in field_mapping.items():
                column_name = parsing_config.get(config_key)
                if column_name:
                    column_mapping[column_name] = form_field
            
            profile_column = parsing_config.get("profile_name", "Название")
            
            # Шукаємо рядок з відповідним Profile Name (використовуємо часткове співпадіння)
            found_rows = []
            for row in csv_reader:
                name = row.get(profile_column, "").strip()
                if profile_name.lower() in name.lower():
                    found_rows.append((name, row))
            
            if not found_rows:
                messagebox.showwarning("Не знайдено", f"Профілі з назвою що містить '{profile_name}' не знайдено в CSV")
                return
            
            # Якщо знайдено кілька - показуємо список для вибору
            if len(found_rows) > 1:
                self.show_search_results(found_rows, column_mapping)
            else:
                # Якщо знайдено тільки один - одразу заповнюємо
                name, row = found_rows[0]
                self.fill_form_data(row, column_mapping, name)
                messagebox.showinfo("Успіх", f"Дані для профілю '{name}' завантажені!")
                
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка ручного пошуку: {str(e)}")
    
    def show_search_results(self, found_rows, column_mapping):
        """Показує діалог з результатами пошуку для вибору"""
        # Створюємо діалогове вікно
        dialog = ctk.CTkToplevel(self)
        dialog.title("Результати пошуку")
        dialog.geometry("400x300")
        dialog.grab_set()  # Модальне вікно
        
        ctk.CTkLabel(dialog, text="Знайдено кілька профілів:", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=10)
        
        # Список результатів
        results_frame = ctk.CTkScrollableFrame(dialog)
        results_frame.pack(fill="both", expand=True, padx=20, pady=10)
        
        for i, (name, row) in enumerate(found_rows):
            btn = ctk.CTkButton(results_frame, text=name, 
                               command=lambda r=row, n=name: self.select_from_results(r, column_mapping, n, dialog),
                               width=350, height=35, font=self.font)
            btn.pack(pady=2, fill="x")
        
        # Кнопка скасування
        ctk.CTkButton(dialog, text="Скасувати", 
                     command=dialog.destroy, width=100, height=30).pack(pady=10)
    
    def select_from_results(self, row, column_mapping, name, dialog):
        """Вибирає профіль з результатів пошуку"""
        self.fill_form_data(row, column_mapping, name)
        dialog.destroy()
        messagebox.showinfo("Успіх", f"Дані для профілю '{name}' завантажені!")
    
    def fill_form_data(self, row, column_mapping, profile_name):
        """Заповнює поля форми даними з CSV рядка"""
        try:
            for csv_column, form_field in column_mapping.items():
                if form_field != "profile_name" and csv_column in row:
                    value = row[csv_column].strip() if row[csv_column] else ""
                    if form_field in self.entries_dict:
                        self.entries_dict[form_field].delete(0, "end")
                        self.entries_dict[form_field].insert(0, value)
            
            self.current_profile_name = profile_name
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка заповнення форми: {str(e)}")
    
    def profile_up(self):
        """Переміщує вибір на одну позицію вверх"""
        try:
            current_value = self.profile_dropdown.get()
            current_index = self.profile_names.index(current_value)
            if current_index > 0:
                new_value = self.profile_names[current_index - 1]
                self.profile_dropdown.set(new_value)
                self.on_profile_selected(new_value)
        except (ValueError, IndexError):
            pass
    
    def profile_down(self):
        """Переміщує вибір на одну позицію вниз"""
        try:
            current_value = self.profile_dropdown.get()
            current_index = self.profile_names.index(current_value)
            if current_index < len(self.profile_names) - 1:
                new_value = self.profile_names[current_index + 1]
                self.profile_dropdown.set(new_value)
                self.on_profile_selected(new_value)
        except (ValueError, IndexError):
            pass
    
    # Стара функція видалена - використовуємо select_account_by_offset_from_current
    
    def update_quick_select_buttons(self, selected_profile):
        """Оновлює підписи кнопок швидкого вибору відносно вибраного аккаунта"""
        try:
            if not hasattr(self, 'quick_btn_1'):
                return  # Кнопки ще не створені
                
            current_index = self.profile_names.index(selected_profile)
            
            # Оновлюємо підписи кнопок
            for i, btn in enumerate([self.quick_btn_1, self.quick_btn_2, self.quick_btn_3, self.quick_btn_4]):
                target_index = current_index + i
                if target_index < len(self.profile_names):
                    # Кнопка активна - показуємо номер
                    btn.configure(text=str(i+1), state="normal")
                else:
                    # Кнопка неактивна - показуємо X
                    btn.configure(text="✗", state="disabled")
                    
        except (ValueError, IndexError):
            # Якщо профіль не знайдено, відключаємо всі кнопки
            if hasattr(self, 'quick_btn_1'):
                for btn in [self.quick_btn_1, self.quick_btn_2, self.quick_btn_3, self.quick_btn_4]:
                    btn.configure(text="✗", state="disabled")
    
    def select_account_by_offset_from_current(self, offset):
        """Вибирає аккаунт з певним зміщенням від БАЗОВОГО ВИБРАНОГО аккаунта
        offset 0 = базовий вибраний аккаунт (кнопка 1)
        offset 1 = +1 від базового (кнопка 2)  
        offset 2 = +2 від базового (кнопка 3)
        offset 3 = +3 від базового (кнопка 4)
        """
        try:
            # Використовуємо збережений базовий аккаунт, якщо існує
            if not hasattr(self, 'base_selected_account') or not self.base_selected_account:
                # Якщо базовий аккаунт не збережений, використовуємо поточний як базовий
                current_selected = self.profile_dropdown.get()
                if current_selected == "Nothing":
                    if hasattr(self, 'log_display'):
                        self.log_display.insert("end", "⚠️ Спочатку виберіть аккаунт зі списку\n")
                    return
                self.base_selected_account = current_selected
            
            # Використовуємо базовий аккаунт для розрахунку
            base_index = self.profile_names.index(self.base_selected_account)
            target_index = base_index + offset
            
            # Мінімальне логування для перевірки
            # debug_msg = f"🔍 DEBUG: Кнопка {offset+1}, базовий={self.base_selected_account} (індекс {base_index}), цільовий індекс={target_index}"
            
            if 0 <= target_index < len(self.profile_names):
                new_value = self.profile_names[target_index]
                
                # Мінімальне логування
                # debug_msg2 = f"🔍 DEBUG: Обираємо {new_value} на позиції {target_index}"
                
                # Встановлюємо флаг програмного вибору
                self._programmatic_selection = True
                
                self.profile_dropdown.set(new_value)
                self.on_profile_selected(new_value)
                
                # Скидаємо флаг програмного вибору
                self._programmatic_selection = False
                
                # Коротке повідомлення про швидкий вибір
                timestamp = datetime.datetime.now().strftime("%H:%M:%S")
                success_msg = f"⚡ [{timestamp}] Кнопка {offset+1} → {new_value}"
                if hasattr(self, 'log_display'):
                    self.log_display.insert("end", success_msg + "\n")
                print(success_msg)
                
            else:
                # Якщо індекс виходить за межі
                available_count = len(self.profile_names) - base_index
                error_msg = f"⚠️ Кнопка {offset+1}: недостатньо аккаунтів від {current_selected} (доступно {available_count})"
                if hasattr(self, 'log_display'):
                    self.log_display.insert("end", error_msg + "\n")
                print(error_msg)
                    
        except (ValueError, IndexError) as e:
            error_msg = f"❌ Помилка швидкого вибору кнопки {offset+1}: {str(e)}"
            if hasattr(self, 'log_display'):
                self.log_display.insert("end", error_msg + "\n")
            print(error_msg)
    
    def on_mouse_wheel(self, event):
        """Обробляє прокрутку колесиком миші"""
        if event.delta > 0:
            self.profile_up()
        else:
            self.profile_down()

    def generate_octo_profile(self):
        """Генерує файл OctoProfile в Excel форматі"""
        try:
            # Отримуємо дані з полів
            profile_template = self.octo_profile_entry.get().strip()
            proxy_template = self.octo_proxy_entry.get().strip()
            count_str = self.octo_count_entry.get().strip()
            
            if not profile_template or not proxy_template or not count_str:
                print("Заповніть всі поля для генерації OctoProfile")
                return
            
            try:
                count = int(count_str)
                if count <= 0:
                    count = 20
            except ValueError:
                count = 20
            
            # Створюємо директорію якщо не існує
            import os
            import glob
            octo_dir = get_app_stuff_path("Octo Import")
            
            # Видаляємо всі старі файли профілів з папки
            self.clean_old_profile_files(octo_dir)
            
            # Парсимо числа з шаблонів (підтримка шаблонів з числом або з {i})
            profile_base, profile_num = self.extract_number_from_template(profile_template)
            proxy_base, proxy_num = self.extract_proxy_number(proxy_template)

            # Визначаємо стартові індекси (якщо шаблони не містять числа, починаємо з 1)
            start_profile_index = profile_num if profile_num and profile_num > 0 else 1
            start_proxy_index = proxy_num if proxy_num and proxy_num > 0 else 1

            # Генеруємо пари Name / Proxy
            rows = []
            for i in range(count):
                seq = i + 1

                # Profile name: якщо витягнули базу та число, додаємо інкрементоване число
                if profile_num and profile_num > 0:
                    # profile_base вже містить підкреслення в кінці (наприклад: "Alex_FarmPunch_USA_5_131125_")
                    new_profile_name = f"{profile_base}{start_profile_index + i}"
                elif '{i}' in profile_template:
                    new_profile_name = profile_template.replace('{i}', str(start_profile_index + i))
                else:
                    new_profile_name = f"{profile_template}_{start_profile_index + i}"

                # Proxy: якщо є sessionid-номер, замінюємо його; якщо використовується {i}, замінюємо; інакше додаємо sessionid
                if proxy_num and proxy_num > 0:
                    new_proxy = proxy_base.replace(f"sessionid-{proxy_num}", f"sessionid-{start_proxy_index + i}")
                elif '{i}' in proxy_template:
                    new_proxy = proxy_template.replace('{i}', str(start_proxy_index + i))
                else:
                    # Якщо у проксі взагалі немає маркера, просто додаємо sessionid в кінець (без гарантій)
                    new_proxy = f"{proxy_template};sessionid-{start_proxy_index + i}"

                rows.append({
                    "Name": new_profile_name,
                    "Proxy": new_proxy
                })

            # Спробуємо створити Excel файл з двома колонками Name та Proxy
            try:
                import pandas as pd
                df = pd.DataFrame(rows)
                file_path = os.path.join(octo_dir, f"names_proxies_{start_profile_index}-{start_profile_index + count - 1}.xlsx")
                df.to_excel(file_path, index=False)
                msg = f"Файл з іменами та проксями створено: {file_path}"
                print(msg)
                if hasattr(self, 'log_display'):
                    self.log_display.insert('end', msg + "\n")
            except ImportError:
                # Якщо pandas відсутній, зберігаємо CSV і конвертуємо
                import csv
                csv_path = os.path.join(octo_dir, f"names_proxies_{start_profile_index}-{start_profile_index + count - 1}.csv")
                with open(csv_path, 'w', newline='', encoding='utf-8') as csvfile:
                    fieldnames = ["Name", "Proxy"]
                    writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
                    writer.writeheader()
                    writer.writerows(rows)

                # Спробуємо конвертувати CSV в Excel
                try:
                    self.convert_csv_to_excel(csv_path)
                except Exception:
                    print(f"CSV файл створено: {csv_path}. Встановіть pandas або openpyxl для створення Excel файлів")

            # Зберігаємо останні номера для функції Choose Last +1
            last_profile_index = start_profile_index + count - 1
            last_proxy_index = start_proxy_index + count - 1
            self.save_last_generated_profile(profile_template, proxy_template, last_profile_index, last_proxy_index)
                
        except Exception as e:
            print(f"Помилка генерації OctoProfile: {str(e)}")
    
    def extract_number_from_template(self, template):
        """Витягує останнє число після останнього підкреслення з шаблону Profile name
        Формат: Alex_FarmPunch_USA_5_131125_1 -> витягує '1'
        """
        import re
        # Шукаємо останнє число після останнього підкреслення
        match = re.search(r'_(\d+)$', template)
        if match:
            last_number = int(match.group(1))
            # Зберігаємо базу без останнього числа (все до останнього _число)
            base = template[:match.start()] + '_'
            return base, last_number
        return template, 0
    
    def extract_proxy_number(self, proxy_template):
        """Витягує номер sessionid з proxy"""
        import re
        # Шукаємо sessionid-номер в проксі
        match = re.search(r'sessionid-(\d+)', proxy_template)
        if match:
            sessionid_num = int(match.group(1))
            base_proxy = proxy_template
            return base_proxy, sessionid_num
        return proxy_template, 0
    
    def convert_csv_to_excel(self, csv_path):
        """Конвертує CSV в Excel якщо можливо"""
        try:
            # Спробуємо використати openpyxl для створення Excel
            import openpyxl
            from openpyxl import Workbook
            import csv
            
            wb = Workbook()
            ws = wb.active
            ws.title = "OctoProfiles"
            
            # Читаємо CSV та записуємо в Excel
            with open(csv_path, 'r', encoding='utf-8') as csvfile:
                csv_reader = csv.reader(csvfile)
                for row_index, row in enumerate(csv_reader, 1):
                    for col_index, value in enumerate(row, 1):
                        ws.cell(row=row_index, column=col_index, value=value)
            
            # Зберігаємо Excel файл
            excel_path = csv_path.replace('.csv', '.xlsx')
            wb.save(excel_path)
            
            # Видаляємо CSV файл
            os.remove(csv_path)
            print(f"OctoProfile файл створено: {excel_path}")
            
        except ImportError:
            print(f"CSV файл створено: {csv_path}")
            print("Встановіть pandas або openpyxl для створення Excel файлів")
    
    def open_octo_folder(self):
        """Відкриває папку Octo Import"""
        try:
            import os
            import subprocess
            import sys
            
            octo_dir = get_app_stuff_path("Octo Import")
            
            # Відкриваємо папку залежно від ОС
            if os.name == 'nt':  # Windows
                os.startfile(octo_dir)
            elif sys.platform == 'darwin':  # macOS
                subprocess.call(['open', octo_dir])
            else:  # Linux
                subprocess.call(['xdg-open', octo_dir])
                
        except Exception as e:
            print(f"Помилка відкриття папки: {str(e)}")

    def clean_old_profile_files(self, octo_dir):
        """Видаляє всі старі файли профілів з папки Octo Import"""
        try:
            import os
            import glob
            
            # Шаблони файлів для видалення
            patterns = [
                "octo_profiles_*.xlsx",
                "octo_profiles_*.csv",
                "octo_profiles*.xlsx", 
                "octo_profiles*.csv"
            ]
            
            deleted_files = []
            
            # Проходимо по всіх шаблонах та видаляємо відповідні файли
            for pattern in patterns:
                file_pattern = os.path.join(octo_dir, pattern)
                files_to_delete = glob.glob(file_pattern)
                
                for file_path in files_to_delete:
                    try:
                        os.remove(file_path)
                        deleted_files.append(os.path.basename(file_path))
                        print(f"Видалено старий файл: {os.path.basename(file_path)}")
                    except Exception as e:
                        print(f"Не вдалося видалити файл {os.path.basename(file_path)}: {str(e)}")
            
            if deleted_files:
                print(f"Очищено {len(deleted_files)} старих файлів профілів")
            else:
                print("Старих файлів профілів не знайдено")
                
        except Exception as e:
            print(f"Помилка при очищенні старих файлів: {str(e)}")

    # ================== Organisation Control Methods ==================
    def __init_org_control__(self):
        """Ініціалізуємо змінні для Organisation Control"""
        self.org_config_file = get_config_path("organisation_config.json")
        self.current_folder_path = ""
        self.save_org_directory = True  # За замовчуванням зберігаємо шлях
        self.organisation_folders = []  # Список знайдених папок-об'єктів
        self.current_selected_object = None  # Поточно обраний об'єкт
        self.load_org_config()
    
    def load_org_config(self):
        """Завантажує налаштування Organisation Control"""
        try:
            if os.path.exists(self.org_config_file):
                with open(self.org_config_file, "r", encoding="utf-8") as f:
                    config = json.load(f)
                    self.save_org_directory = config.get("save_directory", True)
                    if self.save_org_directory and "last_folder" in config:
                        saved_folder = config["last_folder"]
                        if os.path.exists(saved_folder):
                            self.org_folder_path.insert(0, saved_folder)
                            self.current_folder_path = saved_folder
                            self.scan_organisation_folder()
        except Exception as e:
            print(f"Помилка завантаження Organisation config: {str(e)}")
    
    def save_org_config(self):
        """Зберігає налаштування Organisation Control"""
        try:
            config = {
                "save_directory": self.save_org_directory,
                "last_folder": self.current_folder_path if self.save_org_directory else ""
            }
            with open(self.org_config_file, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Помилка збереження Organisation config: {str(e)}")
    
    def select_organisation_folder(self):
        """Відкриває діалог вибору папки для Organisation Control"""
        try:
            folder_path = filedialog.askdirectory(title="Оберіть папку з організацією")
            if folder_path:
                self.org_folder_path.delete(0, 'end')
                self.org_folder_path.insert(0, folder_path)
                self.current_folder_path = folder_path
                self.save_org_config()
                self.scan_organisation_folder()
        except Exception as e:
            messagebox.showerror("Помилка", f"Не вдалося обрати папку: {str(e)}")
    
    def scan_organisation_folder(self):
        """Сканує обрану папку та відображає підпапки як об'єкти"""
        try:
            if not self.current_folder_path or not os.path.exists(self.current_folder_path):
                return
            
            # Отримуємо всі підпапки в обраній директорії
            folders = []
            
            for item in os.listdir(self.current_folder_path):
                item_path = os.path.join(self.current_folder_path, item)
                if os.path.isdir(item_path):
                    # Перевіряємо наявність company.txt у папці
                    company_txt_path = os.path.join(item_path, "company.txt")
                    has_company_txt = os.path.exists(company_txt_path)
                    
                    folders.append({
                        'name': item,
                        'path': item_path,
                        'has_company': has_company_txt
                    })
            
            # Зберігаємо інформацію про папки
            self.organisation_folders = sorted(folders, key=lambda x: x['name'])
            
            # Оновлюємо dropdown список (тільки якщо не в режимі x4)
            if not hasattr(self, 'multi_org_mode') or not self.multi_org_mode:
                self.update_object_list()
                
                # Очищаємо відображення (поки що нічого не обрано)
                if hasattr(self, 'org_content_display') and self.org_content_display.winfo_exists():
                    try:
                        self.org_content_display.delete("0.0", "end")
                        self.org_content_display.insert("0.0", "Оберіть об'єкт для перегляду...")
                    except tk.TclError:
                        pass
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка сканування папки: {str(e)}")
    
    def parse_company_files(self, company_files):
        """Парсить company.txt файли з кожної папки та відображає їх вміст"""
        try:
            # Очищаємо область відображення
            self.org_content_display.delete("0.0", "end")
            
            if not company_files:
                self.org_content_display.insert("0.0", "📄 company.txt файли не знайдені")
                return
            
            all_content = []
            
            for company_file in company_files:
                try:
                    # Отримуємо назву папки (об'єкта)
                    folder_name = os.path.basename(os.path.dirname(company_file))
                    
                    with open(company_file, 'r', encoding='utf-8') as f:
                        content = f.read().strip()
                        if content:
                            # Парсимо поштовий код
                            parsed_content, postal_code = self.parse_postal_code(content)
                            
                            all_content.append(f"📁 Об'єкт: {folder_name}")
                            all_content.append(f"📄 company.txt:")
                            all_content.append(parsed_content)
                            all_content.append("=" * 50)  # Розділювач між об'єктами
                        else:
                            all_content.append(f"📁 Об'єкт: {folder_name}")
                            all_content.append("📄 company.txt порожній")
                            all_content.append("=" * 50)
                except Exception as e:
                    folder_name = os.path.basename(os.path.dirname(company_file)) if company_file else "Unknown"
                    all_content.append(f"📁 Об'єкт: {folder_name}")
                    all_content.append(f"❌ Помилка читання: {str(e)}")
                    all_content.append("=" * 50)
            
            if all_content:
                display_text = "\n".join(all_content)
                self.org_content_display.insert("0.0", display_text)
            else:
                self.org_content_display.insert("0.0", "📄 company.txt файли порожні або не читаються")
                
        except Exception as e:
            print(f"Помилка парсингу company.txt файлів: {str(e)}")
    
    def copy_org_content(self):
        """Копіює весь вміст .txt файлів в буфер обміну"""
        try:
            content = self.org_content_display.get("0.0", "end-1c").strip()
            if content and content != "📄 .txt файли не знайдені" and content != "📄 .txt файли порожні або не читаються":
                safe_clipboard_operation("set", content)
                # Логування в Account Manager
                if hasattr(self, 'log_display'):
                    self.log_display.configure(state='normal')
                    self.log_display.insert('end', f"📋 Скопійовано Organisation content\n")
                    self.log_display.configure(state='disabled')
                    self.log_display.see('end')
            else:
                messagebox.showwarning("Попередження", "Немає даних для копіювання")
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка копіювання: {str(e)}")
    
    def mark_folder_as_done(self):
        """Переміщує поточно обраний об'єкт в підпапку 'used'"""
        try:
            if not self.current_folder_path or not os.path.exists(self.current_folder_path):
                messagebox.showwarning("Попередження", "Не обрано робочу директорію")
                return
            
            # Перевіряємо чи обрано об'єкт
            if not self.current_selected_object:
                messagebox.showwarning("Попередження", "Не обрано об'єкт для переміщення")
                return
            
            # Отримуємо інформацію про обраний об'єкт
            folder_path = self.current_selected_object['path']
            folder_name = self.current_selected_object['name']
            
            # Створюємо папку 'used' в App Stuff якщо не існує
            used_dir = get_app_stuff_path("used")
            print(f"Створено папку: {used_dir}")
            
            # Шлях для переміщення
            destination_path = os.path.join(used_dir, folder_name)
            
            # Якщо папка з такою назвою вже існує в used, додаємо число
            counter = 1
            original_destination = destination_path
            while os.path.exists(destination_path):
                new_name = f"{folder_name}_{counter}"
                destination_path = os.path.join(used_dir, new_name)
                counter += 1
            
            # Переміщуємо папку
            shutil.move(folder_path, destination_path)
            
            # Логування
            self._safe_log_to_display(f"✅ Об'єкт '{folder_name}' переміщено в used/\n")
            
            # Оновлюємо список папок після переміщення
            self.scan_organisation_folder()
            
            # Скидаємо вибір
            self.current_selected_object = None
            self._safe_widget_update(self.object_dropdown, 'set', "Не обрано")
            self._safe_widget_update(self.object_status_label, 'configure', text="")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка переміщення об'єкта: {str(e)}")
    
    def mark_first_4_as_done(self):
        """Переміщує перші 4 об'єкти зі списку в підпапку 'used'"""
        print("\n" + "="*80)
        print("🔧 DEBUG: mark_first_4_as_done ВИКЛИКАНО")
        print("="*80)
        
        try:
            # DEBUG: Перевірка current_folder_path
            print(f"🔍 DEBUG: current_folder_path = {getattr(self, 'current_folder_path', 'НЕ ІСНУЄ')}")
            
            if not self.current_folder_path or not os.path.exists(self.current_folder_path):
                print("❌ DEBUG: Папка не обрана або не існує")
                messagebox.showwarning("Попередження", "Не обрано робочу директорію")
                return
            
            # DEBUG: Перевірка organisation_folders
            print(f"🔍 DEBUG: organisation_folders існує = {hasattr(self, 'organisation_folders')}")
            print(f"🔍 DEBUG: кількість organisation_folders = {len(self.organisation_folders) if hasattr(self, 'organisation_folders') else 0}")
            
            # Перевіряємо чи є об'єкти в списку
            if not self.organisation_folders or len(self.organisation_folders) == 0:
                print("❌ DEBUG: Немає об'єктів для переміщення")
                messagebox.showwarning("Попередження", "Немає об'єктів для переміщення")
                return
            
            # DEBUG: Перевірка режиму x4
            print(f"🔍 DEBUG: multi_org_mode = {getattr(self, 'multi_org_mode', False)}")
            print(f"🔍 DEBUG: multi_org_start_index = {getattr(self, 'multi_org_start_index', 0)}")
            
            # В режимі x4 беремо об'єкти з поточної групи
            if hasattr(self, 'multi_org_mode') and self.multi_org_mode:
                start_idx = self.multi_org_start_index
                end_idx = min(start_idx + 4, len(self.organisation_folders))
                folders_to_move = self.organisation_folders[start_idx:end_idx]
                print(f"📦 DEBUG: Режим x4 - беремо індекси {start_idx} до {end_idx}")
            else:
                # В звичайному режимі беремо перші 4
                folders_to_move = self.organisation_folders[:4]
                print(f"📦 DEBUG: Звичайний режим - беремо перші 4")
            
            print(f"📦 DEBUG: Папок для переміщення: {len(folders_to_move)}")
            for i, f in enumerate(folders_to_move):
                print(f"   {i+1}. {f['name']} -> {f['path']}")
            
            moved_count = 0
            moved_folders = []  # Зберігаємо інфо для можливого Undo
            
            # Створюємо папку 'used' в App Stuff якщо не існує
            used_dir = get_app_stuff_path("used")
            print(f"📁 DEBUG: used_dir = {used_dir}")
            
            # Логування початку
            self._safe_log_to_display(f"🚀 Переміщення {len(folders_to_move)} об'єктів...\n")
            
            # Переміщуємо кожну папку
            for folder_info in folders_to_move:
                try:
                    folder_path = folder_info['path']
                    folder_name = folder_info['name']
                    
                    print(f"\n🔄 DEBUG: Переміщення '{folder_name}'")
                    print(f"   Джерело: {folder_path}")
                    
                    # Шлях для переміщення
                    destination_path = os.path.join(used_dir, folder_name)
                    
                    # Якщо папка з такою назвою вже існує в used, додаємо число
                    counter = 1
                    original_destination = destination_path
                    while os.path.exists(destination_path):
                        new_name = f"{folder_name}_{counter}"
                        destination_path = os.path.join(used_dir, new_name)
                        counter += 1
                    
                    print(f"   Призначення: {destination_path}")
                    
                    # Переміщуємо папку
                    shutil.move(folder_path, destination_path)
                    moved_count += 1
                    
                    # Зберігаємо інфо для Undo
                    moved_folders.append({
                        'name': folder_name,
                        'from': self.current_folder_path,
                        'to': destination_path
                    })
                    
                    # Логування кожного переміщення
                    self._safe_log_to_display(f"  ✅ {moved_count}. '{folder_name}'\n")
                    
                    print(f"✅ DEBUG: Успішно переміщено {moved_count}/{len(folders_to_move)}: {folder_name}")
                    
                except Exception as e:
                    print(f"❌ DEBUG: Помилка переміщення {folder_name}: {e}")
                    import traceback
                    traceback.print_exc()
                    self._safe_log_to_display(f"  ⚠️ Помилка: {folder_name}\n")
            
            # Зберігаємо інфо про переміщені папки для Undo
            if not hasattr(self, 'last_moved_folders'):
                self.last_moved_folders = []
            self.last_moved_folders = moved_folders
            print(f"\n💾 DEBUG: Збережено {len(moved_folders)} папок для Undo")
            
            # Фінальне логування
            self._safe_log_to_display(f"✅ Завершено! Переміщено {moved_count} об'єктів\n")
            
            print(f"\n🔄 DEBUG: Оновлюємо список організацій...")
            # Оновлюємо список папок після переміщення
            self.scan_organisation_folder()
            
            print(f"🔍 DEBUG: Після scan - кількість organisation_folders = {len(self.organisation_folders)}")
            
            # Якщо в режимі x4, оновлюємо відображення
            if hasattr(self, 'multi_org_mode') and self.multi_org_mode:
                print(f"🔄 DEBUG: Оновлюємо x4 відображення")
                # Якщо переміщені всі орги з поточної групи, повертаємось назад
                if self.multi_org_start_index >= len(self.organisation_folders):
                    self.multi_org_start_index = max(0, len(self.organisation_folders) - 4)
                    print(f"   Коригуємо multi_org_start_index = {self.multi_org_start_index}")
                self.display_multi_org_group()
            else:
                # Скидаємо вибір
                self.current_selected_object = None
                self._safe_widget_update(self.object_dropdown, 'set', "Не обрано")
                self._safe_widget_update(self.object_status_label, 'configure', text="")
            
            print(f"\n✅ DEBUG: mark_first_4_as_done ЗАВЕРШЕНО")
            print("="*80 + "\n")
            
            # messagebox.showinfo("Успіх", f"Переміщено {moved_count} об'єктів в used/")
            
        except Exception as e:
            print(f"\n❌ DEBUG: КРИТИЧНА ПОМИЛКА в mark_first_4_as_done:")
            print(f"   {str(e)}")
            import traceback
            traceback.print_exc()
            print("="*80 + "\n")
            messagebox.showerror("Помилка", f"Помилка переміщення об'єктів: {str(e)}")
    
    def undo_last_move(self):
        """Відновлює останні переміщені папки назад з used"""
        print("\n" + "="*80)
        print("🔧 DEBUG: undo_last_move ВИКЛИКАНО")
        print("="*80)
        
        try:
            # Перевіряємо чи є збережені переміщені папки
            if not hasattr(self, 'last_moved_folders') or not self.last_moved_folders:
                print("⚠️ DEBUG: Немає збережених переміщень для Undo")
                messagebox.showinfo("Інформація", "Немає операцій для відміни")
                return
            
            print(f"📦 DEBUG: Знайдено {len(self.last_moved_folders)} папок для відновлення:")
            for i, folder in enumerate(self.last_moved_folders):
                print(f"   {i+1}. {folder['name']}")
                print(f"      З: {folder['to']}")
                print(f"      В: {folder['from']}")
            
            restored_count = 0
            
            # Відновлюємо кожну папку
            for folder_info in self.last_moved_folders:
                try:
                    folder_name = folder_info['name']
                    source_path = folder_info['to']  # Де зараз (в used)
                    dest_folder = folder_info['from']  # Куди повернути (робоча папка)
                    
                    # Повний шлях для відновлення
                    dest_path = os.path.join(dest_folder, folder_name)
                    
                    print(f"\n🔄 DEBUG: Відновлення '{folder_name}'")
                    print(f"   Джерело: {source_path}")
                    print(f"   Призначення: {dest_path}")
                    
                    # Перевіряємо чи існує папка в used
                    if not os.path.exists(source_path):
                        print(f"⚠️ DEBUG: Папка не знайдена в used: {source_path}")
                        continue
                    
                    # Перевіряємо чи не існує вже папка в робочій директорії
                    if os.path.exists(dest_path):
                        print(f"⚠️ DEBUG: Папка вже існує в робочій директорії: {dest_path}")
                        # Додаємо суфікс _restored
                        counter = 1
                        while os.path.exists(dest_path):
                            dest_path = os.path.join(dest_folder, f"{folder_name}_restored_{counter}")
                            counter += 1
                        print(f"   Нове призначення: {dest_path}")
                    
                    # Переміщуємо папку назад
                    shutil.move(source_path, dest_path)
                    restored_count += 1
                    
                    print(f"✅ DEBUG: Успішно відновлено {restored_count}: {folder_name}")
                    self._safe_log_to_display(f"  ↩️ Відновлено: '{folder_name}'\n")
                    
                except Exception as e:
                    print(f"❌ DEBUG: Помилка відновлення {folder_name}: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Очищаємо список переміщених папок
            self.last_moved_folders = []
            print(f"\n💾 DEBUG: Очищено список last_moved_folders")
            
            # Оновлюємо список організацій
            print(f"🔄 DEBUG: Оновлюємо список організацій...")
            self.scan_organisation_folder()
            
            print(f"🔍 DEBUG: Після scan - кількість organisation_folders = {len(self.organisation_folders)}")
            
            # Якщо в режимі x4, оновлюємо відображення
            if hasattr(self, 'multi_org_mode') and self.multi_org_mode:
                print(f"🔄 DEBUG: Оновлюємо x4 відображення")
                self.display_multi_org_group()
            
            print(f"\n✅ DEBUG: undo_last_move ЗАВЕРШЕНО")
            print(f"   Відновлено папок: {restored_count}")
            print("="*80 + "\n")
            
            # messagebox.showinfo("Успіх", f"Відновлено {restored_count} об'єктів з used/")
            
        except Exception as e:
            print(f"\n❌ DEBUG: КРИТИЧНА ПОМИЛКА в undo_last_move:")
            print(f"   {str(e)}")
            import traceback
            traceback.print_exc()
            print("="*80 + "\n")
            messagebox.showerror("Помилка", f"Помилка відновлення об'єктів: {str(e)}")
    
    def open_selected_organisation_folder(self):
        """Відкриває папку поточно обраного об'єкта в провіднику"""
        try:
            if not self.current_selected_object:
                messagebox.showwarning("Увага", "Спочатку оберіть об'єкт для відкриття папки")
                return
            
            folder_path = self.current_selected_object['path']
            folder_name = self.current_selected_object['name']
            
            if not os.path.exists(folder_path):
                messagebox.showerror("Помилка", f"Папка '{folder_name}' не існує або була переміщена")
                return
            
            # Відкриваємо папку залежно від ОС
            if os.name == 'nt':  # Windows
                os.startfile(folder_path)
            elif sys.platform == 'darwin':  # macOS
                subprocess.call(['open', folder_path])
            else:  # Linux
                subprocess.call(['xdg-open', folder_path])
            
            # Логування
            if hasattr(self, 'log_display'):
                self.log_display.configure(state='normal')
                self.log_display.insert('end', f"📂 Відкрито папку: {folder_name}\n")
                self.log_display.configure(state='disabled')
                self.log_display.see('end')
                
        except Exception as e:
            messagebox.showerror("Помилка", f"Не вдалося відкрити папку: {str(e)}")
    
    def update_object_list(self):
        """Оновлює dropdown список об'єктів"""
        try:
            # В режимі x4 не оновлюємо dropdown (його немає)
            if hasattr(self, 'multi_org_mode') and self.multi_org_mode:
                print("⚠️ DEBUG: update_object_list пропущено (режим x4)")
                return
            
            if not self.organisation_folders:
                object_names = ["Не обрано"]
            else:
                object_names = ["Не обрано"] + [folder['name'] for folder in self.organisation_folders]
            
            self.object_names = object_names
            
            # Безпечно оновлюємо dropdown якщо він існує
            if hasattr(self, 'object_dropdown') and self.object_dropdown.winfo_exists():
                try:
                    self.object_dropdown.configure(values=object_names)
                    self.object_dropdown.set("Не обрано")
                except tk.TclError as e:
                    print(f"⚠️ DEBUG: object_dropdown недоступний: {e}")
            
        except Exception as e:
            print(f"Помилка оновлення списку об'єктів: {str(e)}")
    
    def on_object_selected(self, selected_object):
        """Обробляє вибір об'єкта з dropdown"""
        if selected_object == "Не обрано":
            self.current_selected_object = None
            self.object_status_label.configure(text="")
            self.org_content_display.delete("0.0", "end")
            self.org_content_display.insert("0.0", "Оберіть об'єкт для перегляду...")
        else:
            # Знаходимо обраний об'єкт
            for folder in self.organisation_folders:
                if folder['name'] == selected_object:
                    self.current_selected_object = folder
                    status_text = "✅ З company.txt" if folder['has_company'] else "❌ Без company.txt"
                    self.object_status_label.configure(text=status_text)
                    self.display_selected_object()
                    break
    
    def parse_postal_code(self, content):
        """Парсить поштовий код з останнього рядка адреси"""
        import re
        
        lines = content.strip().split('\n')
        if not lines:
            return content, None
            
        # Беремо останній рядок
        last_line = lines[-1].strip()
        
        # Шукаємо поштовий код в форматі 5 цифр в кінці рядка
        postal_code_match = re.search(r'\b(\d{5})\b', last_line)
        
        if postal_code_match:
            postal_code = postal_code_match.group(1)
            # Видаляємо поштовий код з останнього рядка
            new_last_line = last_line.replace(postal_code, '').strip().rstrip(',').strip()
            
            # Створюємо новий контент з окремим рядком для поштового коду
            new_lines = lines[:-1] + [new_last_line, postal_code]
            new_content = '\n'.join(new_lines)
            
            return new_content, postal_code
        
        return content, None
    
    def display_selected_object(self):
        """Відображає вміст обраного об'єкта"""
        try:
            if not self.current_selected_object:
                return
            
            # Очищаємо область відображення
            self.org_content_display.delete("0.0", "end")
            
            folder_name = self.current_selected_object['name']
            company_txt_path = os.path.join(self.current_selected_object['path'], "company.txt")
            
            if self.current_selected_object['has_company']:
                try:
                    with open(company_txt_path, 'r', encoding='utf-8') as f:
                        content = f.read().strip()
                        if content:
                            # Парсимо поштовий код
                            parsed_content, postal_code = self.parse_postal_code(content)
                            
                            display_text = f"📁 Об'єкт: {folder_name}\n📄 company.txt:\n\n{parsed_content}"
                        else:
                            display_text = f"📁 Об'єкт: {folder_name}\n📄 company.txt порожній"
                except Exception as e:
                    display_text = f"📁 Об'єкт: {folder_name}\n❌ Помилка читання company.txt: {str(e)}"
            else:
                display_text = f"📁 Об'єкт: {folder_name}\n❌ company.txt відсутній"
            
            self.org_content_display.insert("0.0", display_text)
            
        except Exception as e:
            print(f"Помилка відображення об'єкта: {str(e)}")
    
    def object_up(self):
        """Переміщує вибір на одну позицію вверх"""
        try:
            current_value = self.object_dropdown.get()
            current_index = self.object_names.index(current_value)
            if current_index > 0:
                new_index = current_index - 1
                new_value = self.object_names[new_index]
                self.object_dropdown.set(new_value)
                self.on_object_selected(new_value)
        except (ValueError, IndexError):
            pass
    
    def object_down(self):
        """Переміщує вибір на одну позицію вниз"""
        try:
            current_value = self.object_dropdown.get()
            current_index = self.object_names.index(current_value)
            if current_index < len(self.object_names) - 1:
                new_index = current_index + 1
                new_value = self.object_names[new_index]
                self.object_dropdown.set(new_value)
                self.on_object_selected(new_value)
        except (ValueError, IndexError):
            pass
    
    def _setup_textbox_click_to_copy(self, textbox):
        """Універсальна функція для налаштування click-to-copy для будь-якого textbox"""
        try:
            # Отримуємо внутрішній Tkinter віджет
            if hasattr(textbox, '_textbox'):
                inner_text = textbox._textbox
                
                # Прив'язуємо подію кліку
                inner_text.bind("<Button-1>", lambda event: self._textbox_on_click_copy(event, inner_text))
                inner_text.bind('<Key>', lambda event: 'break')  # Блокуємо редагування
                
        except Exception as e:
            print(f"Помилка налаштування click-to-copy: {e}")

    def _textbox_on_click_copy(self, event, inner_text):
        """Обробка кліку - копіює рядок під курсором"""
        try:
            # Отримуємо позицію кліку
            click_index = inner_text.index(f"@{event.x},{event.y}")
            line_start = inner_text.index(f"{click_index} linestart")
            line_end = inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка
            line_text = inner_text.get(line_start, line_end).strip()
            
            if line_text:
                # Копіюємо рядок
                safe_clipboard_operation("set", line_text)
                
                # Логування скопійованого
                print(f"📋 Скопійовано (x4): {line_text}")
                
                # Візуальний фідбек
                inner_text.tag_remove('copied_line', '1.0', 'end')
                inner_text.tag_add('copied_line', line_start, line_end)
                inner_text.tag_config('copied_line', background='#404040', foreground='#90EE90')
                
                # Прибираємо підсвітку через 500мс
                inner_text.after(500, lambda: inner_text.tag_remove('copied_line', '1.0', 'end'))
                
        except Exception as e:
            print(f"Помилка копіювання рядка: {e}")

    def _setup_org_click_to_copy(self):
        """Налаштування click-to-copy для Organisation Control company.txt області"""
        try:
            # Отримуємо внутрішній Tkinter віджет
            if hasattr(self.org_content_display, '_textbox'):
                inner_text = self.org_content_display._textbox
                
                # Прив'язуємо події
                inner_text.bind("<Button-1>", self._org_on_click_copy)
                inner_text.bind("<Button-3>", self._org_show_context_menu)
                inner_text.bind('<Key>', lambda event: 'break')  # Блокуємо редагування
                
                # Створюємо контекстне меню для Organisation Control
                self.org_context_menu = tk.Menu(inner_text, tearoff=0, bg='#2e2e2e', fg='white', 
                                              activebackground='#404040', activeforeground='white')
                self.org_context_menu.add_command(label="📋 Копіювати рядок", command=self._org_copy_current_line)
                self.org_context_menu.add_command(label="📋 Копіювати все", command=self.copy_org_content)
                self.org_context_menu.add_separator()
                self.org_context_menu.add_command(label="🔤 Вибрати все", command=self._org_select_all)
                
                print("Click-to-copy налаштовано для Organisation Control!")
            
        except Exception as e:
            print(f"Помилка налаштування click-to-copy для Organisation Control: {e}")
    
    def _org_on_click_copy(self, event):
        """Обробка кліку - копіює рядок під курсором в company.txt області"""
        try:
            # Отримуємо внутрішній текстовий віджет
            inner_text = self.org_content_display._textbox
            
            # Отримуємо позицію кліку
            click_index = inner_text.index(f"@{event.x},{event.y}")
            line_start = inner_text.index(f"{click_index} linestart")
            line_end = inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка
            line_text = inner_text.get(line_start, line_end).strip()
            
            if line_text:
                # Копіюємо рядок
                safe_clipboard_operation("set", line_text)
                
                # Візуальний фідбек
                inner_text.tag_remove('org_copied_line', '1.0', 'end')
                inner_text.tag_add('org_copied_line', line_start, line_end)
                inner_text.tag_config('org_copied_line', background='#404040', foreground='#90EE90')
                
                # Прибираємо підсвітку через 500мс
                inner_text.after(500, lambda: inner_text.tag_remove('org_copied_line', '1.0', 'end'))
                
                # Логуємо копіювання в Account Manager лог
                if hasattr(self, 'log_display'):
                    self.log_display.configure(state='normal')
                    self.log_display.insert('end', f"📋 Скопійовано: {line_text}\n")
                    self.log_display.configure(state='disabled')
                    self.log_display.see('end')
                
                print(f"📋 Скопійовано з company.txt: {line_text}")
                
        except Exception as e:
            print(f"Помилка копіювання в Organisation Control: {e}")
            
    def _org_show_context_menu(self, event):
        """Показати контекстне меню для Organisation Control"""
        try:
            self.org_context_menu.post(event.x_root, event.y_root)
        except:
            pass
            
    def _org_copy_current_line(self):
        """Копіювати поточний рядок з company.txt"""
        try:
            inner_text = self.org_content_display._textbox
            current_pos = inner_text.index(tk.INSERT)
            line_start = inner_text.index(f"{current_pos} linestart")
            line_end = inner_text.index(f"{current_pos} lineend")
            line_text = inner_text.get(line_start, line_end).strip()
            
            if line_text:
                safe_clipboard_operation("set", line_text)
                
                # Логуємо в Account Manager
                if hasattr(self, 'log_display'):
                    self.log_display.configure(state='normal')
                    self.log_display.insert('end', f"📋 Скопійовано рядок: {line_text}\n")
                    self.log_display.configure(state='disabled')
                    self.log_display.see('end')
        except:
            pass
            
    def _org_select_all(self):
        """Вибрати весь текст в company.txt області"""
        try:
            inner_text = self.org_content_display._textbox
            inner_text.tag_add(tk.SEL, "1.0", tk.END)
            inner_text.mark_set(tk.INSERT, "1.0")
            inner_text.see(tk.INSERT)
        except:
            pass
    
    def copy_line_on_click(self, event, textbox):
        """Обробка кліку для копіювання рядка з textbox (для режиму x4)"""
        try:
            # Перевіряємо чи textbox існує
            if not textbox.winfo_exists():
                return
            
            # Отримуємо внутрішній текстовий віджет
            if hasattr(textbox, '_textbox'):
                inner_text = textbox._textbox
            else:
                inner_text = textbox
            
            # Перевіряємо чи inner_text існує
            if not hasattr(inner_text, 'winfo_exists') or not inner_text.winfo_exists():
                return
            
            # Отримуємо позицію кліку
            click_index = inner_text.index(f"@{event.x},{event.y}")
            line_start = inner_text.index(f"{click_index} linestart")
            line_end = inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка
            line_text = inner_text.get(line_start, line_end).strip()
            
            if line_text:
                # Копіюємо рядок
                safe_clipboard_operation("set", line_text)
                
                # Візуальний фідбек
                try:
                    inner_text.tag_remove('copied_line', '1.0', 'end')
                    inner_text.tag_add('copied_line', line_start, line_end)
                    inner_text.tag_config('copied_line', background='#404040', foreground='#90EE90')
                    
                    # Прибираємо підсвітку через 500мс
                    def remove_highlight():
                        try:
                            if inner_text.winfo_exists():
                                inner_text.tag_remove('copied_line', '1.0', 'end')
                        except:
                            pass
                    inner_text.after(500, remove_highlight)
                except:
                    pass
                
                # Логуємо копіювання
                self._safe_log_to_display(f"📋 Скопійовано: {line_text}\n")
                print(f"📋 Скопійовано: {line_text}")
                
        except (tk.TclError, AttributeError):
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"Помилка копіювання: {e}")
    
    def generate_2fa_code(self):
        """Генерує 2FA код на основі секретного ключа в полі 2FA та копіює в буфер"""
        try:
            # Отримуємо значення з поля 2FA
            field_value = self.entries_dict["2FA"].get().strip()
            
            if not field_value:
                # Тихо повертаємося, якщо поле порожнє
                return
            
            # Якщо це секретний ключ
            secret_key = field_value
            
            # Очищаємо ключ від зайвих символів та приводимо до верхнього регістру
            secret_key_clean = re.sub(r'[\s\-=]+', '', secret_key).upper()
            
            # Перевіряємо мінімальну довжину (зазвичай Base32 ключі від 16 символів)
            if len(secret_key_clean) < 16:
                return
            
            # Перевіряємо формат Base32 (допускаємо тільки A-Z та 2-7)
            if not re.match(r'^[A-Z2-7]+$', secret_key_clean):
                return
            
            # Генеруємо 2FA код через lazy import
            pyotp_module = get_pyotp()
            if not pyotp_module:
                print("⚠️ PyOTP не доступний")
                return
                
            totp = pyotp_module.TOTP(secret_key_clean)
            current_code = totp.now()
            
            # Копіюємо код в буфер обміну
            success = copy_to_clipboard(current_code)
            if success:
                print(f"🔑 2FA код {current_code} скопійовано в буфер")
            else:
                print("❌ Не вдалося скопіювати 2FA код")
            
            # Змінюємо колір кнопки на короткий час для візуального фідбеку
            self.provide_visual_feedback()
            
        except Exception as e:
            # Тихо ігноруємо помилки, щоб не засмічувати інтерфейс
            pass
    
    def provide_visual_feedback(self):
        """Надає візуальний фідбек після копіювання 2FA коду"""
        try:
            if hasattr(self, 'gen_2fa_btn'):
                # Змінюємо колір кнопки на синій на 300мс для підтвердження копіювання
                self.gen_2fa_btn.configure(fg_color="blue")
                self.after(300, lambda: self.gen_2fa_btn.configure(fg_color="green"))
        except:
            pass

    def save_last_generated_profile(self, profile_template, proxy_template, last_profile_num, last_proxy_num):
        """Зберігає останній згенерований профіль в конфіг"""
        try:
            last_profile_data = {
                "profile_template": profile_template,
                "proxy_template": proxy_template,
                "last_profile_num": last_profile_num,
                "last_proxy_num": last_proxy_num,
                "timestamp": time.time()
            }
            
            with open(self.octo_profile_config_file, "w", encoding="utf-8") as f:
                json.dump(last_profile_data, f, ensure_ascii=False, indent=2)
            
            # Оновлюємо стан кнопки після збереження
            self.update_choose_last_button_state()
                
        except Exception as e:
            print(f"Помилка збереження останнього профілю: {str(e)}")
    
    def load_last_generated_profile(self):
        """Завантажує останній згенерований профіль з конфігу"""
        try:
            if os.path.exists(self.octo_profile_config_file):
                with open(self.octo_profile_config_file, "r", encoding="utf-8") as f:
                    return json.load(f)
            return None
        except Exception as e:
            print(f"Помилка завантаження останнього профілю: {str(e)}")
            return None
    
    def choose_last_plus_one(self):
        """Завантажує останній профіль +1 в поля генератора"""
        try:
            last_profile = self.load_last_generated_profile()
            
            if not last_profile:
                messagebox.showinfo("Інформація", "Спочатку згенеруйте хоча б один профіль!")
                return
            
            # Отримуємо дані останнього профілю
            profile_template = last_profile.get("profile_template", "")
            proxy_template = last_profile.get("proxy_template", "")
            last_profile_num = last_profile.get("last_profile_num", 0)
            last_proxy_num = last_profile.get("last_proxy_num", 0)
            
            # Створюємо нові шаблони з +1
            new_profile_template = self.increment_profile_number(profile_template, last_profile_num + 1)
            new_proxy_template = self.increment_proxy_sessionid(proxy_template, last_proxy_num + 1)
            
            # Заповнюємо поля
            safe_text_input(self.octo_profile_entry, new_profile_template)
            safe_text_input(self.octo_proxy_entry, new_proxy_template)
            
            print(f"Завантажено профіль +1: {new_profile_template}")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Не вдалося завантажити останній профіль +1: {str(e)}")
    
    def increment_profile_number(self, profile_template, new_number):
        """Змінює останнє число після останнього підкреслення на нове
        Alex_FarmPunch_USA_5_131125_1 + new_number=100 -> Alex_FarmPunch_USA_5_131125_100
        """
        import re
        # Шукаємо останнє число після останнього підкреслення
        match = re.search(r'_(\d+)$', profile_template)
        if match:
            # Замінюємо тільки останнє число
            return profile_template[:match.start()] + f'_{new_number}'
        # Якщо не знайшли, додаємо число в кінець
        return f"{profile_template}_{new_number}"
    
    def increment_proxy_sessionid(self, proxy_template, new_sessionid):
        """Змінює sessionid у шаблоні проксі на новий"""
        import re
        # Знаходимо sessionid-номер в проксі
        match = re.search(r'sessionid-(\d+)', proxy_template)
        if match:
            old_sessionid = match.group(1)
            return proxy_template.replace(f"sessionid-{old_sessionid}", f"sessionid-{new_sessionid}")
        return proxy_template
    
    def update_choose_last_button_state(self):
        """Оновлює стан кнопки Choose Last +1 в залежності від наявності збереженого профілю"""
        try:
            if hasattr(self, 'choose_last_btn'):
                last_profile = self.load_last_generated_profile()
                if last_profile:
                    self.choose_last_btn.configure(state="normal", fg_color="orange")
                else:
                    self.choose_last_btn.configure(state="disabled", fg_color="gray")
        except:
            pass

    # ===== CSV CONFIG METHODS =====
    def init_csv_config(self):
        """Ініціалізує локальний CSV конфіг для збереження даних аккаунтів"""
        try:
            if not os.path.exists(self.LOCAL_CSV_CONFIG):
                with open(self.LOCAL_CSV_CONFIG, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    writer.writerow([
                        'Timestamp', 'Email/Login', 'Password', '2FA', 'Backup Codes', 
                        'Reserve Mail', 'Credit Card', 'Конверсія', 
                        'API Cloudflare', 'Cloudflare ID', 'Cloudflare Password', 'Status'
                    ])
                print(f"✅ CSV конфіг створено: {self.LOCAL_CSV_CONFIG}")
        except Exception as e:
            print(f"❌ Помилка при створенні CSV конфігу: {e}")

    def save_account_to_csv(self):
        """Зберігає дані поточного аккаунту безпосередньо в CSV з sheets_textarea"""
        try:
            # Збираємо дані з усіх полів
            account_data = {
                'Email/Login': self.entries_dict["Email/Login"].get().strip(),
                'Password': self.entries_dict["Password"].get().strip(),
                '2FA': self.entries_dict["2FA"].get().strip(),
                'Backup Codes': self.entries_dict["Backup Codes"].get().strip(),
                'Reserve Mail': self.entries_dict["Reserve Mail"].get().strip(),
                'Credit Card': self.entries_dict["Credit Card"].get().strip(),
                'Конверсія': self.entries_dict["Конверсія"].get().strip(),
                'API Cloudflare': self.entries_dict["API Cloudflare"].get().strip(),
                'Cloudflare ID': self.entries_dict["Cloudflare ID"].get().strip(),
                'Cloudflare Password': self.entries_dict["Cloudflare Password"].get().strip(),
            }
            
            # Перевіряємо на порожні поля
            if not account_data['Email/Login']:
                # Тихо пропускаємо якщо немає email
                return
            
            # Отримуємо CSV з textarea
            csv_content = self.sheets_textarea.get("0.0", "end-1c").strip()
            if not csv_content:
                print("⚠️ CSV дані відсутні в Profile Search")
                return
            
            # Визначаємо роздільник
            delimiter = ';' if ';' in csv_content.split('\n')[0] else ','
            
            # Парсимо CSV
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            headers = csv_reader.fieldnames or []
            
            # Додаємо нові колонки якщо їх немає
            new_columns = ['API Cloudflare', 'Cloudflare ID', 'Cloudflare Password']
            for col in new_columns:
                if col not in headers:
                    headers.append(col)
            
            # Визначаємо колонку профілю
            profile_column = None
            for header in headers:
                if any(name.lower() in header.lower() for name in ["название", "назва", "profile", "профіль", "name", "ім'я"]):
                    profile_column = header
                    break
            
            # Шукаємо колонку пошти
            email_column = None
            for header in headers:
                if any(name.lower() in header.lower() for name in ["почта", "пошта", "email", "mail", "login"]):
                    email_column = header
                    break
            
            if not email_column:
                email_column = "Почта"
            
            # Читаємо всі рядки
            rows = []
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            found = False
            
            for row in csv_reader:
                # Шукаємо рядок з нашим email
                if row.get(email_column, '').strip().lower() == account_data['Email/Login'].lower():
                    # Оновлюємо існуючий рядок
                    row['API Cloudflare'] = account_data['API Cloudflare']
                    row['Cloudflare ID'] = account_data['Cloudflare ID']
                    row['Cloudflare Password'] = account_data['Cloudflare Password']
                    # Оновлюємо також основні поля якщо вони змінились
                    if 'Пароль' in row:
                        row['Пароль'] = account_data['Password']
                    if '2фа' in row:
                        row['2фа'] = account_data['2FA']
                    if 'Бэкап' in row:
                        row['Бэкап'] = account_data['Backup Codes']
                    if 'Резервка' in row:
                        row['Резервка'] = account_data['Reserve Mail']
                    if 'Карта' in row:
                        row['Карта'] = account_data['Credit Card']
                    if 'конверсія' in row:
                        row['конверсія'] = account_data['Конверсія']
                    found = True
                rows.append(row)
            
            if not found:
                print(f"⚠️ Email '{account_data['Email/Login']}' не знайдено в CSV, пропускаємо збереження")
                return
            
            # Формуємо оновлений CSV
            output = io.StringIO()
            writer = csv.DictWriter(output, fieldnames=headers, delimiter=delimiter)
            writer.writeheader()
            writer.writerows(rows)
            
            # Оновлюємо textarea
            new_csv_content = output.getvalue().strip()
            self.sheets_textarea.delete("0.0", "end")
            self.sheets_textarea.insert("0.0", new_csv_content)
            
            # Зберігаємо в локальний файл
            with open(self.LOCAL_PROFILES_CSV, 'w', encoding='utf-8', newline='') as f:
                f.write(new_csv_content)
            
            self.log_action(f"✅ Дані для '{account_data['Email/Login']}' оновлено в CSV")
            print(f"💾 Дані збережені в локальний файл: {self.LOCAL_PROFILES_CSV}")
            
        except Exception as e:
            print(f"❌ Помилка при збереженні до CSV: {e}")
            self.log_action(f"❌ Помилка при збереженні до CSV: {e}")

    def load_account_from_csv(self):
        """Завантажує збережені аккаунти з CSV конфігу"""
        try:
            if not os.path.exists(self.LOCAL_CSV_CONFIG):
                messagebox.showinfo("Інформація", "Локальний конфіг ще не створено")
                return
            
            with open(self.LOCAL_CSV_CONFIG, 'r', encoding='utf-8') as f:
                reader = csv.DictReader(f)
                accounts = list(reader)
            
            if not accounts:
                messagebox.showinfo("Інформація", "У конфігу немає збережених аккаунтів")
                return
            
            # Показуємо вибір останнього аккаунту
            last_account = accounts[-1]
            self.entries_dict["Email/Login"].delete(0, tk.END)
            self.entries_dict["Email/Login"].insert(0, last_account.get('Email/Login', ''))
            self.entries_dict["Password"].delete(0, tk.END)
            self.entries_dict["Password"].insert(0, last_account.get('Password', ''))
            self.entries_dict["2FA"].delete(0, tk.END)
            self.entries_dict["2FA"].insert(0, last_account.get('2FA', ''))
            self.entries_dict["API Cloudflare"].delete(0, tk.END)
            self.entries_dict["API Cloudflare"].insert(0, last_account.get('API Cloudflare', ''))
            self.entries_dict["Cloudflare ID"].delete(0, tk.END)
            self.entries_dict["Cloudflare ID"].insert(0, last_account.get('Cloudflare ID', ''))
            self.entries_dict["Cloudflare Password"].delete(0, tk.END)
            self.entries_dict["Cloudflare Password"].insert(0, last_account.get('Cloudflare Password', ''))
            
            self.log_action(f"✅ Дані завантажені для: {last_account.get('Email/Login', 'Unknown')}")
            messagebox.showinfo("Успіх", f"Завантажено дані для: {last_account.get('Email/Login', '')}")
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка при завантаженні з CSV: {e}")

    # ===== DATABASE UPLOAD METHODS =====
    def upload_to_database(self):
        """Відправляє дані аккаунту до датабази (одиночний або масово)"""
        try:
            import requests
            
            email = self.entries_dict["Email/Login"].get().strip()
            api_cf = self.entries_dict["API Cloudflare"].get().strip()
            
            # Якщо поля заповнені - додаємо конкретний акаунт
            if email or api_cf:
                self._upload_single_account()
            else:
                # Якщо поля порожні - масове додавання з CSV
                self._upload_bulk_accounts()
                
        except Exception as e:
            self.log_action(f"❌ Помилка при відправці: {e}")
            messagebox.showerror("Помилка", f"Помилка при відправці: {e}")
    
    def _upload_single_account(self):
        """Додає один акаунт до датабази"""
        try:
            import requests
            
            email = self.entries_dict["Email/Login"].get().strip()
            api_cf = self.entries_dict["API Cloudflare"].get().strip()
            cf_id = self.entries_dict["Cloudflare ID"].get().strip()
            cf_password = self.entries_dict["Cloudflare Password"].get().strip()
            
            # Формуємо дані для відправки (назви полів як на сайті)
            data = {
                'mail': email,              # Mail на сайті = Email/Login
                'api_key': api_cf,          # Api Key на сайті = API Cloudflare
                'account_id': cf_id,        # Accout id на сайті = Cloudflare ID
                'password': cf_password     # Password на сайті = Cloudflare Password
            }
            
            # Показуємо статус
            self.log_action(f"⬆️ Відправляю дані до датабази: {email or api_cf}...")
            self.log_to_db_panel(f"⬆️ Відправка: {email or api_cf}")
            
            # Відправляємо запит
            try:
                response = requests.post(self.DATABASE_URL, data=data, timeout=10)
                
                # Парсимо JSON відповідь
                try:
                    result = response.json()
                    status = result.get('status', '')
                    message = result.get('message', '')
                    
                    print(f"Response: Status={status}, Message='{message}'")
                    
                    if status == 'success':
                        self.log_action(f"✅ Успішно додано до датабази: {email or api_cf}")
                        self.log_to_db_panel(f"✅ ДОДАНО: {email or api_cf}")
                        
                        # Зберігаємо до локального CSV
                        self.save_account_to_csv()
                        
                        # Очищуємо поля
                        for entry in self.entries_dict.values():
                            entry.delete(0, tk.END)
                    
                    elif status == 'error':
                        self.log_action(f"⏭️ Помилка: {message} ({email or api_cf})")
                        self.log_to_db_panel(f"⏭️ {message[:30]}: {email or api_cf}")
                        if "існує" in message.lower():
                            messagebox.showinfo("Інформація", message)
                        else:
                            messagebox.showwarning("Помилка", message)
                    
                    else:
                        self.log_action(f"⚠️ Невідома відповідь: {message}")
                        self.log_to_db_panel(f"⚠️ Невідомий статус")
                        messagebox.showwarning("Попередження", f"Відповідь:\n{message}")
                    
                except json.JSONDecodeError:
                    response_text = response.text[:200]
                    self.log_action(f"⚠️ Сервер повернув не-JSON: {response_text}")
                    self.log_to_db_panel(f"❌ Не-JSON відповідь")
                    messagebox.showerror("Помилка", f"Сервер повернув некоректну відповідь:\n{response_text}")
                    
                if response.status_code != 200:
                    self.log_action(f"⚠️ Помилка датабази: HTTP {response.status_code}")
                    self.log_to_db_panel(f"❌ HTTP {response.status_code}: {email or api_cf}")
                    messagebox.showwarning("Попередження", f"Помилка датабази: {response.status_code}")
                    
            except requests.exceptions.Timeout:
                self.log_action(f"❌ Помилка: Час очікування вичерпаний")
                self.log_to_db_panel(f"⏱️ Timeout: {email or api_cf}")
                messagebox.showerror("Помилка", "Час очікування вичерпаний!")
            except requests.exceptions.ConnectionError:
                self.log_action(f"❌ Помилка: Проблема з з'єднанням")
                self.log_to_db_panel(f"🔌 Connection Error")
                messagebox.showerror("Помилка", "Проблема з з'єднанням з датабазою!")
                
        except Exception as e:
            self.log_action(f"❌ Помилка при відправці: {e}")
            messagebox.showerror("Помилка", f"Помилка при відправці: {e}")
    
    def _upload_bulk_accounts(self):
        """Масово додає всі акаунти з CSV, які ще не додані до датабази"""
        try:
            import requests
            import csv
            import io
            
            # Перевіряємо наявність CSV даних
            csv_content = self.sheets_textarea.get("0.0", "end").strip()
            if not csv_content:
                messagebox.showwarning("Попередження", "Немає даних для масового додавання!\nВставте CSV дані в поле Profile Search.")
                return
            
            # Підтвердження від користувача
            confirm = messagebox.askyesno(
                "Масове додавання", 
                "Розпочати масове додавання всіх акаунтів з CSV до датабази?\n\n"
                "Акаунти, які вже є в локальному CSV конфігу, будуть пропущені."
            )
            if not confirm:
                return
            
            # Читаємо локальний CSV для перевірки дублікатів
            existing_emails = set()
            existing_api_keys = set()
            if os.path.exists(self.LOCAL_CSV_CONFIG):
                with open(self.LOCAL_CSV_CONFIG, 'r', encoding='utf-8') as f:
                    reader = csv.DictReader(f)
                    for row in reader:
                        if row.get('Email/Login'):
                            existing_emails.add(row['Email/Login'].strip().lower())
                        if row.get('API Cloudflare'):
                            existing_api_keys.add(row['API Cloudflare'].strip())
            
            # Визначаємо роздільник
            delimiter = ';' if ';' in csv_content.split('\n')[0] else ','
            
            # Парсимо CSV дані
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            
            # Лічильники
            total = 0
            uploaded = 0
            skipped = 0
            errors = 0
            
            self.log_action(f"🚀 Розпочато масове додавання акаунтів...")
            self.log_to_db_panel(f"🚀 Розпочато масове додавання...")
            
            for row in csv_reader:
                total += 1
                
                # Отримуємо дані
                email = row.get('Почта', '').strip()
                password = row.get('Пароль', '').strip()
                api_cf = ''  # API Cloudflare може бути в іншій колонці
                cf_id = ''
                
                # Пропускаємо якщо немає email
                if not email:
                    skipped += 1
                    continue
                
                # Перевіряємо чи акаунт вже є в локальному CSV
                if email.lower() in existing_emails or api_cf in existing_api_keys:
                    self.log_action(f"⏭️ Пропускаємо дублікат: {email}")
                    self.log_to_db_panel(f"⏭️ Дублікат: {email}")
                    skipped += 1
                    continue
                
                # Генеруємо Cloudflare Password
                cf_password = f"{password}@" if password else ""
                
                # Формуємо дані для відправки
                data = {
                    'email': email,
                    'api_key': api_cf,
                    'cloudflare_id': cf_id,
                    'cloudflare_password': cf_password,
                    'timestamp': datetime.datetime.now().isoformat(),
                    'source': 'Gmail Hacks Bulk Upload'
                }
                
                try:
                    response = requests.post(self.DATABASE_URL, data=data, timeout=10)
                    
                    if response.status_code == 200:
                        uploaded += 1
                        self.log_action(f"✅ [{uploaded}/{total}] Додано: {email}")
                        self.log_to_db_panel(f"✅ [{uploaded}] {email}")
                        
                        # Додаємо до локального CSV
                        existing_emails.add(email.lower())
                        if api_cf:
                            existing_api_keys.add(api_cf)
                    else:
                        errors += 1
                        self.log_action(f"⚠️ [{total}] Помилка HTTP {response.status_code}: {email}")
                        self.log_to_db_panel(f"❌ HTTP {response.status_code}: {email}")
                        
                except Exception as e:
                    errors += 1
                    self.log_action(f"❌ [{total}] Помилка відправки {email}: {str(e)}")
                    self.log_to_db_panel(f"❌ Error: {email}")
            
            # Підсумок
            summary = (
                f"📊 Масове додавання завершено!\n\n"
                f"Всього оброблено: {total}\n"
                f"✅ Додано: {uploaded}\n"
                f"⏭️ Пропущено (дублікати): {skipped}\n"
                f"❌ Помилки: {errors}"
            )
            
            self.log_action(summary.replace('\n', ' | '))
            self.log_to_db_panel(f"📊 Завершено: {uploaded} додано, {skipped} пропущено, {errors} помилок")
            messagebox.showinfo("Завершено", summary)
            
        except Exception as e:
            self.log_action(f"❌ Критична помилка масового додавання: {e}")
            messagebox.showerror("Помилка", f"Критична помилка: {e}")
    
    def upload_cloudflare_accounts(self):
        """Знаходить та завантажує всі акаунти з Cloudflare даними до датабази"""
        try:
            import requests
            import csv
            import io
            import time
            import re
            
            # Отримуємо CSV дані
            csv_content = self.sheets_textarea.get("0.0", "end").strip()
            if not csv_content:
                messagebox.showwarning("Попередження", "Немає CSV даних!\nЗавантажте дані з Google Sheets.")
                return
            
            # Визначаємо роздільник
            delimiter = ';' if ';' in csv_content.split('\n')[0] else ','
            
            # Парсимо CSV
            csv_reader = csv.DictReader(io.StringIO(csv_content), delimiter=delimiter)
            
            # Збираємо акаунти з Cloudflare даними (БЕЗ дублікатів + ВАЛІДАЦІЯ)
            cloudflare_accounts = []
            seen_emails = set()
            email_pattern = re.compile(r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$')
            
            for row in csv_reader:
                email = row.get('Почта', '').strip().lower()
                api_cf = row.get('API Cloudflare', '').strip()
                cf_id = row.get('Cloudflare ID', '').strip()
                cf_password = row.get('Cloudflare Password', '').strip()
                
                # Пропускаємо якщо немає email або вже бачили цей email
                if not email or email in seen_emails:
                    continue
                
                # Валідація email
                if not email_pattern.match(email):
                    self.log_to_db_panel(f"⚠️ Невалідний email: {email}")
                    continue
                
                # Перевіряємо що ВСІ 3 Cloudflare поля заповнені (API, ID, Password)
                if api_cf and cf_id and cf_password:
                    # Додаткова валідація довжини
                    if len(api_cf) < 10 or len(cf_id) < 10:
                        self.log_to_db_panel(f"⚠️ Короткі дані CF: {email}")
                        continue
                    
                    cloudflare_accounts.append({
                        'email': email,
                        'api_cf': api_cf,
                        'cf_id': cf_id,
                        'cf_password': cf_password
                    })
                    seen_emails.add(email)
            
            if not cloudflare_accounts:
                self.log_to_db_panel("⚠️ Не знайдено акаунтів з Cloudflare даними")
                messagebox.showinfo("Інформація", "Не знайдено акаунтів з Cloudflare даними в CSV")
                return
            
            # Підтвердження
            confirm = messagebox.askyesno(
                "Cloudflare Upload",
                f"Знайдено {len(cloudflare_accounts)} акаунтів з Cloudflare даними.\n\n"
                f"Розпочати завантаження до датабази?\n"
                f"(Затримка між запитами: 0.5 сек)"
            )
            if not confirm:
                return
            
            # Лічильники
            total = len(cloudflare_accounts)
            uploaded = 0
            errors = 0
            skipped = 0
            self.upload_cancelled = False
            
            self.log_to_db_panel(f"☁️ Розпочато завантаження {total} Cloudflare акаунтів...")
            
            # Завантажуємо по черзі
            for idx, account in enumerate(cloudflare_accounts, 1):
                # Перевірка на скасування
                if getattr(self, 'upload_cancelled', False):
                    self.log_to_db_panel(f"⏹️ Завантаження скасовано користувачем")
                    break
                email = account['email']
                
                # Формуємо дані (назви полів як на сайті)
                data = {
                    'mail': email,              # Mail на сайті = Email/Login
                    'api_key': account['api_cf'],          # Api Key на сайті = API Cloudflare
                    'account_id': account['cf_id'],        # Accout id на сайті = Cloudflare ID
                    'password': account['cf_password']     # Password на сайті = Cloudflare Password
                }
                
                # Оновлюємо прогрес у вікні
                try:
                    self.master.title(f"PunchITNow 9.0 Octopus - Upload Progress: {idx}/{total} ({uploaded} success)")
                except:
                    pass
                
                # Retry логіка (до 3 спроб)
                max_retries = 3
                retry_count = 0
                success = False
                
                while retry_count < max_retries and not success:
                    try:
                        response = requests.post(self.DATABASE_URL, data=data, timeout=10)
                        
                        # Парсимо JSON відповідь
                        try:
                            result = response.json()
                            status = result.get('status', '')
                            message = result.get('message', '')
                            
                            if status == 'success':
                                uploaded += 1
                                self.log_action(f"✅ [{uploaded}/{total}] Cloudflare: {email}")
                                self.log_to_db_panel(f"✅ [{uploaded}/{total}] {email}")
                                success = True
                            elif status == 'error':
                                if "існує" in message.lower():
                                    skipped += 1
                                    self.log_action(f"⏭️ [{idx}/{total}] Вже існує: {email}")
                                    self.log_to_db_panel(f"⏭️ Дублікат: {email}")
                                    success = True  # Не retry для дублікатів
                                else:
                                    errors += 1
                                    self.log_action(f"❌ [{idx}/{total}] Помилка: {message[:50]}")
                                    self.log_to_db_panel(f"❌ {message[:30]}: {email}")
                                    success = True  # Не retry для помилок валідації
                            else:
                                raise Exception(f"Невідомий статус: {status}")
                        
                        except json.JSONDecodeError:
                            raise Exception("Не-JSON відповідь від сервера")
                        
                        if response.status_code != 200:
                            raise Exception(f"HTTP {response.status_code}")
                            
                    except requests.exceptions.Timeout:
                        retry_count += 1
                        if retry_count < max_retries:
                            self.log_to_db_panel(f"⏱️ Timeout (retry {retry_count}/{max_retries}): {email}")
                            time.sleep(2)
                        else:
                            errors += 1
                            self.log_to_db_panel(f"⏱️ Timeout (final): {email}")
                    except requests.exceptions.ConnectionError:
                        retry_count += 1
                        if retry_count < max_retries:
                            self.log_to_db_panel(f"🔌 Connection Error (retry {retry_count}/{max_retries}): {email}")
                            time.sleep(2)
                        else:
                            errors += 1
                            self.log_to_db_panel(f"🔌 Connection Error (final): {email}")
                    except Exception as e:
                        retry_count += 1
                        if retry_count < max_retries:
                            self.log_to_db_panel(f"⚠️ Error (retry {retry_count}/{max_retries}): {str(e)[:30]}")
                            time.sleep(2)
                        else:
                            errors += 1
                            self.log_to_db_panel(f"❌ Error (final): {email}")
                
                # Затримка між запитами (0.5 секунди)
                if idx < total:
                    time.sleep(0.5)
            
            # Відновлюємо заголовок вікна
            try:
                self.master.title("PunchITNow 9.0 Octopus")
            except:
                pass
            
            # Підсумок
            if getattr(self, 'upload_cancelled', False):
                summary = (
                    f"⛔ Cloudflare Upload скасовано!\n\n"
                    f"Оброблено: {idx}/{total}\n"
                    f"✅ Завантажено: {uploaded}\n"
                    f"❌ Помилки: {errors}\n"
                    f"⏭️ Пропущено: {skipped}"
                )
                self.log_to_db_panel(f"⛔ Скасовано: {uploaded}/{idx} завантажено")
            else:
                summary = (
                    f"☁️ Cloudflare Upload завершено!\n\n"
                    f"Всього оброблено: {total}\n"
                    f"✅ Завантажено: {uploaded}\n"
                    f"❌ Помилки: {errors}\n"
                    f"⏭️ Пропущено: {skipped}"
                )
                self.log_to_db_panel(f"☁️ Завершено: {uploaded}/{total} завантажено")
            
            self.log_action(summary.replace('\n', ' | '))
            
        except Exception as e:
            self.log_action(f"❌ Критична помилка Cloudflare Upload: {e}")
            self.log_to_db_panel(f"❌ Критична помилка: {str(e)[:50]}")
            messagebox.showerror("Помилка", f"Критична помилка: {e}")

    def log_action(self, message):
        """Записує дію до логу"""
        try:
            timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_message = f"[{timestamp}] {message}\n"
            
            # Записуємо до файлу логу
            with open(self.LOG_FILE, 'a', encoding='utf-8') as f:
                f.write(log_message)
            
            # Виводимо до консолі
            print(log_message.strip())
            
        except Exception as e:
            print(f"Помилка при запису до логу: {e}")
    
    def log_to_db_panel(self, message):
        """Додає повідомлення в DB лог панель"""
        try:
            if hasattr(self, 'db_log_display'):
                self.db_log_display.configure(state="normal")
                timestamp = datetime.datetime.now().strftime('%H:%M:%S')
                self.db_log_display.insert("end", f"[{timestamp}] {message}\n")
                self.db_log_display.see("end")  # Прокручуємо до кінця
                self.db_log_display.configure(state="disabled")
        except Exception as e:
            print(f"Помилка логування в DB панель: {e}")

    def open_log_file(self):
        """Відкриває файл логу"""
        try:
            if os.name == 'nt':  # Windows
                os.startfile(self.LOG_FILE)
            else:
                import subprocess
                subprocess.call(['open', self.LOG_FILE] if sys.platform == 'darwin' else ['xdg-open', self.LOG_FILE])
        except Exception as e:
            messagebox.showerror("Помилка", f"Не вдалося відкрити лог: {e}")

    def open_csv_config(self):
        """Відкриває локальний CSV конфіг"""
        try:
            if os.name == 'nt':  # Windows
                os.startfile(self.LOCAL_CSV_CONFIG)
            else:
                import subprocess
                subprocess.call(['open', self.LOCAL_CSV_CONFIG] if sys.platform == 'darwin' else ['xdg-open', self.LOCAL_CSV_CONFIG])
        except Exception as e:
            messagebox.showerror("Помилка", f"Не вдалося відкрити CSV: {e}")


# ================== SMS Checker ==================
class SmsCheckerTab(ctk.CTkFrame):
    def __init__(self, master, font=None):
        super().__init__(master)
        self.font = font
        self.pack(fill="both", expand=True)
        self.root = master.winfo_toplevel()
        self.root.title("PunchITNow 9.0 Octopus - Global Key Binding")

        # API key for Daisysms
        self.API_KEY = self.load_api_key()  # Load from config or use default

        # Services
        self.SERVICES = {
            'go': 'Google (отримання коду підтвердження)',
            'ds': 'Discord (реєстрація нового акаунта)',
            'wa': 'WhatsApp (підтвердження номера)',
            'ig': 'Instagram (реєстрація/відновлення)',
            'fb': 'Facebook (підтвердження номера)',
            'tg': 'Telegram (реєстрація)',
            'am': 'Amazon (підтвердження покупок)',
            'tw': 'Twitter (підтвердження номера)'
        }

        # Variables for 4 numbers system
        self.numbers_data = {}  # Dictionary to store data for each number slot
        for i in range(1, 5):  # Numbers 1-4
            self.numbers_data[i] = {
                'activation': None,
                'number': None,
                'id': None,
                'sms_code_copied': False
            }
        
        self.current_service = "go"

        # Create UI
        self.create_widgets()
    
    def load_api_key(self):
        """Завантажує API ключ з конфігу або повертає дефолтний"""
        config_dir = os.path.join(os.path.expanduser("~"), ".punchnow")
        config_file = os.path.join(config_dir, "daisysms_config.json")
        
        try:
            if os.path.exists(config_file):
                with open(config_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    return config.get('api_key', 'qRptHWNir0haRqH5o3sVVe2XrOqtqi')
        except Exception as e:
            print(f"Error loading API key: {e}")
        
        return 'qRptHWNir0haRqH5o3sVVe2XrOqtqi'  # Default key
    
    def save_api_key_to_config(self, api_key):
        """Зберігає API ключ у конфіг файл"""
        config_dir = os.path.join(os.path.expanduser("~"), ".punchnow")
        config_file = os.path.join(config_dir, "daisysms_config.json")
        
        try:
            # Створюємо директорію якщо не існує
            os.makedirs(config_dir, exist_ok=True)
            
            # Зберігаємо конфіг
            config = {'api_key': api_key}
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(config, f, indent=4)
            
            return True
        except Exception as e:
            print(f"Error saving API key: {e}")
            return False
    
    def show_instruction(self):
        instruction_text = """
📱 SMS CHECKER - ІНСТРУКЦІЯ З ВИКОРИСТАННЯ

🔑 API НАЛАШТУВАННЯ:
• Введіть ваш API ключ від DaisySMS
• Натисніть "Save" для збереження ключа
• Перевірте баланс кнопкою "Check Balance"

🌐 ВИБІР СЕРВІСУ:
• Google - для отримання кодів підтвердження
• Discord, WhatsApp, Instagram, Facebook, Telegram
• Amazon, Twitter - підтримуються всі популярні сервіси
• Max Price - встановіть максимальну ціну за номер

📞 УПРАВЛІННЯ НОМЕРАМИ:
• Get Number - отримати новий номер
• Copy Number - скопіювати номер в буфер
• Cancel Number - скасувати активацію
• Номери автоматично активуються

📨 SMS УПРАВЛІННЯ:  
• Refresh SMS - оновити статус повідомлення
• Copy Code - скопіювати код з SMS
• Автоматичний моніторинг вхідних SMS
• Коди автоматично копіюються в буфер

� ЗВУКОВІ СПОВІЩЕННЯ:
• При отриманні SMS програма відтворює звук
• Кнопка 🧪 - тестування поточного звуку
• Кнопка 🎵 - вибір власного звукового файлу
• Підтримує: MP3, WAV, OGG формати
• Для роботи звуків встановіть: pip install pygame
• Статус: 🔊 ON (назва файлу) / 🔇 OFF

�💡 ПРИНЦИП РОБОТИ:
1. Налаштуйте API ключ
2. Оберіть сервіс та ціну
3. Отримайте номер
4. Використайте номер для реєстрації
5. Очікуйте SMS код (з звуковим сповіщенням!)
        """
        messagebox.showinfo("📱 SMS Checker - Інструкція", instruction_text)

        # Auto start SMS polling if ID exists
        if self.current_id:
            self.start_sms_polling()

    def create_widgets(self):
        main_frame = ctk.CTkFrame(self)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        # Заголовок і кнопка інструкції
        header_frame = ctk.CTkFrame(main_frame)
        header_frame.pack(fill=tk.X, padx=5, pady=5)
        
        ctk.CTkLabel(header_frame, text="DAISYSMS - Служба SMS активації", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(side=tk.LEFT, padx=10)
        
        # Звукові кнопки та індикатор
        sound_frame = ctk.CTkFrame(header_frame)
        sound_frame.pack(side=tk.RIGHT, padx=5)
        
        # Індикатор звукової системи
        sound_file = get_current_sound_file()
        sound_status = "🔊 ON" if SOUND_AVAILABLE and sound_file else "🔇 OFF"
        sound_color = "green" if "ON" in sound_status else "gray"
        
        sound_name = ""
        if sound_file:
            sound_name = f" ({os.path.basename(sound_file)})"
        
        ctk.CTkLabel(sound_frame, text=f"Sound: {sound_status}{sound_name}", 
                    font=self.font, text_color=sound_color).pack(side=tk.LEFT, padx=2)
        
        # Кнопка тестування звуку
        ctk.CTkButton(sound_frame, text="🧪", 
                     command=self.test_sound_alert, width=25, height=24, 
                     corner_radius=4, font=ctk.CTkFont(size=10)).pack(side=tk.LEFT, padx=1)
        
        # Кнопка вибору звуку
        ctk.CTkButton(sound_frame, text="🎵", 
                     command=self.choose_sound_file, width=25, height=24, 
                     corner_radius=4, font=ctk.CTkFont(size=10)).pack(side=tk.LEFT, padx=1)
        
        ctk.CTkButton(header_frame, text="Інструкція", 
                     command=self.show_instruction, width=100, height=28, 
                     corner_radius=6, font=self.font).pack(side=tk.RIGHT, padx=5)

        # API Key section
        api_frame = ctk.CTkFrame(main_frame)
        api_frame.pack(fill=tk.X, pady=5, padx=5)
        
        ctk.CTkLabel(api_frame, text="API Settings", font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(5,0))

        api_content = ctk.CTkFrame(api_frame)
        api_content.pack(fill=tk.X, padx=10, pady=10)
        
        ctk.CTkLabel(
    api_content,
    text="API Key:", font=self.font).grid(
        row=0,
        column=0,
         sticky=tk.W, padx=5)
        self.api_entry = ctk.CTkEntry(api_content, width=300, show="*", font=self.font)
        self.api_entry.grid(row=0, column=1, padx=5)
        self.api_entry.insert(0, self.API_KEY)

        save_btn = ctk.CTkButton(
    api_content,
    text="Save",
    command=self.save_api_key,
    width=80,
    height=32,
    corner_radius=6,
    font=self.font)
        save_btn.grid(row=0, column=2, padx=5)

        # Global Controls (Balance & Actions) - MOVED UP
        global_frame = ctk.CTkFrame(main_frame)
        global_frame.pack(fill=tk.X, pady=5, padx=5)
        
        ctk.CTkLabel(global_frame, text="Global Controls & Balance", font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(5,0))
        
        global_btn_frame = ctk.CTkFrame(global_frame)
        global_btn_frame.pack(padx=10, pady=10)
        
        self.check_balance_btn = ctk.CTkButton(
            global_btn_frame,
            text="💰 Check Balance",
            command=self.check_balance,
            width=130, height=32,
            corner_radius=6, font=self.font)
        self.check_balance_btn.pack(side=tk.LEFT, padx=5)
        
        self.get_all_numbers_btn = ctk.CTkButton(
            global_btn_frame,
            text="📞 Get All Numbers",
            command=self.get_all_numbers,
            width=130, height=32,
            corner_radius=6, font=self.font)
        self.get_all_numbers_btn.pack(side=tk.LEFT, padx=5)
        
        self.cancel_all_numbers_btn = ctk.CTkButton(
            global_btn_frame,
            text="❌ Cancel All",
            command=self.cancel_all_numbers,
            width=130, height=32,
            corner_radius=6, font=self.font)
        self.cancel_all_numbers_btn.pack(side=tk.LEFT, padx=5)
        
        # Balance display label
        self.balance_label = ctk.CTkLabel(
            global_frame,
            text="Balance: Not checked",
            font=ctk.CTkFont(size=13, weight="bold"))
        self.balance_label.pack(pady=5)

        # Service selection
        service_frame = ctk.CTkFrame(main_frame)
        service_frame.pack(fill=tk.X, pady=5, padx=5)
        
        ctk.CTkLabel(service_frame, text="Service Selection", font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(5,0))
        
        service_content = ctk.CTkFrame(service_frame)
        service_content.pack(fill=tk.X, padx=10, pady=10)

        ctk.CTkLabel(
    service_content,
    text="Service:", font=self.font).grid(
        row=0,
        column=0,
         sticky=tk.W, padx=5)
        
        self.service_btn = ctk.CTkButton(
    service_content,
    text="Select Service...",
    command=self.show_service_popup,
    width=150, height=28,
     corner_radius=6, font=self.font)
        self.service_btn.grid(row=0, column=1, padx=5, sticky=tk.W)

        self.service_label = ctk.CTkLabel(
    service_content, text="Google (отримання коду підтвердження)", font=self.font)
        self.service_label.grid(row=0, column=2, padx=5, sticky=tk.W)

        ctk.CTkLabel(
    service_content,
    text="Max Price ($):", font=self.font).grid(
        row=1,
        column=0,
        padx=5,
             sticky=tk.W)
        self.max_price_entry = ctk.CTkEntry(service_content, width=100, font=self.font)
        self.max_price_entry.grid(row=1, column=1, padx=5, sticky=tk.W)
        self.max_price_entry.insert(0, "0.45")

        # Google search field
        self.google_frame = ctk.CTkFrame(service_content)
        self.google_frame.grid(
    row=2,
    column=0,
    columnspan=4,
    sticky=tk.W,
    pady=5, padx=5)
        ctk.CTkLabel(
    self.google_frame,
    text="Google Search Query:", font=self.font).pack(
        side=tk.LEFT, padx=5)
        self.google_search_entry = ctk.CTkEntry(self.google_frame, width=200, placeholder_text="Enter search query", font=self.font)
        self.google_search_entry.pack(side=tk.LEFT, padx=5)

        # Hide Google fields by default
        self.google_frame.grid_remove()

        # 4 Numbers Management in 2x2 Grid
        numbers_frame = ctk.CTkFrame(main_frame)
        numbers_frame.pack(fill=tk.BOTH, expand=True, pady=5, padx=5)
        
        ctk.CTkLabel(numbers_frame, text="Multi Numbers Management (4 Slots)", font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(5,0))

        # Create grid container for 2x2 layout
        grid_container = ctk.CTkFrame(numbers_frame)
        grid_container.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # Configure grid weights for equal distribution
        grid_container.grid_columnconfigure(0, weight=1)
        grid_container.grid_columnconfigure(1, weight=1)
        grid_container.grid_rowconfigure(0, weight=1)
        grid_container.grid_rowconfigure(1, weight=1)

        # Initialize storage for number widgets
        self.number_widgets = {}

        # Create 4 number slots in 2x2 grid
        slot_positions = [
            (1, 0, 0),  # Slot 1: row 0, col 0
            (2, 0, 1),  # Slot 2: row 0, col 1
            (3, 1, 0),  # Slot 3: row 1, col 0
            (4, 1, 1),  # Slot 4: row 1, col 1
        ]
        
        for slot_num, row, col in slot_positions:
            self.create_number_slot_grid(grid_container, slot_num, row, col)

        # Global SMS management
        global_sms_frame = ctk.CTkFrame(main_frame)
        global_sms_frame.pack(fill=tk.X, pady=5, padx=5)
        
        ctk.CTkLabel(global_sms_frame, text="Global SMS Management", font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(5,0))

        global_sms_btn_frame = ctk.CTkFrame(global_sms_frame)
        global_sms_btn_frame.pack(padx=10, pady=10)

        self.clear_log_btn = ctk.CTkButton(
            global_sms_btn_frame,
            text="Clear All Logs",
            command=self.clear_log,
            width=120, height=32,
            corner_radius=6, font=self.font)
        self.clear_log_btn.pack(side=tk.LEFT, padx=5)

        # Status bar
        self.status_var = tk.StringVar()
        self.status_var.set("Ready")
        self.status_label = ctk.CTkLabel(
    main_frame,
    textvariable=self.status_var,
    font=self.font)
        self.status_label.pack(fill=tk.X, pady=5)
        
        # Auto-check balance on startup (in background)
        def auto_check_balance():
            time.sleep(1)  # Wait for UI to fully load
            try:
                self.check_balance()
            except:
                pass  # Ignore errors on startup
        
        Thread(target=auto_check_balance, daemon=True).start()

    def toggle_google_fields(self):
        if self.current_service == "go":
            self.google_frame.grid()
        else:
            self.google_frame.grid_remove()

    def show_service_popup(self):
        popup = tk.Toplevel(self.root)
        popup.title("Select Service")
        popup.geometry("400x300")

        frame = ttk.Frame(popup, padding="10")
        frame.pack(fill=tk.BOTH, expand=True)

        ttk.Label(frame, text="Select service:").pack(pady=5)

        # Service list
        for code, desc in self.SERVICES.items():
            if HAS_CTK:
                btn = ctk.CTkButton(
                    frame,
                    text=f"{code.upper()} - {desc}",
                    command=lambda c=code, d=desc: self.set_service(
                        c, d, popup),
                    corner_radius=10
                )
            else:
                btn = ttk.Button(
                    frame,
                    text=f"{code.upper()} - {desc}",
                    command=lambda c=code, d=desc: self.set_service(
                        c, d, popup)
                )
            btn.pack(fill=tk.X, pady=2)

    def set_service(self, code, desc, popup):
        self.current_service = code
        self.service_label.configure(text=desc)
        self.toggle_google_fields()
        popup.destroy()

    def create_number_slot(self, parent, slot_number):
        """Створює слот для одного номера"""
        # Frame for this number slot
        slot_frame = ctk.CTkFrame(parent, corner_radius=10)
        slot_frame.pack(fill=tk.X, pady=10, padx=5)
        
        # Header with slot number
        header_frame = ctk.CTkFrame(slot_frame)
        header_frame.pack(fill=tk.X, padx=10, pady=(10, 5))
        
        ctk.CTkLabel(header_frame, text=f"📱 Slot {slot_number}", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(side=tk.LEFT)
        
        # Status indicator
        status_label = ctk.CTkLabel(header_frame, text="🔴 Empty", 
                                   font=ctk.CTkFont(size=12))
        status_label.pack(side=tk.RIGHT)
        
        # Number display
        number_display = ctk.CTkLabel(slot_frame, text="No number assigned", 
                                     font=ctk.CTkFont(size=12, weight="bold"))
        number_display.pack(pady=5)
        
        # Buttons frame
        btn_frame = ctk.CTkFrame(slot_frame)
        btn_frame.pack(pady=10, padx=10)
        
        # Get Number button
        get_btn = ctk.CTkButton(btn_frame, text="Get Number", 
                               command=lambda: self.get_number_for_slot(slot_number),
                               width=100, height=28, corner_radius=6, font=self.font)
        get_btn.pack(side=tk.LEFT, padx=2)
        
        # Copy Number button
        copy_num_btn = ctk.CTkButton(btn_frame, text="Copy Number", 
                                command=lambda: self.copy_number_for_slot(slot_number),
                                width=100, height=28, corner_radius=6, font=self.font,
                                state=tk.DISABLED)
        copy_num_btn.pack(side=tk.LEFT, padx=2)
        
        # Copy Code button (NEW!)
        copy_code_btn = ctk.CTkButton(btn_frame, text="Copy Code", 
                                  command=lambda: self.copy_code_for_slot(slot_number),
                                  width=90, height=28, corner_radius=6, font=self.font,
                                  state=tk.DISABLED)
        copy_code_btn.pack(side=tk.LEFT, padx=2)
        
        # Cancel Number button
        cancel_btn = ctk.CTkButton(btn_frame, text="Cancel", 
                                  command=lambda: self.cancel_number_for_slot(slot_number),
                                  width=80, height=28, corner_radius=6, font=self.font,
                                  state=tk.DISABLED)
        cancel_btn.pack(side=tk.LEFT, padx=2)
        
        # SMS text area for this slot
        sms_text = ctk.CTkTextbox(slot_frame, height=80, wrap="word", 
                                 state="disabled", font=self.font)
        sms_text.pack(fill=tk.X, padx=10, pady=(5, 10))
        
        # Store widget references
        self.number_widgets[slot_number] = {
            'frame': slot_frame,
            'status_label': status_label,
            'number_display': number_display,
            'get_btn': get_btn,
            'copy_num_btn': copy_num_btn,
            'copy_code_btn': copy_code_btn,
            'cancel_btn': cancel_btn,
            'sms_text': sms_text
        }

    def create_number_slot_grid(self, parent, slot_number, row, col):
        """Створює слот для одного номера в grid layout"""
        # Frame for this number slot
        slot_frame = ctk.CTkFrame(parent, corner_radius=10)
        slot_frame.grid(row=row, column=col, padx=5, pady=5, sticky="nsew")
        
        # Header with slot number
        header_frame = ctk.CTkFrame(slot_frame)
        header_frame.pack(fill=tk.X, padx=10, pady=(10, 5))
        
        ctk.CTkLabel(header_frame, text=f"📱 Slot {slot_number}", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(side=tk.LEFT)
        
        # Status indicator
        status_label = ctk.CTkLabel(header_frame, text="🔴 Empty", 
                                   font=ctk.CTkFont(size=12))
        status_label.pack(side=tk.RIGHT)
        
        # Number display
        number_display = ctk.CTkLabel(slot_frame, text="No number", 
                                     font=ctk.CTkFont(size=11, weight="bold"))
        number_display.pack(pady=5)
        
        # Buttons frame (2 rows)
        btn_frame_top = ctk.CTkFrame(slot_frame)
        btn_frame_top.pack(pady=5, padx=10)
        
        btn_frame_bottom = ctk.CTkFrame(slot_frame)
        btn_frame_bottom.pack(pady=2, padx=10)
        
        # Top row buttons
        get_btn = ctk.CTkButton(btn_frame_top, text="Get Number", 
                               command=lambda: self.get_number_for_slot(slot_number),
                               width=95, height=26, corner_radius=6, 
                               font=ctk.CTkFont(size=11))
        get_btn.pack(side=tk.LEFT, padx=2)
        
        copy_num_btn = ctk.CTkButton(btn_frame_top, text="Copy Num", 
                                command=lambda: self.copy_number_for_slot(slot_number),
                                width=95, height=26, corner_radius=6,
                                font=ctk.CTkFont(size=11),
                                state=tk.DISABLED)
        copy_num_btn.pack(side=tk.LEFT, padx=2)
        
        # Bottom row buttons
        copy_code_btn = ctk.CTkButton(btn_frame_bottom, text="Copy Code", 
                                  command=lambda: self.copy_code_for_slot(slot_number),
                                  width=95, height=26, corner_radius=6,
                                  font=ctk.CTkFont(size=11),
                                  state=tk.DISABLED)
        copy_code_btn.pack(side=tk.LEFT, padx=2)
        
        cancel_btn = ctk.CTkButton(btn_frame_bottom, text="Cancel", 
                                  command=lambda: self.cancel_number_for_slot(slot_number),
                                  width=95, height=26, corner_radius=6,
                                  font=ctk.CTkFont(size=11),
                                  state=tk.DISABLED)
        cancel_btn.pack(side=tk.LEFT, padx=2)
        
        # SMS text area for this slot (smaller height for grid)
        sms_text = ctk.CTkTextbox(slot_frame, height=100, wrap="word", 
                                 state="disabled", font=ctk.CTkFont(size=10))
        sms_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=(5, 10))
        
        # Store widget references
        self.number_widgets[slot_number] = {
            'frame': slot_frame,
            'status_label': status_label,
            'number_display': number_display,
            'get_btn': get_btn,
            'copy_num_btn': copy_num_btn,
            'copy_code_btn': copy_code_btn,
            'cancel_btn': cancel_btn,
            'sms_text': sms_text
        }

    def save_api_key(self):
        self.API_KEY = self.api_entry.get()
        if self.save_api_key_to_config(self.API_KEY):
            self.status_var.set("API Key saved to config")
        else:
            self.status_var.set("API Key saved (config save failed)")

    def get_number_for_slot(self, slot_number):
        """Отримує номер для конкретного слота"""
        import requests
        
        if not self.API_KEY:
            messagebox.showerror("Error", "Please enter API key")
            return

        max_price = self.max_price_entry.get()

        try:
            float(max_price)
        except ValueError:
            messagebox.showerror("Error", "Invalid max price")
            return

        google_search = ""
        if self.current_service == "go":
            google_search = self.google_search_entry.get().strip()
            if google_search:
                google_search = f"&google_search={google_search}"

        url = f"https://daisysms.com/stubs/handler_api.php?api_key={self.API_KEY}&action=getNumber&service={self.current_service}&max_price={max_price}{google_search}"

        self.status_var.set(f"Requesting number for Slot {slot_number}...")
        self.root.update()

        try:
            response = requests.get(url)
            widgets = self.number_widgets[slot_number]
            
            if response.status_code == 200 or response.status_code == 400:
                content = response.text.strip()
                if content.startswith("ACCESS_NUMBER"):
                    parts = content.split(":")
                    number_id = parts[1]
                    number = parts[2]
                    
                    # Save to slot data
                    self.numbers_data[slot_number] = {
                        'activation': {"id": number_id, "number": number, "service": self.current_service},
                        'number': number,
                        'id': number_id,
                        'sms_code_copied': False
                    }

                    # Update UI
                    widgets['status_label'].configure(text="🟢 Active")
                    widgets['number_display'].configure(text=f"Number: +{number}")
                    widgets['get_btn'].configure(state=tk.DISABLED)
                    widgets['copy_num_btn'].configure(state=tk.NORMAL)
                    widgets['cancel_btn'].configure(state=tk.NORMAL)

                    self.start_sms_polling_for_slot(slot_number)
                    self.update_status_statistics()
                    
                else:
                    # Handle errors in slot's SMS text area
                    error_message = self.format_error_message(content, max_price)
                    widgets['sms_text'].configure(state="normal")
                    widgets['sms_text'].delete("1.0", "end")
                    widgets['sms_text'].insert("1.0", error_message)
                    widgets['sms_text'].configure(state="disabled")
                    self.status_var.set(f"Error for Slot {slot_number}: {content}")
            else:
                widgets['sms_text'].configure(state="normal")
                widgets['sms_text'].delete("1.0", "end")
                widgets['sms_text'].insert("1.0", f"🌐 Connection Error\n\nHTTP {response.status_code}\nCheck internet connection")
                widgets['sms_text'].configure(state="disabled")
                self.status_var.set(f"HTTP error for Slot {slot_number}: {response.status_code}")
                
        except Exception as e:
            widgets = self.number_widgets[slot_number]
            widgets['sms_text'].configure(state="normal")
            widgets['sms_text'].delete("1.0", "end")
            widgets['sms_text'].insert("1.0", f"💥 Execution Error\n\n{str(e)}\nCheck network settings")
            widgets['sms_text'].configure(state="disabled")
            self.status_var.set(f"Error for Slot {slot_number}: {str(e)}")

    def format_error_message(self, content, max_price):
        """Форматує повідомлення про помилку"""
        if "MAX_PRICE_EXCEEDED" in content:
            return f"⚠️ Price Exceeded\n\nMax price: ${max_price}\nTry increasing the price."
        elif "NO_NUMBERS" in content or content == "NO_NUMBERS":
            return f"❌ No Numbers Available\n\nService: {self.SERVICES.get(self.current_service, self.current_service)}\nMax price: ${max_price}\n\nSuggestions:\n• Increase max price\n• Try different service\n• Retry later"
        elif "TOO_MANY_ACTIVE_RENTALS" in content:
            return f"⚠️ Too Many Active Rentals\n\nFinish current activations before creating new ones."
        elif "NO_MONEY" in content:
            return f"� Insufficient Funds\n\nTop up your DaisySMS balance"
        elif "BAD_KEY" in content:
            return f"🔑 Invalid API Key\n\nCheck your API key"
        else:
            return f"❓ Unknown Server Response\n\n{content}\n\nContact DaisySMS support"

    def copy_number_for_slot(self, slot_number):
        """Копіює номер для конкретного слота"""
        if slot_number in self.numbers_data and self.numbers_data[slot_number]['number']:
            safe_clipboard_operation("set", self.numbers_data[slot_number]['number'])
            self.status_var.set(f"Number copied from Slot {slot_number}")

    def copy_code_for_slot(self, slot_number):
        """Копіює код з SMS для конкретного слота"""
        if slot_number not in self.number_widgets:
            return
            
        widgets = self.number_widgets[slot_number]
        sms_content = widgets['sms_text'].get("1.0", tk.END).strip()
        
        if sms_content and "Code:" in sms_content:
            # Витягуємо код з тексту
            lines = sms_content.split('\n')
            for line in lines:
                if "Code:" in line:
                    code = line.split("Code:")[1].strip()
                    safe_clipboard_operation("set", code)
                    self.status_var.set(f"Code copied from Slot {slot_number}: {code}")
                    return
        
        self.status_var.set(f"No code available in Slot {slot_number}")

    def cancel_number_for_slot(self, slot_number):
        """Скасовує номер для конкретного слота"""
        import requests
        
        if slot_number not in self.numbers_data or not self.numbers_data[slot_number]['id']:
            return

        number_id = self.numbers_data[slot_number]['id']
        url = f"https://daisysms.com/stubs/handler_api.php?api_key={self.API_KEY}&action=setStatus&id={number_id}&status=8"

        try:
            response = requests.get(url)
            if response.status_code == 200:
                if response.text == "ACCESS_CANCEL":
                    self.clear_slot_activation(slot_number)
                    self.status_var.set(f"Slot {slot_number} number cancelled")
                else:
                    self.status_var.set(f"Cancel error for Slot {slot_number}: {response.text}")
            else:
                self.status_var.set(f"HTTP error for Slot {slot_number}: {response.status_code}")
        except Exception as e:
            self.status_var.set(f"Error for Slot {slot_number}: {str(e)}")

    def clear_slot_activation(self, slot_number):
        """Очищає активацію для конкретного слота"""
        # Reset data
        self.numbers_data[slot_number] = {
            'activation': None,
            'number': None,
            'id': None,
            'sms_code_copied': False
        }

        # Update UI
        widgets = self.number_widgets[slot_number]
        widgets['status_label'].configure(text="🔴 Empty")
        widgets['number_display'].configure(text="No number assigned")
        widgets['sms_text'].configure(state="normal")
        widgets['sms_text'].delete("1.0", "end")
        widgets['sms_text'].configure(state="disabled")
        widgets['get_btn'].configure(state=tk.NORMAL)
        widgets['copy_num_btn'].configure(state=tk.DISABLED)
        widgets['copy_code_btn'].configure(state=tk.DISABLED)
        widgets['cancel_btn'].configure(state=tk.DISABLED)
        
        # Update statistics
        self.update_status_statistics()

    def update_status_statistics(self):
        """Оновлює статистику активних номерів у статус бар"""
        active_count = 0
        waiting_count = 0
        received_count = 0
        
        for slot_number in range(1, 5):
            if (slot_number in self.numbers_data and 
                self.numbers_data[slot_number]['id']):
                active_count += 1
                if self.numbers_data[slot_number]['sms_code_copied']:
                    received_count += 1
                else:
                    waiting_count += 1
        
        if active_count == 0:
            status = "Ready - No active numbers"
        else:
            status = f"Active: {active_count} | Waiting: {waiting_count} | Received: {received_count}"
        
        # Зберігаємо поточний статус якщо він не є статистикою
        current_status = self.status_var.get()
        if not ("Active:" in current_status or "Ready -" in current_status):
            # Якщо поточний статус - це не статистика, додаємо статистику
            status = f"{current_status} | {status}"
        
        self.status_var.set(status)

    def check_balance(self):
        import requests  # Важливо: імпортуємо requests локально
        
        if not self.API_KEY:
            messagebox.showerror("Error", "Please enter API key")
            return

        url = f"https://daisysms.com/stubs/handler_api.php?api_key={
    self.API_KEY}&action=getBalance"

        try:
            response = requests.get(url)
            if response.status_code == 200:
                if response.text.startswith("ACCESS_BALANCE"):
                    balance = response.text.split(":")[1]
                    self.balance_label.configure(text=f"💰 Balance: ${balance}")
                    self.status_var.set(f"Balance updated: ${balance}")
                elif response.text == "BAD_KEY":
                    self.balance_label.configure(text="❌ Bad API key")
                    self.status_var.set("Bad API key")
                else:
                    self.balance_label.configure(text=f"⚠️ Error: {response.text}")
                    self.status_var.set(f"Error: {response.text}")
            else:
                self.balance_label.configure(text=f"🌐 HTTP error: {response.status_code}")
                self.status_var.set(f"HTTP error: {response.status_code}")
        except Exception as e:
            self.balance_label.configure(text=f"💥 Error checking balance")
            self.status_var.set(f"Error: {str(e)}")

    def start_sms_polling_for_slot(self, slot_number):
        """Запускає моніторинг SMS для конкретного слота"""
        if slot_number not in self.numbers_data or not self.numbers_data[slot_number]['id']:
            return

        def polling_thread():
            while (slot_number in self.numbers_data and 
                   self.numbers_data[slot_number]['id']):
                self.poll_sms_for_slot(slot_number)
                time.sleep(3)

        Thread(target=polling_thread, daemon=True).start()

    def poll_sms_for_slot(self, slot_number):
        """Перевіряє SMS для конкретного слота"""
        import requests
        
        if slot_number not in self.numbers_data or not self.numbers_data[slot_number]['id']:
            return

        number_id = self.numbers_data[slot_number]['id']
        url = f"https://daisysms.com/stubs/handler_api.php?api_key={self.API_KEY}&action=getStatus&id={number_id}&text=1"

        try:
            response = requests.get(url)
            widgets = self.number_widgets[slot_number]
            
            if response.status_code == 200:
                if response.text.startswith("STATUS_OK"):
                    code = response.text.split(":")[1]
                    text = response.headers.get("X-Text", "Message text unavailable")

                    # Update slot's SMS text area
                    widgets['sms_text'].configure(state="normal")
                    widgets['sms_text'].delete("1.0", "end")
                    widgets['sms_text'].insert("1.0", f"✅ SMS Received!\n\nCode: {code}\n\nMessage:\n{text}")
                    widgets['sms_text'].configure(state="disabled")
                    
                    # Enable Copy Code button
                    widgets['copy_code_btn'].configure(state=tk.NORMAL)

                    if not self.numbers_data[slot_number]['sms_code_copied']:
                        safe_clipboard_operation("set", code)
                        self.numbers_data[slot_number]['sms_code_copied'] = True
                        # Play sound notification
                        play_alert_sound()
                        # Update statistics
                        self.update_status_statistics()
                        
                elif response.text == "STATUS_WAIT_CODE":
                    # Show waiting status in slot
                    widgets['sms_text'].configure(state="normal")
                    current_text = widgets['sms_text'].get("1.0", "end").strip()
                    if not current_text or "Waiting for SMS" not in current_text:
                        widgets['sms_text'].delete("1.0", "end")
                        widgets['sms_text'].insert("1.0", f"⏳ Waiting for SMS...\n\nSlot {slot_number} is active and waiting for verification code.")
                    widgets['sms_text'].configure(state="disabled")
                    
                elif response.text == "STATUS_CANCEL":
                    self.clear_slot_activation(slot_number)
                    self.status_var.set(f"Slot {slot_number} activation cancelled")
                    
                elif response.text == "NO_ACTIVATION":
                    self.clear_slot_activation(slot_number)
                    self.status_var.set(f"Slot {slot_number} activation not found")
                    
                else:
                    widgets['sms_text'].configure(state="normal")
                    widgets['sms_text'].delete("1.0", "end")
                    widgets['sms_text'].insert("1.0", f"❓ Unknown Status\n\n{response.text}")
                    widgets['sms_text'].configure(state="disabled")
                    
        except Exception as e:
            widgets = self.number_widgets[slot_number]
            widgets['sms_text'].configure(state="normal")
            widgets['sms_text'].delete("1.0", "end")
            widgets['sms_text'].insert("1.0", f"💥 Poll Error\n\n{str(e)}")
            widgets['sms_text'].configure(state="disabled")

    def clear_log(self):
        """Очищує логи у всіх слотах"""
        for slot_number in range(1, 5):
            if slot_number in self.number_widgets:
                widgets = self.number_widgets[slot_number]
                widgets['sms_text'].configure(state="normal")
                widgets['sms_text'].delete("1.0", "end")
                widgets['sms_text'].configure(state="disabled")
        self.status_var.set("All logs cleared")

    def get_all_numbers(self):
        """Отримує номери для всіх 4 слотів одночасно"""
        import threading
        
        def get_numbers_threaded():
            for slot_number in range(1, 5):
                # Запускаємо в окремих потоках з затримкою
                def get_for_slot(slot):
                    time.sleep(slot * 0.5)  # Затримка для уникнення перевантаження API
                    self.get_number_for_slot(slot)
                
                thread = threading.Thread(target=lambda s=slot_number: get_for_slot(s), daemon=True)
                thread.start()
        
        # Запускаємо всі запити в фоновому режимі
        main_thread = threading.Thread(target=get_numbers_threaded, daemon=True)
        main_thread.start()
        
        self.status_var.set("Getting numbers for all slots...")

    def cancel_all_numbers(self):
        """Скасовує всі активні номери"""
        import threading
        
        def cancel_numbers_threaded():
            cancelled_count = 0
            for slot_number in range(1, 5):
                if (slot_number in self.numbers_data and 
                    self.numbers_data[slot_number]['id']):
                    self.cancel_number_for_slot(slot_number)
                    cancelled_count += 1
                    time.sleep(0.2)  # Коротка затримка між скасуваннями
            
            if cancelled_count > 0:
                self.status_var.set(f"Cancelled {cancelled_count} active numbers")
            else:
                self.status_var.set("No active numbers to cancel")
        
        thread = threading.Thread(target=cancel_numbers_threaded, daemon=True)
        thread.start()

    def test_sound_alert(self):
        """Тестує звукове сповіщення SMS"""
        try:
            success = test_sound()
            if success:
                self.status_var.set("🔊 Sound test: OK")
            else:
                self.status_var.set("🔇 Sound test: Failed")
        except Exception as e:
            self.status_var.set(f"Sound test error: {e}")
    
    def choose_sound_file(self):
        """Дозволяє вибрати власний звуковий файл"""
        try:
            success = choose_custom_sound()
            if success:
                self.status_var.set("🎵 Custom sound installed!")
                # Оновлюємо інтерфейс
                self.update_sound_indicator()
            else:
                self.status_var.set("Sound selection cancelled")
        except Exception as e:
            self.status_var.set(f"Sound selection error: {e}")
    
    def update_sound_indicator(self):
        """Оновлює індикатор звуку в інтерфейсі"""
        try:
            # Потрібно перезавантажити вкладку для оновлення індикатора
            # Або можна динамічно оновити, але це складніше
            pass
        except Exception as e:
            print(f"Sound indicator update error: {e}")


# ================== File Generator Tab ==================
class FileGeneratorTab:
    """File Generator - Generate various file formats with random meaningful names"""
    
    def __init__(self, master, font=None):
        self.master = master
        self.font = font or ctk.CTkFont()
        self.tab_frame = None
        self.output_folder = None
        self.current_generation_folder = None  # Поточна папка з файлами
        self.generated_files = []
        self.is_generating = False
        
        # Словники для генерації назв
        self.adjectives = [
            "Annual", "Monthly", "Weekly", "Daily", "Quarterly", "Final", "Draft", "Updated",
            "Revised", "Complete", "Summary", "Detailed", "Brief", "Extended", "Special",
            "Important", "Urgent", "Confidential", "Public", "Internal", "External", "Personal",
            "Professional", "Business", "Financial", "Technical", "Creative", "Strategic"
        ]
        
        self.nouns = [
            "Report", "Document", "Analysis", "Presentation", "Proposal", "Invoice", "Contract",
            "Agreement", "Statement", "Summary", "Review", "Plan", "Budget", "Forecast",
            "Schedule", "Agenda", "Minutes", "Memo", "Letter", "Certificate", "Form",
            "Application", "Request", "Approval", "Record", "File", "Data", "Info"
        ]
        
        self.years = ["2023", "2024", "2025"]
        self.months = ["January", "February", "March", "April", "May", "June", 
                      "July", "August", "September", "October", "November", "December"]
        
        # Словники для генерації англійського тексту
        self.words = [
            "the", "be", "to", "of", "and", "a", "in", "that", "have", "I",
            "it", "for", "not", "on", "with", "he", "as", "you", "do", "at",
            "this", "but", "his", "by", "from", "they", "we", "say", "her", "she",
            "or", "an", "will", "my", "one", "all", "would", "there", "their",
            "business", "company", "project", "management", "development", "report",
            "analysis", "strategy", "planning", "implementation", "evaluation",
            "process", "system", "technology", "innovation", "solution", "service",
            "quality", "performance", "efficiency", "productivity", "growth",
            "revenue", "profit", "market", "customer", "client", "partner",
            "team", "employee", "organization", "department", "operation",
            "budget", "financial", "investment", "resource", "asset", "value"
        ]
        
        self.sentences_templates = [
            "The {noun} was successfully completed in {month} {year}.",
            "Our team has achieved significant progress in {noun} development.",
            "This {noun} demonstrates high quality and efficiency standards.",
            "The analysis shows positive trends in {noun} performance.",
            "Strategic planning for {noun} is essential for business growth.",
            "Implementation of new {noun} procedures will improve productivity.",
            "The {noun} evaluation report indicates successful outcomes.",
            "Management decided to expand the {noun} operations.",
            "Our company prioritizes {noun} innovation and development.",
            "The {noun} strategy aligns with long-term business objectives.",
        ]
    
    def generate_random_text(self, min_paragraphs=3, max_paragraphs=8):
        """Generate random English text"""
        import random
        
        num_paragraphs = random.randint(min_paragraphs, max_paragraphs)
        paragraphs = []
        
        for _ in range(num_paragraphs):
            num_sentences = random.randint(3, 6)
            sentences = []
            
            for _ in range(num_sentences):
                # Вибираємо шаблон або генеруємо випадкове речення
                if random.choice([True, False]):
                    template = random.choice(self.sentences_templates)
                    sentence = template.format(
                        noun=random.choice(self.nouns),
                        month=random.choice(self.months),
                        year=random.choice(self.years)
                    )
                else:
                    # Генеруємо випадкове речення
                    length = random.randint(8, 15)
                    words = [random.choice(self.words) for _ in range(length)]
                    words[0] = words[0].capitalize()
                    sentence = " ".join(words) + "."
                
                sentences.append(sentence)
            
            paragraph = " ".join(sentences)
            paragraphs.append(paragraph)
        
        return "\n\n".join(paragraphs)
    
    def create_tab(self, tabview, title):
        """Create tab in the tabview"""
        self.tab_frame = tabview.add(title)
        self.create_widgets()
    
    def generate_unique_filename(self, extension):
        """Generate unique meaningful filename"""
        import random
        import datetime
        
        while True:
            # Різні патерни для назв
            pattern = random.choice([1, 2, 3, 4, 5])
            
            if pattern == 1:
                # Pattern: Adjective_Noun_Year
                adj = random.choice(self.adjectives)
                noun = random.choice(self.nouns)
                year = random.choice(self.years)
                filename = f"{adj}_{noun}_{year}{extension}"
            
            elif pattern == 2:
                # Pattern: Noun_Month_Year
                noun = random.choice(self.nouns)
                month = random.choice(self.months)
                year = random.choice(self.years)
                filename = f"{noun}_{month}_{year}{extension}"
            
            elif pattern == 3:
                # Pattern: Adjective_Noun_Date
                adj = random.choice(self.adjectives)
                noun = random.choice(self.nouns)
                day = random.randint(1, 28)
                month = random.choice(self.months)
                filename = f"{adj}_{noun}_{month}_{day}{extension}"
            
            elif pattern == 4:
                # Pattern: Project_Noun_Version
                adj = random.choice(self.adjectives)
                noun = random.choice(self.nouns)
                version = f"v{random.randint(1,5)}.{random.randint(0,9)}"
                filename = f"{adj}_{noun}_{version}{extension}"
            
            else:
                # Pattern: Noun_Number
                noun = random.choice(self.nouns)
                num = random.randint(1000, 9999)
                filename = f"{noun}_{num}{extension}"
            
            # Перевіряємо унікальність
            if filename not in self.generated_files:
                self.generated_files.append(filename)
                return filename
    
    def create_widgets(self):
        """Create File Generator interface"""
        if not self.tab_frame:
            print("⚠️ Tab frame not initialized")
            return
            
        main_frame = ctk.CTkScrollableFrame(self.tab_frame)
        main_frame.pack(fill="both", expand=True, padx=10, pady=10)
        
        # ==================== Header ====================
        header_frame = ctk.CTkFrame(main_frame)
        header_frame.pack(fill="x", pady=(10, 20))
        
        ctk.CTkLabel(header_frame, text="📁 File Generator", 
                    font=ctk.CTkFont(size=20, weight="bold")).pack(side="left", padx=10, pady=10)
        
        ctk.CTkButton(header_frame, text="❓ Інструкція", 
                     command=self.show_instruction,
                     width=120, height=32, corner_radius=6,
                     fg_color="#2B5278", hover_color="#1e3a52").pack(side="right", padx=10)
        
        # ==================== File Types Selection ====================
        types_frame = ctk.CTkFrame(main_frame)
        types_frame.pack(fill="x", pady=10, padx=5)
        
        ctk.CTkLabel(types_frame, text="📂 Типи файлів для генерації:", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        
        # Checkboxes для типів файлів
        self.file_types_frame = ctk.CTkFrame(types_frame)
        self.file_types_frame.pack(fill="x", padx=10, pady=10)
        
        self.file_type_vars = {}
        file_types = [
            ("📄 DOCX (Word Documents)", ".docx"),
            ("📊 XLSX (Excel Spreadsheets)", ".xlsx"),
            ("📽️ PPTX (PowerPoint)", ".pptx"),
            ("🖼️ PNG (Images)", ".png"),
            ("🎵 MP3 (Audio)", ".mp3"),
            ("🎬 MP4 (Video)", ".mp4"),
            ("📝 TXT (Text Files)", ".txt"),
            ("📋 PDF (Documents)", ".pdf")
        ]
        
        for i, (label, ext) in enumerate(file_types):
            var = ctk.BooleanVar(value=True)
            self.file_type_vars[ext] = var
            row = i // 2
            col = i % 2
            ctk.CTkCheckBox(self.file_types_frame, text=label, variable=var,
                           font=self.font).grid(row=row, column=col, sticky="w", padx=20, pady=5)
        
        # ==================== Settings ====================
        settings_frame = ctk.CTkFrame(main_frame)
        settings_frame.pack(fill="x", pady=10, padx=5)
        
        ctk.CTkLabel(settings_frame, text="⚙️ Налаштування:", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        
        # Кількість файлів
        count_frame = ctk.CTkFrame(settings_frame)
        count_frame.pack(fill="x", padx=10, pady=5)
        
        ctk.CTkLabel(count_frame, text="Кількість файлів кожного типу:", 
                    font=self.font).pack(side="left", padx=10)
        
        self.count_entry = ctk.CTkEntry(count_frame, width=100, font=self.font,
                                        placeholder_text="10")
        self.count_entry.pack(side="left", padx=5)
        self.count_entry.insert(0, "10")
        
        # Інформація про папку (автоматично створюється)
        folder_frame = ctk.CTkFrame(settings_frame)
        folder_frame.pack(fill="x", padx=10, pady=5)
        
        info_text = "📁 Файли будуть збережені в папці:\n'generated_files' (в директорії програми)"
        ctk.CTkLabel(folder_frame, text=info_text, 
                    font=self.font, text_color="gray", justify="left").pack(padx=10, pady=5)
        
        # ==================== Progress ====================
        self.progress_frame = ctk.CTkFrame(main_frame)
        self.progress_frame.pack(fill="x", pady=10, padx=5)
        
        self.progress_label = ctk.CTkLabel(self.progress_frame, text="", 
                                          font=self.font)
        self.progress_label.pack(pady=5)
        
        self.progress_bar = ctk.CTkProgressBar(self.progress_frame)
        self.progress_bar.pack(fill="x", padx=20, pady=5)
        self.progress_bar.set(0)
        
        # ==================== Action Buttons ====================
        button_frame = ctk.CTkFrame(main_frame)
        button_frame.pack(fill="x", pady=20, padx=5)
        
        center_frame = ctk.CTkFrame(button_frame)
        center_frame.pack(expand=True)
        
        self.generate_btn = ctk.CTkButton(center_frame, text="🚀 Генерувати файли", 
                                         command=self.start_generation,
                                         width=180, height=40, corner_radius=6,
                                         font=ctk.CTkFont(size=14, weight="bold"),
                                         fg_color="#1f6aa5", hover_color="#144870")
        self.generate_btn.pack(side="left", padx=10)
        
        self.open_folder_btn = ctk.CTkButton(center_frame, text="📂 Відкрити папку", 
                                            command=self.open_output_folder,
                                            width=160, height=40, corner_radius=6,
                                            font=ctk.CTkFont(size=14),
                                            state="disabled")
        self.open_folder_btn.pack(side="left", padx=10)
        
        # ==================== Results ====================
        results_frame = ctk.CTkFrame(main_frame)
        results_frame.pack(fill="both", expand=True, pady=10, padx=5)
        
        ctk.CTkLabel(results_frame, text="📊 Результати:", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        
        self.results_text = ctk.CTkTextbox(results_frame, height=200, font=self.font)
        self.results_text.pack(fill="both", expand=True, padx=10, pady=10)
    
    def show_instruction(self):
        """Show instruction popup"""
        instruction_text = """
📁 FILE GENERATOR - ГЕНЕРАТОР ФАЙЛІВ

🚀 МОЖЛИВОСТІ:
• Генерація файлів різних форматів
• Унікальні осмислені назви файлів
• Вибір кількості файлів для кожного типу
• Автоматичне створення папки для файлів

📂 ПІДТРИМУВАНІ ФОРМАТИ:
• DOCX - документи Word
• XLSX - таблиці Excel
• PPTX - презентації PowerPoint
• PNG - зображення
• MP3 - аудіо файли
• MP4 - відео файли
• TXT - текстові файли
• PDF - документи PDF

⚙️ ЯК ВИКОРИСТОВУВАТИ:
1. Виберіть типи файлів (галочки)
2. Вкажіть кількість файлів кожного типу
3. Натисніть "Генерувати файли"
4. Дочекайтесь завершення
5. Відкрийте папку з файлами

📁 ПАПКА:
• Автоматично створюється 'generated_files'
• Знаходиться в директорії програми
• При повторній генерації старі файли видаляються

📝 ФОРМАТ НАЗВ:
• Annual_Report_2024.docx
• Budget_January_2025.xlsx
• Summary_Document_v2.3.pptx
• Contract_5847.pdf
• і інші осмислені комбінації

✅ ФАЙЛИ НЕ ДУБЛЮЮТЬСЯ!
"""
        messagebox.showinfo("Інструкція - File Generator", instruction_text)
    
    def start_generation(self):
        """Start file generation in separate thread"""
        if self.is_generating:
            messagebox.showwarning("Увага", "Генерація вже виконується!")
            return
        
        # Перевірка кількості
        try:
            count = int(self.count_entry.get())
            if count <= 0 or count > 1000:
                messagebox.showwarning("Увага", "Кількість файлів має бути від 1 до 1000!")
                return
        except ValueError:
            messagebox.showwarning("Увага", "Введіть коректну кількість файлів!")
            return
        
        # Перевірка вибраних типів
        selected_types = [ext for ext, var in self.file_type_vars.items() if var.get()]
        if not selected_types:
            messagebox.showwarning("Увага", "Виберіть хоча б один тип файлів!")
            return
        
        # Запускаємо генерацію в окремому потоці
        from threading import Thread
        Thread(target=self.generate_files, args=(count, selected_types), daemon=True).start()
    
    def generate_files(self, count, selected_types):
        """Generate files (runs in separate thread)"""
        self.is_generating = True
        self.generated_files = []
        
        # Вимикаємо кнопку генерації
        self.generate_btn.configure(state="disabled", text="⏳ Генерування...")
        
        # Очищаємо результати
        self.results_text.delete("1.0", "end")
        
        total_files = count * len(selected_types)
        generated_count = 0
        
        import os
        from datetime import datetime
        
        # Визначаємо директорію програми
        script_dir = os.path.dirname(os.path.abspath(__file__))
        
        # Створюємо папку 'generated_files' в директорії програми
        output_dir = os.path.join(script_dir, "generated_files")
        
        # Якщо папка існує - видаляємо старі файли
        if os.path.exists(output_dir):
            self.add_result("🗑️ Видалення старих файлів...")
            import shutil
            try:
                shutil.rmtree(output_dir)
                self.add_result("✅ Старі файли видалено")
            except Exception as e:
                self.add_result(f"⚠️ Помилка видалення: {e}")
        
        # Створюємо папку заново
        os.makedirs(output_dir, exist_ok=True)
        self.add_result(f"📁 Створено папку: {output_dir}\n")
        
        # Зберігаємо шлях до поточної папки
        self.current_generation_folder = output_dir
        
        for ext in selected_types:
            ext_name = ext.replace(".", "").upper()
            self.update_progress(f"Генерація {ext_name} файлів...", generated_count / total_files)
            
            for i in range(count):
                filename = self.generate_unique_filename(ext)
                filepath = os.path.join(output_dir, filename)
                
                # Генеруємо файл в залежності від типу
                self.create_file(filepath, ext)
                
                generated_count += 1
                self.update_progress(f"Згенеровано: {generated_count}/{total_files}", 
                                   generated_count / total_files)
                
                # Додаємо в результати
                self.add_result(f"✅ {filename}")
        
        # Завершення
        self.update_progress(f"✅ Завершено! Згенеровано {total_files} файлів", 1.0)
        self.is_generating = False
        self.generate_btn.configure(state="normal", text="🚀 Генерувати файли")
        self.open_folder_btn.configure(state="normal")
        
        # Показуємо повідомлення про успіх (БЕЗ popup)
        self.add_result(f"\n{'='*50}")
        self.add_result(f"🎉 УСПІХ! Згенеровано {total_files} файлів")
        self.add_result(f"📂 Папка: {output_dir}")
        self.add_result(f"{'='*50}")
    
    def create_file(self, filepath, ext):
        """Create file based on extension"""
        import random
        
        if ext == ".docx":
            self.create_docx(filepath)
        elif ext == ".xlsx":
            self.create_xlsx(filepath)
        elif ext == ".pptx":
            self.create_pptx(filepath)
        elif ext == ".png":
            self.create_png(filepath)
        elif ext == ".mp3":
            self.create_mp3(filepath)
        elif ext == ".mp4":
            self.create_mp4(filepath)
        elif ext == ".txt":
            self.create_txt(filepath)
        elif ext == ".pdf":
            self.create_pdf(filepath)
    
    def create_docx(self, filepath):
        """Create DOCX file"""
        try:
            from docx import Document
            from datetime import datetime
            import os
            
            doc = Document()
            # Заголовок з назви файлу
            filename = os.path.basename(filepath).replace('.docx', '').replace('_', ' ')
            doc.add_heading(filename, 0)
            
            # Дата створення
            doc.add_paragraph(f'Created: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
            doc.add_paragraph('')  # Порожній рядок
            
            # Додаємо рандомний текст
            random_text = self.generate_random_text(min_paragraphs=4, max_paragraphs=8)
            for paragraph in random_text.split('\n\n'):
                doc.add_paragraph(paragraph)
            
            doc.save(filepath)
        except ImportError:
            # Якщо бібліотека недоступна, створюємо порожній файл
            with open(filepath, 'wb') as f:
                f.write(b'PK\x03\x04')  # ZIP header (DOCX is ZIP)
    
    def create_xlsx(self, filepath):
        """Create XLSX file"""
        try:
            from openpyxl import Workbook
            from datetime import datetime
            import random
            import os
            
            wb = Workbook()
            ws = wb.active
            
            # Заголовок
            filename = os.path.basename(filepath).replace('.xlsx', '').replace('_', ' ')
            ws['A1'] = filename
            ws['A2'] = f'Created: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}'
            
            # Заголовки колонок
            ws['A4'] = 'Item'
            ws['B4'] = 'Category'
            ws['C4'] = 'Value'
            ws['D4'] = 'Status'
            ws['E4'] = 'Date'
            
            # Генеруємо рандомні дані
            for row in range(5, 25):  # 20 рядків даних
                ws[f'A{row}'] = f'Item {row-4}'
                ws[f'B{row}'] = random.choice(self.nouns)
                ws[f'C{row}'] = random.randint(100, 9999)
                ws[f'D{row}'] = random.choice(['Active', 'Pending', 'Completed', 'In Progress'])
                ws[f'E{row}'] = f'{random.choice(self.months)} {random.randint(1, 28)}, {random.choice(self.years)}'
            
            wb.save(filepath)
        except ImportError:
            # Якщо бібліотека недоступна, створюємо порожній файл
            with open(filepath, 'wb') as f:
                f.write(b'PK\x03\x04')
    
    def create_pptx(self, filepath):
        """Create PPTX file"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import os
            import random
            
            prs = Presentation()
            
            # Слайд 1 - Титульний
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide1.shapes.title
            subtitle = slide1.placeholders[1]
            
            filename = os.path.basename(filepath).replace('.pptx', '').replace('_', ' ')
            title.text = filename
            subtitle.text = f"Presentation {random.choice(self.years)}"
            
            # Слайд 2 - Контент
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            title2 = slide2.shapes.title
            content2 = slide2.placeholders[1]
            title2.text = random.choice(self.adjectives) + " " + random.choice(self.nouns)
            
            # Генеруємо bullet points
            tf = content2.text_frame
            bullet_text = self.generate_random_text(min_paragraphs=3, max_paragraphs=5)
            for paragraph in bullet_text.split('\n\n')[:5]:
                p = tf.add_paragraph()
                p.text = paragraph[:100] + "..." if len(paragraph) > 100 else paragraph
                p.level = 0
            
            # Слайд 3 - Ще контент
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            title3 = slide3.shapes.title
            content3 = slide3.placeholders[1]
            title3.text = "Key Points"
            
            tf3 = content3.text_frame
            for i in range(4):
                p = tf3.add_paragraph()
                p.text = f"{random.choice(self.adjectives)} {random.choice(self.nouns)} in {random.choice(self.years)}"
                p.level = 0
            
            prs.save(filepath)
        except ImportError:
            with open(filepath, 'wb') as f:
                f.write(b'PK\x03\x04')
    
    def create_png(self, filepath):
        """Create PNG image"""
        try:
            from PIL import Image, ImageDraw, ImageFont
            from datetime import datetime
            import random
            import os
            
            # Створюємо зображення з випадковим кольором
            width, height = 800, 600
            color = (random.randint(50, 200), random.randint(50, 200), random.randint(50, 200))
            img = Image.new('RGB', (width, height), color)
            
            draw = ImageDraw.Draw(img)
            
            # Назва файлу
            filename = os.path.basename(filepath).replace('.png', '').replace('_', ' ')
            
            # Заголовок
            try:
                title_font = ImageFont.truetype("arial.ttf", 24)
                text_font = ImageFont.truetype("arial.ttf", 14)
            except:
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()
            
            # Заголовок по центру
            title_bbox = draw.textbbox((0, 0), filename, font=title_font)
            title_width = title_bbox[2] - title_bbox[0]
            draw.text(((width - title_width) // 2, 50), filename, fill=(255, 255, 255), font=title_font)
            
            # Дата
            date_text = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            date_bbox = draw.textbbox((0, 0), date_text, font=text_font)
            date_width = date_bbox[2] - date_bbox[0]
            draw.text(((width - date_width) // 2, 90), date_text, fill=(255, 255, 255), font=text_font)
            
            # Рандомні слова
            y_pos = 140
            for i in range(15):
                text = f"{random.choice(self.adjectives)} {random.choice(self.nouns)}"
                text_bbox = draw.textbbox((0, 0), text, font=text_font)
                text_width = text_bbox[2] - text_bbox[0]
                draw.text(((width - text_width) // 2, y_pos), text, fill=(255, 255, 255), font=text_font)
                y_pos += 25
            
            img.save(filepath)
        except ImportError:
            # Мінімальний PNG файл
            with open(filepath, 'wb') as f:
                f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde')
    
    def create_mp3(self, filepath):
        """Create MP3 file"""
        # Створюємо мінімальний MP3 header
        with open(filepath, 'wb') as f:
            f.write(b'\xff\xfb\x90\x00' * 100)  # MP3 frame header
    
    def create_mp4(self, filepath):
        """Create MP4 file"""
        # Створюємо мінімальний MP4 header
        with open(filepath, 'wb') as f:
            f.write(b'\x00\x00\x00\x20ftypisom\x00\x00\x02\x00isomiso2mp41')
    
    def create_txt(self, filepath):
        """Create TXT file"""
        from datetime import datetime
        import os
        
        filename = os.path.basename(filepath).replace('.txt', '').replace('_', ' ')
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(f"{filename}\n")
            f.write(f"{'='*60}\n")
            f.write(f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"{'='*60}\n\n")
            
            # Додаємо рандомний текст
            random_text = self.generate_random_text(min_paragraphs=5, max_paragraphs=10)
            f.write(random_text)
    
    def create_pdf(self, filepath):
        """Create PDF file"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from datetime import datetime
            import os
            
            c = canvas.Canvas(filepath, pagesize=letter)
            width, height = letter
            
            # Заголовок
            filename = os.path.basename(filepath).replace('.pdf', '').replace('_', ' ')
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 50, filename)
            
            # Дата
            c.setFont("Helvetica", 10)
            c.drawString(50, height - 70, f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            
            # Лінія
            c.line(50, height - 80, width - 50, height - 80)
            
            # Генеруємо текст
            random_text = self.generate_random_text(min_paragraphs=3, max_paragraphs=6)
            
            # Виводимо текст з переносами
            c.setFont("Helvetica", 11)
            y_position = height - 110
            
            for paragraph in random_text.split('\n\n'):
                # Розбиваємо параграф на рядки по 80 символів
                words = paragraph.split()
                lines = []
                current_line = ""
                
                for word in words:
                    if len(current_line + " " + word) < 80:
                        current_line += (" " + word) if current_line else word
                    else:
                        lines.append(current_line)
                        current_line = word
                
                if current_line:
                    lines.append(current_line)
                
                # Виводимо рядки
                for line in lines:
                    if y_position < 50:  # Якщо місця немає, створюємо нову сторінку
                        c.showPage()
                        c.setFont("Helvetica", 11)
                        y_position = height - 50
                    
                    c.drawString(50, y_position, line)
                    y_position -= 15
                
                y_position -= 10  # Відступ між параграфами
            
            c.save()
        except ImportError:
            # Мінімальний PDF
            with open(filepath, 'wb') as f:
                f.write(b'%PDF-1.4\n')
    
    def update_progress(self, text, value):
        """Update progress bar and label"""
        self.progress_label.configure(text=text)
        self.progress_bar.set(value)
    
    def add_result(self, text):
        """Add result to textbox"""
        self.results_text.insert("end", text + "\n")
        self.results_text.see("end")
    
    def open_output_folder(self):
        """Open output folder in file explorer"""
        if self.current_generation_folder:
            import os
            import subprocess
            if os.path.exists(self.current_generation_folder):
                subprocess.Popen(f'explorer "{self.current_generation_folder}"')
            else:
                messagebox.showerror("Помилка", "Папка не існує!")
        else:
            messagebox.showwarning("Увага", "Спочатку згенеруйте файли!")


# ================== Mail.TM Tab (Temporary Email Service) ==================
class RegistrationTab(ctk.CTkFrame):
    """Mail.TM - Temporary Email Service Integration"""
    
    def __init__(self, master, font=None):
        super().__init__(master)
        self.font = font or ctk.CTkFont()
        self.pack(fill="both", expand=True)
        
        # Mail.TM API Base URL
        self.API_BASE = "https://api.mail.tm"
        
        # Storage for accounts
        self.accounts = []  # List of {email, password, token, id}
        self.selected_account = None
        self.messages = []  # Messages for selected account
        self.stop_creation = False  # Flag to stop account creation
        
        self.create_widgets()
    
    def show_instruction(self):
        instruction_text = """
📧 MAIL.TM - ТИМЧАСОВА ПОШТА

🚀 МОЖЛИВОСТІ:
• Створення тимчасових email адрес
• Автоматична реєстрація на mail.tm
• Вхід в існуючі акаунти
• Перегляд вхідних повідомлень
• Копіювання: логін, пароль, login:pass

📧 СТВОРЕННЯ ПОШТИ:
• Натисніть "Create New Email"
• Система автоматично створить акаунт
• Email та пароль буде збережено
• Можна створити декілька адрес

🔐 ВХІД В АКАУНТ:
• Виберіть email зі списку
• Натисніть "Login"  
• Токен автоматично оновиться
• Можна переглядати повідомлення

� ПОВІДОМЛЕННЯ:
• Список всіх вхідних листів
• Автоматичне оновлення (30 сек)
• Перегляд повного тексту
• Копіювання важливої інформації

� КОПІЮВАННЯ:
• Copy Email - тільки адреса
• Copy Password - тільки пароль
• Copy Login:Pass - формат email:password
• Зручно для реєстрацій

� ЗБЕРЕЖЕННЯ:
• Всі акаунти зберігаються автоматично
• Файл: ~/.punchnow/mailtm_accounts.json
• Завантажуються при наступному запуску

� АВТООНОВЛЕННЯ:
• Повідомлення оновлюються кожні 30 сек
• Можна вимкнути/увімкнути
• Ручне оновлення кнопкою "Refresh"
        """
        messagebox.showinfo("📧 Mail.TM - Інструкція", instruction_text)
    
    def create_widgets(self):
        """Створює інтерфейс Mail.TM"""
        main_frame = ctk.CTkScrollableFrame(self)
        main_frame.pack(fill="both", expand=True, padx=10, pady=10)
        
        # Header
        header_frame = ctk.CTkFrame(main_frame)
        header_frame.pack(fill="x", pady=(10, 20))
        
        ctk.CTkLabel(header_frame, text="📧 Mail.TM - Temporary Email", 
                    font=ctk.CTkFont(size=18, weight="bold")).pack(side="left", padx=10)
        
        ctk.CTkButton(header_frame, text="📋 View Logs", 
                     command=self.show_logs_window, width=100, height=28).pack(side="right", padx=5)
        
        # Current account log
        current_acc_frame = ctk.CTkFrame(main_frame)
        current_acc_frame.pack(fill="x", pady=(0, 10), padx=5)
        
        ctk.CTkLabel(current_acc_frame, text="🔐 Current Account:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(side="left", padx=10, pady=10)
        
        self.current_account_label = ctk.CTkLabel(current_acc_frame, 
                                                  text="Not logged in", 
                                                  font=ctk.CTkFont(size=11),
                                                  text_color="gray")
        self.current_account_label.pack(side="left", padx=5, pady=10)
        
        # Controls Frame
        controls_frame = ctk.CTkFrame(main_frame)
        controls_frame.pack(fill="x", pady=10, padx=5)
        
        ctk.CTkLabel(controls_frame, text="Account Management", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(10,5))
        
        # Bulk creation frame
        bulk_frame = ctk.CTkFrame(controls_frame)
        bulk_frame.pack(pady=5)
        
        ctk.CTkLabel(bulk_frame, text="Create:", 
                    font=ctk.CTkFont(size=11)).pack(side="left", padx=5)
        
        self.bulk_count_entry = ctk.CTkEntry(bulk_frame, width=70, 
                                            placeholder_text="1-1000")
        self.bulk_count_entry.pack(side="left", padx=5)
        self.bulk_count_entry.insert(0, "1")
        
        ctk.CTkLabel(bulk_frame, text="emails at once", 
                    font=ctk.CTkFont(size=11)).pack(side="left", padx=5)
        
        ctk.CTkLabel(bulk_frame, text="| Delay:", 
                    font=ctk.CTkFont(size=11)).pack(side="left", padx=(20, 5))
        
        self.delay_entry = ctk.CTkEntry(bulk_frame, width=60, 
                                       placeholder_text="sec")
        self.delay_entry.pack(side="left", padx=5)
        self.delay_entry.insert(0, "15")
        
        ctk.CTkLabel(bulk_frame, text="sec/account", 
                    font=ctk.CTkFont(size=11)).pack(side="left", padx=5)
        
        btn_frame = ctk.CTkFrame(controls_frame)
        btn_frame.pack(pady=10)
        
        self.create_email_btn = ctk.CTkButton(btn_frame, text="✉️ Create Email(s)", 
                                             command=self.create_new_email,
                                             width=140, height=32)
        self.create_email_btn.pack(side="left", padx=5)
        
        self.stop_btn = ctk.CTkButton(btn_frame, text="⏹️ Stop", 
                                     command=self.stop_creation_process,
                                     width=80, height=32, state="disabled",
                                     fg_color="#8B0000", hover_color="#A52A2A")
        self.stop_btn.pack(side="left", padx=5)
        
        self.refresh_btn = ctk.CTkButton(btn_frame, text="🔄 Refresh", 
                                        command=self.refresh_accounts_list,
                                        width=100, height=32)
        self.refresh_btn.pack(side="left", padx=5)
        
        self.login_btn = ctk.CTkButton(btn_frame, text="🔑 Login", 
                                      command=self.login_account,
                                      width=100, height=32, state="disabled")
        self.login_btn.pack(side="left", padx=5)
        
        self.delete_btn = ctk.CTkButton(btn_frame, text="🗑️ Delete", 
                                       command=self.delete_account,
                                       width=100, height=32, state="disabled")
        self.delete_btn.pack(side="left", padx=5)
        
        self.save_btn = ctk.CTkButton(btn_frame, text="💾 Save All", 
                                     command=self.save_accounts,
                                     width=100, height=32)
        self.save_btn.pack(side="left", padx=5)
        
        self.open_file_btn = ctk.CTkButton(btn_frame, text="📂 Open File", 
                                          command=self.open_accounts_file,
                                          width=100, height=32)
        self.open_file_btn.pack(side="left", padx=5)
        
        # Accounts List Frame
        accounts_frame = ctk.CTkFrame(main_frame)
        accounts_frame.pack(fill="both", expand=True, pady=10, padx=5)
        
        acc_header = ctk.CTkFrame(accounts_frame)
        acc_header.pack(fill="x", pady=(10,5), padx=10)
        
        ctk.CTkLabel(acc_header, text="📬 Email Accounts", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(side="left")
        
        self.accounts_count_var = tk.StringVar(value="(0)")
        ctk.CTkLabel(acc_header, textvariable=self.accounts_count_var, 
                    font=ctk.CTkFont(size=12), 
                    text_color="gray").pack(side="left", padx=5)
        
        # List and buttons container
        list_container = ctk.CTkFrame(accounts_frame)
        list_container.pack(fill="both", expand=True, padx=10, pady=10)
        
        # Accounts listbox
        self.accounts_listbox = tk.Listbox(list_container, height=8, 
                                          font=("Consolas", 10))
        self.accounts_listbox.pack(side="left", fill="both", expand=True)
        self.accounts_listbox.bind('<<ListboxSelect>>', self.on_account_select)
        
        scrollbar = tk.Scrollbar(list_container, command=self.accounts_listbox.yview)
        scrollbar.pack(side="right", fill="y")
        self.accounts_listbox.config(yscrollcommand=scrollbar.set)
        
        # Copy buttons
        copy_frame = ctk.CTkFrame(accounts_frame)
        copy_frame.pack(pady=5)
        
        ctk.CTkButton(copy_frame, text="📧 Copy Email", 
                     command=self.copy_email, width=110, height=28).pack(side="left", padx=3)
        ctk.CTkButton(copy_frame, text="🔑 Copy Password", 
                     command=self.copy_password, width=110, height=28).pack(side="left", padx=3)
        ctk.CTkButton(copy_frame, text="📋 Copy Login:Pass", 
                     command=self.copy_login_pass, width=130, height=28).pack(side="left", padx=3)
        
        # Messages Frame
        messages_frame = ctk.CTkFrame(main_frame)
        messages_frame.pack(fill="both", expand=True, pady=10, padx=5)
        
        msg_header = ctk.CTkFrame(messages_frame)
        msg_header.pack(fill="x", pady=(10,5), padx=10)
        
        ctk.CTkLabel(msg_header, text="📨 Messages", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(side="left")
        
        self.auto_refresh_var = tk.BooleanVar(value=True)
        ctk.CTkCheckBox(msg_header, text="Auto-refresh (30s)", 
                       variable=self.auto_refresh_var).pack(side="right", padx=5)
        
        ctk.CTkButton(msg_header, text="🔄 Refresh", 
                     command=self.refresh_messages, width=90, height=26).pack(side="right", padx=5)
        
        # Messages listbox
        msg_list_container = ctk.CTkFrame(messages_frame)
        msg_list_container.pack(fill="both", expand=True, padx=10, pady=5)
        
        self.messages_listbox = tk.Listbox(msg_list_container, height=6, 
                                          font=("Consolas", 9))
        self.messages_listbox.pack(side="left", fill="both", expand=True)
        self.messages_listbox.bind('<<ListboxSelect>>', self.on_message_select)
        
        msg_scrollbar = tk.Scrollbar(msg_list_container, command=self.messages_listbox.yview)
        msg_scrollbar.pack(side="right", fill="y")
        self.messages_listbox.config(yscrollcommand=msg_scrollbar.set)
        
        # Message content
        self.message_text = ctk.CTkTextbox(messages_frame, height=150, 
                                          wrap="word", state="disabled")
        self.message_text.pack(fill="both", expand=True, padx=10, pady=(5,10))
        
        # Status bar
        self.status_var = tk.StringVar(value="Ready - Create or login to an account")
        self.status_label = ctk.CTkLabel(main_frame, textvariable=self.status_var)
        self.status_label.pack(fill="x", pady=5)
        
        # Load saved accounts
        self.load_accounts()
        
        # Start auto-refresh thread
        self.start_auto_refresh()
    
    def get_config_path(self):
        """Get config file path"""
        import os
        config_dir = os.path.join(os.path.expanduser("~"), ".punchnow")
        os.makedirs(config_dir, exist_ok=True)
        return os.path.join(config_dir, "mailtm_accounts.json")
    
    def load_accounts(self):
        """Load saved accounts from config"""
        import os
        import json
        try:
            config_file = self.get_config_path()
            if os.path.exists(config_file):
                with open(config_file, 'r', encoding='utf-8') as f:
                    self.accounts = json.load(f)
                self.update_accounts_list()
                
                # Set first logged-in account as current
                for acc in self.accounts:
                    if acc.get('token'):
                        self.selected_account = acc
                        self.current_account_label.configure(
                            text=f"🟢 {acc['email']}", 
                            text_color="green"
                        )
                        break
                
                self.status_var.set(f"Loaded {len(self.accounts)} saved accounts")
        except Exception as e:
            print(f"Error loading accounts: {e}")
            self.accounts = []
    
    def save_accounts(self):
        """Save accounts to config"""
        import json
        try:
            config_file = self.get_config_path()
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(self.accounts, f, indent=4)
            self.status_var.set(f"Saved {len(self.accounts)} accounts")
        except Exception as e:
            self.status_var.set(f"Error saving: {str(e)}")
    
    def update_accounts_list(self):
        """Update accounts listbox"""
        import tkinter as tk
        self.accounts_listbox.delete(0, tk.END)
        for acc in self.accounts:
            email = acc.get('email', 'Unknown')
            status = "🟢" if acc.get('token') else "🔴"
            self.accounts_listbox.insert(tk.END, f"{status} {email}")
        
        # Update counter
        logged_in = sum(1 for acc in self.accounts if acc.get('token'))
        self.accounts_count_var.set(f"({len(self.accounts)} total, {logged_in} logged in)")
    
    def create_new_email(self):
        """Create new temporary email(s)"""
        import requests
        import string
        import random
        from threading import Thread
        import time
        
        # Get count from entry
        try:
            count = int(self.bulk_count_entry.get().strip())
            if count < 1 or count > 1000:
                self.status_var.set("⚠️ Please enter a number between 1 and 1000")
                return
        except ValueError:
            self.status_var.set("⚠️ Please enter a valid number")
            return
        
        # Get delay from entry
        try:
            delay = float(self.delay_entry.get().strip())
            if delay < 0:
                delay = 15
        except ValueError:
            delay = 15
        
        self.status_var.set(f"Creating {count} email(s)...")
        self.create_email_btn.configure(state="disabled")
        self.stop_btn.configure(state="normal")
        self.stop_creation = False  # Reset flag
        
        def create_thread():
            created = 0
            failed = 0
            
            try:
                # Get available domains
                domains_resp = requests.get(f"{self.API_BASE}/domains")
                if domains_resp.status_code != 200:
                    error_msg = f"Error: Cannot get domains (status {domains_resp.status_code})"
                    self.status_var.set(error_msg)
                    print(f"DEBUG: {error_msg}")
                    return
                
                domains_data = domains_resp.json()
                domains = domains_data.get('hydra:member', domains_data if isinstance(domains_data, list) else [])
                
                if not domains:
                    self.status_var.set("Error: No domains available")
                    return
                
                domain = domains[0].get('domain') if isinstance(domains[0], dict) else domains[0]
                print(f"DEBUG: Using domain: {domain}")
                
                # Create multiple accounts
                for i in range(count):
                    # Check if stop was requested
                    if self.stop_creation:
                        self.status_var.set(f"⏹️ Stopped! Created: {created}, Failed: {failed}")
                        print(f"⏹️ Creation stopped by user")
                        break
                    
                    try:
                        self.status_var.set(f"Creating email {i+1}/{count}...")
                        
                        # Generate random email
                        username = ''.join(random.choices(string.ascii_lowercase + string.digits, k=10))
                        email = f"{username}@{domain}"
                        password = ''.join(random.choices(string.ascii_letters + string.digits, k=16))
                        
                        # Create account
                        account_data = {
                            "address": email,
                            "password": password
                        }
                        
                        create_resp = requests.post(f"{self.API_BASE}/accounts", 
                                                  json=account_data,
                                                  headers={"Content-Type": "application/json"})
                        
                        if create_resp.status_code in [200, 201]:
                            result = create_resp.json()
                            account_id = result.get('id')
                            
                            # Login to get token
                            login_resp = requests.post(f"{self.API_BASE}/token", 
                                                     json=account_data,
                                                     headers={"Content-Type": "application/json"})
                            
                            token = None
                            if login_resp.status_code == 200:
                                token_data = login_resp.json()
                                token = token_data.get('token')
                            
                            # Save account
                            new_account = {
                                "email": email,
                                "password": password,
                                "token": token,
                                "id": account_id
                            }
                            self.accounts.append(new_account)
                            
                            # Update current account label for the first/last created account
                            if created == 0:  # First account
                                self.selected_account = new_account
                                self.current_account_label.configure(
                                    text=f"🟢 {email}", 
                                    text_color="green"
                                )
                            
                            created += 1
                            
                            # Update UI every 5 accounts or on last
                            if created % 5 == 0 or i == count - 1:
                                self.update_accounts_list()
                                self.save_accounts()
                            
                            print(f"✅ Created: {email}")
                        else:
                            failed += 1
                            print(f"❌ Failed to create account {i+1}: {create_resp.status_code}")
                        
                        # Delay between accounts
                        if i < count - 1:
                            time.sleep(delay)
                            
                    except Exception as e:
                        failed += 1
                        print(f"❌ Error creating account {i+1}: {e}")
                
                # Final save and update
                self.update_accounts_list()
                self.save_accounts()
                
                if not self.stop_creation:
                    self.status_var.set(f"✅ Done! Created: {created}, Failed: {failed}")
                    
            except Exception as e:
                error_msg = f"Error: {str(e)}"
                self.status_var.set(error_msg)
                print(f"DEBUG EXCEPTION: {error_msg}")
                import traceback
                traceback.print_exc()
            finally:
                self.create_email_btn.configure(state="normal")
                self.stop_btn.configure(state="disabled")
        
        Thread(target=create_thread, daemon=True).start()
    
    def stop_creation_process(self):
        """Stop account creation process"""
        self.stop_creation = True
        self.status_var.set("⏹️ Stopping creation...")
        self.stop_btn.configure(state="disabled")
    
    def refresh_accounts_list(self):
        """Refresh accounts list from file"""
        self.load_accounts()
        self.update_accounts_list()
        self.status_var.set("🔄 List refreshed!")
    
    def on_account_select(self, event):
        """Handle account selection"""
        selection = self.accounts_listbox.curselection()
        if selection:
            self.login_btn.configure(state="normal")
            self.delete_btn.configure(state="normal")
    
    def login_account(self):
        """Login to selected account"""
        import requests
        from threading import Thread
        
        selection = self.accounts_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        account = self.accounts[idx]
        
        self.status_var.set(f"Logging in to {account['email']}...")
        
        def login_thread():
            try:
                login_data = {
                    "address": account['email'],
                    "password": account['password']
                }
                
                resp = requests.post(f"{self.API_BASE}/token", 
                                   json=login_data,
                                   headers={"Content-Type": "application/json"})
                
                if resp.status_code == 200:
                    token_data = resp.json()
                    account['token'] = token_data.get('token')
                    account['id'] = token_data.get('id')
                    
                    self.selected_account = account
                    self.update_accounts_list()
                    self.save_accounts()
                    
                    # Update current account label
                    self.current_account_label.configure(
                        text=f"🟢 {account['email']}", 
                        text_color="green"
                    )
                    
                    self.status_var.set(f"✅ Logged in: {account['email']}")
                    self.refresh_messages()
                else:
                    self.status_var.set(f"Login failed: {resp.status_code}")
            except Exception as e:
                self.status_var.set(f"Login error: {str(e)}")
        
        Thread(target=login_thread, daemon=True).start()
    
    def delete_account(self):
        """Delete selected account"""
        selection = self.accounts_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        account = self.accounts[idx]
        
        if messagebox.askyesno("Delete Account", 
                               f"Delete {account['email']}?"):
            self.accounts.pop(idx)
            self.update_accounts_list()
            self.save_accounts()
            self.status_var.set(f"Deleted: {account['email']}")
            
            if self.selected_account == account:
                self.selected_account = None
                self.messages_listbox.delete(0, tk.END)
                # Update current account label
                self.current_account_label.configure(
                    text="Not logged in", 
                    text_color="gray"
                )
    
    def open_accounts_file(self):
        """Open the accounts JSON file in default editor"""
        import os
        import subprocess
        
        try:
            config_file = self.get_config_path()
            
            # Make sure file exists
            if not os.path.exists(config_file):
                self.save_accounts()  # Create file if doesn't exist
            
            # Open with default application
            if os.name == 'nt':  # Windows
                os.startfile(config_file)
            elif os.name == 'posix':  # Linux/Mac
                opener = 'open' if os.uname().sysname == 'Darwin' else 'xdg-open'
                subprocess.call([opener, config_file])
            
            self.status_var.set(f"📂 Opened: {config_file}")
        except Exception as e:
            self.status_var.set(f"Error opening file: {str(e)}")
    
    def show_logs_window(self):
        """Show logs window with all generated emails"""
        import tkinter as tk
        from tkinter import ttk
        
        # Create new window
        logs_window = tk.Toplevel(self)
        logs_window.title("📋 Mail.TM - Generated Emails Log")
        logs_window.geometry("900x500")
        logs_window.configure(bg="#2b2b2b")
        
        # Header
        header_frame = tk.Frame(logs_window, bg="#1e1e1e", height=60)
        header_frame.pack(fill="x", padx=10, pady=10)
        header_frame.pack_propagate(False)
        
        title_label = tk.Label(header_frame, text=f"📧 Total Accounts: {len(self.accounts)}", 
                              font=("Segoe UI", 14, "bold"), fg="white", bg="#1e1e1e")
        title_label.pack(side="left", padx=10, pady=15)
        
        logged_in = sum(1 for acc in self.accounts if acc.get('token'))
        status_label = tk.Label(header_frame, text=f"🟢 Logged In: {logged_in}  |  🔴 Not Logged: {len(self.accounts) - logged_in}", 
                               font=("Segoe UI", 11), fg="#aaaaaa", bg="#1e1e1e")
        status_label.pack(side="left", padx=20, pady=15)
        
        # Buttons frame
        btn_frame = tk.Frame(header_frame, bg="#1e1e1e")
        btn_frame.pack(side="right", padx=10, pady=10)
        
        refresh_btn = tk.Button(btn_frame, text="🔄 Refresh", command=lambda: self.refresh_logs(logs_window),
                               bg="#2d5f2d", fg="white", font=("Segoe UI", 10), relief="flat", 
                               padx=15, pady=5, cursor="hand2")
        refresh_btn.pack(side="left", padx=5)
        
        copy_all_btn = tk.Button(btn_frame, text="📋 Copy All", command=self.copy_all_accounts,
                                bg="#2d4d5f", fg="white", font=("Segoe UI", 10), relief="flat", 
                                padx=15, pady=5, cursor="hand2")
        copy_all_btn.pack(side="left", padx=5)
        
        # Search frame
        search_frame = tk.Frame(logs_window, bg="#2b2b2b")
        search_frame.pack(fill="x", padx=10, pady=(0, 10))
        
        tk.Label(search_frame, text="🔍 Search:", font=("Segoe UI", 10), 
                fg="white", bg="#2b2b2b").pack(side="left", padx=5)
        
        search_entry = tk.Entry(search_frame, font=("Segoe UI", 11), width=40)
        search_entry.pack(side="left", padx=5, fill="x", expand=True)
        
        def on_search(*args):
            search_text = search_entry.get().lower()
            # Clear tree
            for item in tree.get_children():
                tree.delete(item)
            # Filter and insert
            for acc in self.accounts:
                email = acc.get('email', '')
                if search_text in email.lower():
                    status = "🟢 Online" if acc.get('token') else "🔴 Offline"
                    password = acc.get('password', 'N/A')
                    acc_id = acc.get('id', 'N/A')
                    tag = "logged_in" if acc.get('token') else "not_logged"
                    tree.insert("", "end", values=(status, email, password, acc_id), tags=(tag,))
        
        search_entry.bind('<KeyRelease>', on_search)
        
        clear_search_btn = tk.Button(search_frame, text="✖ Clear", 
                                     command=lambda: (search_entry.delete(0, tk.END), on_search()),
                                     bg="#5f2d2d", fg="white", font=("Segoe UI", 9), relief="flat", 
                                     padx=10, pady=3, cursor="hand2")
        clear_search_btn.pack(side="left", padx=5)
        
        # Create Treeview with scrollbar
        tree_frame = tk.Frame(logs_window, bg="#2b2b2b")
        tree_frame.pack(fill="both", expand=True, padx=10, pady=(0, 10))
        
        # Scrollbars
        vsb = ttk.Scrollbar(tree_frame, orient="vertical")
        hsb = ttk.Scrollbar(tree_frame, orient="horizontal")
        
        # Treeview
        columns = ("status", "email", "password", "id")
        tree = ttk.Treeview(tree_frame, columns=columns, show="headings", 
                           yscrollcommand=vsb.set, xscrollcommand=hsb.set)
        
        vsb.config(command=tree.yview)
        hsb.config(command=tree.xview)
        
        # Define columns
        tree.heading("status", text="Status")
        tree.heading("email", text="Email")
        tree.heading("password", text="Password")
        tree.heading("id", text="Account ID")
        
        tree.column("status", width=80, anchor="center")
        tree.column("email", width=300, anchor="w")
        tree.column("password", width=200, anchor="w")
        tree.column("id", width=250, anchor="w")
        
        # Configure tags for colors
        tree.tag_configure("logged_in", background="#1e3a1e", foreground="white")
        tree.tag_configure("not_logged", background="#3a1e1e", foreground="white")
        
        # Insert data
        for idx, acc in enumerate(self.accounts):
            status = "🟢 Online" if acc.get('token') else "🔴 Offline"
            email = acc.get('email', 'N/A')
            password = acc.get('password', 'N/A')
            acc_id = acc.get('id', 'N/A')
            
            tag = "logged_in" if acc.get('token') else "not_logged"
            tree.insert("", "end", values=(status, email, password, acc_id), tags=(tag,))
        
        # Pack scrollbars and tree
        vsb.pack(side="right", fill="y")
        hsb.pack(side="bottom", fill="x")
        tree.pack(fill="both", expand=True)
        
        # Context menu for copying
        def on_right_click(event):
            item = tree.identify_row(event.y)
            if item:
                tree.selection_set(item)
                menu = tk.Menu(logs_window, tearoff=0)
                menu.add_command(label="📧 Copy Email", command=lambda: copy_field(1))
                menu.add_command(label="🔑 Copy Password", command=lambda: copy_field(2))
                menu.add_command(label="📋 Copy Email:Password", command=lambda: copy_login_pass())
                menu.add_command(label="🆔 Copy ID", command=lambda: copy_field(3))
                menu.post(event.x_root, event.y_root)
        
        def copy_field(col_idx):
            selected = tree.selection()
            if selected:
                item = tree.item(selected[0])
                value = item['values'][col_idx]
                safe_clipboard_operation("set", str(value))
                self.status_var.set(f"✅ Copied: {value}")
        
        def copy_login_pass():
            selected = tree.selection()
            if selected:
                item = tree.item(selected[0])
                email = item['values'][1]
                password = item['values'][2]
                safe_clipboard_operation("set", f"{email}:{password}")
                self.status_var.set(f"✅ Copied: {email}:{password}")
        
        tree.bind("<Button-3>", on_right_click)
        
        # Footer info
        footer_frame = tk.Frame(logs_window, bg="#1e1e1e", height=40)
        footer_frame.pack(fill="x", padx=10, pady=(0, 10))
        footer_frame.pack_propagate(False)
        
        info_label = tk.Label(footer_frame, 
                             text="💡 Right-click on any row to copy individual fields", 
                             font=("Segoe UI", 9), fg="#888888", bg="#1e1e1e")
        info_label.pack(pady=10)
        
        # Store tree reference for refresh
        logs_window.tree = tree
    
    def refresh_logs(self, logs_window):
        """Refresh logs window data"""
        if hasattr(logs_window, 'tree'):
            tree = logs_window.tree
            # Clear existing items
            for item in tree.get_children():
                tree.delete(item)
            
            # Re-insert data
            for acc in self.accounts:
                status = "🟢 Online" if acc.get('token') else "🔴 Offline"
                email = acc.get('email', 'N/A')
                password = acc.get('password', 'N/A')
                acc_id = acc.get('id', 'N/A')
                
                tag = "logged_in" if acc.get('token') else "not_logged"
                tree.insert("", "end", values=(status, email, password, acc_id), tags=(tag,))
            
            self.status_var.set("✅ Logs refreshed")
    
    def copy_all_accounts(self):
        """Copy all accounts in format email:password"""
        if not self.accounts:
            self.status_var.set("⚠️ No accounts to copy")
            return
        
        all_accounts = "\n".join([f"{acc['email']}:{acc['password']}" for acc in self.accounts])
        safe_clipboard_operation("set", all_accounts)
        self.status_var.set(f"✅ Copied {len(self.accounts)} accounts to clipboard")
    
    def copy_email(self):
        """Copy selected email"""
        selection = self.accounts_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        email = self.accounts[idx]['email']
        safe_clipboard_operation("set", email)
        self.status_var.set(f"Copied email: {email}")
    
    def copy_password(self):
        """Copy selected password"""
        selection = self.accounts_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        password = self.accounts[idx]['password']
        safe_clipboard_operation("set", password)
        self.status_var.set("Copied password")
    
    def copy_login_pass(self):
        """Copy login:pass format"""
        selection = self.accounts_listbox.curselection()
        if not selection:
            return
        
        idx = selection[0]
        acc = self.accounts[idx]
        login_pass = f"{acc['email']}:{acc['password']}"
        safe_clipboard_operation("set", login_pass)
        self.status_var.set(f"Copied: {acc['email']}:****")
    
    def refresh_messages(self):
        """Refresh messages for selected account"""
        import requests
        from threading import Thread
        
        if not self.selected_account or not self.selected_account.get('token'):
            self.status_var.set("Please login to an account first")
            return
        
        def refresh_thread():
            try:
                token = self.selected_account['token']
                headers = {"Authorization": f"Bearer {token}"}
                
                resp = requests.get(f"{self.API_BASE}/messages", headers=headers)
                
                if resp.status_code == 200:
                    messages_data = resp.json()
                    self.messages = messages_data.get('hydra:member', [])
                    
                    self.messages_listbox.delete(0, tk.END)
                    for msg in self.messages:
                        subject = msg.get('subject', 'No Subject')
                        from_addr = msg.get('from', {}).get('address', 'Unknown')
                        date = msg.get('createdAt', '')[:10]
                        self.messages_listbox.insert(tk.END, f"[{date}] {from_addr}: {subject}")
                    
                    self.status_var.set(f"📬 {len(self.messages)} messages")
                elif resp.status_code == 401:
                    self.status_var.set("Session expired - please login again")
                    self.selected_account['token'] = None
                    self.update_accounts_list()
                else:
                    self.status_var.set(f"Error: {resp.status_code}")
            except Exception as e:
                self.status_var.set(f"Refresh error: {str(e)}")
        
        Thread(target=refresh_thread, daemon=True).start()
    
    def on_message_select(self, event):
        """Handle message selection"""
        import requests
        from threading import Thread
        
        selection = self.messages_listbox.curselection()
        if not selection or not self.selected_account:
            return
        
        idx = selection[0]
        message = self.messages[idx]
        message_id = message.get('id')
        
        def load_message_thread():
            try:
                token = self.selected_account['token']
                headers = {"Authorization": f"Bearer {token}"}
                
                resp = requests.get(f"{self.API_BASE}/messages/{message_id}", 
                                  headers=headers)
                
                if resp.status_code == 200:
                    msg_data = resp.json()
                    
                    # Extract message details
                    subject = msg_data.get('subject', 'No Subject')
                    from_addr = msg_data.get('from', {}).get('address', 'Unknown')
                    text = msg_data.get('text', '')
                    
                    # Display message
                    self.message_text.configure(state="normal")
                    self.message_text.delete("1.0", "end")
                    
                    content = f"📧 Subject: {subject}\n"
                    content += f"📤 From: {from_addr}\n"
                    content += f"{'='*50}\n\n"
                    content += text if text else "No text content"
                    
                    self.message_text.insert("1.0", content)
                    self.message_text.configure(state="disabled")
                else:
                    self.status_var.set(f"Error loading message: {resp.status_code}")
            except Exception as e:
                self.status_var.set(f"Error: {str(e)}")
        
        Thread(target=load_message_thread, daemon=True).start()
    
    def start_auto_refresh(self):
        """Start auto-refresh thread"""
        import time
        from threading import Thread
        
        def auto_refresh_loop():
            while True:
                time.sleep(30)  # 30 seconds
                if self.auto_refresh_var.get() and self.selected_account:
                    self.refresh_messages()
        
        Thread(target=auto_refresh_loop, daemon=True).start()


# ================== Gmail Parser Tab ==================

class GmailParserTab(ctk.CTkFrame):
    def __init__(self, master, font=None):
        super().__init__(master)
        self.font = font
        self.pack(fill="both", expand=True)
        self.setup_ui()

    def setup_ui(self):
        # Верхня частина - введення даних
        input_frame = ctk.CTkFrame(self)
        input_frame.pack(padx=10, pady=5, fill="both")
        ctk.CTkLabel(input_frame, text="Введіть акаунти (по одному на рядок):", font=self.font).pack(anchor="w", padx=10, pady=(10,5))
        self.input_textarea = ctk.CTkTextbox(input_frame, width=600, height=150, font=self.font)
        self.input_textarea.pack(padx=10, pady=5, fill="both", expand=True)

        # Секція парсингу з Gmail Hacks
        parse_frame = ctk.CTkFrame(self)
        parse_frame.pack(padx=10, pady=5, fill="x")
        
        ctk.CTkLabel(parse_frame, text="Парсинг з Gmail Hacks:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10,5))
        
        # Контейнер для полів вводу
        parse_input_frame = ctk.CTkFrame(parse_frame)
        parse_input_frame.pack(padx=10, pady=5, fill="x")
        
        # Поле для email
        ctk.CTkLabel(parse_input_frame, text="Profile Name:", font=self.font).grid(row=0, column=0, sticky="w", padx=5, pady=5)
        self.parse_email_entry = ctk.CTkEntry(parse_input_frame, width=200, font=self.font, 
                                            placeholder_text="Введіть profile name...")
        self.parse_email_entry.grid(row=0, column=1, padx=5, pady=5, sticky="ew")
        
        # Поле для кількості
        ctk.CTkLabel(parse_input_frame, text="Кількість:", font=self.font).grid(row=0, column=2, sticky="w", padx=5, pady=5)
        self.parse_count_entry = ctk.CTkEntry(parse_input_frame, width=80, font=self.font, 
                                            placeholder_text="20")
        self.parse_count_entry.grid(row=0, column=3, padx=5, pady=5)
        self.parse_count_entry.insert(0, "20")  # Значення за замовчуванням
        
        # Кнопка вибору профілю з popup
        ctk.CTkButton(parse_input_frame, text="Вибрати профіль", command=self.show_profile_popup, 
                     width=120, height=32, corner_radius=6, font=self.font).grid(row=0, column=4, padx=5, pady=5)
        
        # Налаштовуємо розтягування колонки з email entry
        parse_input_frame.grid_columnconfigure(1, weight=1)

        # Кнопки дій
        button_frame = ctk.CTkFrame(self)
        button_frame.pack(pady=5, padx=10, fill="x")
        
        # Контейнер для центрування кнопок
        center_frame = ctk.CTkFrame(button_frame)
        center_frame.pack(expand=True)

        ctk.CTkButton(center_frame, text="Парсити з CSV", command=self.parse_from_csv, 
                     width=130, height=32, corner_radius=6, font=self.font).pack(side=tk.LEFT, padx=5, pady=10)
        ctk.CTkButton(center_frame, text="Перевірити email-и", command=self.on_check, 
                     width=150, height=32, corner_radius=6, font=self.font).pack(side=tk.LEFT, padx=5, pady=10)
        self.open_csv_button = ctk.CTkButton(center_frame, text="Відкрити CSV", command=self.open_csv_file, 
                                           state=tk.DISABLED, width=120, height=32, corner_radius=6, font=self.font)
        self.open_csv_button.pack(side=tk.LEFT, padx=5, pady=10)

        # Прогрес-бар
        self.progress_frame = ctk.CTkFrame(self)
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ctk.CTkProgressBar(self.progress_frame)
        self.progress_bar.pack(fill="x", padx=10, pady=5)
        self.progress_bar.set(0)
        self.progress_label = ctk.CTkLabel(self.progress_frame, text="", font=self.font)
        self.progress_label.pack(pady=2)
        # Спочатку ховаємо прогрес-бар
        # self.progress_frame.pack_forget()

        # Середня частина - результати перевірки
        result_frame = ctk.CTkFrame(self)
        result_frame.pack(padx=10, pady=5, fill="both", expand=True)

        # Лічильники
        counter_frame = ctk.CTkFrame(result_frame)
        counter_frame.pack(fill="x", pady=5, padx=5)

        self.live_count_label = ctk.CTkLabel(counter_frame, text="Live emails: 0", font=self.font, text_color="green")
        self.live_count_label.pack(side=tk.LEFT, padx=10, pady=5)

        self.die_count_label = ctk.CTkLabel(counter_frame, text="Die emails: 0", font=self.font, text_color="red")
        self.die_count_label.pack(side=tk.LEFT, padx=10, pady=5)

        self.total_count_label = ctk.CTkLabel(counter_frame, text="Total: 0", font=ctk.CTkFont(weight="bold"), text_color="white")
        self.total_count_label.pack(side=tk.RIGHT, padx=10, pady=5)

        # Результати Live/Die
        results_frame = ctk.CTkFrame(result_frame)
        results_frame.pack(fill="both", expand=True, padx=5, pady=5)

        # Ліва колонка - Live emails
        live_frame = ctk.CTkFrame(results_frame)
        live_frame.pack(side=tk.LEFT, fill="both", expand=True, padx=5, pady=5)

        # Заголовок і кнопка копіювання для Live emails
        live_header_frame = ctk.CTkFrame(live_frame)
        live_header_frame.pack(fill="x", padx=5, pady=5)

        ctk.CTkLabel(live_header_frame, text="Live emails:", font=self.font, text_color="green").pack(side=tk.LEFT, anchor="w", padx=5)
        ctk.CTkButton(live_header_frame, text="Копіювати", command=self.copy_live_emails, 
                     width=100, height=28, corner_radius=6, font=self.font).pack(side=tk.RIGHT, padx=5)

        self.live_textarea = ctk.CTkTextbox(live_frame, width=300, height=200, font=self.font)
        self.live_textarea.pack(fill="both", expand=True, padx=5, pady=5)

        # Права колонка - Die emails
        die_frame = ctk.CTkFrame(results_frame)
        die_frame.pack(side=tk.LEFT, fill="both", expand=True, padx=5, pady=5)

        # Заголовок і кнопка копіювання для Die emails
        die_header_frame = ctk.CTkFrame(die_frame)
        die_header_frame.pack(fill="x", padx=5, pady=5)

        ctk.CTkLabel(die_header_frame, text="Die emails:", font=self.font, text_color="red").pack(side=tk.LEFT, anchor="w", padx=5)
        ctk.CTkButton(die_header_frame, text="Копіювати", command=self.copy_die_emails, 
                     width=100, height=28, corner_radius=6, font=self.font).pack(side=tk.RIGHT, padx=5)

        self.die_textarea = ctk.CTkTextbox(die_frame, width=300, height=200, font=self.font)
        self.die_textarea.pack(fill="both", expand=True, padx=5, pady=5)

        # Ініціалізуємо змінну для зберігання результатів
        self.check_results = None

    def update_progress(self, progress, current, total):
        self.progress_bar.set(progress / 100.0)  # CTkProgressBar використовує значення від 0 до 1
        self.progress_label.configure(text=f"Перевірено: {current}/{total} ({int(progress)}%)")
        self.update_idletasks()

    def on_check(self):
        input_text = self.input_textarea.get("1.0", "end")  # CTkTextbox використовує "end" замість tk.END
        if not input_text.strip():
            messagebox.showwarning("Увага", "Введіть email-и для перевірки.")
            return

        # Отримуємо список email-ів
        accounts = parse_accounts(input_text)
        if not accounts:
            messagebox.showwarning("Увага", "Не вдалося розпізнати жодного email.")
            return

        emails = [acc['email'] for acc in accounts]

        # Показуємо прогрес-бар
        self.progress_frame.pack(pady=5, fill="x")
        self.progress_var.set(0)
        self.progress_label.configure(text="Підготовка до перевірки...")

        # Запускаємо перевірку в окремому потоці
        def check_thread():
            live, die = check_emails_api(emails, self.update_progress)

            # Оновлюємо UI в головному потоці
            self.after(0, lambda: self.update_results(live, die))

        Thread(target=check_thread).start()

    def update_results(self, live, die):
        print(f"\n{'='*80}")
        print(f"🖥️  update_results CALLED")
        print(f"   LIVE received: {len(live)} emails")
        print(f"   DIE received: {len(die)} emails")
        print(f"   LIVE sample: {live[:3] if len(live) > 0 else 'empty'}")
        print(f"   DIE sample: {die[:3] if len(die) > 0 else 'empty'}")
        print(f"{'='*80}\n")
        
        # Оновлюємо текстові поля з результатами  
        self.live_textarea.delete("1.0", "end")
        self.die_textarea.delete("1.0", "end")

        self.live_textarea.insert("1.0", "\n".join(live))
        self.die_textarea.insert("1.0", "\n".join(die))

        # Оновлюємо лічильники
        total = len(live) + len(die)
        self.live_count_label.configure(text=f"Live emails: {len(live)}")
        self.die_count_label.configure(text=f"Die emails: {len(die)}")
        self.total_count_label.configure(text=f"Total: {total}")

        # Ховаємо прогрес-бар
        self.progress_label.configure(text=f"Перевірку завершено. Live: {len(live)}, Die: {len(die)}, Total: {total}")

        # Активуємо кнопку відкриття CSV
        self.open_csv_button.configure(state=tk.NORMAL)

        # Зберігаємо результати для подальшого експорту
        self.check_results = (live, die)

        # Автоматично зберігаємо результати в CSV файл
        csv_path = get_csv_path()
        try:
            save_to_csv(csv_path, live, die)
            # Активуємо кнопку відкриття CSV
            self.open_csv_button.configure(state=tk.NORMAL)
        except Exception as e:
            print(f"Помилка при збереженні CSV: {str(e)}")



    def open_csv_file(self):
        csv_path = get_csv_path()
        if not os.path.exists(csv_path):
            messagebox.showwarning("Увага", "CSV файл ще не створено. Спочатку виконайте перевірку email-ів.")
            return

        # Відкриваємо CSV файл у стандартній програмі
        if os.name == 'nt':  # Windows
            os.startfile(csv_path)
        elif os.name == 'posix':  # macOS, Linux
            import subprocess
            subprocess.call(('open', csv_path)) if sys.platform == 'darwin' else subprocess.call(('xdg-open', csv_path))
        else:
            messagebox.showwarning("Увага", "Не вдалося відкрити CSV файл. Спробуйте відкрити його вручну.")

    def copy_live_emails(self):
        live_text = self.live_textarea.get("1.0", "end").strip()
        if live_text:
            safe_clipboard_operation("set", live_text)
            # messagebox.showinfo("Копіювання", "Live emails скопійовано в буфер обміну")

    def copy_die_emails(self):
        die_text = self.die_textarea.get("1.0", "end").strip()
        if die_text:
            safe_clipboard_operation("set", die_text)
            # messagebox.showinfo("Копіювання", "Die emails скопійовано в буфер обміну")

    def parse_from_csv(self):
        """Парсить пошти з CSV файлу Gmail Hacks, починаючи з певного profile name"""
        try:
            profile_name = self.parse_email_entry.get().strip()
            if not profile_name:
                print("⚠️ Введіть profile name для пошуку")
                return
            
            try:
                count = int(self.parse_count_entry.get().strip())
                if count <= 0:
                    count = 20
            except ValueError:
                count = 20
            
            # Шукаємо CSV файл Gmail Hacks
            gmail_hacks_config = get_config_path("gmail_hacks_config.json")
            if not os.path.exists(gmail_hacks_config):
                print("⚠️ Конфігурація Gmail Hacks не знайдена. Спочатку завантажте дані в Gmail Hacks tab.")
                return
            
            # Завантажуємо кешовані CSV дані
            with open(gmail_hacks_config, "r", encoding="utf-8") as f:
                config = json.load(f)
                csv_content = config.get("cached_csv_data", "")
            
            if not csv_content:
                print("⚠️ CSV дані не знайдені. Спочатку завантажте дані в Gmail Hacks tab.")
                return
            
            # Парсимо CSV
            import csv
            import io
            
            csv_reader = csv.DictReader(io.StringIO(csv_content))
            rows = list(csv_reader)
            
            # 🔍 DEBUG: Виводимо назви колонок
            if rows:
                print(f"📊 CSV колонки: {list(rows[0].keys())}")
            
            # Шукаємо початковий рядок по profile name
            # Підтримка різних варіантів назв колонок
            profile_columns = ["Название", "Назва", "Name", "Profile", "Профіль"]
            email_columns = ["Почта", "Email", "Login", "Mail", "E-mail"]
            password_columns = ["Пароль", "Password", "Pass", "Pwd"]
            
            start_index = -1
            for i, row in enumerate(rows):
                # Перевіряємо всі можливі варіанти колонок профілю
                for col in profile_columns:
                    if col in row and row.get(col, "").strip().lower() == profile_name.lower():
                        start_index = i
                        print(f"✅ Знайдено профіль '{profile_name}' на позиції {i} (колонка: {col})")
                        break
                if start_index != -1:
                    break
            
            if start_index == -1:
                available_profiles = []
                for col in profile_columns:
                    if col in rows[0]:
                        available_profiles = [row.get(col, "").strip() for row in rows[:5] if row.get(col, "").strip()]
                        break
                
                print(f"❌ Profile name '{profile_name}' не знайдено в CSV")
                if available_profiles:
                    print(f"Перші профілі:\n" + "\n".join(f"  • {p}" for p in available_profiles[:5]))
                return
            
            # Витягуємо пошти починаючи з знайденого рядка
            # 🎯 ФІЛЬТРУЄМО: тільки @gmail.com БЕЗ паролів
            emails = []
            end_index = min(start_index + count, len(rows))
            
            parsed_count = 0
            skipped_count = {"has_password": 0, "not_gmail": 0, "no_email": 0}
            
            for i in range(start_index, len(rows)):
                # Зупиняємось якщо зібрали достатньо
                if parsed_count >= count:
                    break
                    
                row = rows[i]
                
                # 🔍 Шукаємо email в різних можливих колонках
                email_field = ""
                for col in email_columns:
                    if col in row and row.get(col, "").strip():
                        email_field = row.get(col, "").strip()
                        break
                
                # 🔍 Шукаємо пароль в різних можливих колонках
                password_field = ""
                for col in password_columns:
                    if col in row and row.get(col, "").strip():
                        password_field = row.get(col, "").strip()
                        break
                
                # 🔧 ПАРСИМО: якщо email містить роздільник, розділяємо
                # Формати: email:password, email;password, email|password, email password
                if email_field:
                    # Пробуємо розділити по роздільниках
                    import re
                    # Розділяємо по : ; | пробіл таб
                    parts = re.split(r'[:;|\s\t]+', email_field, maxsplit=1)
                    
                    if len(parts) > 1:
                        # Є роздільник - беремо тільки email частину
                        email = parts[0].strip()
                        embedded_password = parts[1].strip()  # Пароль в email полі
                    else:
                        # Немає роздільника - це чистий email
                        email = email_field
                        embedded_password = ""
                    
                    # Пароль може бути або в окремому полі, або вбудований в email
                    has_password = bool(password_field) or bool(embedded_password)
                    
                    # ✅ Беремо тільки Gmail адреси (незалежно від паролів, але показуємо тільки email)
                    if email and "@gmail.com" in email.lower():
                        # Зберігаємо тільки чистий email (без пароля)
                        emails.append(email)
                        parsed_count += 1
                        if parsed_count <= 3:  # Debug перших 3
                            print(f"  ✅ {parsed_count}. {email}")
                    else:
                        # Debug статистика відхилення
                        if not email:
                            skipped_count["no_email"] += 1
                        elif "@gmail.com" not in email.lower():
                            skipped_count["not_gmail"] += 1
                else:
                    skipped_count["no_email"] += 1
            
            # 📊 Debug статистика
            print(f"\n📊 Статистика парсингу:")
            print(f"  Переглянуто рядків: {len(rows) - start_index}")
            print(f"  ✅ Знайдено Gmail адрес: {len(emails)}")
            print(f"  ❌ Відхилено:")
            print(f"     • Не Gmail домен: {skipped_count['not_gmail']}")
            print(f"     • Немає email: {skipped_count['no_email']}")
            
            if emails:
                # Вставляємо знайдені пошти в поле введення
                self.input_textarea.delete("0.0", "end")
                self.input_textarea.insert("0.0", "\n".join(emails))
                print(f"✅ Успіх! Знайдено {len(emails)} Gmail адрес починаючи з профілю '{profile_name}'")
            else:
                # Детальне повідомлення про проблему
                self.input_textarea.delete("0.0", "end")
                print(f"❌ Не знайдено жодної Gmail адреси починаючи з профілю '{profile_name}'")
                
        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            print(f"❌ Помилка парсингу:\n{error_details}")

    def show_profile_popup(self):
        """Показує popup вікно для вибору профілю з списку"""
        try:
            # Завантажуємо CSV дані з Gmail Hacks
            gmail_hacks_config = get_config_path("gmail_hacks_config.json")
            if not os.path.exists(gmail_hacks_config):
                print("Gmail Hacks конфігурація не знайдена")
                return
            
            with open(gmail_hacks_config, "r", encoding="utf-8") as f:
                config = json.load(f)
                csv_content = config.get("cached_csv_data", "")
            
            if not csv_content:
                print("CSV дані не знайдені")
                return
            
            # Парсимо CSV та отримуємо список профілів
            import csv
            import io
            
            csv_reader = csv.DictReader(io.StringIO(csv_content))
            profiles = []
            
            for row in csv_reader:
                profile_name = row.get("Название", "").strip()
                if profile_name:
                    profiles.append(profile_name)
            
            if not profiles:
                print("Профілі не знайдені в CSV")
                return
            
            # Створюємо popup вікно
            popup = ctk.CTkToplevel(self)
            popup.title("Вибір профілю")
            popup.geometry("500x400")
            popup.lift()  # Піднімаємо вікно на передній план
            popup.focus_force()  # Фокус на вікно
            popup.grab_set()  # Модальне вікно
            
            # Центруємо вікно відносно батьківського
            popup.transient(self)
            popup.after(100, lambda: popup.lift())
            
            # Заголовок
            ctk.CTkLabel(popup, text="Виберіть профіль для парсингу:", 
                        font=ctk.CTkFont(size=16, weight="bold")).pack(pady=20)
            
            # Поле пошуку в popup
            search_frame = ctk.CTkFrame(popup)
            search_frame.pack(fill="x", padx=20, pady=10)
            
            ctk.CTkLabel(search_frame, text="Пошук:", font=self.font).pack(anchor="w", padx=10, pady=(10,5))
            search_entry = ctk.CTkEntry(search_frame, placeholder_text="Введіть для пошуку профілю...", font=self.font)
            search_entry.pack(fill="x", padx=10, pady=5)
            
            # Скролюючий список профілів
            profiles_frame = ctk.CTkScrollableFrame(popup, label_text="Профілі:")
            profiles_frame.pack(fill="both", expand=True, padx=20, pady=10)
            
            # Змінна для зберігання кнопок
            profile_buttons = []
            
            def update_profile_list(filter_text=""):
                # Очищуємо попередні кнопки
                for widget in profiles_frame.winfo_children():
                    widget.destroy()
                profile_buttons.clear()
                
                # Фільтруємо профілі
                filtered_profiles = [p for p in profiles if filter_text.lower() in p.lower()]
                
                # Створюємо кнопки для відфільтрованих профілів
                for profile in filtered_profiles:
                    btn = ctk.CTkButton(profiles_frame, text=profile, 
                                       command=lambda p=profile: self.select_profile_from_popup(p, popup),
                                       width=440, height=35, font=self.font)
                    btn.pack(pady=2, fill="x")
                    profile_buttons.append(btn)
            
            # Bind для пошуку в реальному часі
            def on_search_change(*args):
                update_profile_list(search_entry.get())
            
            search_entry.bind('<KeyRelease>', on_search_change)
            
            # Початково показуємо всі профілі
            update_profile_list()
            
            # Кнопка закриття
            ctk.CTkButton(popup, text="Закрити", command=popup.destroy, 
                         width=100, height=30, font=self.font).pack(pady=20)
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка відкриття popup: {str(e)}")
    
    def select_profile_from_popup(self, profile_name, popup):
        """Вибирає профіль з popup та автоматично парсить дані"""
        try:
            # Встановлюємо вибраний профіль в поле
            self.parse_email_entry.delete(0, "end")
            self.parse_email_entry.insert(0, profile_name)
            
            # Закриваємо popup
            popup.destroy()
            
            # Автоматично парсимо дані
            self.parse_from_csv()
            
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка вибору профілю: {str(e)}")


# ================== Settings Tab ==================
class SettingsTab(ctk.CTkFrame):
    """Вкладка з налаштуваннями програми"""
    
    def __init__(self, master, app, font=None):
        super().__init__(master)
        self.app = app  # Посилання на головний App
        self.font = font or ctk.CTkFont(family="Segoe UI", size=11)
        self.pack(fill="both", expand=True)
        
        # Створюємо скролювальну область
        main_scroll = ctk.CTkScrollableFrame(self)
        main_scroll.pack(fill="both", expand=True, padx=20, pady=20)
        
        # Заголовок
        title_label = ctk.CTkLabel(main_scroll, text="⚙️ Налаштування програми", 
                                   font=ctk.CTkFont(size=24, weight="bold"))
        title_label.pack(pady=(0, 30))
        
        # === СЕКЦІЯ 1: Управління вкладками ===
        tabs_section = ctk.CTkFrame(main_scroll)
        tabs_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(tabs_section, text="📑 Управління вкладками", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        ctk.CTkLabel(tabs_section, text="Налаштуйте видимість та порядок вкладок", 
                    font=self.font, text_color="gray").pack(anchor="w", padx=20, pady=(0, 10))
        
        tabs_btn_frame = ctk.CTkFrame(tabs_section)
        tabs_btn_frame.pack(padx=20, pady=(0, 15))
        
        ctk.CTkButton(tabs_btn_frame, text="📑 Налаштування вкладок", 
                     command=self.app.tab_manager.show_tab_settings,
                     width=250, height=40, corner_radius=8,
                     font=ctk.CTkFont(size=13, weight="bold")).pack(pady=5)
        
        # === СЕКЦІЯ 2: Інтерфейс ===
        ui_section = ctk.CTkFrame(main_scroll)
        ui_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(ui_section, text="🎨 Налаштування інтерфейсу", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        ctk.CTkLabel(ui_section, text="Змініть тему, швидкість скролу та інші параметри UI", 
                    font=self.font, text_color="gray").pack(anchor="w", padx=20, pady=(0, 10))
        
        ui_btn_frame = ctk.CTkFrame(ui_section)
        ui_btn_frame.pack(padx=20, pady=(0, 15))
        
        ctk.CTkButton(ui_btn_frame, text="🎨 UI Settings", 
                     command=self.app.open_ui_settings,
                     width=250, height=40, corner_radius=8,
                     font=ctk.CTkFont(size=13, weight="bold")).pack(pady=5)
        
        # === СЕКЦІЯ 3: Гарячі клавіші ===
        hotkeys_section = ctk.CTkFrame(main_scroll)
        hotkeys_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(hotkeys_section, text="⌨️ Гарячі клавіші", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        ctk.CTkLabel(hotkeys_section, text="Налаштуйте комбінації клавіш для швидких дій", 
                    font=self.font, text_color="gray").pack(anchor="w", padx=20, pady=(0, 10))
        
        hotkeys_btn_frame = ctk.CTkFrame(hotkeys_section)
        hotkeys_btn_frame.pack(padx=20, pady=(0, 15))
        
        ctk.CTkButton(hotkeys_btn_frame, text="⌨️ Налаштування Hotkeys", 
                     command=self.app.open_hotkeys_settings,
                     width=250, height=40, corner_radius=8,
                     font=ctk.CTkFont(size=13, weight="bold")).pack(pady=5)
        
        # === СЕКЦІЯ 3.5: Парсинг таблиць ===
        sheets_parsing_section = ctk.CTkFrame(main_scroll)
        sheets_parsing_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(sheets_parsing_section, text="📊 Парсинг Google Sheets", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        ctk.CTkLabel(sheets_parsing_section, text="Налаштуйте назви колонок для автоматичного парсингу таблиць", 
                    font=self.font, text_color="gray").pack(anchor="w", padx=20, pady=(0, 10))
        
        sheets_btn_frame = ctk.CTkFrame(sheets_parsing_section)
        sheets_btn_frame.pack(padx=20, pady=(0, 15))
        
        ctk.CTkButton(sheets_btn_frame, text="📊 Налаштування колонок", 
                     command=self.app.open_sheets_parsing_settings,
                     width=250, height=40, corner_radius=8,
                     font=ctk.CTkFont(size=13, weight="bold")).pack(pady=5)
        
        # === СЕКЦІЯ 4: Оновлення програми ===
        update_section = ctk.CTkFrame(main_scroll)
        update_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(update_section, text="🔄 Оновлення програми", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        ctk.CTkLabel(update_section, text="Перевірте наявність оновлень та встановіть нову версію", 
                    font=self.font, text_color="gray").pack(anchor="w", padx=20, pady=(0, 10))
        
        update_btn_frame = ctk.CTkFrame(update_section)
        update_btn_frame.pack(padx=20, pady=(0, 15))
        
        ctk.CTkButton(update_btn_frame, text="🔄 Перевірити оновлення", 
                     command=self.app.check_for_updates,
                     width=250, height=40, corner_radius=8,
                     font=ctk.CTkFont(size=13, weight="bold"),
                     fg_color="#4CAF50", hover_color="#45a049").pack(pady=5)
        
        # === СЕКЦІЯ 5: Конфігурація ===
        config_section = ctk.CTkFrame(main_scroll)
        config_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(config_section, text="📁 Файли конфігурації", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        ctk.CTkLabel(config_section, text="Перегляньте шляхи до файлів конфігурації та CSV", 
                    font=self.font, text_color="gray").pack(anchor="w", padx=20, pady=(0, 10))
        
        config_btn_frame = ctk.CTkFrame(config_section)
        config_btn_frame.pack(padx=20, pady=(0, 15))
        
        ctk.CTkButton(config_btn_frame, text="📁 Показати конфігурацію", 
                     command=self.app.show_config_info,
                     width=250, height=40, corner_radius=8,
                     font=ctk.CTkFont(size=13, weight="bold")).pack(pady=5)
        
        # === СЕКЦІЯ 6: Про програму ===
        about_section = ctk.CTkFrame(main_scroll)
        about_section.pack(fill="x", pady=10)
        
        ctk.CTkLabel(about_section, text="ℹ️ Про програму", 
                    font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(15, 10))
        
        version_label = ctk.CTkLabel(about_section, 
                                     text="Punch IT Now 9.2\nВерсія: 9.2.0\nАвтор: @Alex_FarmPunch", 
                                     font=self.font, justify="left")
        version_label.pack(anchor="w", padx=20, pady=(0, 10))
        
        telegram_btn = ctk.CTkButton(about_section, text="💬 Відкрити Telegram", 
                                     command=self.app.open_telegram_link,
                                     width=250, height=40, corner_radius=8,
                                     font=ctk.CTkFont(size=13, weight="bold"),
                                     fg_color="#0088cc", hover_color="#006699")
        telegram_btn.pack(padx=20, pady=(0, 15))


# ================== ChatGPT Integration Tab ==================
class ChatGPTTab(ctk.CTkFrame):
    def __init__(self, master, font=None):
        super().__init__(master)
        self.font = font or ctk.CTkFont(family="Segoe UI", size=11)
        self.pack(fill="both", expand=True)
        
        # API ключ та конфігурація
        self.config_file = get_config_path("chatgpt_config.json")
        self.api_key = ""
        self.conversation_history = []
        
        # Завантажуємо збережені налаштування
        self.load_config()
        
        # Створюємо UI
        self.create_widgets()
    
    def create_widgets(self):
        """Створює інтерфейс для ChatGPT інтеграції"""
        main_frame = ctk.CTkFrame(self)
        main_frame.pack(fill="both", expand=True, padx=10, pady=10)
        
        # Заголовок
        title_label = ctk.CTkLabel(main_frame, text="🤖 ChatGPT Integration", 
                                 font=ctk.CTkFont(size=18, weight="bold"))
        title_label.pack(pady=(10, 20))
        
        # API Key секція
        api_frame = ctk.CTkFrame(main_frame)
        api_frame.pack(fill="x", padx=10, pady=10)
        
        ctk.CTkLabel(api_frame, text="OpenAI API Settings", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(10, 5))
        
        # Контейнер для API ключа
        api_container = ctk.CTkFrame(api_frame)
        api_container.pack(fill="x", padx=10, pady=10)
        
        ctk.CTkLabel(api_container, text="API Key:", font=self.font).pack(anchor="w", padx=5)
        
        api_input_frame = ctk.CTkFrame(api_container)
        api_input_frame.pack(fill="x", padx=5, pady=5)
        
        self.api_entry = ctk.CTkEntry(api_input_frame, width=400, show="*", 
                                     placeholder_text="sk-proj-...", font=self.font)
        self.api_entry.pack(side="left", fill="x", expand=True, padx=(0, 5))
        # Вставляємо збережений API ключ якщо є
        if self.api_key:
            self.api_entry.insert(0, self.api_key)
        
        save_btn = ctk.CTkButton(api_input_frame, text="💾 Save", 
                               command=self.save_api_key, width=80, height=32, 
                               corner_radius=6, font=self.font)
        save_btn.pack(side="right")
        
        test_btn = ctk.CTkButton(api_input_frame, text="🔧 Test", 
                               command=self.test_api_key, width=80, height=32, 
                               corner_radius=6, font=self.font)
        test_btn.pack(side="right", padx=(0, 5))
        
        # Статус API
        self.api_status = ctk.CTkLabel(api_container, text="API Status: Not configured", 
                                     font=self.font, text_color="gray")
        self.api_status.pack(anchor="w", padx=5, pady=(0, 5))
        
        # Chat секція
        chat_frame = ctk.CTkFrame(main_frame)
        chat_frame.pack(fill="both", expand=True, padx=10, pady=10)
        
        ctk.CTkLabel(chat_frame, text="💬 Chat Interface", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(10, 5))
        
        # Історія розмови
        self.chat_history = ctk.CTkTextbox(chat_frame, height=300, font=self.font, 
                                         wrap="word", state="disabled")
        self.chat_history.pack(fill="both", expand=True, padx=10, pady=(5, 10))
        
        # Поле для вводу повідомлення
        input_frame = ctk.CTkFrame(chat_frame)
        input_frame.pack(fill="x", padx=10, pady=(0, 10))
        
        self.message_entry = ctk.CTkTextbox(input_frame, height=80, font=self.font, 
                                          wrap="word")
        self.message_entry.pack(fill="both", expand=True, padx=(0, 5))
        
        # Кнопки управління
        button_frame = ctk.CTkFrame(input_frame)
        button_frame.pack(fill="x", pady=5)
        
        send_btn = ctk.CTkButton(button_frame, text="📤 Send Message", 
                               command=self.send_message, width=120, height=35, 
                               corner_radius=6, font=self.font)
        send_btn.pack(side="left", padx=5)
        
        clear_btn = ctk.CTkButton(button_frame, text="🗑️ Clear Chat", 
                                command=self.clear_chat, width=100, height=35, 
                                corner_radius=6, font=self.font)
        clear_btn.pack(side="left", padx=5)
        
        copy_btn = ctk.CTkButton(button_frame, text="📋 Copy Last Response", 
                               command=self.copy_last_response, width=140, height=35, 
                               corner_radius=6, font=self.font)
        copy_btn.pack(side="left", padx=5)
        
        # Шаблони швидких запитань
        templates_frame = ctk.CTkFrame(main_frame)
        templates_frame.pack(fill="x", padx=10, pady=10)
        
        ctk.CTkLabel(templates_frame, text="⚡ Quick Templates", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(10, 5))
        
        templates_container = ctk.CTkFrame(templates_frame)
        templates_container.pack(fill="x", padx=10, pady=(5, 10))
        
        # Шаблони кнопок
        templates = [
            ("💡 Generate Ideas", "Generate 5 creative ideas for:"),
            ("✍️ Improve Text", "Please improve this text and make it more professional:"),
            ("🔍 Analyze", "Please analyze this content and provide insights:"),
            ("📝 Summarize", "Please summarize the following content:"),
            ("🌐 Translate", "Please translate this to English:"),
            ("🐛 Debug Code", "Please help debug this code and explain the issue:")
        ]
        
        for i, (btn_text, template) in enumerate(templates):
            row = i // 3
            col = i % 3
            
            btn = ctk.CTkButton(templates_container, text=btn_text, 
                              command=lambda t=template: self.insert_template(t),
                              width=180, height=32, corner_radius=6, font=self.font)
            btn.grid(row=row, column=col, padx=5, pady=5, sticky="ew")
        
        # Конфігуруємо колонки для рівномірного розподілу
        for i in range(3):
            templates_container.grid_columnconfigure(i, weight=1)
        
        # Bind Enter для відправки повідомлення
        self.message_entry.bind("<Control-Return>", lambda e: self.send_message())
        
        # Показуємо початкове повідомлення
        self.display_message("System", "👋 Привіт! Я ChatGPT асистент. Налаштуйте API ключ та почніть спілкування!")
    
    def load_config(self):
        """Завантажує збережені налаштування"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, "r", encoding="utf-8") as f:
                    config = json.load(f)
                    self.api_key = config.get("api_key", "")
        except Exception as e:
            print(f"Помилка завантаження конфігурації ChatGPT: {str(e)}")
    
    def save_config(self):
        """Зберігає налаштування"""
        try:
            config = {
                "api_key": self.api_key
            }
            with open(self.config_file, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Помилка збереження конфігурації ChatGPT: {str(e)}")
    
    def save_api_key(self):
        """Зберігає API ключ"""
        self.api_key = self.api_entry.get().strip()
        self.save_config()
        self.api_status.configure(text="API Status: Saved", text_color="green")
        self.display_message("System", "✅ API ключ збережено!")
    
    def test_api_key(self):
        """Тестує API ключ"""
        if not HAS_OPENAI():
            self.api_status.configure(text="API Status: OpenAI library not installed", text_color="red")
            self.display_message("System", "❌ OpenAI бібліотека не встановлена. Встановіть: pip install openai")
            return
        
        api_key = self.api_entry.get().strip()
        if not api_key:
            self.api_status.configure(text="API Status: No API key", text_color="red")
            self.display_message("System", "❌ Введіть API ключ")
            return
        
        try:
            self.api_status.configure(text="API Status: Testing...", text_color="orange")
            self.update()
            
            # Тестуємо API
            client = openai.OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": "Hello, test message"}],
                max_tokens=50
            )
            
            self.api_status.configure(text="API Status: ✅ Working", text_color="green")
            self.display_message("System", "✅ API ключ працює! Готово до використання.")
            
        except Exception as e:
            self.api_status.configure(text="API Status: ❌ Error", text_color="red")
            self.display_message("System", f"❌ Помилка тестування API: {str(e)}")
    
    def send_message(self):
        """Відправляє повідомлення до ChatGPT"""
        if not HAS_OPENAI():
            self.display_message("System", "❌ OpenAI бібліотека не встановлена")
            return
        
        message = self.message_entry.get("0.0", "end-1c").strip()
        if not message:
            return
        
        api_key = self.api_entry.get().strip()
        if not api_key:
            self.display_message("System", "❌ Налаштуйте API ключ")
            return
        
        # Показуємо повідомлення користувача
        self.display_message("You", message)
        
        # Очищаємо поле вводу
        self.message_entry.delete("0.0", "end")
        
        # Додаємо до історії
        self.conversation_history.append({"role": "user", "content": message})
        
        # Показуємо індикатор завантаження
        self.display_message("ChatGPT", "🤔 Думаю...")
        
        # Відправляємо запит в окремому потоці
        Thread(target=self._send_request_thread, args=(api_key,), daemon=True).start()
    
    def _send_request_thread(self, api_key):
        """Відправляє запит до ChatGPT в окремому потоці"""
        try:
            client = openai.OpenAI(api_key=api_key)
            
            # Обмежуємо історію до останніх 10 повідомлень для економії токенів
            recent_history = self.conversation_history[-10:]
            
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=recent_history,
                max_tokens=1000,
                temperature=0.7
            )
            
            reply = response.choices[0].message.content
            
            # Додаємо відповідь до історії
            self.conversation_history.append({"role": "assistant", "content": reply})
            
            # Показуємо відповідь в основному потоці
            self.after(0, lambda: self._update_chat_with_response(reply))
            
        except Exception as e:
            error_msg = f"❌ Помилка: {str(e)}"
            self.after(0, lambda: self._update_chat_with_response(error_msg))
    
    def _update_chat_with_response(self, response):
        """Оновлює чат з відповіддю ChatGPT"""
        # Видаляємо індикатор завантаження
        self.chat_history.configure(state="normal")
        content = self.chat_history.get("0.0", "end")
        lines = content.split("\n")
        
        # Знаходимо та видаляємо останній рядок з "Думаю..."
        for i in range(len(lines)-1, -1, -1):
            if "🤔 Думаю..." in lines[i]:
                # Видаляємо цей рядок
                start_pos = "\n".join(lines[:i]) + ("\n" if i > 0 else "")
                self.chat_history.delete("0.0", "end")
                self.chat_history.insert("0.0", start_pos)
                break
        
        self.chat_history.configure(state="disabled")
        
        # Додаємо реальну відповідь
        self.display_message("ChatGPT", response)
        self.last_response = response
    
    def display_message(self, sender, message):
        """Відображає повідомлення в чаті"""
        self.chat_history.configure(state="normal")
        
        # Форматування відправника
        if sender == "You":
            prefix = "👤 You: "
            color = "lightblue"
        elif sender == "ChatGPT":
            prefix = "🤖 ChatGPT: "
            color = "lightgreen"
        else:
            prefix = "🔧 System: "
            color = "orange"
        
        # Додаємо повідомлення
        current_content = self.chat_history.get("0.0", "end-1c")
        if current_content.strip():
            self.chat_history.insert("end", "\n\n")
        
        self.chat_history.insert("end", f"{prefix}{message}")
        
        # Прокручуємо вниз
        self.chat_history.see("end")
        self.chat_history.configure(state="disabled")
    
    def clear_chat(self):
        """Очищує історію чату"""
        self.chat_history.configure(state="normal")
        self.chat_history.delete("0.0", "end")
        self.chat_history.configure(state="disabled")
        
        self.conversation_history = []
        self.display_message("System", "🗑️ Чат очищено")
    
    def copy_last_response(self):
        """Копіює останню відповідь ChatGPT"""
        if hasattr(self, 'last_response'):
            safe_clipboard_operation("set", self.last_response)
            self.display_message("System", "📋 Останню відповідь скопійовано!")
        else:
            self.display_message("System", "❌ Немає відповіді для копіювання")
    
    def insert_template(self, template):
        """Вставляє шаблон в поле вводу"""
        current = self.message_entry.get("0.0", "end-1c")
        if current.strip():
            self.message_entry.insert("end", f"\n\n{template} ")
        else:
            self.message_entry.insert("0.0", f"{template} ")
        
        self.message_entry.focus()


class App(ctk.CTk):
    def setup_global_hotkeys(self):
        """ВИМКНЕНО: Глобальні гарячі клавіші викликали рекурсивні помилки"""
        
        # ЗАКОМЕНТОВАНО для стабільності програми
        # Використовуйте стандартні системні hotkeys (Ctrl+C, Ctrl+V тощо)
        
        print("⚠️ Глобальні гарячі клавіші вимкнено для стабільності")
        print("✅ Використовуйте стандартні системні hotkeys")
        
        return  # Функція вимкнена
        
        # Старий код нижче закоментовано для стабільності:
        #
        # hotkey_mappings = [
        #     ('<Control-c>', self.global_copy),
        #     ('<Control-v>', self.global_paste), 
        #     ('<Control-a>', self.global_select_all),
        #     ('<Control-s>', self.global_save),
        #     ('<Control-o>', self.global_open),
        #     ('<Control-n>', self.global_new),
        #     ('<Control-f>', self.global_find),
        #     ('<Control-z>', self.global_undo),
        #     ('<Control-g>', self.global_generate_all),
        #     ('<Control-b>', self.global_copy_combined),
        # ]
        # 
        # self.bind_all('<KeyPress>', self.handle_keypress)
        # 
        # for hotkey, callback in hotkey_mappings:
        #     try:
        #         self.bind_all(hotkey, callback)
        #     except Exception as e:
        #         print(f"✗ Помилка реєстрації {hotkey}: {e}")
        # 
        # print("✅ Глобальні гарячі клавіші налаштовано для всіх розкладок")
    
    def setup_tab_hotkeys(self):
        """Налаштування hotkey для швидкого переключення вкладок та Gmail Hacks функцій"""
        try:
            # Локальні hotkeys як fallback (коли глобальні не працюють)
            self.bind_all('<F1>', self.switch_to_generators)
            self.bind_all('<F2>', self.switch_to_gmail_hacks)
            self.bind_all('<F4>', self.copy_2fa_code)
            self.bind_all('<Escape>', self.force_restore_window)
            
            # Локальні Shift комбінації як fallback
            self.bind_all('<Shift-Key-1>', lambda e: self.select_gmail_section(1))
            self.bind_all('<Shift-Key-2>', lambda e: self.select_gmail_section(2))
            self.bind_all('<Shift-Key-3>', lambda e: self.select_gmail_section(3))
            self.bind_all('<Shift-Key-4>', lambda e: self.select_gmail_section(4))
            
            # Налаштовуємо ВСІ глобальні system-wide hotkeys
            self.setup_global_f3_hotkey()
            
            print("⚡ Налаштовано hotkeys:")
            print("   🌍 ГЛОБАЛЬНІ (працюють завжди): F1-F4 + Shift+1-4")
            print("   📱 ЛОКАЛЬНІ (fallback): bind_all для програми")
            print("   💡 Hotkeys працюють незалежно від фокуса віджетів!")
        except Exception as e:
            print(f"❌ Помилка налаштування hotkeys: {e}")
    
    def setup_global_f3_hotkey(self):
        """Налаштовує всі глобальні system-wide hotkeys"""
        try:
            import keyboard
            import threading
            
            # Перевіряємо чи вже є глобальні hotkeys
            if hasattr(self, '_global_hotkeys_registered') and self._global_hotkeys_registered:
                return
            
            # Завантажуємо custom hotkeys
            hotkeys_config = self.load_hotkeys_config()
            minimize_key = hotkeys_config.get("minimize_restore", "f3")
            
            print(f"⌨️ Завантажено custom hotkeys:")
            print(f"   🪟 Minimize/Restore: {minimize_key}")
            
            def global_f1_handler():
                """F1 → Generators"""
                try:
                    if self.winfo_exists():
                        self.after(0, self.switch_to_generators)
                except Exception as e:
                    print(f"❌ Помилка глобального F1: {e}")
                    
            def global_f2_handler():
                """F2 → Gmail Hacks"""
                try:
                    if self.winfo_exists():
                        self.after(0, self.switch_to_gmail_hacks)
                except Exception as e:
                    print(f"❌ Помилка глобального F2: {e}")
                    
            def global_minimize_handler():
                """Custom Key → Minimize/Restore"""
                try:
                    if self.winfo_exists():
                        self.after(0, lambda: self.toggle_window_state())
                except Exception as e:
                    print(f"❌ Помилка глобального {minimize_key}: {e}")
                    
            def global_f4_handler():
                """F4 → Copy 2FA"""
                try:
                    if self.winfo_exists():
                        self.after(0, self.copy_2fa_code)
                except Exception as e:
                    print(f"❌ Помилка глобального F4: {e}")
                    
            def global_shift_1_handler():
                """Shift+1 → Gmail Section 1"""
                try:
                    if self.winfo_exists():
                        self.after(0, lambda: self.select_gmail_section(1))
                except Exception as e:
                    print(f"❌ Помилка глобального Shift+1: {e}")
                    
            def global_shift_2_handler():
                """Shift+2 → Gmail Section 2"""
                try:
                    if self.winfo_exists():
                        self.after(0, lambda: self.select_gmail_section(2))
                except Exception as e:
                    print(f"❌ Помилка глобального Shift+2: {e}")
                    
            def global_shift_3_handler():
                """Shift+3 → Gmail Section 3"""
                try:
                    if self.winfo_exists():
                        self.after(0, lambda: self.select_gmail_section(3))
                except Exception as e:
                    print(f"❌ Помилка глобального Shift+3: {e}")
                    
            def global_shift_4_handler():
                """Shift+4 → Gmail Section 4"""
                try:
                    if self.winfo_exists():
                        self.after(0, lambda: self.select_gmail_section(4))
                except Exception as e:
                    print(f"❌ Помилка глобального Shift+4: {e}")
            
            # 🌍 Реєструємо всі глобальні hotkeys (cross-platform)
            def register_global_hotkeys():
                try:
                    # Базові функціональні клавіші (універсальні)
                    keyboard.add_hotkey('f1', global_f1_handler, suppress=False)
                    keyboard.add_hotkey('f2', global_f2_handler, suppress=False)
                    keyboard.add_hotkey(minimize_key, global_minimize_handler, suppress=False)  # Custom hotkey
                    keyboard.add_hotkey('f4', global_f4_handler, suppress=False)
                    
                    # Платформо-специфічні комбінації
                    if IS_MACOS:
                        # macOS використовує Command замість Shift для системних hotkeys
                        keyboard.add_hotkey('cmd+1', global_shift_1_handler, suppress=False)
                        keyboard.add_hotkey('cmd+2', global_shift_2_handler, suppress=False)
                        keyboard.add_hotkey('cmd+3', global_shift_3_handler, suppress=False)
                        keyboard.add_hotkey('cmd+4', global_shift_4_handler, suppress=False)
                        # Також додаємо Shift варіанти для сумісності
                        keyboard.add_hotkey('shift+1', global_shift_1_handler, suppress=False)
                        keyboard.add_hotkey('shift+2', global_shift_2_handler, suppress=False)
                        keyboard.add_hotkey('shift+3', global_shift_3_handler, suppress=False)
                        keyboard.add_hotkey('shift+4', global_shift_4_handler, suppress=False)
                        hotkey_info = "F1-F4 + Cmd+1-4 + Shift+1-4 (macOS)"
                    else:
                        # Windows/Linux використовують Shift
                        keyboard.add_hotkey('shift+1', global_shift_1_handler, suppress=False)
                        keyboard.add_hotkey('shift+2', global_shift_2_handler, suppress=False)
                        keyboard.add_hotkey('shift+3', global_shift_3_handler, suppress=False)
                        keyboard.add_hotkey('shift+4', global_shift_4_handler, suppress=False)
                        hotkey_info = "F1-F4 + Shift+1-4"
                    
                    self._global_hotkeys_registered = True
                    print(f"⚡ Всі глобальні hotkeys активовано ({PLATFORM}):")
                    print(f"   {hotkey_info} працюють незалежно від фокуса!")
                except Exception as e:
                    print(f"❌ Не вдалося зареєструвати глобальні hotkeys: {e}")
            
            # Запускаємо реєстрацію в окремому thread'і
            hotkey_thread = threading.Thread(target=register_global_hotkeys, daemon=True)
            hotkey_thread.start()
            
        except ImportError:
            print("⚠️ Бібліотека keyboard недоступна - глобальні hotkeys не активовані")
        except Exception as e:
            print(f"❌ Помилка налаштування глобальних hotkeys: {e}")
    

    
    def switch_to_generators(self, event=None):
        """Перемикає на вкладку Generators (F1)"""
        try:
            if hasattr(self, 'tabview') and self.tabview:
                if "Generators" in self.tabview._tab_dict:
                    self.tabview.set("Generators")
                    print("⚡ F1 → Перемкнуто на Generators")
                    return "break"  # Зупинити подальшу обробку
        except Exception as e:
            print(f"❌ Помилка переключення на Generators: {e}")
        return None
    
    def switch_to_gmail_hacks(self, event=None):  
        """Перемикає на вкладку Gmail Hacks (F2)"""
        try:
            if hasattr(self, 'tabview') and self.tabview:
                if "Gmail Hacks" in self.tabview._tab_dict:
                    self.tabview.set("Gmail Hacks")
                    print("⚡ F2 → Перемкнуто на Gmail Hacks")
                    return "break"  # Зупинити подальшу обробку
        except Exception as e:
            print(f"❌ Помилка переключення на Gmail Hacks: {e}")
        return None
    
    def toggle_window_state(self, event=None):
        """Згортає вікно в таскбар або розгортає назад (F3) - працює глобально"""
        try:
            # Перевіряємо чи не було недавнього F3 (захист від подвійного спрацювання)
            current_time = time.time()
            if hasattr(self, '_last_f3_time') and (current_time - self._last_f3_time) < 0.5:
                print("🛡️ F3 ігнорується - занадто швидке натискання")
                return "break" if event else None
            
            self._last_f3_time = current_time
            current_state = self.state()
            print(f"� Поточний стан вікна: {current_state}")
            
            # Перевіряємо флаг мінімізації
            minimized_by_f3 = getattr(self, '_minimized_by_f3', False)
            
            if current_state == 'iconic' or minimized_by_f3:
                # Вікно згорнуто - розгортаємо
                self.restore_window()
                print("⚡ F3 → Вікно розгорнуто з таскбару")
            elif current_state in ['normal', 'zoomed']:
                # Вікно видиме - згортаємо
                self._minimized_by_f3 = True
                self.iconify()
                print("⚡ F3 → Вікно згорнуто в таскбар")
            elif current_state == 'withdrawn':
                # Вікно приховано - розгортаємо
                self.restore_window() 
                print("⚡ F3 → Вікно відновлено з прихованого стану")
            else:
                # Fallback - розгортаємо
                self.restore_window()
                print("⚡ F3 → Вікно розгорнуто (fallback)")
                
            return "break" if event else None
        except Exception as e:
            print(f"❌ Помилка згортання/розгортання: {e}")
        return None
    
    def restore_window(self):
        """Розгортає вікно з правильним фокусом"""
        try:
            print(f"🔄 Розгортаємо вікно зі стану: {self.state()}")
            
            # Розгортаємо вікно
            self.deiconify()
            
            # Даємо час вікну з'явитися
            self.after(50, self._complete_restore)
            
        except Exception as e:
            print(f"❌ Помилка розгортання: {e}")
    
    def _complete_restore(self):
        """Завершує розгортання вікна"""
        try:
            # Піднімаємо на передній план та даємо фокус
            self.lift()
            self.focus_force()
            
            # Тимчасово робимо topmost для гарантії що вікно з'явиться
            try:
                self.attributes('-topmost', True)
                self.after(100, lambda: self.attributes('-topmost', False))
            except:
                pass  # Не критично якщо topmost не працює
                
            # Скидаємо флаг мінімізації
            self._minimized_by_f3 = False
            print("✅ Вікно успішно розгорнуто")
            
        except Exception as e:
            print(f"❌ Помилка завершення розгортання: {e}")
    
    def schedule_restore_check(self):
        """Створює періодичну перевірку для можливості розгортання"""
        def check_and_restore():
            try:
                if hasattr(self, '_minimized_by_f3') and self._minimized_by_f3:
                    if self.state() == 'iconic':
                        # Вікно все ще згорнуто, перевіряємо знову через 1 секунду
                        self.after(1000, check_and_restore)
                    else:
                        # Вікно розгорнуто іншим способом, скидаємо флаг
                        self._minimized_by_f3 = False
            except:
                pass
        
        # Початкова перевірка через 100мс
        self.after(100, check_and_restore)
    
    def force_restore_window(self, event=None):
        """Примусово розгортає вікно (Escape)"""
        try:
            current_state = self.state()
            if current_state == 'iconic':
                self.restore_window()
                print("⚡ Escape → Вікно примусово розгорнуто")
                return "break"
        except Exception as e:
            print(f"❌ Помилка примусового розгортання: {e}")
        return None
    
    def enable_restore_hotkey(self):
        """Дозволяє розгорнути вікно через F3 навіть коли воно згорнуто"""
        try:
            if self.state() == 'iconic':
                # Вікно все ще згорнуто, повторюємо через секунду
                self.after(1000, self.enable_restore_hotkey)
        except:
            pass
    
    def copy_2fa_code(self, event=None):
        """Генерує та копіює 2FA код в буфер (F4) - використовує існуючу кнопку 🔑"""
        try:
            print("🔑 F4 натиснуто - спроба генерації 2FA коду")
            
            # Автоматично переключаємося на Gmail Hacks якщо не там
            if hasattr(self, 'tabview') and self.tabview:
                current_tab = self.tabview.get()
                if current_tab != "Gmail Hacks":
                    print("📋 Переключаємось на Gmail Hacks для F4")
                    self.tabview.set("Gmail Hacks")
                    # Даємо час для переключення
                    self.after(100, lambda: self._execute_f4_action())
                    return "break"
                else:
                    return self._execute_f4_action()
            else:
                print("❌ Tabview не знайдено")
        except Exception as e:
            print(f"❌ Помилка F4: {e}")
        return None
    
    def _execute_f4_action(self):
        """Виконує F4 дію після переключення на Gmail Hacks"""
        try:
            # Шукаємо Gmail Hacks instance та викликаємо generate_2fa_code()
            gmail_hacks_instance = None
            for tab_name, tab_frame in self.tabview._tab_dict.items():
                if tab_name == "Gmail Hacks":
                    # Шукаємо Gmail Hacks instance
                    gmail_hacks_instance = self.find_gmail_hacks_instance(tab_frame)
                    if gmail_hacks_instance:
                        print("✅ Gmail Hacks instance знайдено, генеруємо 2FA")
                        # Викликаємо існуючу функцію генерації 2FA
                        gmail_hacks_instance.generate_2fa_code()
                        return "break"
                    break
                    
            print("⚠️ Gmail Hacks instance не знайдено")
            return None
        except Exception as e:
            print(f"❌ Помилка виконання F4: {e}")
            return None
    
    def find_gmail_hacks_instance(self, parent_widget):
        """Рекурсивно шукає Gmail Hacks instance з функцією generate_2fa_code"""
        try:
            # Перевіряємо чи поточний віджет має функцію generate_2fa_code
            if hasattr(parent_widget, 'generate_2fa_code') and hasattr(parent_widget, 'entries_dict'):
                return parent_widget
            
            # Рекурсивний пошук в дочірніх елементах
            for child in parent_widget.winfo_children():
                if hasattr(child, 'generate_2fa_code') and hasattr(child, 'entries_dict'):
                    return child
                elif hasattr(child, 'winfo_children'):
                    result = self.find_gmail_hacks_instance(child)
                    if result:
                        return result
            return None
        except Exception as e:
            print(f"❌ Помилка пошуку Gmail Hacks instance: {e}")
            return None
    
    def select_gmail_section(self, section_num):
        """Натискає відповідну кнопку швидкого вибору 1-4 (Shift+1-4)"""
        try:
            print(f"🔥 Shift+{section_num} натиснуто! (фокус: {self.focus_get()})")
            
            # Переключаємося на Gmail Hacks якщо не там
            if hasattr(self, 'tabview') and self.tabview:
                current_tab = self.tabview.get()
                if current_tab != "Gmail Hacks":
                    self.tabview.set("Gmail Hacks")
                    print(f"📋 Переключено з '{current_tab}' на Gmail Hacks")
                    # Даємо час для переключення
                    self.after(100, lambda: self._execute_gmail_section_selection(section_num))
                else:
                    self._execute_gmail_section_selection(section_num)
            else:
                print("❌ Tabview не знайдено")
        except Exception as e:
            print(f"❌ Помилка вибору кнопки {section_num}: {e}")
        return "break"  # Завжди повертаємо break щоб запобігти подальшій обробці
    
    def _execute_gmail_section_selection(self, section_num):
        """Виконує вибір секції Gmail після переключення вкладки"""
        try:
            # Шукаємо Gmail Hacks instance та викликаємо відповідну функцію
            gmail_hacks_instance = None
            for tab_name, tab_frame in self.tabview._tab_dict.items():
                if tab_name == "Gmail Hacks":
                    gmail_hacks_instance = self.find_gmail_hacks_instance(tab_frame)
                    if gmail_hacks_instance:
                        # Викликаємо відповідну функцію кнопки (section_num - 1 бо offset починається з 0)
                        offset = section_num - 1
                        gmail_hacks_instance.select_account_by_offset_from_current(offset)
                        print(f"⚡ Shift+{section_num} → Кнопка {section_num} виконана (offset {offset})")
                        return
                    break
                    
            print("⚠️ Gmail Hacks instance не знайдено")
        except Exception as e:
            print(f"❌ Помилка виконання вибору секції {section_num}: {e}")

    def handle_keypress(self, event):
        """ВИМКНЕНО: Обробник викликав рекурсивні помилки"""
        
        # Функція повністю вимкнена для стабільності
        return None
        
        # Старий код закоментовано:
        # if hasattr(self, '_processing_keypress') and self._processing_keypress:
        #     return None
        # 
        # if not (event.state & 0x4):  # Ctrl modifier
        #     return None
        
        # Весь код нижче закоментовано для стабільності:
        #
        # try:
        #     self._processing_keypress = True
        #     
        #     layout_mapping = {
        #         'c': self.global_copy, 'с': self.global_copy,
        #         'v': self.global_paste, 'м': self.global_paste,
        #         'a': self.global_select_all, 'ф': self.global_select_all,
        #         # ... інші mapping
        #     }
        #     
        #     char = event.char
        #     if char in layout_mapping:
        #         try:
        #             result = layout_mapping[char](event)
        #             return result
        #         except Exception as e:
        #             print(f"Помилка обробки hotkey для '{char}': {e}")
        #     
        #     return None
        #     
        # except Exception as e:
        #     print(f"Критична помилка в handle_keypress: {e}")
        #     return None
        # finally:
        #     if hasattr(self, '_processing_keypress'):
        #         self._processing_keypress = False
    
    # ========================================
    # ВСІ ГЛОБАЛЬНІ HOTKEY ФУНКЦІЇ ЗАКОМЕНТОВАНО ДЛЯ СТАБІЛЬНОСТІ
    # ========================================
    
    # def global_copy(self, event=None):
    #     """Глобальне копіювання"""
    #     try:
    #         focused = self.focus_get()
    #         if not focused:
    #             return None
    #             
    #         # Спочатку пробуємо стандартне копіювання
    #         try:
    #             focused.event_generate('<Control-c>')
    #             return "break"
    #         except Exception as e:
    #             pass
    #         
    #         # Якщо стандартне не працює, пробуємо через selection
    #         if hasattr(focused, 'selection_get'):
    #             try:
    #                 selected_text = focused.selection_get()
    #                 safe_clipboard_operation("set", selected_text)
    #                 return "break"
    #             except Exception as e:
    #                 pass
    #         
    #         # Для CustomTkinter Text віджетів
    #         if hasattr(focused, 'get') and hasattr(focused, 'tag_ranges'):
    #             try:
    #                 if focused.tag_ranges("sel"):
    #                     selected_text = focused.get("sel.first", "sel.last")
    #                     safe_clipboard_operation("set", selected_text)
    #                     return "break"
    #             except Exception as e:
    #                 pass
    #                 
    #     except Exception as e:
    #         print(f"Global copy error: {e}")
    #         
    #     return None
    
    # def global_paste(self, event=None):
    #     """Глобальне вставлення"""
    #     try:
    #         focused = self.focus_get()
    #         if not focused:
    #             return None
    #             
    #         # Спочатку пробуємо стандартне вставлення
    #         try:
    #             focused.event_generate('<Control-v>')
    #             return "break"
    #         except Exception as e:
    #             pass
    #         
    #         # Якщо стандартне не працює, пробуємо ручне вставлення
    #         if hasattr(focused, 'insert'):
    #             try:
    #                 clipboard_text = safe_clipboard_operation("get")
    #                 if clipboard_text:
    #                     if hasattr(focused, 'index') and hasattr(focused, 'get'):
    #                         # Для Entry віджетів
    #                         cursor_pos = focused.index('insert')
    #                         focused.insert(cursor_pos, clipboard_text)
    #                     else:
    #                         # Для Text віджетів
    #                         focused.insert('insert', clipboard_text)
    #                     return "break"
    #             except Exception as e:
    #                 pass
    #                 
    #     except Exception as e:
    #         print(f"Global paste error: {e}")
    #         
    #     return None
    
    # def global_select_all(self, event=None):
    #     """Глобальне виділення всього"""
    #     try:
    #         focused = self.focus_get()
    #         if focused:
    #             focused.event_generate('<Control-a>')
    #             return "break"
    #     except Exception as e:
    #         print(f"Global select all error: {e}")
    #     return None
    
    # def global_save(self, event=None):
    #     """Глобальне збереження"""
    #     try:
    #         print("💾 Збереження...")
    #         return "break"
    #     except Exception as e:
    #         print(f"Global save error: {e}")
    #     return None
    
    # def global_open(self, event=None):
    #     """Глобальне відкриття"""
    #     try:
    #         print("📂 Відкриття...")
    #         return "break"
    #     except Exception as e:
    #         print(f"Global open error: {e}")
    #     return None
        
    # def global_new(self, event=None):
    #     """Створення нового"""
    #     try:
    #         print("📄 Створення нового...")
    #         return "break"
    #     except Exception as e:
    #         print(f"Global new error: {e}")
    #     return None
        
    # def global_find(self, event=None):
    #     """Глобальний пошук"""
    #     try:
    #         print("🔍 Пошук...")
    #         return "break"
    #     except Exception as e:
    #         print(f"Global find error: {e}")
    #     return None
        
    # def global_undo(self, event=None):
    #     """Глобальне скасування"""
    #     try:
    #         focused = self.focus_get()
    #         if focused:
    #             focused.event_generate('<Control-z>')
    #             return "break"
    #     except Exception as e:
    #         print(f"Global undo error: {e}")
    #     return None
    
    # def global_generate_all(self, event=None):
    #     """Глобальна генерація всіх секцій (Ctrl+G/П)"""
    #     try:
    #         if hasattr(self, 'generate_all'):
    #             self.generate_all()
    #             print("🎲 Згенеровано всі секції (Ctrl+G)")
    #             return "break"
    #     except Exception as e:
    #         print(f"Global generate all error: {e}")
    #     return None
    
    # def global_copy_combined(self, event=None):
    #     """Глобальне копіювання об'єднаного контенту (Ctrl+B/И)"""
    #     try:
    #         if hasattr(self, 'copy_combined'):
    #             self.copy_combined()
    #             print("📋 Об'єднаний контент скопійовано (Ctrl+B)")
    #             return "break"
    #     except Exception as e:
    #         print(f"Global copy combined error: {e}")
    #     return None

    def __init__(self):
        super().__init__()
        self.font_default = ctk.CTkFont(family="Segoe UI", size=11)
        self.title('Punch IT Now 9.2 - Global Key Binding')
        self.resizable(True, True)
        
        # Початкові розміри (буде змінено в setup_responsive_design)
        self.geometry('1100x750')
        self.minsize(800, 500)
        
        # Очищуємо тимчасові папки при запуску
        cleanup_temp_folders()
        
        # Ініціалізуємо змінні для захисту від рекурсії
        self._processing_keypress = False
        self._saving_farmer_name = False
        self._save_timer = None
        
        # Налаштування швидкості скролу (за замовчуванням 1 = нормально, 3 = швидко)
        load_global_scroll_speed()  # Завантажуємо глобальну швидкість
        self.scroll_speed = get_global_scroll_speed()  # Зберігаємо локальну копію для сумісності
        
        # ⚙️ Features Settings - завантажуємо конфігурацію функцій
        self.features_config = _features_config
        print(f"✅ Features Config завантажено: {self.features_config.config_file}")
        
        # 🖼️ Додаємо іконку для програми (cross-platform)
        # ✨ Іконка тепер лежить в папці config
        try:
            icon_loaded = False
            
            # Визначаємо формат іконки для платформи
            if IS_MACOS:
                icon_filename = "Punch SOFT.icns"  # macOS використовує .icns
                icon_method = "iconphoto"  # macOS підтримує iconphoto
            else:
                icon_filename = "Punch SOFT.ico"   # Windows/Linux використовують .ico
                icon_method = "iconbitmap"
            
            # Спочатку шукаємо в папці config
            icon_path = os.path.join(os.path.dirname(__file__), "config", icon_filename)
            if os.path.exists(icon_path):
                if IS_MACOS and icon_method == "iconphoto":
                    # macOS потребує PhotoImage
                    from PIL import Image, ImageTk
                    img = Image.open(icon_path)
                    photo = ImageTk.PhotoImage(img)
                    self.iconphoto(True, photo)
                    icon_loaded = True
                else:
                    self.iconbitmap(icon_path)
                    icon_loaded = True
                print(f"✅ Icon loaded from config: {icon_path}")
                
            elif os.path.exists(f"config/{icon_filename}"):
                # Якщо файл в config директорії
                if IS_MACOS and icon_method == "iconphoto":
                    from PIL import Image, ImageTk
                    img = Image.open(f"config/{icon_filename}")
                    photo = ImageTk.PhotoImage(img)
                    self.iconphoto(True, photo)
                    icon_loaded = True
                else:
                    self.iconbitmap(f"config/{icon_filename}")
                    icon_loaded = True
                print("✅ Icon loaded from config directory")
                
            else:
                # Для пакованих файлів шукаємо в тимчасовій папці
                if hasattr(sys, '_MEIPASS'):
                    icon_path = os.path.join(sys._MEIPASS, "config", icon_filename)
                    if os.path.exists(icon_path):
                        if IS_MACOS and icon_method == "iconphoto":
                            from PIL import Image, ImageTk
                            img = Image.open(icon_path)
                            photo = ImageTk.PhotoImage(img)
                            self.iconphoto(True, photo)
                            icon_loaded = True
                        else:
                            self.iconbitmap(icon_path)
                            icon_loaded = True
                        print(f"✅ Icon loaded from bundle config: {icon_path}")
                        
            if not icon_loaded:
                print(f"⚠ Icon file not found in config ({icon_filename})")
                
        except Exception as e:
            print(f"❌ Не вдалось завантажити іконку: {e}")
        
        # Додаємо обробник зміни розміру вікна
        self.bind('<Configure>', self.on_window_resize)
        self.fullscreen = False
        self.bind('<F11>', lambda e: self.toggle_fullscreen())
        
        # Гаряча клавіша для налаштувань вкладок (Ctrl+,)
        self.bind('<Control-comma>', lambda e: self.tab_manager.show_tab_settings())
        
        # Глобальні гарячі клавіші вимкнено для стабільності
        # self.setup_global_hotkeys()
        
        self.grid_columnconfigure(0, weight=2)
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        # create logger early
        self.logger = Logger(None)
        
        # Ініціалізуємо менеджер вкладок
        print("🔧 Ініціалізація TabManager...")
        self.tab_manager = TabManager(self)
        print("✅ TabManager ініціалізовано")

        self.theme_file = get_config_path('current_theme.json')
        # Завантажити збережену тему при старті
        self.load_saved_theme()
        
        # Ініціалізуємо CSV шлях з збереженого конфігу
        self.csv_path = self.load_csv_config()
        data = load_csv_columns(self.csv_path)
        # keep data accessible for day-generation popup
        self.data = data

        # Use CTkTabview for tabs
        self.tabview = ctk.CTkTabview(self)
        self.tabview.grid(row=0, column=0, columnspan=2, sticky='nsew')

        # Generators Tab
        self.generators_tab = self.tabview.add("Generators")
        
        # Додаємо кнопку налаштувань до Generators
        self.tab_manager.add_settings_button_to_tab(self.generators_tab)
        
        # Створюємо контейнер для скролювальної області та кнопок
        main_container = ctk.CTkFrame(self.generators_tab)
        main_container.pack(fill="both", expand=True)
        
        # Створюємо скролювальну область для контенту (БЕЗ expand щоб не займала весь простір)
        self.generators_scrollable = ctk.CTkScrollableFrame(main_container)
        self.generators_scrollable.pack(fill="both", expand=True, padx=5, pady=(5, 0))  # pady=(5, 0) - відступ тільки зверху
        
        # Налаштовуємо grid для скролювальної області
        self.generators_scrollable.grid_columnconfigure(0, weight=2, minsize=600)  # Ліва колонка мінімум 600px
        self.generators_scrollable.grid_columnconfigure(1, weight=1, minsize=300)  # Права колонка мінімум 300px
        self.generators_scrollable.grid_rowconfigure(0, weight=0)  # Day Log верхня панель
        self.generators_scrollable.grid_rowconfigure(1, weight=0)  # Контент секцій
        
        # Ініціалізуємо CSV змінні для Generators
        if not hasattr(self, 'csv_files'):
            self.csv_files = []
        if not hasattr(self, 'main_csv_file'):
            self.main_csv_file = None
            
        # Сканируем доступні CSV файли
        try:
            self.csv_files = self.scan_csv_files()
        except:
            self.csv_files = []

        # Top menu
        top_frame = ctk.CTkFrame(self.generators_scrollable)
        top_frame.grid(row=0, column=0, columnspan=2, sticky='ew', padx=5, pady=5)
        
        # Ряд із основними кнопками
        main_buttons_frame = ctk.CTkFrame(top_frame)
        main_buttons_frame.pack(fill='x', pady=5)
        
        self.day_var = tk.StringVar(value='Day 1')
        day_cb = ctk.CTkOptionMenu(main_buttons_frame, variable=self.day_var, values=[f'Day {i}' for i in range(1,6)], font=self.font_default, width=80, height=24)
        day_cb.pack(side='left', padx=3)
        day_gen = ctk.CTkButton(main_buttons_frame, text='Generate Day', command=lambda: self._day_generate_to_list(self.day_var.get()), width=90, height=24, corner_radius=6, font=self.font_default)
        day_gen.pack(side='left', padx=3)
        self.day_next_btn = ctk.CTkButton(main_buttons_frame, text='Next', command=self._day_next, width=60, height=24, corner_radius=6, font=self.font_default)
        self.day_next_btn.pack(side='left', padx=3)
        self.day_copy_btn = ctk.CTkButton(main_buttons_frame, text='Copy', command=self._day_copy, width=60, height=24, corner_radius=6, font=self.font_default)
        self.day_copy_btn.pack(side='left', padx=3)
        
        # Галочка 3 days farm
        self.three_days_var = tk.BooleanVar()
        self.three_days_checkbox = ctk.CTkCheckBox(main_buttons_frame, text="3 days farm", 
                                                   variable=self.three_days_var, 
                                                   command=self.on_three_days_toggle,
                                                   font=self.font_default, checkbox_width=16, checkbox_height=16)
        self.three_days_checkbox.pack(side='left', padx=(10, 3))
        
        # Галочка 4 windows mode
        self.four_windows_var = tk.BooleanVar()
        
        # ЗАВАНТАЖУЄМО конфіг 4-windows ПЕРЕД створенням UI елементів
        try:
            config_path = get_config_path("four_windows_config.json")
            if os.path.exists(config_path):
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                is_4win_enabled = config.get("four_windows_enabled", False)
                self.four_windows_var.set(is_4win_enabled)
                print(f"📖 Завантажено 4-windows конфіг: {'активний' if is_4win_enabled else 'неактивний'}")
        except:
            self.four_windows_var.set(False)
        
        self.four_windows_checkbox = ctk.CTkCheckBox(main_buttons_frame, text="4 windows", 
                                                     variable=self.four_windows_var, 
                                                     command=self.on_four_windows_toggle,
                                                     font=self.font_default, checkbox_width=16, checkbox_height=16)
        self.four_windows_checkbox.pack(side='left', padx=(10, 3))
        
        # Кнопка інструкції
        ctk.CTkButton(main_buttons_frame, text='Інструкція', command=self.show_generators_instruction, width=80, height=24, corner_radius=6, font=self.font_default).pack(side='left', padx=10)

        # Компактний CSV Manager (поруч з інструкцією)
        ctk.CTkLabel(main_buttons_frame, text="📊", font=self.font_default).pack(side="left", padx=(20,3))
        
        # Ініціалізуємо CSV файли для Generators
        if not hasattr(self, 'csv_files'):
            self.csv_files = []
        if not hasattr(self, 'main_csv_file'):
            self.main_csv_file = None
        
        # Компактний dropdown для CSV
        self.gen_main_csv_dropdown = ctk.CTkComboBox(main_buttons_frame, values=self.csv_files or ["Немає файлів"],
                                                   width=150, font=self.font_default, state="readonly",
                                                   command=self.on_gen_main_csv_selected, height=24)
        self.gen_main_csv_dropdown.pack(side="left", padx=3)
        
        # Кнопка оновити список CSV
        ctk.CTkButton(main_buttons_frame, text="🔄", 
                     command=self.refresh_gen_csv_files, width=24, height=24, 
                     corner_radius=6, font=self.font_default).pack(side="left", padx=2)
        
        # Load CSV кнопка
        ctk.CTkButton(main_buttons_frame, text="�", 
                     command=self.load_csv, width=24, height=24, 
                     corner_radius=6, font=self.font_default).pack(side="left", padx=2)

        # Ініціалізація змінних бінда (залишаємо для сумісності)
        self.bound_key = None
        self.binding_mode = False
        self.current_global_hotkey = None  # Для глобальних hotkeys

        # Left column
        left_frame = ctk.CTkFrame(self.generators_scrollable)
        left_frame.grid(row=1, column=0, sticky='nsew', padx=5, pady=5)
        
        # Налаштовуємо weight для row щоб day_frame розтягувався
        left_frame.grid_rowconfigure(0, weight=1)  # Day Log розтягується
        left_frame.grid_rowconfigure(1, weight=0)  # Секції не розтягуються
        
        # Створюємо контейнер для Day Log (може бути 1 або 4 вікна)
        self.day_frame = ctk.CTkFrame(left_frame, height=580)
        self.day_frame.grid(row=0, column=0, columnspan=4, sticky='nsew', pady=3)
        self.day_frame.grid_propagate(False)  # Відключаємо автоматичний розмір
        
        # Налаштовуємо weight для розширення
        self.day_frame.grid_rowconfigure(0, weight=1)
        
        # Стандартний режим - одне вікно Day Log
        self.day_log_box = ctk.CTkTextbox(self.day_frame, wrap="word", font=self.font_default)
        self.day_log_box.pack(fill='both', expand=True, pady=3)
        self.day_log_box.bind("<Button-1>", self._on_day_log_click)
        
        # Створюємо 4 додаткові Day Log вікна (спочатку приховані)
        self.day_log_boxes = [self.day_log_box]  # Перше вікно
        
        for i in range(1, 4):
            day_log = ctk.CTkTextbox(self.day_frame, wrap="word", font=self.font_default)
            day_log.bind("<Button-1>", lambda event, idx=i: self._on_day_log_click(event, window_idx=idx))
            self.day_log_boxes.append(day_log)
        
        print("Click-to-copy налаштовано для Day Log!")
        
        # Конфігуруємо left_frame для 2 колонок (всі колонки однакової ширини)
        left_frame.grid_columnconfigure(0, weight=1, uniform="sections")
        left_frame.grid_columnconfigure(1, weight=1, uniform="sections")
        left_frame.grid_columnconfigure(2, weight=1, uniform="sections")
        left_frame.grid_columnconfigure(3, weight=1, uniform="sections")
        
        # Створюємо фрейми для секцій (спочатку 2, можемо розширити до 4)
        self.sections_frame_1 = ctk.CTkScrollableFrame(left_frame, height=550)
        self.sections_frame_1.grid(row=1, column=0, sticky='nsew', padx=(0, 2), pady=3)
        
        self.sections_frame_2 = ctk.CTkScrollableFrame(left_frame, height=550)
        self.sections_frame_2.grid(row=1, column=1, sticky='nsew', padx=(2, 0), pady=3)
        
        # Додаткові фрейми для 4-windows режиму (спочатку приховані)
        self.sections_frame_3 = ctk.CTkScrollableFrame(left_frame, height=550)
        self.sections_frame_4 = ctk.CTkScrollableFrame(left_frame, height=550)
        
        # Налаштовуємо індивідуальний скрол для кожного фрейма
        self._setup_individual_scroll()
        
        # Зберігаємо старі назви для зворотної сумісності
        self.sections_frame_left = self.sections_frame_1
        self.sections_frame_right = self.sections_frame_2

        # Right column (UI Generators) - тільки якщо увімкнено
        if is_feature_enabled("ui_sections", "action_log"):
            right_frame = ctk.CTkFrame(self.generators_scrollable)
            right_frame.grid(row=1, column=1, sticky='nsew', padx=5, pady=5)
            
            # Налаштовуємо генератори у правій колонці
            self.setup_ui_generators(right_frame)
        else:
            print("⚙️ Лог дій (UI Generators) - вимкнено через Features Settings")
            # Розширюємо ліву колонку на весь простір
            self.generators_scrollable.grid_columnconfigure(0, weight=1, minsize=0)
            self.generators_scrollable.grid_columnconfigure(1, weight=0, minsize=0)

        # Bottom panel - ПОЗА generators_scrollable, завжди внизу!
        bottom_frame = ctk.CTkFrame(main_container)
        bottom_frame.pack(side='bottom', fill='x', padx=5, pady=5)  # side='bottom' - прикріплюється до низу
        # expose as self.bottom so other code can reference it
        self.bottom = bottom_frame
        btn_frame = ctk.CTkFrame(bottom_frame)
        btn_frame.pack(fill='x', pady=5)
        
        # Налаштовуємо кнопки
        self.setup_ui_buttons(btn_frame)

        # storage for flat items for Next/Copy
        self._last_day_items = []
        self._last_day_index = 0
        
        # Авторський напис (БЕЗ кнопок - вони перенесені в Settings)
        author_panel = ctk.CTkFrame(bottom_frame)
        author_panel.pack(fill='x', pady=2)
        
        self.author_lbl = ctk.CTkLabel(author_panel, 
                                     text="Створено @Alex_FarmPunch • По всім питанням та пропозиціям звертатись в особисті • Дякую за користування",
                                     font=ctk.CTkFont(size=10),
                                     text_color=("blue", "lightblue"),
                                     cursor="hand2")
        self.author_lbl.pack(pady=2)
        self.author_lbl.bind("<Button-1>", lambda e: self.open_telegram_link())
        
        # Налаштовуємо підтримку різних розкладок клавіатури (відкладено)
        # setup_keyboard_locale_support(self)

        # Create section frames with required ranges
        ga_vals = data.get('Google Alerts', [])
        gmail_vals = data.get('Gmail answer', [])
        prompts_vals = data.get('Prompts', [])
        yt_vals = data.get('youtube comentary', [])
        # email subscription source and Random Stuff append words
        email_vals = data.get('Email Subscription') or data.get('Email subscription') or data.get('email subscription') or []
        random_stuff = data.get('Random Stuff') or data.get('random stuff') or data.get('Random stuff') or []
        # email list for direct mails
        emails_for_mail = data.get('Email for mail') or data.get('Email for Mail') or data.get('email for mail') or []
        # try common variants for the new column name
        gsign_vals = data.get('google sign') or data.get('Google sign') or data.get('Google Sign') or []
        # weirdo column data
        weirdo_vals = data.get('Weirdo') or data.get('weirdo') or data.get('WEIRDO') or []

        # Create sections in scrollable area
        # Перевіряємо чи активний 4-windows режим для збільшення sampling
        is_four_windows = getattr(self, 'four_windows_var', None) and self.four_windows_var.get()
        sampling_multiplier = 2 if is_four_windows else 1
        
        # Створюємо список секцій тільки для тих, що увімкнені в Features Settings
        sections = []
        
        if is_feature_enabled("generators", "google_alerts"):
            sections.append(('Google Alerts', 3*sampling_multiplier, 5*sampling_multiplier, ga_vals))
        
        if is_feature_enabled("generators", "google_sign"):
            sections.append(('google sign', 5*sampling_multiplier, 10*sampling_multiplier, gsign_vals))
        
        if is_feature_enabled("generators", "email_subscription"):
            sections.append(('Email Subscription', 3*sampling_multiplier, 5*sampling_multiplier, email_vals))
        
        if is_feature_enabled("generators", "email_for_mail"):
            sections.append(('Email for mail', 5*sampling_multiplier, 10*sampling_multiplier, emails_for_mail))
        
        if is_feature_enabled("generators", "gmail_answer"):
            sections.append(('Gmail answer', 5*sampling_multiplier, 8*sampling_multiplier, gmail_vals))
        
        if is_feature_enabled("generators", "prompts"):
            sections.append(('Prompts', 8*sampling_multiplier, 12*sampling_multiplier, prompts_vals))
        
        if is_feature_enabled("generators", "youtube_commentary"):
            sections.append(('youtube comentary', 5*sampling_multiplier, 7*sampling_multiplier, yt_vals))
        
        if is_feature_enabled("generators", "weirdo"):
            sections.append(('Weirdo', 3*sampling_multiplier, 7*sampling_multiplier, weirdo_vals))
        self.sections_objs = []
        
        # Логування відключених секцій
        all_generators = [
            ("google_alerts", "Google Alerts"),
            ("google_sign", "google sign"),
            ("email_subscription", "Email Subscription"),
            ("email_for_mail", "Email for mail"),
            ("gmail_answer", "Gmail answer"),
            ("prompts", "Prompts"),
            ("youtube_commentary", "youtube comentary"),
            ("weirdo", "Weirdo")
        ]
        
        for key, name in all_generators:
            if not is_feature_enabled("generators", key):
                print(f"⚙️ {name} - вимкнено через Features Settings")
        
        for i, (title, minv, maxv, vals) in enumerate(sections):
            # Визначаємо в яку колонку помістити (ліву або праву)
            target_frame = self.sections_frame_left if i % 2 == 0 else self.sections_frame_right
            
            if title == 'Email Subscription' or title.lower().startswith('google sign'):
                sf = SectionFrame(target_frame, title, minv, maxv, vals, self.logger, append_values=random_stuff, font=self.font_default)
            else:
                sf = SectionFrame(target_frame, title, minv, maxv, vals, self.logger, font=self.font_default)
                
            sf.pack(fill='x', pady=4)
            self.sections_objs.append(sf)

        # map named section shortcuts
        for s in self.sections_objs:
            if s.title == 'Google Alerts':
                self.sec_ga = s
            elif s.title == 'Email Subscription':
                self.sec_email = s
            elif s.title == 'Email for mail':
                self.sec_email_for_mail = s
            elif s.title.lower().startswith('google sign'):
                self.sec_gsign = s
            elif s.title == 'Gmail answer':
                self.sec_gmail = s
            elif s.title == 'Prompts':
                self.sec_prompts = s
            elif s.title == 'youtube comentary':
                self.sec_yt = s
            elif s.title == 'Weirdo':
                self.sec_weirdo = s

        # Зберігаємо оригінальні дані секцій для 4-windows режиму
        self._original_sections = []
        for s in self.sections_objs:
            append_vals = getattr(s, 'append_values', None)
            self._original_sections.append((s.title, s.min_items, s.max_items, s.column_values, append_vals))
        
        print(f"💾 Збережено {len(self._original_sections)} оригінальних секцій")
        
        # Застосовуємо 4-windows layout якщо він активний
        if self.four_windows_var.get():
            print("🔄 Застосовуємо 4-windows layout при ініціалізації...")
            try:
                self.switch_layout_mode(True)
            except Exception as e:
                print(f"⚠️ Помилка при застосуванні 4-windows layout: {e}")

    def _setup_individual_scroll(self):
        """Налаштовує індивідуальний скрол для кожної колонки"""
        try:
            print("🔄 Налаштування індивідуального скролу для колонок...")
            
            frames = [self.sections_frame_1, self.sections_frame_2, 
                     self.sections_frame_3, self.sections_frame_4]
            
            for idx, frame in enumerate(frames):
                if frame and hasattr(frame, 'bind'):
                    # Прив'язуємо скрол тільки до конкретного фрейма
                    frame.bind("<Enter>", lambda event, f=frame: self._on_frame_enter(f))
                    frame.bind("<Leave>", lambda event, f=frame: self._on_frame_leave(f))
                    
                    # Також прив'язуємо до всіх дочірніх віджетів
                    self._bind_scroll_to_children(frame, frame)
                    
                    print(f"✅ Налаштовано скрол для колонки {idx + 1}")
                
        except Exception as e:
            print(f"❌ Помилка налаштування скролу: {str(e)}")
    
    def _bind_scroll_to_children(self, widget, target_frame):
        """Прив'язує скрол до всіх дочірніх віджетів"""
        try:
            # Перевіряємо чи віджет існує
            if not widget.winfo_exists():
                return
                
            if hasattr(widget, 'bind'):
                widget.bind("<Enter>", lambda event, f=target_frame: self._on_frame_enter(f))
                widget.bind("<Leave>", lambda event, f=target_frame: self._on_frame_leave(f))
            
            # Рекурсивно для всіх дітей
            if hasattr(widget, 'winfo_children'):
                for child in widget.winfo_children():
                    self._bind_scroll_to_children(child, target_frame)
                    
        except (tk.TclError, AttributeError):
            pass  # Віджет вже знищено
        except Exception as e:
            pass  # Ігноруємо помилки для окремих віджетів
    
    def _on_frame_enter(self, frame):
        """Активує скрол для конкретного фрейма"""
        try:
            # Перевіряємо чи фрейм існує
            if not frame.winfo_exists():
                return
            # Відключаємо скрол від інших фреймів
            self._disable_all_scroll()
            # Увімкнемо скрол для поточного фрейма
            frame.bind_all("<MouseWheel>", lambda event: self._on_frame_scroll(event, frame))
            frame.focus_set()
        except (tk.TclError, AttributeError) as e:
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"❌ Помилка активації скролу: {str(e)}")
    
    def _on_frame_leave(self, frame):
        """Деактивує скрол для фрейма"""
        try:
            # Перевіряємо чи фрейм існує
            if not frame.winfo_exists():
                return
            # Невелика затримка щоб не відключати скрол миттєво
            self.after(100, self._disable_all_scroll)
        except (tk.TclError, AttributeError) as e:
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"❌ Помилка деактивації скролу: {str(e)}")
    
    def _disable_all_scroll(self):
        """Відключає скрол від усіх фреймів"""
        try:
            if self.winfo_exists():
                self.unbind_all("<MouseWheel>")
        except (tk.TclError, AttributeError) as e:
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"❌ Помилка відключення скролу: {str(e)}")
    
    def _on_frame_scroll(self, event, frame):
        """Обробляє скрол для конкретного фрейма"""
        try:
            # Перевіряємо чи фрейм існує
            if not frame.winfo_exists():
                return
                
            # Знаходимо canvas для скролінгу
            canvas = None
            
            # Шукаємо canvas у батьківських віджетах
            current = frame
            while current and not canvas:
                try:
                    if not current.winfo_exists():
                        break
                    if hasattr(current, '_parent_canvas') and current._parent_canvas:
                        canvas = current._parent_canvas
                        break
                    if isinstance(current, tk.Canvas):
                        canvas = current
                        break
                    current = current.master
                except (tk.TclError, AttributeError):
                    break
            
            if canvas:
                try:
                    if canvas.winfo_exists():
                        # Скролимо canvas з урахуванням глобальної швидкості
                        scroll_speed = get_global_scroll_speed()
                        scroll_amount = int(-1*(event.delta/120)) * scroll_speed
                        canvas.yview_scroll(scroll_amount, "units")
                        print(f"🔄 Скрол колонки: {scroll_amount} (швидкість: {scroll_speed}x)")
                except tk.TclError:
                    pass  # Canvas вже знищено
            else:
                print("⚠️ Canvas для скролу не знайдено")
                
        except (tk.TclError, AttributeError) as e:
            pass  # Віджет вже знищено
        except Exception as e:
            print(f"❌ Помилка скролу фрейма: {str(e)}")
    
    def setup_ui_generators(self, right_frame):
        """Налаштовує генератори UI"""
        # Add generators to right frame
        
        # ⚙️ Резервні коди Google - перевірка Features Settings
        if is_feature_enabled("utilities", "google_backup_codes"):
            codes_labelframe = ctk.CTkFrame(right_frame)
            codes_labelframe.pack(fill="x", pady=5, padx=5)
            ctk.CTkLabel(codes_labelframe, text="Резервні коди Google", font=self.font_default).pack(pady=3)
            self.input_area = ctk.CTkTextbox(codes_labelframe, width=250, height=80, font=self.font_default)
            self.input_area.pack(padx=2, pady=2)
            ctk.CTkButton(codes_labelframe, text="Вставити з буфера", command=self.paste_codes, width=120, height=26, corner_radius=6, font=self.font_default).pack(pady=1)
            ctk.CTkButton(codes_labelframe, text="Трансформація кодів", command=self.process_codes, width=120, height=26, corner_radius=6, font=self.font_default).pack(pady=1)
            print("✅ Резервні коди Google - увімкнено")
        else:
            print("⚙️ Резервні коди Google - вимкнено через Features Settings")

        # ⚙️ 2FA Генератор - перевірка Features Settings
        if is_feature_enabled("utilities", "2fa_generator"):
            fa_labelframe = ctk.CTkFrame(right_frame)
            fa_labelframe.pack(fill="x", pady=5, padx=5)
            ctk.CTkLabel(fa_labelframe, text="Генератор 2FA кодів", font=self.font_default).pack(pady=3)
            self.secret_key = ""
            self.secret_label = ctk.CTkLabel(fa_labelframe, text="Секрет: -", font=self.font_default)
            self.secret_label.pack(pady=2)
            ctk.CTkButton(fa_labelframe, text="Вставити секрет", command=self.paste_secret, width=120, height=26, corner_radius=6, font=self.font_default).pack(pady=1)
            ctk.CTkButton(fa_labelframe, text="Згенерувати 2FA", command=self.generate_2fa, width=120, height=26, corner_radius=6, font=self.font_default).pack(pady=1)

            # Вихідні дані
            output_labelframe = ctk.CTkFrame(right_frame)
            output_labelframe.pack(fill="x", pady=5, padx=5)
            ctk.CTkLabel(output_labelframe, text="Результат (клік=копій)", font=self.font_default).pack(pady=3)
            self.output_area = ctk.CTkTextbox(output_labelframe, width=150, height=50, state='disabled', font=self.font_default)
            self.output_area.pack(padx=2, pady=2)
            
            # Налаштовуємо click-to-copy для 2FA результатів
            try:
                inner_text = self.output_area._textbox
                inner_text.bind("<Button-1>", self._on_2fa_click_copy)
                inner_text.bind('<Key>', lambda event: 'break')  # Блокуємо редагування
            except Exception as e:
                print(f"Помилка налаштування click-to-copy для 2FA: {e}")
            print("✅ Генератор 2FA - увімкнено")
        else:
            print("⚙️ Генератор 2FA - вимкнено через Features Settings")

        # ⚙️ Генератор паролів - перевірка Features Settings
        if is_feature_enabled("utilities", "password_generator"):
            password_labelframe = ctk.CTkFrame(right_frame)
            password_labelframe.pack(fill="x", pady=5, padx=5)
            ctk.CTkLabel(password_labelframe, text="Генератор паролів", font=self.font_default).pack(pady=3)
            settings_frame = ctk.CTkFrame(password_labelframe)
            settings_frame.pack(fill="x", padx=2, pady=2)
            ctk.CTkLabel(settings_frame, text="Кількість:", font=self.font_default).grid(row=0, column=0, sticky="w", padx=2, pady=2)
            self.password_count = ctk.CTkEntry(settings_frame, width=40, font=self.font_default)
            self.password_count.grid(row=0, column=1, sticky="w", padx=2, pady=2)
            self.password_count.insert(0, "1")
            
            # Стрілочки для зміни кількості
            arrows_frame = ctk.CTkFrame(settings_frame)
            arrows_frame.grid(row=0, column=2, sticky="w", padx=2)
            up_btn = ctk.CTkButton(arrows_frame, text="▲", width=20, height=12, corner_radius=3, 
                                  command=self.increase_password_count, font=ctk.CTkFont(size=8))
            up_btn.pack(side="top")
            down_btn = ctk.CTkButton(arrows_frame, text="▼", width=20, height=12, corner_radius=3, 
                                    command=self.decrease_password_count, font=ctk.CTkFont(size=8))
            down_btn.pack(side="top")
            
            # Прив'язка колесика миші до entry поля
            self.password_count.bind("<MouseWheel>", self.on_password_count_mousewheel)
            ctk.CTkButton(password_labelframe, text="Згенерувати", command=self.generate_passwords, width=120, height=26, corner_radius=6, font=self.font_default).pack(pady=1)
            self.password_list = ctk.CTkTextbox(password_labelframe, width=250, height=120, font=self.font_default)
            self.password_list.pack(padx=2, pady=2)
            ctk.CTkButton(password_labelframe, text="Копіювати", command=self.copy_passwords, width=120, height=26, corner_radius=6, font=self.font_default).pack(pady=1)
            
            # Налаштовуємо click-to-copy для списку паролів
            try:
                inner_text = self.password_list._textbox
                inner_text.bind("<Button-1>", self._on_password_click_copy)
                inner_text.bind('<Key>', lambda event: 'break')  # Блокуємо редагування
            except Exception as e:
                print(f"Помилка налаштування click-to-copy для паролів: {e}")
            print("✅ Генератор паролів - увімкнено")
        else:
            print("⚙️ Генератор паролів - вимкнено через Features Settings")




        # Add log textbox for generators
        log_frame = ctk.CTkFrame(right_frame)
        log_frame.pack(fill='x', expand=False, pady=5, padx=5)
        
        # Заголовок з додатковою інформацією
        log_header = ctk.CTkFrame(log_frame)
        log_header.pack(fill='x', padx=5, pady=3)
        ctk.CTkLabel(log_header, text="📝 Лог дій", font=self.font_default).pack(side='left')
        ctk.CTkLabel(log_header, text="(1 клік=копій рядок | Право-клік=меню)", 
                    font=ctk.CTkFont(size=9), text_color="gray").pack(side='right')
        
        # Покращене текстове поле з підтримкою інтерактивності
        self.log_text = ctk.CTkTextbox(log_frame, width=250, height=150, 
                                     font=self.font_default,
                                     wrap='word')  # Перенос слів
        self.log_text.pack(fill='x', expand=False, padx=2, pady=2)
        
        # Додаємо стартове повідомлення
        self.log_text.insert('1.0', '🎯 Лог генератора готовий до роботи\n📋 1 клік = копій рядок | Право-клік = меню\n✨ Спробуйте клікнути на цей рядок!\n\n')
        
        # Робимо текст readonly (не можна редагувати, але можна клікати)
        self.log_text.configure(state='normal')
        self.log_text.bind('<Key>', lambda event: 'break')  # Блокуємо введення з клавіатури
        
        # Update logger to use the log textbox
        self.logger = Logger(self.log_text)
        
        # Додатково налаштовуємо click-to-copy для CustomTkinter textbox
        try:
            # Отримуємо внутрішній Tkinter Text віджет
            inner_text = self.log_text._textbox
            
            # Прив'язуємо події до внутрішнього віджету
            inner_text.bind("<Button-1>", self.logger.on_single_click)
            inner_text.bind("<Button-3>", self.logger.show_context_menu)
            inner_text.bind('<Key>', lambda event: 'break')  # Блокуємо редагування
            
            print("Click-to-copy налаштовано для Generator логу!")
        except Exception as e:
            print(f"Помилка налаштування click-to-copy: {e}")
            # Fallback - звичайний лог без click-to-copy
            pass

        # ================== Daily Report Section ==================
        # ⚙️ Перевірка Features Settings
        if is_feature_enabled("daily_report", "enabled"):
            report_frame = ctk.CTkFrame(right_frame)
            report_frame.pack(fill='x', pady=(5, 0), padx=5)
            
            # Заголовок секції
            report_header = ctk.CTkFrame(report_frame)
            report_header.pack(fill='x', padx=5, pady=3)
            ctk.CTkLabel(report_header, text="📊 Daily Report", font=ctk.CTkFont(size=12, weight="bold")).pack(side='left')
            
            # ⚙️ Кнопка копіювання - перевірка Features Settings
            if is_feature_enabled("daily_report", "copy_button"):
                self.copy_report_btn = ctk.CTkButton(report_header, text="📋 Копіювати", 
                                                   command=self.copy_daily_report, width=90, height=24, 
                                                   corner_radius=4, font=ctk.CTkFont(size=10))
                self.copy_report_btn.pack(side='right', padx=5)
            
            # Основний контент звіту
            report_content = ctk.CTkFrame(report_frame)
            report_content.pack(fill='x', padx=5, pady=(0, 5))
            
            # Налаштування grid для красивого розташування
            report_content.grid_columnconfigure(0, weight=1)
            report_content.grid_columnconfigure(1, weight=0)
            
            # ⚙️ Ім'я фармера - перевірка Features Settings
            if is_feature_enabled("daily_report", "farmer_name"):
                name_frame = ctk.CTkFrame(report_content)
                name_frame.grid(row=0, column=0, columnspan=2, sticky='ew', padx=2, pady=2)
                
                ctk.CTkLabel(name_frame, text="👨‍🌾 Фармер:", font=self.font_default, width=80).pack(side='left', padx=(5,2))
                
                # Поле для імені фармера
                self.report_farmer_name = ctk.CTkEntry(name_frame, width=120, font=self.font_default)
                self.report_farmer_name.pack(side='left', padx=2)
                
                # Автоматичне збереження імені при зміні з debounce
                self._save_timer = None
                def on_name_change(*args):
                    # Скасовуємо попередній таймер
                    if self._save_timer:
                        self.after_cancel(self._save_timer)
                    # Встановлюємо новий таймер на 1000мс
                    self._save_timer = self.after(1000, self.save_farmer_name_to_config)
                
                self.report_farmer_name.bind('<KeyRelease>', on_name_change)
                self.report_farmer_name.bind('<FocusOut>', lambda *args: self.save_farmer_name_to_config())
            
            print("✅ Daily Report - увімкнено")
        else:
            # Додаємо placeholder щоб зберегти відступи
            placeholder_frame = ctk.CTkFrame(right_frame, height=10, fg_color="transparent")
            placeholder_frame.pack(fill='x', pady=5)
            print("⚙️ Daily Report - вимкнено через Features Settings")
        
    def setup_ui_buttons(self, btn_frame):
        """Налаштовує кнопки UI"""
        # Buttons - основний ряд
        gen_btn = ctk.CTkButton(btn_frame, text='Generate All', command=self.generate_all, width=90, height=26, corner_radius=6, font=self.font_default)
        copy_all_btn = ctk.CTkButton(btn_frame, text='Copy All', command=self.copy_combined, width=90, height=26, corner_radius=6, font=self.font_default)
        full_screen_btn = ctk.CTkButton(btn_frame, text='FullScreen', command=self.toggle_fullscreen, width=90, height=26, corner_radius=6, font=self.font_default)
        ui_btn = ctk.CTkButton(btn_frame, text='UI Settings', command=self.open_ui_settings, width=90, height=26, corner_radius=6, font=self.font_default)
        
        # Створюємо окремий фрейм для 4-windows кнопок
        self.four_windows_btn_frame = ctk.CTkFrame(btn_frame)
        
        # 4-Windows копіювання кнопки в окремому фреймі
        self.copy_window1_btn = ctk.CTkButton(self.four_windows_btn_frame, text='Copy W1', command=lambda: self.copy_window_data(1), width=70, height=26, corner_radius=6, font=self.font_default, fg_color="orange")
        self.copy_window2_btn = ctk.CTkButton(self.four_windows_btn_frame, text='Copy W2', command=lambda: self.copy_window_data(2), width=70, height=26, corner_radius=6, font=self.font_default, fg_color="orange") 
        self.copy_window3_btn = ctk.CTkButton(self.four_windows_btn_frame, text='Copy W3', command=lambda: self.copy_window_data(3), width=70, height=26, corner_radius=6, font=self.font_default, fg_color="orange")
        self.copy_window4_btn = ctk.CTkButton(self.four_windows_btn_frame, text='Copy W4', command=lambda: self.copy_window_data(4), width=70, height=26, corner_radius=6, font=self.font_default, fg_color="orange")
        
        # Пакуємо основні кнопки
        gen_btn.pack(side='left', padx=2, pady=2)
        copy_all_btn.pack(side='left', padx=2, pady=2)
        full_screen_btn.pack(side='left', padx=2, pady=2)
        ui_btn.pack(side='left', padx=2, pady=2)
        
        # Пакуємо кнопки 4-windows у їхньому фреймі
        self.copy_window1_btn.pack(side='left', padx=1, pady=2)
        self.copy_window2_btn.pack(side='left', padx=1, pady=2)
        self.copy_window3_btn.pack(side='left', padx=1, pady=2)
        self.copy_window4_btn.pack(side='left', padx=1, pady=2)
        
        # Зберігаємо посилання на кнопки і фрейм
        self.four_windows_buttons = [self.copy_window1_btn, self.copy_window2_btn, 
                                   self.copy_window3_btn, self.copy_window4_btn]
        
        # Перевіряємо стан 4-windows режиму при ініціалізації  
        print("🔍 Ініціалізація кнопок 4-windows режиму")



        # Завантаження збереженого стану вікна
        self.load_window_configuration()

        # Налаштовуємо адаптивний дизайн після створення всіх елементів
        self.after(100, self.setup_responsive_design)
        
        # Налаштовуємо hotkeys для вкладок
        self.after(150, self.setup_tab_hotkeys)
        
        # Ініціалізуємо систему бінда клавіш
        try:
            self.load_key_binding()
        except Exception as e:
            print(f"Помилка ініціалізації бінда: {e}")
        
        # Ініціалізуємо систему 3 days farm
        try:
            self.load_three_days_config()
        except Exception as e:
            print(f"Помилка ініціалізації 3 days farm: {e}")
            
        # Ініціалізуємо систему 4 windows mode
        try:
            self.load_four_windows_config()
        except Exception as e:
            print(f"Помилка ініціалізації 4 windows mode: {e}")

        # Налаштування обробника закриття програми
        self.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        # Додаємо можливість розгорнути вікно кліком по іконці
        self.bind("<Map>", self.on_window_map)
        self.bind("<Unmap>", self.on_window_unmap)
        self.bind("<FocusIn>", self.on_window_focus)
        
        # Прив'язка подій для збереження стану
        self.bind("<Configure>", self.on_window_configure)
        
        # Гаряча клавіша для скидання позиції вікна (Ctrl+Shift+R)
        self.bind("<Control-Shift-R>", lambda e: self.reset_window_state())
        
        # Очищення застарілих конфігурацій (через 3 секунди після запуску)
        self.after(3000, self.cleanup_old_configs)
        
        # Створюємо вкладки динамічно
        print("🚀 Готуємось створити вкладки...")
        self.create_tabs()
        print("🏁 Створення вкладок завершено")

    def load_scroll_speed_config(self):
        """Завантажує налаштування швидкості скролу"""
        try:
            config_path = get_config_path('scroll_speed.json')
            if os.path.exists(config_path):
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    speed = config.get('scroll_speed', 1)
                    print(f"📜 Завантажено швидкість скролу: {speed}x")
                    return speed
        except Exception as e:
            print(f"❌ Помилка завантаження швидкості скролу: {e}")
        
        # За замовчуванням
        return 1
    
    def save_scroll_speed_config(self, speed):
        """Зберігає налаштування швидкості скролу"""
        try:
            config_path = get_config_path('scroll_speed.json')
            config = {'scroll_speed': speed}
            with open(config_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, indent=2, ensure_ascii=False)
            print(f"💾 Збережено швидкість скролу: {speed}x")
            return True
        except Exception as e:
            print(f"❌ Помилка збереження швидкості скролу: {e}")
            return False
    
    def update_all_scroll_speeds(self):
        """Оновлює швидкість скролу для всіх існуючих компонентів"""
        try:
            # Оновлюємо локальну копію
            self.scroll_speed = get_global_scroll_speed()
            
            # Перереєструємо скрол для всіх фреймів (якщо вони існують)
            if hasattr(self, 'frame_1') and self.frame_1:
                self._configure_frame_scroll(self.frame_1, 1)
            if hasattr(self, 'frame_2') and self.frame_2:
                self._configure_frame_scroll(self.frame_2, 2) 
            if hasattr(self, 'frame_3') and self.frame_3:
                self._configure_frame_scroll(self.frame_3, 3)
            if hasattr(self, 'frame_4') and self.frame_4:
                self._configure_frame_scroll(self.frame_4, 4)
                
            print(f"🔄 Оновлено швидкість скролу для всіх компонентів: {self.scroll_speed}x")
        except Exception as e:
            print(f"❌ Помилка оновлення швидкості скролу: {e}")

    def load_window_configuration(self):
        """Завантажує збережену конфігурацію вікна"""
        window_state = load_window_state()
        
        if window_state:
            geometry = window_state.get("geometry")
            is_maximized = window_state.get("maximized", False)
            
            if geometry and validate_geometry(geometry):
                self.geometry(geometry)
                print(f"📐 Restored window: {geometry}")
            else:
                print("⚠ Invalid geometry, using defaults")
                self.geometry('1100x750')
                
            if is_maximized:
                try:
                    self.after(100, lambda: self.state('zoomed'))  # Windows equivalent of maximize
                    print("🔄 Window maximized")
                except:
                    pass
        else:
            print("📐 Using default window configuration")

    def on_window_configure(self, event):
        """Обробляє події зміни розміру/позиції вікна"""
        # Зберігаємо стан тільки для головного вікна, не для дочірніх віджетів
        if event.widget == self:
            # Додаємо невелику затримку, щоб не зберігати занадто часто
            if hasattr(self, '_configure_timer'):
                self.after_cancel(self._configure_timer)
            
            self._configure_timer = self.after(1000, self.save_current_window_state)

    def save_current_window_state(self):
        """Зберігає поточний стан вікна"""
        try:
            # Отримуємо поточну геометрію
            geometry = self.geometry()
            
            # Перевіряємо чи вікно розгорнуто
            is_maximized = (self.state() == 'zoomed')
            
            # Не зберігаємо якщо вікно мінімізоване або в неправильному стані
            if self.state() not in ['normal', 'zoomed']:
                return
            
            # Тимчасово змінюємо заголовок для показу збереження
            original_title = self.title()
            self.title('Punch IT Now 9.2 - 💾 Збережено...')
            
            save_window_state(geometry, is_maximized)
            
            # Повертаємо оригінальний заголовок через 1 секунду
            self.after(1000, lambda: self.title('Punch IT Now 9.2 - Global Key Binding'))
            
        except Exception as e:
            print(f"Failed to save window state: {e}")
            # Відновлюємо заголовок при помилці
    
    def on_window_map(self, event):
        """Викликається коли вікно стає видимим (розгортається)"""
        if event.widget == self:
            print("⚡ Вікно розгорнуто")
    
    def on_window_unmap(self, event):
        """Викликається коли вікно приховується (згортається)"""
        if event.widget == self:
            print("⚡ Вікно згорнуто")
    
    def on_window_focus(self, event):
        """Викликається коли вікно отримує фокус (клік по іконці)"""
        if event.widget == self:
            # Якщо вікно було згорнуто F3 і тепер отримало фокус - можливо користувач кликнув по іконці
            if hasattr(self, '_minimized_by_f3') and self._minimized_by_f3:
                self._minimized_by_f3 = False
                print("⚡ Вікно розгорнуто кліком по іконці")
            self.after(100, lambda: self.title('Punch IT Now 9.2 - Global Key Binding'))

    def reset_window_state(self, info_window=None):
        """Скидає збережений стан вікна до значень за замовчуванням"""
        try:
            # Видаляємо файл конфігурації вікна
            config_path = get_config_path("window_state.json")
            if os.path.exists(config_path):
                os.remove(config_path)
                print("🔄 Window state config removed")
            
            # Застосовуємо стандартні налаштування
            self.geometry('1100x750')
            self.state('normal')
            
            # Центруємо вікно на екрані
            self.update_idletasks()
            width = self.winfo_width()
            height = self.winfo_height()
            x = (self.winfo_screenwidth() // 2) - (width // 2)
            y = (self.winfo_screenheight() // 2) - (height // 2)
            self.geometry(f'{width}x{height}+{x}+{y}')
            
            # Показуємо повідомлення
            message = '🔄 Позиція скинута' if info_window is None else '🔄 Позиція скинута'
            self.title(f'Punch IT Now 9.2 - {message}')
            self.after(2000, lambda: self.title('Punch IT Now 9.2 - Octopus'))
            
            # Закриваємо інформаційне вікно якщо відкрите
            if info_window:
                info_window.destroy()
                
            print("🔄 Window position reset to defaults")
            
        except Exception as e:
            print(f"Failed to reset window state: {e}")
            messagebox.showerror("Помилка", f"Не вдалося скинути позицію: {str(e)}")

    def cleanup_old_configs(self):
        """Видаляє застарілі файли конфігурації"""
        try:
            config_dir = get_config_path("")
            if not os.path.exists(config_dir):
                return
            
            current_time = time.time()
            old_configs = 0
            
            # Шукаємо файли старіше 60 днів
            for filename in os.listdir(config_dir):
                if filename.endswith('.json') and filename != 'window_state.json':
                    file_path = os.path.join(config_dir, filename)
                    try:
                        file_time = os.path.getmtime(file_path)
                        if current_time - file_time > 60 * 24 * 3600:  # 60 днів
                            os.remove(file_path)
                            old_configs += 1
                    except:
                        continue
            
            if old_configs > 0:
                print(f"🧹 Cleaned {old_configs} old config files")
                
        except Exception as e:
            print(f"Config cleanup error: {e}")

    # ================== Key Binding System ==================
    
    def load_key_binding(self):
        """Завантажує збережений бінд клавіш"""
        try:
            bind_file = get_config_path("key_binding.json")
            print(f"🔧 Завантаження бінда з {bind_file}")
            if os.path.exists(bind_file):
                with open(bind_file, "r", encoding="utf-8") as f:
                    bind_data = json.load(f)
                    self.bound_key = bind_data.get("bound_key", None)
                    if self.bound_key:
                        print(f"✅ Завантажено бінд: {self.bound_key}")
                        self.update_bind_status(f"Активний: {self.bound_key}")
                        self.setup_key_listener()
                        # ВІДКЛЮЧАЄМО кнопку при завантаженні, якщо є активний бінд
                        if hasattr(self, 'setup_hotkey_button'):
                            self.setup_hotkey_button.configure(state="disabled", text="🔒 Hotkey активний")
                    else:
                        print("❌ Немає збереженого бінда")
                        self.update_bind_status("Не активний")
            else:
                print("❌ Файл бінда не існує")
                self.bound_key = None
                self.update_bind_status("Не активний")
        except Exception as e:
            print(f"❌ Помилка завантаження бінда: {str(e)}")
            self.bound_key = None
            self.update_bind_status("Не активний")
    
    def save_key_binding(self):
        """Зберігає бінд клавіш"""
        try:
            bind_file = get_config_path("key_binding.json")
            bind_data = {"bound_key": self.bound_key}
            with open(bind_file, "w", encoding="utf-8") as f:
                json.dump(bind_data, f, ensure_ascii=False, indent=2)
            print(f"Збережено бінд: {self.bound_key}")
        except Exception as e:
            print(f"Помилка збереження бінда: {str(e)}")
    
    def update_bind_status(self, status):
        """Оновлює статус бінда в інтерфейсі"""
        # Перевіряємо чи існує UI елемент (після оптимізації може не існувати)
        if hasattr(self, 'bind_status_label') and self.bind_status_label:
            self.bind_status_label.configure(text=f"Статус: {status}")
            if "Активний" in status:
                self.bind_status_label.configure(text_color="green")
            else:
                self.bind_status_label.configure(text_color="gray")
    
    def start_key_binding(self):
        """Запускає режим прослуховування клавіш"""
        # Функція відключена після оптимізації UI
        if hasattr(self, 'setup_hotkey_button') and self.setup_hotkey_button:
            self.binding_mode = True
            self.setup_hotkey_button.configure(text="Натисніть клавішу...", fg_color="orange")
            self.update_bind_status("Очікування клавіші...")
        else:
            print("🔇 Key binding UI було видалено для оптимізації простору")
        
        # Фокусуємося на головне вікно для перехоплення клавіш
        self.focus_set()
        self.bind("<Key>", self.on_key_pressed)
        self.bind("<Button-1>", self.on_mouse_clicked)
    
    def on_key_pressed(self, event):
        """Обробляє натискання клавіші в режимі бінда"""
        if not self.binding_mode:
            return
        
        # Формуємо назву клавіші
        key_name = event.keysym
        if key_name == "??":  # Невідома клавіша
            key_name = f"Key_{event.keycode}"
        
        self.bound_key = key_name
        self.binding_mode = False
        
        # Відновлюємо кнопку (якщо існує)
        if hasattr(self, 'setup_hotkey_button') and self.setup_hotkey_button:
            self.setup_hotkey_button.configure(text="⚙️ Налаштувати Hotkey", fg_color="transparent")
        self.update_bind_status(f"Активний: {key_name}")
        
        # Відв'язуємо тимчасові обробники
        self.unbind("<Key>")
        self.unbind("<Button-1>")
        
        # Зберігаємо та налаштовуємо слухач
        self.save_key_binding()
        self.setup_key_listener()
    
    def on_mouse_clicked(self, event):
        """Обробляє клік миші в режимі бінда"""
        if not self.binding_mode:
            return
        
        # Формуємо назву кнопки миші
        button_name = f"Mouse_{event.num}"
        
        self.bound_key = button_name
        self.binding_mode = False
        
        # Відновлюємо кнопку (якщо існує)
        if hasattr(self, 'setup_hotkey_button') and self.setup_hotkey_button:
            self.setup_hotkey_button.configure(text="⚙️ Налаштувати Hotkey", fg_color="transparent")
        self.update_bind_status(f"Активний: {button_name}")
        
        # Відв'язуємо тимчасові обробники
        self.unbind("<Key>")
        self.unbind("<Button-1>")
        
        # Зберігаємо та налаштовуємо слухач
        self.save_key_binding()
        self.setup_key_listener()
    
    def setup_key_listener(self):
        """Налаштовує глобальний постійний слухач для збереженої клавіші"""
        if not self.bound_key:
            print("❌ Немає ключа для прив'язки")
            return
        
        try:
            # Очищуємо попередні глобальні hotkeys
            self.clear_global_hotkeys()
            
            # Налаштовуємо глобальний hotkey через keyboard модуль
            import keyboard
            
            if self.bound_key.startswith("Mouse_"):
                # Для миші використовуємо стандартні клавіші як альтернативу
                print(f"🖱️ Мишка не підтримується глобально, використовуйте клавіатуру")
                self.update_bind_status("Помилка: Мишка не підтримується глобально")
                return
            else:
                # Для клавіатури - глобальний hotkey
                # Конвертуємо tkinter keysym в keyboard format
                keyboard_key = self.convert_tkinter_key_to_keyboard(self.bound_key)
                
                if keyboard_key:
                    # Реєструємо глобальний hotkey
                    keyboard.add_hotkey(keyboard_key, self.trigger_ctrl_shift_e)
                    self.current_global_hotkey = keyboard_key
                    print(f"🌍 Налаштовано ГЛОБАЛЬНИЙ слухач: {keyboard_key} → Ctrl+Shift+E")
                    self.update_bind_status(f"Активний глобально: {self.bound_key}")
                    
                    # ВІДКЛЮЧАЄМО кнопку setup hotkey, щоб не було конфліктів (якщо існує)
                    if hasattr(self, 'setup_hotkey_button') and self.setup_hotkey_button:
                        self.setup_hotkey_button.configure(state="disabled", text="🔒 Hotkey активний")
                    
                else:
                    print(f"❌ Неможливо конвертувати клавішу: {self.bound_key}")
                    self.update_bind_status("Помилка конвертації клавіші")
                    
        except ImportError:
            print("❌ Модуль keyboard не встановлено! Встановіть: pip install keyboard")
            self.update_bind_status("Помилка: немає модуля keyboard")
        except Exception as e:
            print(f"❌ Помилка налаштування глобального слухача: {str(e)}")
            self.update_bind_status(f"Помилка: {str(e)}")
    
    def convert_tkinter_key_to_keyboard(self, tkinter_key):
        """Конвертує tkinter keysym в keyboard module format"""
        # Мапінг популярних клавіш
        key_mapping = {
            'F1': 'f1', 'F2': 'f2', 'F3': 'f3', 'F4': 'f4', 'F5': 'f5', 
            'F6': 'f6', 'F7': 'f7', 'F8': 'f8', 'F9': 'f9', 'F10': 'f10',
            'F11': 'f11', 'F12': 'f12',
            'Insert': 'insert', 'Delete': 'delete', 'Home': 'home', 'End': 'end',
            'Page_Up': 'page up', 'Page_Down': 'page down',
            'Up': 'up', 'Down': 'down', 'Left': 'left', 'Right': 'right',
            'Return': 'enter', 'BackSpace': 'backspace', 'Tab': 'tab',
            'Escape': 'esc', 'space': 'space',
            'Shift_L': 'left shift', 'Shift_R': 'right shift',
            'Control_L': 'left ctrl', 'Control_R': 'right ctrl',
            'Alt_L': 'left alt', 'Alt_R': 'right alt',
            'Super_L': 'left windows', 'Super_R': 'right windows',
            'Menu': 'menu', 'Pause': 'pause', 'Scroll_Lock': 'scroll lock',
            'Num_Lock': 'num lock', 'Caps_Lock': 'caps lock'
        }
        
        # Спочатку перевіряємо точний збіг
        if tkinter_key in key_mapping:
            return key_mapping[tkinter_key]
        
        # Для звичайних літер і цифр
        if len(tkinter_key) == 1 and tkinter_key.isalnum():
            return tkinter_key.lower()
        
        # Для цифрових клавіш на numpad
        if tkinter_key.startswith('KP_'):
            num = tkinter_key.replace('KP_', '')
            if num.isdigit():
                return f'num {num}'
        
        # Якщо не знайдено, повертаємо як є (нижній регістр)
        return tkinter_key.lower()
    
    def clear_global_hotkeys(self):
        """Очищає всі глобальні hotkeys"""
        try:
            import keyboard
            # Очищаємо всі глобальні hotkeys
            keyboard.unhook_all_hotkeys()
            self._global_hotkeys_registered = False
            print("🗑️ Всі глобальні hotkeys очищено")
        except ImportError:
            pass
        except Exception as e:
            print(f"Помилка очищення глобальних hotkeys: {e}")
    
    def trigger_ctrl_shift_e(self):
        """Емулює натискання Ctrl+Shift+E глобально в будь-якому додатку"""
        try:
            # Спробуємо через keyboard (краще для глобального використання)
            import keyboard
            keyboard.send('ctrl+shift+e')
            print("🚀 ПРЯМА ЕМУЛЯЦІЯ: Ctrl+Shift+E → antidetect browser")
            
            # Додаткове логування для відстеження
            timestamp = datetime.datetime.now().strftime('%H:%M:%S')
            self.logger.log(f"🎯 [{timestamp}] ПРЯМА ЕМУЛЯЦІЯ: Ctrl+Shift+E (кнопка натиснута)")
            
        except ImportError:
            try:
                # Альтернативно через pyautogui
                import pyautogui
                pyautogui.hotkey('ctrl', 'shift', 'e')
                print("🚀 ПРЯМА ЕМУЛЯЦІЯ: Ctrl+Shift+E (pyautogui)")
                timestamp = datetime.datetime.now().strftime('%H:%M:%S')
                self.logger.log(f"🎯 [{timestamp}] ПРЯМА ЕМУЛЯЦІЯ: Ctrl+Shift+E (pyautogui)")
            except ImportError:
                print("❌ КРИТИЧНО: Немає модулів для глобальної емуляції клавіш!")
                print("💡 Встановіть: pip install keyboard pyautogui")
                self.logger.log("❌ Помилка: Немає модулів для емуляції клавіш")
        except Exception as e:
            print(f"❌ Помилка глобальної емуляції клавіш: {str(e)}")
            self.logger.log(f"❌ Помилка емуляції: {str(e)}")
    
    def clear_key_binding(self):
        """Очищує глобальний бінд клавіш"""
        # Відв'язуємо існуючий локальний бінд
        if self.bound_key:
            try:
                if self.bound_key.startswith("Mouse_"):
                    button_num = self.bound_key.split("_")[1]
                    self.unbind(f"<Button-{button_num}>")
                else:
                    self.unbind(f"<{self.bound_key}>")
            except:
                pass
        
        # Очищуємо глобальні hotkeys
        self.clear_global_hotkeys()
        
        # ВКЛЮЧАЄМО назад кнопку setup hotkey (якщо існує)
        if hasattr(self, 'setup_hotkey_button') and self.setup_hotkey_button:
            self.setup_hotkey_button.configure(state="normal", text="⚙️ Налаштувати Hotkey")
        
        self.bound_key = None
        self.update_bind_status("Не активний")
        self.save_key_binding()
        print("🗑️ Глобальний бінд клавіш очищено")

    # ================== 3 Days Farm System ==================
    
    def load_three_days_config(self):
        """Завантажує налаштування 3 days farm"""
        try:
            config_file = get_config_path("three_days_config.json")
            if os.path.exists(config_file):
                with open(config_file, "r", encoding="utf-8") as f:
                    config = json.load(f)
                    enabled = config.get("enabled", False)
                    self.three_days_var.set(enabled)
                    if enabled:
                        self.update_day_options_for_three_days()
                    print(f"✅ Завантажено 3 days farm: {'активно' if enabled else 'неактивно'}")
            else:
                print("❌ Файл конфігурації 3 days farm не існує")
        except Exception as e:
            print(f"❌ Помилка завантаження 3 days farm: {str(e)}")
    
    def save_three_days_config(self):
        """Зберігає налаштування 3 days farm"""
        try:
            config_file = get_config_path("three_days_config.json")
            config = {"enabled": self.three_days_var.get()}
            with open(config_file, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            print(f"💾 Збережено 3 days farm: {'активно' if config['enabled'] else 'неактивно'}")
        except Exception as e:
            print(f"❌ Помилка збереження 3 days farm: {str(e)}")
    
    def on_three_days_toggle(self):
        """Обробляє зміну стану галочки 3 days farm"""
        is_enabled = self.three_days_var.get()
        print(f"🔄 Переключено 3 days farm: {'активно' if is_enabled else 'неактивно'}")
        
        if is_enabled:
            self.update_day_options_for_three_days()
        else:
            self.update_day_options_for_five_days()
        
        self.save_three_days_config()
    
    def update_day_options_for_three_days(self):
        """Оновлює опції днів для 3-денного режиму"""
        try:
            # Знаходимо dropdown днів
            day_cb = None
            for widget in self.generators_scrollable.winfo_children():
                if hasattr(widget, 'winfo_children'):
                    for child in widget.winfo_children():
                        if hasattr(child, 'winfo_children'):
                            for grandchild in child.winfo_children():
                                if isinstance(grandchild, ctk.CTkOptionMenu) and hasattr(grandchild, '_variable'):
                                    if grandchild._variable == self.day_var:
                                        day_cb = grandchild
                                        break
            
            if day_cb:
                day_cb.configure(values=[f'Day {i}' for i in range(1, 4)])  # Days 1-3
                if self.day_var.get() not in ['Day 1', 'Day 2', 'Day 3']:
                    self.day_var.set('Day 1')
                print("✅ Оновлено опції днів: 1-3")
        except Exception as e:
            print(f"❌ Помилка оновлення опцій днів для 3 days: {str(e)}")
    
    def update_day_options_for_five_days(self):
        """Оновлює опції днів для 5-денного режиму"""
        try:
            # Знаходимо dropdown днів
            day_cb = None
            for widget in self.generators_scrollable.winfo_children():
                if hasattr(widget, 'winfo_children'):
                    for child in widget.winfo_children():
                        if hasattr(child, 'winfo_children'):
                            for grandchild in child.winfo_children():
                                if isinstance(grandchild, ctk.CTkOptionMenu) and hasattr(grandchild, '_variable'):
                                    if grandchild._variable == self.day_var:
                                        day_cb = grandchild
                                        break
            
            if day_cb:
                day_cb.configure(values=[f'Day {i}' for i in range(1, 6)])  # Days 1-5
                print("✅ Оновлено опції днів: 1-5")
        except Exception as e:
            print(f"❌ Помилка оновлення опцій днів для 5 days: {str(e)}")

    def on_four_windows_toggle(self):
        """Обробляє зміну стану галочки 4 windows режиму"""
        is_enabled = self.four_windows_var.get()
        print(f"🪟 Переключено 4 windows режим: {'активно' if is_enabled else 'неактивно'}")
        
        # Діагностика
        print(f"🔍 Кількість секцій: {len(self.sections_objs) if hasattr(self, 'sections_objs') else 'Немає'}")
        print(f"🔍 Є фрейми: frame_1={hasattr(self, 'sections_frame_1')}, frame_3={hasattr(self, 'sections_frame_3')}")
        
        # Зберігаємо стан у конфіг
        self.save_four_windows_config()
        
        # Оновлюємо sampling ranges для всіх секцій
        try:
            sampling_multiplier = 2 if is_enabled else 1
            
            # Оновлюємо min/max items для всіх секцій
            if hasattr(self, 'sections_objs') and self.sections_objs:
                for section in self.sections_objs:
                    if hasattr(section, 'title'):
                        # Визначаємо базові значення для кожної секції
                        if section.title == 'Google Alerts':
                            section.min_items = 3 * sampling_multiplier
                            section.max_items = 5 * sampling_multiplier
                        elif section.title.lower().startswith('google sign'):
                            section.min_items = 5 * sampling_multiplier
                            section.max_items = 10 * sampling_multiplier
                        elif section.title == 'Email Subscription':
                            section.min_items = 3 * sampling_multiplier
                            section.max_items = 5 * sampling_multiplier
                        elif section.title == 'Email for mail':
                            section.min_items = 5 * sampling_multiplier
                            section.max_items = 10 * sampling_multiplier
                        elif section.title == 'Gmail answer':
                            section.min_items = 5 * sampling_multiplier
                            section.max_items = 8 * sampling_multiplier
                        elif section.title == 'Prompts':
                            section.min_items = 8 * sampling_multiplier
                            section.max_items = 12 * sampling_multiplier
                        elif section.title == 'youtube comentary':
                            section.min_items = 5 * sampling_multiplier
                            section.max_items = 7 * sampling_multiplier
                        elif section.title == 'Weirdo':
                            section.min_items = 3 * sampling_multiplier
                            section.max_items = 7 * sampling_multiplier
            
            print(f"✅ Оновлено sampling ranges для {'4-х вікон' if is_enabled else 'стандартного'} режиму")
            
        except Exception as e:
            print(f"❌ Помилка оновлення sampling ranges: {str(e)}")
        
        # Перемикаємо layout між 2 і 4 колонками
        try:
            self.switch_layout_mode(is_enabled)
        except Exception as e:
            print(f"❌ Помилка зміни layout: {str(e)}")
        
        # Оновлюємо інтерфейс
        if is_enabled:
            print("🪟 Увімкнуто режим 4-х вікон - розділено на 4 колонки з окремими кнопками копіювання")
        else:
            print("🪟 Вимкнуто режим 4-х вікон - стандартні 2 колонки")
    
    def save_four_windows_config(self):
        """Зберігає налаштування 4 windows режиму"""
        try:
            config_path = get_config_path("four_windows_config.json")
            config = {
                "four_windows_enabled": self.four_windows_var.get(),
                "last_updated": datetime.datetime.now().isoformat()
            }
            with open(config_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            print(f"💾 Збережено 4 windows config: {'активно' if config['four_windows_enabled'] else 'неактивно'}")
        except Exception as e:
            print(f"❌ Помилка збереження 4 windows config: {str(e)}")
    
    def load_four_windows_config(self):
        """Завантажує налаштування 4 windows режиму"""
        try:
            config_path = get_config_path("four_windows_config.json")
            if os.path.exists(config_path):
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                is_enabled = config.get("four_windows_enabled", False)
                
                # Тимчасово відключаємо callback щоб не викликати toggle при завантаженні
                self.four_windows_checkbox.configure(command=lambda: None)
                
                # Встановлюємо стан галочки
                self.four_windows_var.set(is_enabled)
                
                # Повертаємо callback назад
                self.four_windows_checkbox.configure(command=self.on_four_windows_toggle)
                
                # Оновлюємо sampling ranges відповідно до завантаженого стану
                if is_enabled and hasattr(self, 'sections_objs') and self.sections_objs:
                    sampling_multiplier = 2
                    for section in self.sections_objs:
                        if hasattr(section, 'title'):
                            # Визначаємо базові значення для кожної секції
                            if section.title == 'Google Alerts':
                                section.min_items = 3 * sampling_multiplier
                                section.max_items = 5 * sampling_multiplier
                            elif section.title.lower().startswith('google sign'):
                                section.min_items = 5 * sampling_multiplier
                                section.max_items = 10 * sampling_multiplier
                            elif section.title == 'Email Subscription':
                                section.min_items = 3 * sampling_multiplier
                                section.max_items = 5 * sampling_multiplier
                            elif section.title == 'Email for mail':
                                section.min_items = 5 * sampling_multiplier
                                section.max_items = 10 * sampling_multiplier
                            elif section.title == 'Gmail answer':
                                section.min_items = 5 * sampling_multiplier
                                section.max_items = 8 * sampling_multiplier
                            elif section.title == 'Prompts':
                                section.min_items = 8 * sampling_multiplier
                                section.max_items = 12 * sampling_multiplier
                            elif section.title == 'youtube comentary':
                                section.min_items = 5 * sampling_multiplier
                                section.max_items = 7 * sampling_multiplier
                            elif section.title == 'Weirdo':
                                section.min_items = 3 * sampling_multiplier
                                section.max_items = 7 * sampling_multiplier
                
                # Застосовуємо layout відповідно до завантаженого стану
                try:
                    self.switch_layout_mode(is_enabled)
                    # Також показуємо кнопки якщо режим активний
                    if is_enabled:
                        self.show_four_windows_buttons()
                except Exception as e:
                    print(f"❌ Помилка застосування layout при завантаженні: {str(e)}")
                
                print(f"📖 Завантажено 4 windows config: {'активно' if is_enabled else 'неактивно'}")
            else:
                self.four_windows_var.set(False)
                print("📖 4 windows config файл не знайдено, використовуємо стандартні налаштування")
        except Exception as e:
            print(f"❌ Помилка завантаження 4 windows config: {str(e)}")
            self.four_windows_var.set(False)

    def switch_layout_mode(self, four_windows_mode):
        """Перемикає між 2-колонковим та 4-колонковим режимом"""
        try:
            print(f"🔧 switch_layout_mode викликано: режим={'4-вікна' if four_windows_mode else '2-вікна'}")
            
            if not hasattr(self, 'sections_objs') or not self.sections_objs:
                print("❌ Немає секцій для перемикання")
                return
            
            # Отримуємо батьківський фрейм
            left_frame = self.sections_frame_1.master
            print(f"🎯 Батьківський фрейм: {left_frame}")
            
            if four_windows_mode:
                print("🔄 Переключення на 4-колонковий режим...")
                
                # Налаштовуємо сітку для 4 колонок (всі однакової ширини через uniform)
                left_frame.grid_columnconfigure(0, weight=1, uniform="sections")
                left_frame.grid_columnconfigure(1, weight=1, uniform="sections") 
                left_frame.grid_columnconfigure(2, weight=1, uniform="sections")
                left_frame.grid_columnconfigure(3, weight=1, uniform="sections")
                
                # Показуємо всі 4 колонки
                self.sections_frame_1.grid(row=1, column=0, sticky='nsew', padx=(0, 1), pady=3)
                self.sections_frame_2.grid(row=1, column=1, sticky='nsew', padx=(1, 1), pady=3)
                self.sections_frame_3.grid(row=1, column=2, sticky='nsew', padx=(1, 1), pady=3)
                self.sections_frame_4.grid(row=1, column=3, sticky='nsew', padx=(1, 0), pady=3)
                
                print(f"📊 Фрейми розміщені: 1={self.sections_frame_1.winfo_children()}")
                print(f"📊 Фрейм 1 має {len(self.sections_frame_1.winfo_children())} дочірніх елементів")
                
                print("✅ Показано 4 колонки")
                
                # Перерозподіляємо секції
                self.redistribute_sections_to_four_columns()
                    
            else:
                print("🔄 Переключення на 2-колонковий режим...")
                
                # Приховуємо колонки 3 і 4
                self.sections_frame_3.grid_remove()
                self.sections_frame_4.grid_remove()
                
                # Налаштовуємо сітку для 2 колонок (всі однакової ширини через uniform)
                left_frame.grid_columnconfigure(2, weight=0)
                left_frame.grid_columnconfigure(3, weight=0)
                left_frame.grid_columnconfigure(0, weight=1, uniform="sections")
                left_frame.grid_columnconfigure(1, weight=1, uniform="sections")
                
                # Повертаємо стандартний вигляд для колонок 1 і 2
                self.sections_frame_1.grid(row=1, column=0, sticky='nsew', padx=(0, 2), pady=3)
                self.sections_frame_2.grid(row=1, column=1, sticky='nsew', padx=(2, 0), pady=3)
                
                print("✅ Приховано колонки 3-4, показано 2 колонки")
                
                # Перерозподіляємо секції назад до 2 колонок
                self.redistribute_sections_to_two_columns()
            
            print(f"✅ Переключено на {'4-колонковий' if four_windows_mode else '2-колонковий'} режим")
            
            # Показуємо/приховуємо відповідні кнопки
            if four_windows_mode:
                self.show_four_windows_buttons()
            else:
                self.hide_four_windows_buttons()
            
            # Перемикаємо Day Log між 1 і 4 вікнами
            self.switch_day_log_mode(four_windows_mode)
            
            # Принудово оновлюємо інтерфейс
            self.update_idletasks()
            
            # Додаткова перевірка через 100ms
            self.after(100, lambda: self.verify_layout_state(four_windows_mode))
            
        except Exception as e:
            print(f"❌ Помилка перемикання layout: {str(e)}")
            import traceback
            traceback.print_exc()

    def redistribute_sections_to_four_columns(self):
        """Створює копії ВСІХ секцій у кожній з 4 колонок"""
        try:
            frames = [self.sections_frame_1, self.sections_frame_2, 
                     self.sections_frame_3, self.sections_frame_4]
            
            # Очищуємо всі фрейми
            for frame in frames:
                for widget in frame.winfo_children():
                    widget.destroy()
            
            # Створюємо копії всіх секцій у кожному фреймі
            self.sections_objs = []  # Скидаємо список секцій
            
            original_sections_data = []
            
            # Зберігаємо дані оригінальних секцій
            if hasattr(self, '_original_sections'):
                original_sections_data = self._original_sections
            else:
                # Якщо це перший раз, зберігаємо дані з CSV
                csv_data = load_csv_columns(CSV_PATH)
                original_sections_data = [
                    ('Google Alerts', 1, 3, csv_data.get('Google Alerts', []), None),
                    ('google sign', 1, 3, csv_data.get('google sign', []), None), 
                    ('Email Subscription', 1, 2, csv_data.get('Email Subscription', []), None),
                    ('Email for mail', 1, 2, csv_data.get('Email for mail', []), None),
                    ('Gmail answer', 1, 2, csv_data.get('Gmail answer', []), None),
                    ('Prompts', 1, 2, csv_data.get('Prompts', []), None),
                    ('youtube comentary', 1, 2, csv_data.get('youtube comentary', []), None),
                    ('Weirdo', 1, 2, csv_data.get('Weirdo', []), None)
                ]
                self._original_sections = original_sections_data
            
            # Створюємо ВСІ секції у КОЖНІЙ колонці
            for col_index, frame in enumerate(frames):
                print(f"🏗️ Створюємо секції для колонки {col_index + 1}")
                
                for title, min_v, max_v, values, append_vals in original_sections_data:
                    # Створюємо секцію з унікальним ідентифікатором для кожної колонки
                    section_title = f"{title} (Col {col_index + 1})" if col_index > 0 else title
                    
                    section = SectionFrame(frame, section_title, min_v, max_v, values, 
                                         self.logger, append_values=append_vals, font=self.font_default)
                    section.pack(fill='x', pady=1)
                    section._original_title = title  # Зберігаємо оригінальну назву для посилань
                    
                    # Тільки секції з першої колонки додаємо до основного списку (для посилань)
                    if col_index == 0:
                        self.sections_objs.append(section)
                    
                    print(f"📦 Створено секцію '{section_title}' у колонці {col_index + 1}")
            
            # Оновлюємо посилання на секції (тільки з першої колонки)
            self._update_section_references()
            
            # Переконуємося що логер працює після перебудови
            if hasattr(self, 'log_text') and self.log_text:
                self.logger = Logger(self.log_text)
                self.logger.log("🔄 Логер відновлено після перебудови 4-windows режиму")
                
        except Exception as e:
            print(f"❌ Помилка перерозподілу на 4 колонки: {str(e)}")

    def redistribute_sections_to_two_columns(self):
        """Перерозподіляє секції по 2 колонках"""
        try:
            print(f"🔄 Перерозподіл на 2 колонки: {len(self.sections_objs)} секцій")
            
            # ВАЖЛИВО: Очищаємо ВСІ 4 фрейми перед перерозподілом
            for frame in [self.sections_frame_1, self.sections_frame_2, 
                         self.sections_frame_3, self.sections_frame_4]:
                for widget in frame.winfo_children():
                    widget.destroy()
            
            # Тепер створюємо секції тільки в перших 2 колонках
            new_sections = []
            
            for i, section in enumerate(self.sections_objs):
                target_frame = self.sections_frame_1 if i % 2 == 0 else self.sections_frame_2
                
                # Зберігаємо дані старої секції
                old_title = getattr(section, '_original_title', section.title)  # Використовуємо оригінальну назву
                old_min = section.min_items
                old_max = section.max_items
                old_values = section.column_values
                old_append_values = getattr(section, 'append_values', None)
                old_items = getattr(section, 'items', [])
                old_index = getattr(section, 'index', 0)
                
                # Створюємо нову секцію без (Col X) суфіксу
                new_section = SectionFrame(target_frame, old_title, old_min, old_max, 
                                         old_values, self.logger, append_values=old_append_values, font=self.font_default)
                new_section.pack(fill='x', pady=3)
                new_section._original_title = old_title  # Зберігаємо оригінальну назву
                
                # Відновлюємо згенеровані елементи
                if old_items:
                    new_section.items = old_items
                    new_section.index = old_index
                    new_section._refresh_listbox()
                
                new_sections.append(new_section)
                print(f"📦 Секція '{old_title}' створена в колонці {(i % 2) + 1}")
            
            # Замінюємо список секцій
            self.sections_objs = new_sections
            
            # Оновлюємо посилання на секції
            self._update_section_references()
            
            # Переконуємося що логер працює після перебудови
            if hasattr(self, 'log_text') and self.log_text:
                self.logger = Logger(self.log_text)
                self.logger.log("🔄 Логер відновлено після перебудови 2-колонкового режиму")
            
            print(f"✅ Перерозподіл завершено: {len(self.sections_objs)} секцій у 2 колонках")
                
        except Exception as e:
            print(f"❌ Помилка перерозподілу на 2 колонки: {str(e)}")
            import traceback
            traceback.print_exc()

    def verify_layout_state(self, four_windows_mode):
        """Перевіряє і виправляє стан layout через затримку"""
        try:
            print(f"🔍 Перевірка стану layout через 100ms...")
            
            if four_windows_mode:
                # Перевіряємо чи відображені 4 колонки
                frame3_visible = self.sections_frame_3.winfo_viewable()
                frame4_visible = self.sections_frame_4.winfo_viewable()
                
                print(f"📊 Колонка 3 видима: {frame3_visible}")
                print(f"📊 Колонка 4 видима: {frame4_visible}")
                
                if not frame3_visible or not frame4_visible:
                    print("🔧 Виправляю відображення 4 колонок...")
                    # Принудово показуємо колонки знову
                    self.sections_frame_3.grid(row=1, column=2, sticky='nsew', padx=(1, 1), pady=3)
                    self.sections_frame_4.grid(row=1, column=3, sticky='nsew', padx=(1, 0), pady=3)
                    
                    # Оновлюємо сітку
                    left_frame = self.sections_frame_1.master
                    left_frame.grid_columnconfigure(2, weight=1)
                    left_frame.grid_columnconfigure(3, weight=1)
                    
                    print("✅ Виправлено відображення 4 колонок")
            
        except Exception as e:
            print(f"❌ Помилка перевірки layout: {str(e)}")

    def show_four_windows_buttons(self):
        """Показує кнопки для 4-windows режиму"""
        try:
            if hasattr(self, 'four_windows_btn_frame'):
                print("👁️ Показую фрейм кнопок 4-windows режиму")
                # Показуємо весь фрейм з кнопками як новий ряд
                self.four_windows_btn_frame.pack(fill='x', pady=2)
                print("  ✅ Показано фрейм з кнопками Copy W1-W4")
        except Exception as e:
            print(f"❌ Помилка показу кнопок 4-windows: {str(e)}")

    def hide_four_windows_buttons(self):
        """Приховує кнопки для 4-windows režиму"""
        try:
            if hasattr(self, 'four_windows_btn_frame'):
                print("👻 Приховую фрейм кнопок 4-windows режиму")
                # Приховуємо весь фрейм з кнопками
                self.four_windows_btn_frame.pack_forget()
                print("  🫥 Приховано фрейм з кнопками Copy W1-W4")
        except Exception as e:
            print(f"❌ Помилка приховування кнопок 4-windows: {str(e)}")

    def copy_window_data(self, window_number):
        """Копіює дані для конкретного вікна (1-4)"""
        try:
            if not hasattr(self, 'sections_objs') or not self.sections_objs:
                self.logger.log("❌ Немає секцій для копіювання")
                return
            
            # Збираємо дані з секцій для конкретного вікна
            # У 4-windows режимі секції розподілені по колонках: 0->1, 1->2, 2->3, 3->4, 4->1, 5->2...
            window_sections = []
            for i, section in enumerate(self.sections_objs):
                section_window = (i % 4) + 1  # Визначаємо до якого вікна належить секція
                if section_window == window_number:
                    window_sections.append(section)
            
            if not window_sections:
                self.logger.log(f"❌ Немає секцій для вікна {window_number}")
                return
                
            # Збираємо вибрані дані з секцій цього вікна
            parts = []
            for section in window_sections:
                selected = section.get_selected() or ''
                if selected:
                    parts.append(f"{section.title}:\n{selected}")
            
            if not parts:
                self.logger.log(f"❌ Немає згенерованих даних для вікна {window_number}")
                return
                
            # Об'єднуємо дані
            combined = '\n\n'.join(parts)
            
            # Копіюємо в буфер обміну
            try:
                self.clipboard_clear()
                self.clipboard_append(combined)
                self.logger.log(f"📋 Скопійовано дані вікна {window_number} ({len(window_sections)} секцій)")
            except Exception as e:
                self.logger.log(f"❌ Помилка копіювання вікна {window_number}: {e}")
                
        except Exception as e:
            print(f"❌ Помилка copy_window_data: {str(e)}")
            self.logger.log(f"❌ Помилка копіювання вікна {window_number}: {e}")

    def restore_four_windows_layout(self):
        """Відновлює 4-windows layout після дій, що могли його скинути"""
        try:
            if hasattr(self, 'four_windows_var') and self.four_windows_var.get():
                print("🔧 restore_four_windows_layout: відновлення...")
                self.switch_layout_mode(True)
                print("✅ 4-windows layout відновлено")
        except Exception as e:
            print(f"❌ Помилка відновлення 4-windows layout: {str(e)}")

    def switch_day_log_mode(self, four_windows_mode):
        """Перемикає Day Log між 1 вікном та 4 вікнами"""
        try:
            print(f"📝 Перемикання Day Log: {'4 вікна' if four_windows_mode else '1 вікно'}")
            
            if four_windows_mode:
                # 4-вікна режим для Day Log
                
                # Приховуємо стандартне вікно
                self.day_log_box.pack_forget()
                
                # Налаштовуємо сітку для 4 колонок у day_frame
                self.day_frame.grid_columnconfigure(0, weight=1)
                self.day_frame.grid_columnconfigure(1, weight=1)
                self.day_frame.grid_columnconfigure(2, weight=1)
                self.day_frame.grid_columnconfigure(3, weight=1)
                
                # Показуємо всі 4 Day Log вікна
                for i, day_log in enumerate(self.day_log_boxes):
                    day_log.grid(row=0, column=i, sticky='nsew', padx=3, pady=3)
                    # Додаємо заголовок для кожного вікна
                    self._set_day_log_header(day_log, f"Account {i+1}")
                
                print("✅ Day Log: показано 4 вікна")
                
            else:
                # Стандартний режим - 1 вікно
                
                # Приховуємо всі 4 вікна з сітки
                for day_log in self.day_log_boxes:
                    day_log.grid_forget()
                
                # Скидаємо grid конфігурацію day_frame до 1 колонки
                self.day_frame.grid_columnconfigure(0, weight=1)
                self.day_frame.grid_columnconfigure(1, weight=0)
                self.day_frame.grid_columnconfigure(2, weight=0)
                self.day_frame.grid_columnconfigure(3, weight=0)
                
                # Показуємо тільки основне вікно через pack
                self.day_log_box.pack(fill='both', expand=True, pady=3)
                
                print("✅ Day Log: показано 1 вікно")
                
        except Exception as e:
            print(f"❌ Помилка перемикання Day Log: {str(e)}")

    def _set_day_log_header(self, day_log, header_text):
        """Встановлює заголовок для Day Log вікна ТІЛЬКИ якщо вікно порожнє"""
        try:
            current_content = day_log.get('1.0', tk.END).strip()
            if not current_content:  # Тільки якщо вікно порожнє
                day_log.configure(state='normal')
                day_log.delete('1.0', tk.END)
                day_log.insert('1.0', f"=== {header_text} ===\n\n")
                day_log.configure(state='disabled')
        except Exception as e:
            print(f"❌ Помилка встановлення заголовку Day Log: {str(e)}")

    def clear_global_hotkeys(self):
        """Очищує глобальні hotkeys"""
        try:
            import keyboard
            keyboard.clear_all_hotkeys()
            print("🗑️ Глобальні hotkeys очищено")
        except:
            pass
    
    def on_closing(self):
        """Обробляє закриття програми"""
        try:
            # Очищуємо глобальні hotkeys перед закриттям
            self.clear_global_hotkeys()
            print("🗑️ Global hotkeys cleared")
            
            # Зберігаємо останній стан вікна перед закриттям
            self.save_current_window_state()
            print("💾 Final window state saved")
            
            # Зберігаємо бінд клавіш
            if hasattr(self, 'bound_key'):
                self.save_key_binding()
                print("💾 Key binding saved")
            
            # Зберігаємо 3 days farm конфігурацію
            if hasattr(self, 'three_days_var'):
                self.save_three_days_config()
                print("💾 3 days farm config saved")
        except Exception as e:
            print(f"Error saving final state: {e}")
        finally:
            # Закриваємо програму
            self.destroy()

    def setup_responsive_design(self):
        """Налаштовує адаптивний дизайн під розмір екрану"""
        try:
            # Отримуємо розмір екрану
            screen_width = self.winfo_screenwidth()
            screen_height = self.winfo_screenheight()
            
            # Розраховуємо оптимальний розмір вікна (80% від екрану)
            window_width = min(1400, int(screen_width * 0.8))
            window_height = min(900, int(screen_height * 0.8))
            
            # Центруємо вікно на екрані
            x = (screen_width - window_width) // 2
            y = (screen_height - window_height) // 2
            
            self.geometry(f'{window_width}x{window_height}+{x}+{y}')
            
            # Встановлюємо адаптивний мінімальний розмір
            min_width = max(800, int(screen_width * 0.4))
            min_height = max(500, int(screen_height * 0.4))
            self.minsize(min_width, min_height)
            
            print(f"✅ Responsive design applied: {window_width}x{window_height}")
            
        except Exception as e:
            # Fallback для випадків помилок
            self.geometry('1100x750')
            self.minsize(800, 500)
            print(f"❌ Responsive design fallback: {e}")

    def create_tabs(self):
        """Створює всі вкладки на основі налаштувань TabManager"""
        print("🔄 Початок створення вкладок...")
        
        # Зберігаємо посилання на вкладки
        if not hasattr(self, 'tab_objects'):
            self.tab_objects = {}
        
        visible_tabs = self.tab_manager.get_visible_tabs()
        print(f"📋 Знайдено {len(visible_tabs)} видимих вкладок для створення")
        
        for tab_config in visible_tabs:
            tab_name = tab_config["name"]
            tab_title = tab_config["title"]
            
            try:
                print(f"🔧 Створення вкладки: {tab_name} -> {tab_title}")
                
                if tab_name == "Generators":
                    # Перевіряємо чи Generators вже існує (уникаємо дублювання)
                    if "Generators" in self.tabview._tab_dict:
                        print(f"✅ Вкладка Generators вже існує - додаємо в tab_objects")
                        self.tab_objects[tab_name] = self.generators_tab
                        continue
                    
                    # Якщо Generators не існує (після видалення в recreate_tabs)
                    # Треба повторно викликати весь код створення з __init__
                    print(f"🔧 Перестворюємо Generators вкладку після видалення...")
                    
                    # Створюємо вкладку Generators
                    self.generators_tab = self.tabview.add("Generators")
                    self.tab_manager.add_settings_button_to_tab(self.generators_tab)
                    
                    # ВАЖЛИВО: Викликаємо метод ініціалізації UI для Generators
                    # Цей код має бути винесений в окремий метод, але поки просто позначаємо що треба перезавантажити
                    print(f"⚠️ УВАГА: Generators UI не може бути повністю перестворений без перезапуску")
                    print(f"⚠️ Рекомендується перезапустити програму для коректної роботи Generators")
                    
                    # Додаємо інформативне повідомлення
                    info_frame = ctk.CTkFrame(self.generators_tab)
                    info_frame.pack(fill="both", expand=True, padx=50, pady=50)
                    
                    icon_label = ctk.CTkLabel(info_frame, text="⚠️", 
                                             font=ctk.CTkFont(size=72))
                    icon_label.pack(pady=20)
                    
                    title_label = ctk.CTkLabel(info_frame, 
                                              text="Generators вкладка відновлена",
                                              font=ctk.CTkFont(size=20, weight="bold"))
                    title_label.pack(pady=10)
                    
                    message_label = ctk.CTkLabel(info_frame, 
                                                text="Для повного відновлення всього функціоналу\n(секції, кнопки, генератори)\n\nбудь ласка, перезапустіть програму",
                                                font=ctk.CTkFont(size=14),
                                                justify="center")
                    message_label.pack(pady=20)
                    
                    restart_btn = ctk.CTkButton(info_frame, text="🔄 Закрити програму (перезапустіть вручну)",
                                               command=self.quit,
                                               width=300, height=50,
                                               font=ctk.CTkFont(size=14, weight="bold"),
                                               fg_color="#ff6b6b", hover_color="#ee5a5a")
                    restart_btn.pack(pady=20)
                    
                    self.tab_objects[tab_name] = self.generators_tab
                    continue
                elif tab_name == "Gmail Hacks":
                    tab_frame = self.tabview.add(tab_title)
                    self.tab_manager.add_settings_button_to_tab(tab_frame)
                    self.gmail_tab = GmailHacksTab(tab_frame, font=self.font_default)
                    self.tab_objects[tab_name] = self.gmail_tab
                    print(f"✅ Створено вкладку Gmail Hacks")
                elif tab_name == "Gmail Parser":
                    tab_frame = self.tabview.add(tab_title)
                    self.tab_manager.add_settings_button_to_tab(tab_frame)
                    self.gmail_parser_tab = GmailParserTab(tab_frame, font=self.font_default)
                    self.tab_objects[tab_name] = self.gmail_parser_tab
                elif tab_name == "Settings":
                    tab_frame = self.tabview.add(tab_title)
                    # Для Settings НЕ додаємо кнопку налаштувань (вона сама є налаштуваннями)
                    self.settings_tab = SettingsTab(tab_frame, self, font=self.font_default)
                    self.tab_objects[tab_name] = self.settings_tab
                    print(f"✅ Створено вкладку Settings")
                elif tab_name == "SMS Checker":
                    tab_frame = self.tabview.add(tab_title)
                    self.tab_manager.add_settings_button_to_tab(tab_frame)
                    self.sms_tab = SmsCheckerTab(tab_frame, font=self.font_default)
                    self.tab_objects[tab_name] = self.sms_tab
                elif tab_name == "Registration":
                    tab_frame = self.tabview.add(tab_title)
                    self.tab_manager.add_settings_button_to_tab(tab_frame)
                    self.registration_tab = RegistrationTab(tab_frame, font=self.font_default)
                    self.tab_objects[tab_name] = self.registration_tab
                elif tab_name == "ChatGPT":
                    if HAS_OPENAI():
                        tab_frame = self.tabview.add(tab_title)
                        self.tab_manager.add_settings_button_to_tab(tab_frame)
                        self.chatgpt_tab = ChatGPTTab(tab_frame, font=self.font_default)
                        self.tab_objects[tab_name] = self.chatgpt_tab
                    else:
                        print(f"⚠️ ChatGPT вкладка пропущена - OpenAI бібліотека недоступна")
                elif tab_name == "File Generator":
                    if FILE_GENERATOR_AVAILABLE():
                        try:
                            tab_frame = self.tabview.add(tab_title)
                            self.tab_manager.add_settings_button_to_tab(tab_frame)
                            self.file_generator_tab = FileGeneratorTab(self, self.font_default)
                            # Встановлюємо tab_frame безпосередньо, не викликаємо create_tab
                            self.file_generator_tab.tab_frame = tab_frame
                            self.file_generator_tab.create_widgets()
                            self.tab_objects[tab_name] = self.file_generator_tab
                            print(f"OK: Вкладка {tab_title} успішно додана!")
                        except Exception as e:
                            print(f"ERROR: Помилка створення File Generator: {e}")
                    else:
                        print(f"WARNING: File Generator вкладка пропущена - модуль недоступний")
                
            except Exception as e:
                        print(f"❌ ERROR: Помилка створення вкладки {tab_title}: {e}")
                        import traceback
                        traceback.print_exc()
        
        print(f"✅ Завершено створення вкладок. Всього створено: {len(self.tab_objects)} вкладок")
    
    def recreate_tabs(self):
        """Перестворює всі вкладки після зміни налаштувань"""
        try:
            # Зберігаємо поточну активну вкладку
            current_tab = None
            try:
                current_tab = self.tabview.get()
            except:
                pass
            
            # Видаляємо ВСІ вкладки (включаючи Generators)
            tab_names = list(self.tabview._tab_dict.keys())
            for tab_name in tab_names:
                try:
                    self.tabview.delete(tab_name)
                    print(f"🗑️ Видалено вкладку: {tab_name}")
                except Exception as e:
                    print(f"⚠️ Не вдалося видалити вкладку {tab_name}: {e}")
            
            # Очищуємо посилання на старі об'єкти
            self.tab_objects.clear()
            
            # НЕ створюємо Generators тут - дозволяємо create_tabs() зробити це
            # Це уникає дублювання вкладки
            
            # Створюємо ВСІ вкладки включаючи Generators
            self.create_tabs()
            
            # Спробуємо відновити активну вкладку
            if current_tab and current_tab in self.tabview._tab_dict:
                try:
                    self.tabview.set(current_tab)
                except:
                    # Якщо попередня вкладка прихована, переключаємось на першу видиму
                    visible_tabs = list(self.tabview._tab_dict.keys())
                    if visible_tabs:
                        self.tabview.set(visible_tabs[0])
            
            print("✅ Вкладки успішно оновлені!")
            
        except Exception as e:
            print(f"❌ Помилка перестворення вкладок: {e}")
            import traceback
            traceback.print_exc()

    def open_telegram_link(self):
        """Відкриває Telegram посилання автора"""
        try:
            import webbrowser
            webbrowser.open("https://t.me/Alex_FarmPunch")
        except Exception as e:
            print(f"Помилка відкриття посилання: {e}")

    def open_hotkeys_settings(self):
        """Відкриває діалог налаштування гарячих клавіш"""
        popup = ctk.CTkToplevel(self)
        popup.title('⌨️ Hotkeys Settings')
        popup.geometry('600x500')
        popup.attributes('-topmost', True)
        
        # Центруємо вікно
        popup.transient(self)
        popup.grab_set()
        
        popup.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (600 // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (500 // 2)
        popup.geometry(f'600x500+{x}+{y}')
        
        # Заголовок
        title_frame = ctk.CTkFrame(popup, fg_color="transparent")
        title_frame.pack(pady=15, padx=20, fill="x")
        
        title = ctk.CTkLabel(
            title_frame,
            text="⌨️ Налаштування гарячих клавіш",
            font=ctk.CTkFont(size=22, weight="bold")
        )
        title.pack()
        
        subtitle = ctk.CTkLabel(
            title_frame,
            text="Налаштуйте комбінації клавіш для швидких дій",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        subtitle.pack(pady=(5, 0))
        
        # Scrollable Frame для налаштувань
        scroll_frame = ctk.CTkScrollableFrame(popup, width=550, height=300)
        scroll_frame.pack(pady=10, padx=20, fill="both", expand=True)
        
        # Завантажуємо поточні hotkeys
        hotkeys_config = self.load_hotkeys_config()
        
        # Словник для зберігання entry widgets
        hotkey_entries = {}
        
        # === HOTKEY 1: Minimize/Restore (F3) ===
        f3_frame = ctk.CTkFrame(scroll_frame)
        f3_frame.pack(fill="x", pady=10, padx=10)
        
        ctk.CTkLabel(f3_frame, text="🪟 Згорнути/Розгорнути вікно", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(anchor="w", padx=10, pady=(10, 5))
        
        ctk.CTkLabel(f3_frame, text="Поточне значення:", 
                    font=ctk.CTkFont(size=11), text_color="gray").pack(anchor="w", padx=10)
        
        f3_entry = ctk.CTkEntry(f3_frame, placeholder_text="Натисніть клавішу...", width=200)
        f3_entry.pack(anchor="w", padx=10, pady=5)
        f3_entry.insert(0, hotkeys_config.get("minimize_restore", "f3"))
        hotkey_entries["minimize_restore"] = f3_entry
        
        ctk.CTkLabel(f3_frame, text="💡 Приклади: f3, f4, ctrl+m, alt+space", 
                    font=ctk.CTkFont(size=10), text_color="gray").pack(anchor="w", padx=10, pady=(0, 10))
        
        # === HOTKEY 2: Octo Browser (F) ===
        octo_frame = ctk.CTkFrame(scroll_frame)
        octo_frame.pack(fill="x", pady=10, padx=10)
        
        ctk.CTkLabel(octo_frame, text="🐙 Octo Browser Toggle", 
                    font=ctk.CTkFont(size=14, weight="bold")).pack(anchor="w", padx=10, pady=(10, 5))
        
        ctk.CTkLabel(octo_frame, text="Поточне значення:", 
                    font=ctk.CTkFont(size=11), text_color="gray").pack(anchor="w", padx=10)
        
        octo_entry = ctk.CTkEntry(octo_frame, placeholder_text="Натисніть клавішу...", width=200)
        octo_entry.pack(anchor="w", padx=10, pady=5)
        octo_entry.insert(0, hotkeys_config.get("octo_browser", "f"))
        hotkey_entries["octo_browser"] = octo_entry
        
        ctk.CTkLabel(octo_frame, text="💡 Приклади: f, o, ctrl+o, alt+f", 
                    font=ctk.CTkFont(size=10), text_color="gray").pack(anchor="w", padx=10, pady=(0, 10))
        
        # Інформаційна панель
        info_frame = ctk.CTkFrame(scroll_frame, fg_color="#2b2b2b")
        info_frame.pack(fill="x", pady=10, padx=10)
        
        info_text = """ℹ️ Підтримувані формати:
• Одна клавіша: a, b, f3, f4, space, enter
• Комбінації: ctrl+a, alt+f, shift+f3
• Модифікатори: ctrl, alt, shift

⚠️ Зверніть увагу:
• Після зміни hotkeys потрібен перезапуск програми
• Деякі комбінації можуть бути зайняті системою
• F1-F4 та Shift+1-4 використовуються для перемикання"""
        
        ctk.CTkLabel(info_frame, text=info_text, 
                    font=ctk.CTkFont(size=10), 
                    justify="left", 
                    text_color="#cccccc").pack(padx=15, pady=15, anchor="w")
        
        # Кнопки
        buttons_frame = ctk.CTkFrame(popup, fg_color="transparent")
        buttons_frame.pack(pady=15, padx=20, fill="x")
        
        def save_hotkeys():
            """Зберігає нові hotkeys"""
            new_config = {
                "minimize_restore": hotkey_entries["minimize_restore"].get().strip().lower(),
                "octo_browser": hotkey_entries["octo_browser"].get().strip().lower()
            }
            
            # Валідація
            if not new_config["minimize_restore"] or not new_config["octo_browser"]:
                messagebox.showwarning("Помилка", "Всі поля мають бути заповнені!")
                return
            
            # Перевірка на дублікати
            if new_config["minimize_restore"] == new_config["octo_browser"]:
                messagebox.showwarning("Помилка", "Hotkeys не можуть співпадати!")
                return
            
            # Зберігаємо конфігурацію
            self.save_hotkeys_config(new_config)
            
            messagebox.showinfo("Успіх", 
                              "Hotkeys збережено!\n\n"
                              "⚠️ Для застосування змін потрібен перезапуск програми.")
            
            popup.destroy()
        
        def reset_defaults():
            """Скидає hotkeys до значень за замовчуванням"""
            if messagebox.askyesno("Підтвердження", "Скинути hotkeys до значень за замовчуванням?"):
                hotkey_entries["minimize_restore"].delete(0, tk.END)
                hotkey_entries["minimize_restore"].insert(0, "f3")
                hotkey_entries["octo_browser"].delete(0, tk.END)
                hotkey_entries["octo_browser"].insert(0, "f")
        
        # Кнопка Reset
        reset_btn = ctk.CTkButton(
            buttons_frame, 
            text='🔄 За замовчуванням', 
            command=reset_defaults,
            width=180, 
            height=35, 
            corner_radius=8,
            fg_color="#666666",
            hover_color="#555555"
        )
        reset_btn.pack(side='left', padx=5)
        
        # Кнопка Cancel
        cancel_btn = ctk.CTkButton(
            buttons_frame, 
            text='❌ Скасувати', 
            command=popup.destroy,
            width=120, 
            height=35, 
            corner_radius=8,
            fg_color="#d32f2f",
            hover_color="#b71c1c"
        )
        cancel_btn.pack(side='left', padx=5)
        
        # Кнопка Save
        save_btn = ctk.CTkButton(
            buttons_frame, 
            text='💾 Зберегти', 
            command=save_hotkeys,
            width=120, 
            height=35, 
            corner_radius=8,
            fg_color="#4caf50",
            hover_color="#388e3c"
        )
        save_btn.pack(side='right', padx=5)

    def load_hotkeys_config(self):
        """Завантажує конфігурацію hotkeys"""
        config_file = get_config_path("hotkeys_config.json")
        try:
            if os.path.exists(config_file):
                with open(config_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception as e:
            print(f"Помилка завантаження hotkeys config: {e}")
        
        # Значення за замовчуванням
        return {
            "minimize_restore": "f3",
            "octo_browser": "f"
        }
    
    def save_hotkeys_config(self, config):
        """Зберігає конфігурацію hotkeys"""
        config_file = get_config_path("hotkeys_config.json")
        try:
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            print(f"✅ Hotkeys config збережено: {config}")
        except Exception as e:
            print(f"Помилка збереження hotkeys config: {e}")
    
    def open_sheets_parsing_settings(self):
        """Відкриває діалог налаштувань парсингу Google Sheets"""
        dialog = ctk.CTkToplevel(self)
        dialog.title("📊 Налаштування парсингу таблиць")
        dialog.geometry("650x700")
        dialog.transient(self)
        dialog.grab_set()
        
        # Центруємо вікно
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (650 // 2)
        y = (dialog.winfo_screenheight() // 2) - (700 // 2)
        dialog.geometry(f"+{x}+{y}")
        
        main_scroll = ctk.CTkScrollableFrame(dialog)
        main_scroll.pack(fill="both", expand=True, padx=20, pady=20)
        
        ctk.CTkLabel(main_scroll, text="📊 Налаштування назв колонок Google Sheets",
                    font=ctk.CTkFont(size=18, weight="bold")).pack(pady=(0, 10))
        
        ctk.CTkLabel(main_scroll, 
                    text="Вкажіть назви колонок у вашій таблиці для автоматичного парсингу.\n"
                         "Наприклад, якщо колонка з email називається 'Почта', введіть 'Почта'.",
                    font=ctk.CTkFont(size=11), text_color="gray",
                    justify="left").pack(pady=(0, 20), padx=10)
        
        # Завантажуємо поточну конфігурацію
        config = self.load_sheets_parsing_config()
        
        entries = {}
        
        # Створюємо поля для кожної колонки
        fields = [
            ("profile_name", "📝 Назва профілю", "Название"),
            ("email", "📧 Email/Login", "Почта"),
            ("password", "🔒 Password", "Пароль"),
            ("2fa", "🔐 2FA", "2фа"),
            ("backup_codes", "📋 Backup Codes", "Бэкап"),
            ("reserve_mail", "📮 Reserve Mail", "Резервка"),
            ("credit_card", "💳 Credit Card", "Карта"),
            ("conversion", "📊 Конверсія", "конверсія"),
            ("api_cf", "🔑 API Cloudflare", "API CF"),
            ("cf_id", "🆔 Cloudflare ID", "CF ID"),
            ("cf_password", "🔐 Cloudflare Password", "CF Pass")
        ]
        
        for field_key, label_text, default_value in fields:
            field_frame = ctk.CTkFrame(main_scroll)
            field_frame.pack(fill="x", pady=8, padx=10)
            
            ctk.CTkLabel(field_frame, text=label_text,
                        font=ctk.CTkFont(size=12, weight="bold"),
                        width=200, anchor="w").pack(side="left", padx=(10, 10))
            
            entry = ctk.CTkEntry(field_frame, width=300, placeholder_text=default_value)
            entry.pack(side="left", padx=5, expand=True, fill="x")
            entry.insert(0, config.get(field_key, default_value))
            entries[field_key] = entry
        
        # Інформаційна підказка
        info_frame = ctk.CTkFrame(main_scroll)
        info_frame.pack(fill="x", pady=15, padx=10)
        
        ctk.CTkLabel(info_frame, 
                    text="💡 Порада: Софт шукаєці назви в першому рядку таблиці (заголовки).\n"
                         "Якщо колонка не знайдена, вона буде пропущена при імпорті.",
                    font=ctk.CTkFont(size=10), text_color="#FFA500",
                    justify="left").pack(pady=10, padx=10)
        
        # Кнопки
        buttons_frame = ctk.CTkFrame(main_scroll)
        buttons_frame.pack(pady=20)
        
        def save_and_close():
            # Збираємо дані з полів
            new_config = {}
            for field_key, entry in entries.items():
                value = entry.get().strip()
                if value:
                    new_config[field_key] = value
            
            # Зберігаємо конфігурацію
            self.save_sheets_parsing_config(new_config)
            
            messagebox.showinfo("Успіх", 
                              "Налаштування збережено!\n\n"
                              "Нові налаштування будуть використані при наступному імпорті.")
            dialog.destroy()
        
        def reset_defaults():
            # Скидаємо до дефолтних значень
            for field_key, label_text, default_value in fields:
                entries[field_key].delete(0, tk.END)
                entries[field_key].insert(0, default_value)
        
        ctk.CTkButton(buttons_frame, text="💾 Зберегти",
                     command=save_and_close,
                     width=140, height=35,
                     font=ctk.CTkFont(size=12, weight="bold")).pack(side="left", padx=5)
        
        ctk.CTkButton(buttons_frame, text="🔄 За замовчуванням",
                     command=reset_defaults,
                     width=160, height=35,
                     font=ctk.CTkFont(size=12, weight="bold"),
                     fg_color="#FF9800", hover_color="#F57C00").pack(side="left", padx=5)
        
        ctk.CTkButton(buttons_frame, text="❌ Скасувати",
                     command=dialog.destroy,
                     width=120, height=35,
                     font=ctk.CTkFont(size=12, weight="bold"),
                     fg_color="gray", hover_color="darkgray").pack(side="left", padx=5)
    
    def load_sheets_parsing_config(self):
        """Завантажує конфігурацію парсингу Google Sheets"""
        try:
            config_file = get_config_path("sheets_parsing_config.json")
            if os.path.exists(config_file):
                with open(config_file, "r", encoding="utf-8") as f:
                    return json.load(f)
        except Exception as e:
            print(f"Помилка завантаження sheets parsing config: {e}")
        
        # Дефолтна конфігурація
        return {
            "profile_name": "Название",
            "email": "Почта",
            "password": "Пароль",
            "2fa": "2фа",
            "backup_codes": "Бэкап",
            "reserve_mail": "Резервка",
            "credit_card": "Карта",
            "conversion": "конверсія",
            "api_cf": "API CF",
            "cf_id": "CF ID",
            "cf_password": "CF Pass"
        }
    
    def save_sheets_parsing_config(self, config):
        """Зберігає конфігурацію парсингу Google Sheets"""
        try:
            config_file = get_config_path("sheets_parsing_config.json")
            with open(config_file, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Помилка збереження sheets parsing config: {e}")
    
    def check_for_updates(self):
        """Перевіряє наявність оновлень на GitHub"""
        try:
            import requests
            from packaging import version
            import json
            
            # URL до version.json на GitHub
            VERSION_URL = "https://raw.githubusercontent.com/qbyLive1/PunchItNow/main/version.json"
            
            # Поточна версія
            CURRENT_VERSION = "9.3.1"
            
            # Показуємо діалог з прогресом
            progress_dialog = ctk.CTkToplevel(self)
            progress_dialog.title("🔄 Перевірка оновлень")
            progress_dialog.geometry("400x200")
            progress_dialog.transient(self)
            
            # Центруємо діалог
            progress_dialog.update_idletasks()
            x = (progress_dialog.winfo_screenwidth() // 2) - (400 // 2)
            y = (progress_dialog.winfo_screenheight() // 2) - (200 // 2)
            progress_dialog.geometry(f"+{x}+{y}")
            
            status_label = ctk.CTkLabel(progress_dialog, 
                                       text="🔍 Перевірка оновлень...",
                                       font=ctk.CTkFont(size=14))
            status_label.pack(pady=30)
            
            progress = ctk.CTkProgressBar(progress_dialog, width=300)
            progress.pack(pady=20)
            progress.set(0.3)
            
            # Оновлюємо GUI
            progress_dialog.update()
            self.update()
            
            try:
                # Завантажуємо інформацію про версію (з обходом кешу)
                import time
                cache_bust_url = f"{VERSION_URL}?t={int(time.time())}"
                response = requests.get(cache_bust_url, timeout=10)
                response.raise_for_status()
                
                # Читаємо контент і видаляємо BOM якщо є
                content = response.content
                if content.startswith(b'\xef\xbb\xbf'):
                    content = content[3:]  # Видаляємо UTF-8 BOM
                
                version_info = json.loads(content.decode('utf-8'))
                
                latest_version = version_info.get("version", "0.0.0")
                download_url = version_info.get("download_url", "")
                changelog = version_info.get("changelog", "")
                config_url = version_info.get("config_url", "")
                
                progress.set(1.0)
                progress_dialog.update()
                progress_dialog.destroy()
                
                # Порівнюємо версії
                if version.parse(latest_version) > version.parse(CURRENT_VERSION):
                    # Є оновлення
                    self.show_update_dialog(latest_version, CURRENT_VERSION, download_url, changelog, config_url)
                else:
                    messagebox.showinfo("✅ Оновлення", 
                                      f"У вас встановлена остання версія!\n\n"
                                      f"Поточна версія: {CURRENT_VERSION}")
                    
            except requests.exceptions.RequestException as e:
                progress_dialog.destroy()
                messagebox.showerror("❌ Помилка", 
                                   f"Не вдалося перевірити оновлення:\n{str(e)}\n\n"
                                   f"Перевірте інтернет-з'єднання.")
                
        except ImportError:
            messagebox.showerror("❌ Помилка", 
                               "Відсутній модуль 'requests' або 'packaging'.\n\n"
                               "Встановіть: pip install requests packaging")
        except Exception as e:
            messagebox.showerror("❌ Помилка", f"Помилка перевірки оновлень:\n{str(e)}")
    
    def show_update_dialog(self, latest_version, current_version, download_url, changelog, config_url):
        """Показує діалог з інформацією про оновлення"""
        dialog = ctk.CTkToplevel(self)
        dialog.title("🎉 Доступне оновлення!")
        dialog.geometry("600x500")
        dialog.transient(self)
        dialog.grab_set()
        
        # Центруємо діалог
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() // 2) - (600 // 2)
        y = (dialog.winfo_screenheight() // 2) - (500 // 2)
        dialog.geometry(f"+{x}+{y}")
        
        main_frame = ctk.CTkFrame(dialog)
        main_frame.pack(fill="both", expand=True, padx=20, pady=20)
        
        # Заголовок
        ctk.CTkLabel(main_frame, 
                    text="🎉 Доступне оновлення!",
                    font=ctk.CTkFont(size=20, weight="bold")).pack(pady=(10, 5))
        
        # Версії
        version_frame = ctk.CTkFrame(main_frame)
        version_frame.pack(pady=10, padx=10, fill="x")
        
        ctk.CTkLabel(version_frame, 
                    text=f"Поточна версія: {current_version}",
                    font=ctk.CTkFont(size=13)).pack(pady=5)
        
        ctk.CTkLabel(version_frame, 
                    text=f"Нова версія: {latest_version}",
                    font=ctk.CTkFont(size=13, weight="bold"),
                    text_color="#4CAF50").pack(pady=5)
        
        # Changelog
        ctk.CTkLabel(main_frame, 
                    text="📝 Що нового:",
                    font=ctk.CTkFont(size=14, weight="bold")).pack(pady=(15, 5), anchor="w", padx=10)
        
        changelog_box = ctk.CTkTextbox(main_frame, height=200, width=540)
        changelog_box.pack(pady=5, padx=10)
        changelog_box.insert("1.0", changelog or "Інформація відсутня")
        changelog_box.configure(state="disabled")
        
        # Кнопки
        button_frame = ctk.CTkFrame(main_frame)
        button_frame.pack(pady=20)
        
        def download_and_install():
            dialog.destroy()
            self.download_and_install_update(download_url, config_url)
        
        ctk.CTkButton(button_frame, 
                     text="⬇️ Завантажити та встановити",
                     command=download_and_install,
                     width=250, height=40,
                     font=ctk.CTkFont(size=13, weight="bold"),
                     fg_color="#4CAF50", hover_color="#45a049").pack(side="left", padx=5)
        
        ctk.CTkButton(button_frame, 
                     text="❌ Пізніше",
                     command=dialog.destroy,
                     width=120, height=40,
                     font=ctk.CTkFont(size=13, weight="bold"),
                     fg_color="gray", hover_color="darkgray").pack(side="left", padx=5)
    
    def download_and_install_update(self, download_url, config_url):
        """Завантажує та встановлює оновлення з .zip архіву"""
        try:
            import requests
            import tempfile
            import subprocess
            import zipfile
            import shutil
            
            # Створюємо діалог прогресу
            progress_dialog = ctk.CTkToplevel(self)
            progress_dialog.title("⬇️ Завантаження оновлення")
            progress_dialog.geometry("500x250")
            progress_dialog.transient(self)
            progress_dialog.grab_set()
            
            # Центруємо діалог
            progress_dialog.update_idletasks()
            x = (progress_dialog.winfo_screenwidth() // 2) - (500 // 2)
            y = (progress_dialog.winfo_screenheight() // 2) - (250 // 2)
            progress_dialog.geometry(f"+{x}+{y}")
            
            status_label = ctk.CTkLabel(progress_dialog, 
                                       text="⬇️ Завантаження архіву...",
                                       font=ctk.CTkFont(size=14))
            status_label.pack(pady=20)
            
            progress_bar = ctk.CTkProgressBar(progress_dialog, width=400)
            progress_bar.pack(pady=10)
            progress_bar.set(0)
            
            percent_label = ctk.CTkLabel(progress_dialog, text="0%", font=ctk.CTkFont(size=12))
            percent_label.pack(pady=5)
            
            self.update()
            
            # Завантажуємо .zip архів
            response = requests.get(download_url, stream=True, timeout=30)
            response.raise_for_status()
            
            total_size = int(response.headers.get('content-length', 0))
            downloaded = 0
            
            # Зберігаємо в тимчасовий файл
            temp_zip = tempfile.NamedTemporaryFile(delete=False, suffix=".zip")
            temp_zip_path = temp_zip.name
            
            with temp_zip as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
                        downloaded += len(chunk)
                        
                        if total_size > 0:
                            progress = downloaded / total_size
                            progress_bar.set(progress)
                            percent_label.configure(text=f"{int(progress * 100)}%")
                            self.update()
            
            status_label.configure(text="📦 Розпаковування архіву...")
            progress_bar.set(0.5)
            self.update()
            
            # Створюємо тимчасову папку для розпакування
            temp_extract_dir = tempfile.mkdtemp(prefix="punchinow_update_")
            
            try:
                # Розпаковуємо архів
                with zipfile.ZipFile(temp_zip_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_extract_dir)
                
                status_label.configure(text="✅ Архів розпаковано!")
                progress_bar.set(1.0)
                self.update()
                time.sleep(1)
                
                progress_dialog.destroy()
                
                # Визначаємо поточну директорію програми
                if getattr(sys, 'frozen', False):
                    current_exe = sys.executable
                    current_dir = os.path.dirname(current_exe)
                else:
                    current_exe = os.path.abspath(__file__)
                    current_dir = os.path.dirname(current_exe)
                
                # Створюємо bat-скрипт для оновлення
                bat_content = f"""@echo off
chcp 65001 >nul
echo 🔄 Оновлення Punch IT Now...
timeout /t 2 /nobreak >nul

cd /d "{current_dir}"

REM Створюємо backup папку
if not exist "backup" mkdir "backup"

REM Зберігаємо важливі файли в backup
if exist "PunchITNow.exe" (
    echo 💾 Створення резервної копії...
    copy /Y "PunchITNow.exe" "backup\\PunchITNow_backup.exe" >nul
)
if exist "config" (
    xcopy /E /I /Y "config" "backup\\config" >nul
)

REM Копіюємо нові файли з тимчасової папки
echo 📂 Копіювання нових файлів...
xcopy /E /I /Y "{temp_extract_dir}\\*" "{current_dir}" >nul

REM Видаляємо тимчасовий архів
del "{temp_zip_path}" >nul

REM Видаляємо тимчасову папку
rmdir /S /Q "{temp_extract_dir}" >nul

echo ✅ Оновлення завершено!
echo ⏳ Очікування завершення старого процесу...
timeout /t 5 /nobreak >nul

REM Запускаємо оновлену програму
echo 🚀 Запуск оновленої програми...
start "" "{current_exe}"

REM Видаляємо bat-скрипт
(goto) 2>nul & del "%~f0"
"""
                
                bat_path = os.path.join(tempfile.gettempdir(), "update_punchinow.bat")
                with open(bat_path, "w", encoding="utf-8") as f:
                    f.write(bat_content)
                
                # Показуємо повідомлення
                result = messagebox.askyesno("🔄 Перезапуск", 
                                            "Оновлення завантажено!\n\n"
                                            "Програма буде перезапущена для встановлення оновлення.\n"
                                            "Резервна копія буде збережена в папці 'backup'.\n\n"
                                            "Продовжити?")
                
                if result:
                    # Запускаємо bat-скрипт і закриваємо програму
                    subprocess.Popen([bat_path], shell=True, creationflags=subprocess.CREATE_NO_WINDOW)
                    self.quit()
                else:
                    # Видаляємо тимчасові файли
                    os.unlink(temp_zip_path)
                    os.unlink(bat_path)
                    shutil.rmtree(temp_extract_dir, ignore_errors=True)
                    messagebox.showinfo("ℹ️", "Оновлення відкладено.")
                    
            except zipfile.BadZipFile:
                progress_dialog.destroy()
                messagebox.showerror("❌ Помилка", 
                                   "Завантажений файл не є коректним ZIP архівом.")
            finally:
                # Очищення тимчасових файлів у разі помилки
                if os.path.exists(temp_zip_path):
                    os.unlink(temp_zip_path)
                
        except requests.exceptions.RequestException as e:
            messagebox.showerror("❌ Помилка", 
                               f"Не вдалося завантажити оновлення:\n{str(e)}")
        except Exception as e:
            messagebox.showerror("❌ Помилка", 
                               f"Помилка встановлення оновлення:\n{str(e)}")

    def show_config_info(self):
        """Показує інформацію про розташування файлів конфігурації"""
        try:
            info_text = f"""📁 Розташування файлів програми Punch Gmail GOD 5.0:

🏠 Робоча папка: {WORKDIR}

⚙️ Конфігурації (папка config/):
• Налаштування Gmail Hacks: gmail_hacks_config.json
• Автозбереження полів: account_info_autosave.json  
• Збережена тема: current_theme.json
• Позиція та розмір вікна: window_state.json

📊 Файли даних:
• CSV для Generator: Account generation.csv
• Результати перевірки: checker results.csv
• Згенеровані акаунти: accounts.txt (Registration)

🎯 Гарячі клавіші (працюють на всіх розкладках EN/UA/RU):
• Ctrl+C/С - копіювання
• Ctrl+V/М - вставлення  
• Ctrl+A/Ф - виділити все
• Ctrl+S/І/Ы - збереження
• Ctrl+O/Щ - відкрити
• Ctrl+F/А - пошук
• Ctrl+Z/Я - скасувати
• Ctrl+G/П - згенерувати всі секції
• Ctrl+B/И - копіювати об'єднаний контент
• Ctrl+Shift+R - швидко скинути позицію вікна
• F11 - повноекранний режим

ℹ️ При створенні .exe файлу всі налаштування
будуть портативними і збережуться поруч з програмою.

💾 Позиція та розмір вікна автоматично зберігаються
при кожній зміні і відновлюються при наступному запуску."""
            
            # Створюємо інформаційне вікно
            info_window = ctk.CTkToplevel(self)
            info_window.title("Інформація про конфігурацію")
            info_window.geometry("500x400")
            info_window.resizable(False, False)
            info_window.attributes('-topmost', True)
            
            # Текстове поле з інформацією
            text_widget = ctk.CTkTextbox(info_window, font=ctk.CTkFont(size=11))
            text_widget.pack(fill='both', expand=True, padx=20, pady=20)
            text_widget.insert('1.0', info_text)
            text_widget.configure(state='disabled')
            
            # Кнопки
            btn_frame = ctk.CTkFrame(info_window)
            btn_frame.pack(fill='x', padx=20, pady=10)
            
            ctk.CTkButton(btn_frame, text="📂 Відкрити папку config", 
                         command=lambda: self.open_folder(CONFIG_DIR), width=150).pack(side='left', padx=5)
            
            ctk.CTkButton(btn_frame, text="📂 Відкрити робочу папку", 
                         command=lambda: self.open_folder(WORKDIR), width=150).pack(side='left', padx=5)
            
            ctk.CTkButton(btn_frame, text="🔄 Скинути розмір вікна", 
                         command=lambda: self.reset_window_state(info_window), width=150).pack(side='left', padx=5)
            
            ctk.CTkButton(btn_frame, text="✅ Закрити", 
                         command=info_window.destroy, width=100).pack(side='right', padx=5)
                         
        except Exception as e:
            messagebox.showerror("Помилка", f"Помилка показу конфігурації: {str(e)}")

    def open_folder(self, path):
        """Відкриває папку в провіднику"""
        try:
            import subprocess
            import os
            
            if os.name == 'nt':  # Windows
                os.startfile(path)
            elif sys.platform == 'darwin':  # macOS
                subprocess.call(['open', path])
            else:  # Linux
                subprocess.call(['xdg-open', path])
        except Exception as e:
            messagebox.showerror("Помилка", f"Не вдалося відкрити папку: {str(e)}")

    def load_saved_theme(self):
        """Завантажити збережену тему при старті програми"""
        try:
            if os.path.exists(self.theme_file):
                with open(self.theme_file, 'r', encoding='utf-8') as f:
                    theme_info = json.load(f)
                    if 'theme_path' in theme_info and os.path.exists(theme_info['theme_path']):
                        ctk.set_default_color_theme(theme_info['theme_path'])
                        self.logger.log(f"Loaded saved theme: {os.path.basename(theme_info['theme_path'])}")
                    elif 'theme_name' in theme_info:
                        ctk.set_default_color_theme(theme_info['theme_name'])
                        self.logger.log(f"Loaded saved theme: {theme_info['theme_name']}")
        except Exception as e:
            self.logger.log(f"Failed to load saved theme: {str(e)}")

    def save_current_theme(self, theme_info):
        """Зберегти поточну тему"""
        try:
            with open(self.theme_file, 'w', encoding='utf-8') as f:
                json.dump(theme_info, f, indent=2, ensure_ascii=False)
        except Exception as e:
            self.logger.log(f"Failed to save theme: {str(e)}")

    def load_theme(self):
        try:
            if os.path.exists(self.theme_file):
                with open(self.theme_file, 'r', encoding='utf-8') as f:
                    theme = json.load(f)
                # Check if it's a CustomTkinter theme and convert to app format
                if "CTk" in theme:
                    root_bg = theme["CTk"]["fg_color"][0]
                    btn_bg = theme["CTkButton"]["fg_color"][0]
                    btn_fg = theme["CTkButton"]["text_color"][0]
                    label_fg = theme["CTkLabel"]["text_color"][0]
                    entry_bg = theme["CTkEntry"]["fg_color"][0]
                    entry_fg = theme["CTkEntry"]["text_color"][0]
                    text_bg = theme["CTkTextbox"]["fg_color"][0]
                    text_fg = theme["CTkTextbox"]["text_color"][0]
                    scrollbar_bg = theme["CTkScrollbar"]["button_color"][0]
                    frame_top_bg = theme["CTkFrame"]["top_fg_color"][0]
                    border_color = theme["CTkButton"]["border_color"][0]
                    theme = {
                        "name": "Converted CTk Theme",
                        "root_bg": root_bg,
                        "styles": {
                            "TButton": {"background": btn_bg, "foreground": btn_fg, "bordercolor": border_color, "relief": "flat"},
                            "TLabel": {"background": root_bg, "foreground": label_fg},
                            "TFrame": {"background": root_bg},
                            "TNotebook": {"background": root_bg, "tabbackground": frame_top_bg},
                            "TCombobox": {"fieldbackground": entry_bg, "background": root_bg, "foreground": entry_fg},
                            "TEntry": {"fieldbackground": entry_bg, "background": root_bg, "foreground": entry_fg},
                            "TText": {"background": text_bg, "foreground": text_fg},
                            "TScrollbar": {"background": scrollbar_bg, "troughcolor": root_bg}
                        }
                    }
                self.apply_theme(theme)
                self.current_theme = theme
        except Exception as e:
            print(f"Failed to load theme: {e}")

    def apply_theme(self, theme):
        root_bg = theme.get('root_bg', '#ffffff')
        fg = theme.get('styles', {}).get('TLabel', {}).get('foreground', '#000000')
        self.configure(bg=root_bg)
        style = ttk.Style()
        styles = theme.get('styles', {})
        for widget, props in styles.items():
            try:
                style.configure(widget, **props)
            except Exception as e:
                print(f"Failed to configure {widget}: {e}")
        # Also configure custom button styles
        if 'TButton' in styles:
            try:
                style.configure('Paste.TButton', **styles['TButton'])
                style.configure('Copy.TButton', **styles['TButton'])
            except Exception as e:
                print(f"Failed to configure custom button styles: {e}")
        # Apply tk widget options for better theme consistency
        self.option_add("*Listbox.background", root_bg)
        self.option_add("*Listbox.foreground", fg)
        self.option_add("*Text.background", root_bg)
        self.option_add("*Text.foreground", fg)
        self.option_add("*Label.background", root_bg)
        self.option_add("*Label.foreground", fg)
        self.option_add("*Entry.background", root_bg)
        self.option_add("*Entry.foreground", fg)
        # Update existing tk widgets
        self.update_tk_widgets(self, root_bg, fg)
        # Update UI settings to match the theme
        brightness = self.get_brightness(root_bg)
        self.ui_theme = 'dark' if brightness < 128 else 'light'
        if 'TButton' in styles and 'background' in styles['TButton']:
            self.button_color = styles['TButton']['background']
        self.logger.log(f"Applied theme: {theme.get('name', 'Unknown')}")

    def update_tk_widgets(self, widget, bg, fg):
        try:
            if isinstance(widget, (tk.Text, tk.Listbox, tk.Entry, tk.Label)):
                widget.configure(bg=bg, fg=fg)
        except tk.TclError:
            pass
        for child in widget.winfo_children():
            self.update_tk_widgets(child, bg, fg)

    def save_theme(self):
        if self.current_theme:
            try:
                with open(self.theme_file, 'w', encoding='utf-8') as f:
                    json.dump(self.current_theme, f, indent=2)
            except Exception as e:
                print(f"Failed to save theme: {e}")

    def on_close(self):
        self.save_theme()
        self.destroy()

    def generate_all(self):
        try:
            # Запам'ятовуємо стан 4-windows режиму перед генерацією
            four_windows_active = getattr(self, 'four_windows_var', None) and self.four_windows_var.get()
            print(f"🔄 generate_all: 4-windows режим {'активний' if four_windows_active else 'неактивний'}")
            
            # Переконуємося, що посилання на секції оновлені
            if four_windows_active and hasattr(self, 'sections_objs'):
                print("🔧 Оновлюємо посилання на секції перед генерацією...")
                self._update_section_references()
            
            # Відладочна інформація про стан секцій
            print(f"🔍 Секції: ga={self.sec_ga is not None}, gsign={self.sec_gsign is not None}, email={self.sec_email is not None}")
            print(f"🔍 Секції: email_mail={self.sec_email_for_mail is not None}, gmail={self.sec_gmail is not None}, prompts={self.sec_prompts is not None}")
            print(f"🔍 Секції: yt={self.sec_yt is not None}, weirdo={self.sec_weirdo is not None}")
            print(f"🔍 Всього секцій в sections_objs: {len(getattr(self, 'sections_objs', []))}")
            
            # Переконуємося що логер працює
            if hasattr(self, 'logger') and self.logger:
                self.logger.log("🔄 Початок генерації всіх секцій")
            
            # В 4-windows режимі генеруємо секції у всіх колонках
            if four_windows_active:
                # Знаходимо всі секції у всіх фреймах
                all_sections = []
                frames = [self.sections_frame_1, self.sections_frame_2, 
                         self.sections_frame_3, self.sections_frame_4]
                
                for frame_idx, frame in enumerate(frames):
                    frame_sections = []
                    # Отримуємо внутрішній скроловий фрейм
                    try:
                        scrollable_frame = frame._scrollable_frame if hasattr(frame, '_scrollable_frame') else frame
                        for widget in scrollable_frame.winfo_children():
                            if hasattr(widget, 'title') and hasattr(widget, 'generate_items'):
                                frame_sections.append(widget)
                                all_sections.append(widget)
                        print(f"📊 Фрейм {frame_idx + 1}: знайдено {len(frame_sections)} секцій")
                    except Exception as e:
                        print(f"❌ Помилка сканування фрейму {frame_idx + 1}: {str(e)}")
                
                print(f"🔄 Загалом знайдено {len(all_sections)} секцій у всіх 4 колонках")
                if hasattr(self, 'logger') and self.logger:
                    self.logger.log(f"🔄 Генерація {len(all_sections)} секцій у 4-windows режимі")
                
                generated_count = 0
                for section in all_sections:
                    try:
                        section.generate_items()
                        generated_count += 1
                        print(f"✅ Згенеровано: {section.title}")
                        if hasattr(self, 'logger') and self.logger:
                            self.logger.log(f"✅ Згенеровано: {section.title}")
                    except Exception as e:
                        error_msg = f"❌ Помилка генерації {getattr(section, 'title', 'Unknown')}: {str(e)}"
                        print(error_msg)
                        if hasattr(self, 'logger') and self.logger:
                            self.logger.log(error_msg)
                
                if hasattr(self, 'logger') and self.logger:
                    self.logger.log(f"✅ Загалом згенеровано {generated_count} секцій у 4-windows режимі")
                        
            else:
                # Стандартний режим - генеруємо тільки основні секції
                if hasattr(self, 'sections_objs') and self.sections_objs:
                    for section in self.sections_objs:
                        try:
                            section.generate_items()
                            print(f"✅ Згенеровано секцію: {section.title}")
                            if hasattr(self, 'logger') and self.logger:
                                self.logger.log(f"✅ Згенеровано секцію: {section.title}")
                        except Exception as e:
                            error_msg = f"❌ Помилка генерації секції {getattr(section, 'title', 'Unknown')}: {str(e)}"
                            print(error_msg)
                            if hasattr(self, 'logger') and self.logger:
                                self.logger.log(error_msg)
                else:
                    print("❌ Немає доступних секцій для генерації")
                    if hasattr(self, 'logger') and self.logger:
                        self.logger.log("❌ Немає доступних секцій для генерації")
                
            # Резервний спосіб через посилання (якщо sections_objs недоступний)  
            if not hasattr(self, 'sections_objs') or not self.sections_objs:
                if hasattr(self, 'sec_ga') and self.sec_ga:
                    self.sec_ga.generate_items()
                # generate google sign
                if hasattr(self, 'sec_gsign') and self.sec_gsign:
                    self.sec_gsign.generate_items()
                # generate email subscription
                if hasattr(self, 'sec_email') and self.sec_email:
                    self.sec_email.generate_items()
                # generate email for mail list
                if hasattr(self, 'sec_email_for_mail') and self.sec_email_for_mail:
                    self.sec_email_for_mail.generate_items()
                if hasattr(self, 'sec_gmail') and self.sec_gmail:
                    self.sec_gmail.generate_items()
                if hasattr(self, 'sec_prompts') and self.sec_prompts:
                    self.sec_prompts.generate_items()
                if hasattr(self, 'sec_yt') and self.sec_yt:
                    self.sec_yt.generate_items()
                # generate weirdo
                if hasattr(self, 'sec_weirdo') and self.sec_weirdo:
                    self.sec_weirdo.generate_items()
            
            self.logger.log('Generated all sections')
            
            # НЕ відновлюємо 4-windows режим після генерації, щоб не очищувати згенерований контент
            # if four_windows_active:
            #     print("🔧 Відновлюємо 4-windows режим після генерації")
            #     self.after(50, lambda: self.restore_four_windows_layout())
            
        except Exception as e:
            self.logger.log(f'Error generating sections: {e}')
            print(f"Error in generate_all: {e}")

    def copy_combined(self):
        parts = []
        for s in (self.sec_ga, getattr(self, 'sec_gsign', None), getattr(self, 'sec_email', None), getattr(self, 'sec_email_for_mail', None), self.sec_gmail, self.sec_prompts, self.sec_yt, getattr(self, 'sec_weirdo', None)):
            if s is None:
                continue
            v = s.get_selected() or ''
            parts.append(v)
        combined = '\n'.join([p for p in parts if p])
        try:
            self.clipboard_clear()
            self.clipboard_append(combined)
            self.logger.log('Copied combined selection to clipboard')
        except Exception:
            self.logger.log('Failed to copy combined')

    def copy_daily_report(self):
        """Копіює Daily Report в буфер обміну"""
        try:
            # Отримуємо значення з полів
            farmer_name = self.report_farmer_name.get().strip() or "Alex"
            day = self.report_day.get().strip() or "00"
            month = self.report_month.get().strip() or "00"
            
            accounts_ready = self.report_accounts_ready.get().strip() or "0"
            emails_registered = self.report_emails_registered.get().strip() or "0"
            in_progress = self.report_in_progress.get().strip() or "0"
            ads = self.report_ads.get().strip() or "0"
            unlocks = self.report_unlocks.get().strip() or "0"
            appeals = self.report_appeals.get().strip() or "0"
            problems = self.report_problems.get().strip() or "0"
            additional = self.report_additional.get().strip() or "0"
            
            # Зберігаємо конфіг якщо галочка активна
            if self.save_report_config.get():
                self.save_daily_report_config()
            
            # Формуємо звіт
            report = f"""Daily report:
👨‍🌾 Фармер: {farmer_name}
📅 Дата – {day}.{month}

✅ Аккаунтів готово – {accounts_ready}
📧 Зареєстровано пошт – {emails_registered}
🛠️ У процесі фарму – {in_progress}
📢 ADS – {ads}

🔓 Пройдено розлогів – {unlocks}
📝 Написано апеляцій – {appeals}
⚠️ Проблеми – {problems}
➕ Дофарм – {additional}"""
            
            # Копіюємо в буфер
            safe_clipboard_operation("set", report)
            self.logger.log('📊 Daily Report скопійовано в буфер обміну')
            
        except Exception as e:
            self.logger.log(f'❌ Помилка копіювання звіту: {e}')

    def save_daily_report_config(self):
        """Зберігає налаштування Daily Report в конфіг"""
        try:
            config = {
                "farmer_name": self.report_farmer_name.get().strip() or "Alex",
                "accounts_ready": self.report_accounts_ready.get().strip() or "0",
                "emails_registered": self.report_emails_registered.get().strip() or "0", 
                "in_progress": self.report_in_progress.get().strip() or "0",
                "ads": self.report_ads.get().strip() or "0",
                "unlocks": self.report_unlocks.get().strip() or "0",
                "appeals": self.report_appeals.get().strip() or "0",
                "problems": self.report_problems.get().strip() or "0",
                "additional": self.report_additional.get().strip() or "0"
            }
            
            config_path = get_config_path("daily_report_config.json")
            with open(config_path, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            
            self.logger.log('💾 Налаштування Daily Report збережені')
            
        except Exception as e:
            self.logger.log(f'❌ Помилка збереження конфігу: {e}')

    def save_farmer_name_to_config(self):
        """Зберігає ім'я фармера в конфіг"""
        # Захист від рекурсії
        if hasattr(self, '_saving_farmer_name') and self._saving_farmer_name:
            return
            
        try:
            self._saving_farmer_name = True
            config_path = get_config_path("daily_report_config.json")
            
            # Завантажуємо існуючий конфіг або створюємо новий
            config = {}
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    config = json.load(f)
            
            # Оновлюємо ім'я фармера
            if hasattr(self, 'report_farmer_name'):
                farmer_name = self.report_farmer_name.get().strip()
                if farmer_name:  # Зберігаємо тільки якщо не пусте
                    # Перевіряємо, чи змінилося ім'я
                    if config.get("farmer_name") != farmer_name:
                        config["farmer_name"] = farmer_name
                        
                        # Зберігаємо конфіг з правильним кодуванням
                        with open(config_path, "w", encoding="utf-8") as f:
                            json.dump(config, f, indent=2, ensure_ascii=False)
                        
        except Exception as e:
            if hasattr(self, 'logger'):
                self.logger.log(f'❌ Помилка збереження імені: {e}')
        finally:
            self._saving_farmer_name = False

    def load_daily_report_config(self):
        """Завантажує збережені налаштування Daily Report"""
        try:
            config_path = get_config_path("daily_report_config.json")
            if not os.path.exists(config_path):
                # Встановлюємо значення за замовчуванням
                if hasattr(self, 'report_farmer_name'):
                    self.report_farmer_name.delete(0, 'end')
                    self.report_farmer_name.insert(0, "Alex")
                return
            
            with open(config_path, "r", encoding="utf-8") as f:
                config = json.load(f)
            
            # Завантажуємо значення в поля
            if hasattr(self, 'report_farmer_name'):
                self.report_farmer_name.delete(0, 'end')
                self.report_farmer_name.insert(0, config.get("farmer_name", "Alex"))
            
            if hasattr(self, 'report_accounts_ready'):
                self.report_accounts_ready.delete(0, 'end') 
                self.report_accounts_ready.insert(0, config.get("accounts_ready", "0"))
            
            if hasattr(self, 'report_emails_registered'):
                self.report_emails_registered.delete(0, 'end')
                self.report_emails_registered.insert(0, config.get("emails_registered", "0"))
            
            if hasattr(self, 'report_in_progress'):
                self.report_in_progress.delete(0, 'end')
                self.report_in_progress.insert(0, config.get("in_progress", "0"))
            
            if hasattr(self, 'report_ads'):
                self.report_ads.delete(0, 'end')
                self.report_ads.insert(0, config.get("ads", "0"))
            
            if hasattr(self, 'report_unlocks'):
                self.report_unlocks.delete(0, 'end')
                self.report_unlocks.insert(0, config.get("unlocks", "0"))
            
            if hasattr(self, 'report_appeals'):
                self.report_appeals.delete(0, 'end')
                self.report_appeals.insert(0, config.get("appeals", "0"))
            
            if hasattr(self, 'report_problems'):
                self.report_problems.delete(0, 'end')
                self.report_problems.insert(0, config.get("problems", "0"))
            
            if hasattr(self, 'report_additional'):
                self.report_additional.delete(0, 'end')
                self.report_additional.insert(0, config.get("additional", "0"))
            
            self.logger.log('📁 Налаштування Daily Report завантажені')
            
        except Exception as e:
            self.logger.log(f'❌ Помилка завантаження конфігу: {e}')

    def load_csv_config(self):
        """Завантажує збережений шлях до CSV файлу або використовує за замовчуванням"""
        try:
            config_path = get_config_path("csv_config.json")
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    config = json.load(f)
                last_csv_path = config.get("last_csv_path", "")
                
                # Перевіряємо чи існує збережений файл
                if last_csv_path and os.path.exists(last_csv_path):
                    self.logger.log(f'📁 Завантажено збережений CSV: {os.path.basename(last_csv_path)}')
                    return last_csv_path
            
            # Шукаємо файл за замовчуванням в директорії програми
            default_path = CSV_PATH
            if os.path.exists(default_path):
                self.logger.log(f'📁 Використано CSV за замовчуванням: {os.path.basename(default_path)}')
                return default_path
            
            # Якщо файл за замовчуванням не знайдено, повертаємо його шлях все одно
            self.logger.log(f'⚠️ CSV файл не знайдено, використано шлях: {os.path.basename(default_path)}')
            return default_path
            
        except Exception as e:
            self.logger.log(f'❌ Помилка завантаження CSV конфігу: {e}')
            return CSV_PATH

    def save_csv_config(self, csv_path):
        """Зберігає шлях до останнього вибраного CSV файлу"""
        try:
            config = {
                "last_csv_path": csv_path,
                "last_directory": os.path.dirname(csv_path)
            }
            
            config_path = get_config_path("csv_config.json")
            with open(config_path, "w", encoding="utf-8") as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            
            self.logger.log(f'💾 Збережено шлях CSV: {os.path.basename(csv_path)}')
            
        except Exception as e:
            self.logger.log(f'❌ Помилка збереження CSV конфігу: {e}')

    def get_csv_initial_directory(self):
        """Повертає початкову директорію для файлового діалогу"""
        try:
            # Спочатку перевіряємо директорію програми
            program_dir = WORKDIR
            if os.path.exists(program_dir):
                csv_files_in_program_dir = [f for f in os.listdir(program_dir) if f.endswith('.csv')]
                if csv_files_in_program_dir:
                    return program_dir
            
            # Потім перевіряємо збережену директорію
            config_path = get_config_path("csv_config.json")
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    config = json.load(f)
                last_directory = config.get("last_directory", "")
                if last_directory and os.path.exists(last_directory):
                    return last_directory
            
            # За замовчуванням повертаємо директорію програми
            return program_dir
            
        except Exception as e:
            self.logger.log(f'❌ Помилка визначення початкової директорії: {e}')
            return WORKDIR

    def open_csv_folder(self):
        """Відкриває папку з поточним CSV файлом"""
        try:
            if hasattr(self, 'csv_path') and self.csv_path:
                csv_folder = os.path.dirname(self.csv_path)
                if os.path.exists(csv_folder):
                    # Відкриваємо папку залежно від ОС
                    if os.name == 'nt':  # Windows
                        os.startfile(csv_folder)
                    elif sys.platform == 'darwin':  # macOS
                        subprocess.call(('open', csv_folder))
                    else:  # Linux
                        subprocess.call(('xdg-open', csv_folder))
                    
                    self.logger.log(f'📁 Відкрито папку: {os.path.basename(csv_folder)}')
                else:
                    self.logger.log(f'❌ Папка не існує: {csv_folder}')
            else:
                self.logger.log('❌ CSV файл не обрано')
                
        except Exception as e:
            self.logger.log(f'❌ Помилка відкриття папки: {e}')

    def save_profile_all(self):
        profile = {}
        for s in (self.sec_ga, getattr(self, 'sec_gsign', None), getattr(self, 'sec_email', None), getattr(self, 'sec_email_for_mail', None), self.sec_gmail, self.sec_prompts, self.sec_yt, getattr(self, 'sec_weirdo', None)):
            if s is None:
                continue
            profile[s.title] = s.get_selected()
        path = filedialog.asksaveasfilename(defaultextension='.json', filetypes=[('JSON','*.json')])
        if not path:
            return
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(profile, f, ensure_ascii=False, indent=2)
        self.logger.log(f'Saved profile to {os.path.basename(path)}')

    # -------------------- Day log functionality --------------------

    def _sample(self, col_name, min_n, max_n):
        vals = (self.data.get(col_name) or
                self.data.get(col_name.lower()) or
                self.data.get(col_name.title()) or
                self.data.get(col_name.upper()) or [])
        
        # Генеруємо кількість елементів для показу
        n = random.randint(min_n, max_n)
        
        if not vals:
            # Якщо дані не знайдені, повертаємо один запис з загальною кількістю
            return [f"{col_name} {n} штуки"]
        
        if len(vals) >= n:
            vals_shuf = list(vals)
            random.shuffle(vals_shuf)
            return vals_shuf[:n]
        else:
            return [random.choice(vals) for _ in range(n)]

    def _add_typos_to_text(self, text):
        """Додає 2-5 граматичних помилок в текст для унікальності з випадковими символами a-z, 0-9"""
        if not text or len(text) < 10:
            return text
        
        words = text.split()
        if len(words) < 2:
            return text
        
        # Збільшуємо кількість помилок: 2-5 залежно від довжини тексту
        max_typos = min(5, max(2, len(words) // 3))
        num_typos = random.randint(2, max_typos)
        
        if num_typos == 0:
            return text
        
        # Вибираємо випадкові слова для помилок (не перше і не останнє)
        modifiable_indices = list(range(1, len(words) - 1)) if len(words) > 2 else list(range(len(words)))
        if not modifiable_indices:
            return text
            
        selected_indices = random.sample(modifiable_indices, min(num_typos, len(modifiable_indices)))
        
        for idx in selected_indices:
            word = words[idx]
            if len(word) > 3:  # Тільки довгі слова
                typo_type = random.choice([1, 2, 3, 4])  # 4 типи помилок
                
                if typo_type == 1:  # Випадкова вставка символів a-z, 0-9
                    char_idx = random.randint(1, len(word) - 1)
                    chars = list(word)
                    
                    # Випадково обираємо тип символа для вставки
                    insert_type = random.choice(['letter', 'number'])
                    if insert_type == 'letter':
                        random_char = random.choice('abcdefghijklmnopqrstuvwxyz')
                    else:
                        random_char = random.choice('0123456789')
                    
                    chars.insert(char_idx, random_char)
                    words[idx] = ''.join(chars)
                
                elif typo_type == 2:  # Перестановка сусідніх символів
                    if len(word) > 3:
                        char_idx = random.randint(1, len(word) - 3)
                        chars = list(word)
                        chars[char_idx], chars[char_idx + 1] = chars[char_idx + 1], chars[char_idx]
                        words[idx] = ''.join(chars)
                
                elif typo_type == 3:  # Заміна символа на випадковий a-z, 0-9
                    char_idx = random.randint(1, len(word) - 2)
                    chars = list(word)
                    
                    # Випадково заміняємо символ на літеру або цифру
                    replace_type = random.choice(['letter', 'number'])
                    if replace_type == 'letter':
                        chars[char_idx] = random.choice('abcdefghijklmnopqrstuvwxyz')
                    else:
                        chars[char_idx] = random.choice('0123456789')
                    
                    words[idx] = ''.join(chars)
                
                elif typo_type == 4:  # Подвоєння символа
                    char_idx = random.randint(1, len(word) - 2)
                    chars = list(word)
                    chars.insert(char_idx, chars[char_idx])
                    words[idx] = ''.join(chars)
        
        return ' '.join(words)

    def _sample_with_typos(self, col_name, min_n, max_n):
        """Версія _sample з граматичними помилками для Prompts"""
        items = self._sample(col_name, min_n, max_n)
        
        # Додаємо помилки тільки для Prompts
        if col_name.lower() == 'prompts':
            return [self._add_typos_to_text(item) for item in items]
        
        return items

    def _build_day_entries(self, day_label):
        """Return grouped entries [(title, [items...]), ...] for the given day_label."""
        entries = []
        
        # Перевіряємо чи активний режим 3 days farm
        is_three_days_mode = hasattr(self, 'three_days_var') and self.three_days_var.get()
        
        if is_three_days_mode:
            return self._build_three_days_entries(day_label)
        
        if day_label == 'Day 1':
            # Day 1: авторизації через Google + реєстрації + підписки на newsletters
            entries.append(('🔐 Авторизації через Google', self._sample_with_typos('google sign', 3, 4)))
            entries.append(('📧 Підписки на newsletters', self._sample('Email Subscription', 2, 3)))
            
            # Створити 3 Google документи
            entries.append(('📄 Створити Google документи', [
                "Створити Google Sheets документ",
                "Створити Google Docs документ", 
                "Створити Google Slides презентацію"
            ]))
            
            # Відправити листи
            entries.append(('✉️ Надсилання листів', [
                "Відправити лист на іншу пошту з таблиці",
                "Відправити другий лист на іншу пошту з таблиці",
                "Відповісти на вхідний лист"
            ]))
            
            # Google News підписки
            entries.append(('📰 Google News підписки', [
                'Підписатися на ключове слово "war"',
                'Підписатися на ключове слово "Trump"',
                'Підписатися на ключове слово "football"'
            ]))
            
            # Google Business Profile
            entries.append(('🏢 Google Business Profile', [
                "Створити Google Business Profile (рандомна USA інформація)"
            ]))

        elif day_label == 'Day 2':
            # Day 2: авторизації + YouTube + документи + календар + диск + листи + Google Ads
            entries.append(('🔐 Авторизації через Google', self._sample_with_typos('google sign', 2, 3)))
            entries.append(('📧 Підписки на newsletters', self._sample('Email Subscription', 1, 2)))
            
            # YouTube активність
            entries.append(('🎬 YouTube активність', [
                "Створити YouTube канал",
                "Переглянути 1 відео (144p, 7+ хв на фоні)",
                "Поставити лайки на відео",
                "Підписатися на канал"
            ]))
            
            # Google документи
            entries.append(('📄 Створити Google документи', [
                "Створити Google Sheets документ",
                "Створити Google Docs документ",
                "Створити Google Slides презентацію"
            ]))
            
            # Google Calendar
            entries.append(('📅 Google Calendar', [
                "Створити подію в календарі на завтра",
                "Створити другу подію в календарі на завтра"
            ]))
            
            # Google Drive
            entries.append(('💾 Google Drive', [
                "Завантажити документ (ORG формат) на Drive",
                "Завантажити другий документ на Drive"
            ]))
            
            # Листи
            entries.append(('✉️ Надсилання листів', [
                "Відправити лист на іншу пошту",
                "Відповісти на вхідний лист"
            ]))
            
            # Google Ads
            entries.append(('💰 Google Ads', [
                "Створити Google Ads через сторонній сайт (YouTube реклама або пошук)",
                "Прив'язати карту до Google Ads",
                "Подати документи USA організації",
                "Поповнити рахунок на $10"
            ]))

        elif day_label == 'Day 3':
            # Day 3: сайти + авторизації + YouTube + документи + листи + Ads + Sites + Analytics
            entries.append(('🌐 Відкриття сайтів', self._sample_with_typos('Prompts', 5, 6)))
            
            entries.append(('🔐 Авторизації і реєстрації', [
                "Реєстрація в інтернет-магазині #1",
                "Реєстрація в інтернет-магазині #2",
                "Реєстрація в інтернет-магазині #3"
            ]))
            
            # YouTube
            entries.append(('🎬 YouTube активність', [
                "Переглянути 1 відео (144p, 7+ хв на фоні)",
                "Поставити лайки та підписки"
            ]))
            
            # Google документи
            entries.append(('📄 Створити Google документи', [
                "Створити Google Sheets документ",
                "Створити Google Docs документ",
                "Створити Google Slides презентацію"
            ]))
            
            # Листи (загалом 5+ відповідей)
            entries.append(('✉️ Надсилання листів', [
                "Відправити лист на іншу пошту",
                "Відповісти на вхідний лист",
                "Загалом потрібно 5+ відповідей"
            ]))
            
            # Перевірка Ads
            entries.append(('💰 Перевірка Google Ads', [
                "Перевірити: Верифікація пройшла?",
                "Перевірити: Немає бану?", 
                "Перевірити: Документи OK?"
            ]))
            
            # Google Sites
            entries.append(('🌍 Google Sites', [
                "Створити простий Google сайт"
            ]))
            
            # Google Analytics
            entries.append(('📊 Google Analytics', [
                "Створити Google Analytics",
                "Створити конверсію в Google ADS",
                "Додати код конверсії в опис і таблицю"
            ]))

        elif day_label == 'Day 4':
            # Day 4: кукі + авторизації + YouTube + Forms + Keep + листи + Ads
            entries.append(('🍪 Нагулювання кукі', self._sample_with_typos('Prompts', 5, 6)))
            
            entries.append(('🔐 Авторизації через Google', self._sample_with_typos('google sign', 3, 4)))
            entries.append(('📧 Підписки на newsletters', self._sample('Email Subscription', 2, 3)))
            
            # YouTube активність розширена
            entries.append(('🎬 YouTube активність', [
                "Переглянути 1 відео (144p, 7+ хв на фоні)",
                "Поставити лайки",
                "Підписатися на канал",
                "Написати коментар до відео"
            ]))
            
            # Google Forms
            entries.append(('📝 Google Forms', [
                "Створити першу форму (довільна тематика)",
                "Створити другу форму (довільна тематика)"
            ]))
            
            # Google Keep
            entries.append(('📒 Google Keep', [
                "Створити нотатку зі списком завдань",
                "Створити нотатку з ідеями",
                "Додати нагадування до нотаток"
            ]))
            
            # Листи та поділитися документами
            entries.append(('✉️ Надсилання листів', [
                "Відправити лист на іншу пошту",
                "Відповісти на вхідний лист",
                "Поділитися Google документом з іншою поштою з таблиці"
            ]))
            
            # Перевірка Ads
            entries.append(('💰 Перевірка Google Ads', [
                "Перевірити стан рекламних кампаній"
            ]))

        elif day_label == 'Day 5':
            # Day 5: кукі + авторизації + Maps + документи + листи + фінальна перевірка
            entries.append(('🍪 Нагулювання кукі', self._sample_with_typos('Prompts', 5, 6)))
            entries.append(('🔮 Дивні запити', self._sample_with_typos('Weirdo', 1, 1)))
            
            entries.append(('🔐 Авторизації через Google', self._sample_with_typos('google sign', 2, 3)))
            
            # Google Maps
            entries.append(('🗺️ Google Maps', [
                "Додати в улюблене місце в USA (стадіон)",
                "Додати в улюблене друге місце в USA",
                "Прокласти маршрут між локаціями"
            ]))
            
            # Документи та календар
            entries.append(('📄 Створити Google документи', [
                "Створити 2 Google Sheets документи",
                "Створити Google Docs документ",
                "Створити Google Slides презентацію",
                "Додати події в Calendar"
            ]))
            
            # Листи та обмін
            entries.append(('✉️ Надсилання листів', [
                "Відправити лист на іншу пошту",
                "Відповісти на вхідний лист",
                "Поділитися альбомом з Photos",
                "Поділитися нотаткою з Keep"
            ]))
            
            # Додаємо обов'язковий пункт для створення ADS
            entries.append(('Create ADS Full', 
                          ['После создания и активации добавляем в адс разные сервисы (ютуб, гугл серч, мерчант и другие)']))
        
        # shuffle order of blocks
        random.shuffle(entries)
        
        # Додаємо фінальну перевірку тільки для Day 5 в КІНЦІ (після shuffle)
        if day_label == 'Day 5':
            entries.append(('🔍 ФІНАЛЬНА ПЕРЕВІРКА', [
                "Перевірити всі сервіси (Forms, Keep, Sites, Maps)",
                "Дофармити листи/відповіді до 150+ листів",
                "Переконатися що є 5+ відповідей",
                "Перевірити на бани/підозрілі активності",
                "Переконатися що всі завдання виконані"
            ]))
        
        return entries

    def _build_three_days_entries(self, day_label):
        """Return grouped entries for 3 days farm mode"""
        entries = []
        
        if day_label == 'Day 1':
            # День 1: паралельна генерація + базова активність
            entries.append(('🔐 Додати 2FA', [
                "Додати аутентифікатор (Google Authenticator/Authy)",
                "Згенерувати та зберегти backup коди",
                "Підтвердити активності в листі від Google"
            ]))
            
            entries.append(('👤 Персоналізація', [
                "Встановити аватарку профілю Gmail",
                "Встановити фон на пошті",
                "Налаштувати підпис електронної пошти"
            ]))
            
            entries.append(('🍪 Нагуляти кукі', self._sample_with_typos('Prompts', 5, 6)))
            
            entries.append(('🔐 Авторизації через Google', self._sample('Email Subscription', 2, 3)))
            
            entries.append(('📄 Google документи', [
                "Створити Google Sheets документ",
                "Створити Google Docs документ", 
                "Створити Google Slides презентацію"
            ]))
            
            entries.append(('✉️ Email активність', [
                "Відправити 1-2 листи на іншу пошту з таблиці",
                "Відповісти на вхідні листи"
            ]))
            
            entries.append(('📰 Google News підписки', [
                'Підписатися на "war" (англійською)',
                'Підписатися на "Trump" (англійською)',
                'Підписатися на "football" (англійською)'
            ]))
            
            entries.append(('🏢 Google Business Profile', [
                "Створити Business Profile з рандомною інформацією США"
            ]))
            
        elif day_label == 'Day 2':
            # День 2: поглиблення + YouTube + Ads
            entries.append(('🔐 Авторизації через Google', self._sample('Email Subscription', 2, 3)))
            
            entries.append(('🎥 YouTube активність', [
                "Створити YouTube канал",
                "Переглянути 1 відео (144p, 7+ хв на фоні)",
                "Поставити лайки та підписки"
            ]))
            
            entries.append(('📄 Google документи', [
                "Створити 3 Google документи (Sheets, Docs, Slides)"
            ]))
            
            entries.append(('📅 Google Calendar', [
                "Створити 1-2 події на завтра"
            ]))
            
            entries.append(('💾 Google Drive', [
                "Завантажити документи організації",
                "Завантажити фізичні документи"
            ]))
            
            entries.append(('✉️ Email активність', [
                "Відправити/відповісти на 1-2 листи"
            ]))
            
            entries.append(('💰 Google Ads', [
                "Створити Google Ads через сторонній сайт",
                "Прив'язати карту до аккаунту",
                "Подати документи USA організації",
                "Поповнити рахунок на $10"
            ]))
            
            entries.append(('📊 Google Analytics', [
                "Створити Google Analytics аккаунт"
            ]))
            
            entries.append(('🎯 Конверсія Google Ads', [
                "Створити конверсію в Google ADS",
                "Додати код конверсії в опис профілю",
                "Записати частину коду в таблицю"
            ]))
            
        elif day_label == 'Day 3':
            # День 3: фінальна активність + перевірки
            entries.append(('🍪 Нагуляти кукі', self._sample_with_typos('Prompts', 5, 6)))
            
            entries.append(('🔐 Авторизації', [
                "Зробити 3 авторизації",
                "Реєстрація на 2-3 сайтах"
            ]))
            
            entries.append(('🎥 YouTube активність', [
                "Переглянути 1 відео (144p, 7+ хв на фоні)",
                "Поставити лайки та підписки"
            ]))
            
            entries.append(('📄 Google документи', [
                "Створити 3 Google документи (Sheets, Docs, Slides)"
            ]))
            
            entries.append(('✉️ Email активність', [
                "Відправити/відповісти на 1-2 листи",
                "Загалом має бути 5+ відповідей"
            ]))
            
            entries.append(('💰 Перевірка Ads', [
                "Перевірити чи пройшла верифікація",
                "Переконатися що немає бану",
                "Перевірити статус документів"
            ]))
            
            entries.append(('📅 Calendar події', [
                "Додати події в Calendar, якщо не робили раніше"
            ]))
            
            entries.append(('🔍 ФІНАЛЬНА ПЕРЕВІРКА', [
                "Перевірити листи",
                "Перевірити персоналізацію (аватарка, фон)",
                "Перевірити 2FA налаштування",
                "Дофармити якщо щось бракує",
                "Заповнити таблицю",
                "Підготувати до видачі"
            ]))
        
        # Рандомізуємо порядок секцій (крім фінальної перевірки)
        if day_label == 'Day 3':
            # Залишаємо фінальну перевірку завжди останньою
            final_check = None
            other_entries = []
            
            for entry in entries:
                if '🔍 ФІНАЛЬНА ПЕРЕВІРКА' in entry[0]:
                    final_check = entry
                else:
                    other_entries.append(entry)
            
            # Рандомізуємо порядок всіх секцій крім фінальної
            random.shuffle(other_entries)
            
            # Збираємо назад з фінальною перевіркою в кінці
            entries = other_entries
            if final_check:
                entries.append(final_check)
        else:
            # Для Day 1 і Day 2 також рандомізуємо порядок
            random.shuffle(entries)
        
        return entries

    # Inline day-list helpers for the embedded day listbox
    def _day_generate_to_list(self, day_label):
        # Use grouped entries from _build_day_entries so inline view matches popup
        entries = self._build_day_entries(day_label)
        
        # Перевіряємо чи активний 4-windows режим
        four_windows_active = getattr(self, 'four_windows_var', None) and self.four_windows_var.get()
        
        if four_windows_active:
            # 4-windows режим - заповнюємо всі 4 вікна
            self._fill_four_day_log_windows(day_label, entries)
        else:
            # Стандартний режим - заповнюємо одне вікно
            self._fill_single_day_log_window(day_label, entries)

    def _fill_single_day_log_window(self, day_label, entries):
        """Заповнює одне Day Log вікно в стандартному режимі"""
        flat = []
        line_map = []
        try:
            self.day_log_box.configure(state='normal')
            self.day_log_box.delete('1.0', tk.END)
            self.day_log_box.insert(tk.END, f"{day_label}\n")
            cur_line = 2
            for title, items in entries:
                self.day_log_box.insert(tk.END, f"{title}:\n")
                cur_line += 1
                if items:
                    for it in items:
                        self.day_log_box.insert(tk.END, f" - {it}\n")
                        flat.append(it)
                        line_map.append(cur_line)
                        cur_line += 1
                else:
                    self.day_log_box.insert(tk.END, " (no items)\n")
                    cur_line += 1
                # blank line after block
                self.day_log_box.insert(tk.END, "\n")
                cur_line += 1
            # store flattened items and mapping for Next/Copy
            self._last_day_items = flat
            self._day_line_map = line_map
            self._last_day_index = 0
            # highlight first flat item if exists
            self.day_log_box.tag_configure('highlight', background='#ffff99')
            if line_map:
                ln = line_map[0]
                self.day_log_box.tag_add('highlight', f'{ln}.0', f'{ln}.0 lineend')
                self.day_log_box.see(f'{ln}.0')
            self.day_log_box.configure(state='disabled')
        except Exception:
            pass
        self.logger.log(f'Generated inline {day_label} entries: {len(self._last_day_items)}')
        
        # Автоматично переміщуємо секції відповідно до згенерованого логу
        self._reorder_sections_by_day_log(entries)
        
        # 🔄 Синхронізація: автоматично натискаємо Generate All після Generate Day
        try:
            self.generate_all()
            self.logger.log(f'🔄 Auto-triggered Generate All after {day_label}')
        except Exception as e:
            self.logger.log(f'⚠️ Error auto-triggering Generate All: {e}')

    def _fill_four_day_log_windows(self, day_label, entries):
        """Заповнює 4 Day Log вікна унікальними варіантами завдань"""
        try:
            print(f"🔢 Заповнення 4-х Day Log вікон для {day_label}")
            
            # Для кожного з 4-х вікон створюємо унікальні варіанти
            for window_idx, day_log in enumerate(self.day_log_boxes):
                day_log.configure(state='normal')
                day_log.delete('1.0', tk.END)
                
                # Заголовок для кожного акаунту
                day_log.insert(tk.END, f"=== Account {window_idx + 1} - {day_label} ===\n\n")
                
                # Створюємо унікальні entries для кожного вікна
                window_entries = self._create_unique_entries_for_window(entries, window_idx)
                
                # Заповнюємо завданнями з варіаціями
                for title, items in window_entries:
                    day_log.insert(tk.END, f"{title}:\n")
                    
                    if items:
                        for item in items:
                            day_log.insert(tk.END, f" - {item}\n")
                    else:
                        day_log.insert(tk.END, " (no items)\n")
                    
                    day_log.insert(tk.END, "\n")  # Пустий рядок після блоку
                
                day_log.configure(state='disabled')
                print(f"  ✅ Заповнено вікно {window_idx + 1}")
            
            print("✅ Всі 4 Day Log вікна заповнені")
            
            # НЕ переміщуємо секції щоб не очищати згенерований контент
            # self._reorder_sections_by_day_log(entries)
            
            # Синхронізуємо секції з Day Log завданнями
            self._sync_sections_with_day_log()
            
            # 🔄 Синхронізація: автоматично натискаємо Generate All після Generate Day
            try:
                self.generate_all()
                self.logger.log(f'🔄 Auto-triggered Generate All after {day_label} (4-windows)')
            except Exception as e:
                self.logger.log(f'⚠️ Error auto-triggering Generate All (4-windows): {e}')
            
        except Exception as e:
            print(f"❌ Помилка заповнення 4-х Day Log вікон: {str(e)}")

    def _create_unique_entries_for_window(self, entries, window_idx):
        """Створює унікальні entries (секції) для кожного вікна"""
        try:
            import random
            import copy
            
            # Створюємо глибоку копію entries
            window_entries = copy.deepcopy(entries)
            
            if window_idx == 0:
                # Вікно 1 - оригінальні завдання
                return window_entries
                
            elif window_idx == 1:
                # Вікно 2 - пріоритетні завдання зі змішаним порядком секцій
                random.shuffle(window_entries)
                for i, (title, items) in enumerate(window_entries):
                    if items:
                        priority_suffixes = [" ⭐ Priority", " 🔥 Important", " ⚡ Urgent", " 📌 First", " 🎯 Focus"]
                        modified_items = []
                        for j, item in enumerate(items):
                            suffix = priority_suffixes[j % len(priority_suffixes)]
                            modified_items.append(f"{item}{suffix}")
                        window_entries[i] = (title, modified_items)
                return window_entries
                
            elif window_idx == 2:
                # Вікно 3 - розширені завдання з частковим змішуванням
                # Змішуємо тільки частину секцій
                mid_point = len(window_entries) // 2
                first_half = window_entries[:mid_point]
                second_half = window_entries[mid_point:]
                random.shuffle(second_half)
                window_entries = first_half + second_half
                
                for i, (title, items) in enumerate(window_entries):
                    if items:
                        extended_suffixes = [" + Extended", " + Advanced", " + Plus", " + Enhanced", " + Pro"]
                        modified_items = []
                        for j, item in enumerate(items):
                            suffix = extended_suffixes[j % len(extended_suffixes)]
                            modified_items.append(f"{item}{suffix}")
                        window_entries[i] = (title, modified_items)
                return window_entries
                
            elif window_idx == 3:
                # Вікно 4 - альтернативні завдання зі зворотним порядком
                window_entries.reverse()
                for i, (title, items) in enumerate(window_entries):
                    if items:
                        # Зворотний порядок елементів в кожній секції
                        reversed_items = list(reversed(items))
                        alt_suffixes = [" → Alternative", " → Backup", " → Option B", " → Extra", " → Variant"]
                        modified_items = []
                        for j, item in enumerate(reversed_items):
                            suffix = alt_suffixes[j % len(alt_suffixes)]
                            modified_items.append(f"{item}{suffix}")
                        window_entries[i] = (title, modified_items)
                return window_entries
                
            return window_entries
            
        except Exception as e:
            print(f"❌ Помилка створення унікальних entries: {str(e)}")
            return entries

    def _sync_sections_with_day_log(self):
        """Синхронізує секції в колонках з завданнями з Day Log вікон"""
        try:
            print("🔄 Синхронізація секцій з Day Log завданнями...")
            
            # Витягуємо завдання з кожного Day Log вікна
            for window_idx, day_log in enumerate(self.day_log_boxes):
                content = day_log.get('1.0', tk.END)
                tasks = self._extract_tasks_from_day_log(content)
                
                # Оновлюємо секції в відповідній колонці
                self._update_column_sections_with_tasks(window_idx, tasks)
                
                print(f"✅ Синхронізовано колонку {window_idx + 1} з {len(tasks)} завданнями")
                
        except Exception as e:
            print(f"❌ Помилка синхронізації секцій з Day Log: {str(e)}")
    
    def _extract_tasks_from_day_log(self, content):
        """Витягує завдання з тексту Day Log"""
        try:
            tasks = []
            lines = content.split('\n')
            
            print(f"🔍 Аналізую {len(lines)} рядків з Day Log...")
            
            current_section = None
            for line in lines:
                original_line = line
                line = line.strip()
                if not line or line.startswith('==='):
                    continue
                
                # Якщо рядок закінчується двокрапкою, це назва секції
                if line.endswith(':'):
                    current_section = line[:-1]  # Прибираємо двокрапку
                    print(f"  📌 Знайдено секцію: '{current_section}'")
                    continue
                
                # Якщо рядок починається з " - " або просто з "- ", це завдання
                if current_section:
                    if line.startswith('- '):
                        task = line[2:]  # Прибираємо "- "
                        tasks.append((current_section, task))
                        print(f"    ✓ Завдання: '{task}'")
                    elif original_line.startswith(' - '):  # Перевіряємо оригінальний рядок
                        task = line  # line вже trimmed
                        tasks.append((current_section, task))
                        print(f"    ✓ Завдання: '{task}'")
            
            print(f"📊 Всього витягнуто {len(tasks)} завдань")
            return tasks
            
        except Exception as e:
            print(f"❌ Помилка витягування завдань: {str(e)}")
            return []
    
    def _update_column_sections_with_tasks(self, column_idx, tasks):
        """Оновлює секції в колонці відповідно до завдань"""
        try:
            frames = [self.sections_frame_1, self.sections_frame_2, 
                     self.sections_frame_3, self.sections_frame_4]
            
            if column_idx >= len(frames):
                return
                
            target_frame = frames[column_idx]
            
            # Групуємо завдання по секціям (зберігаємо порядок появи в Day Log)
            section_tasks = {}
            section_order = []  # Зберігаємо порядок секцій з Day Log
            
            for section_name, task in tasks:
                if section_name not in section_tasks:
                    section_tasks[section_name] = []
                    section_order.append(section_name)  # Додаємо в порядку появи
                section_tasks[section_name].append(task)
            
            # Отримуємо всі віджети секцій
            scrollable_frame = target_frame._scrollable_frame if hasattr(target_frame, '_scrollable_frame') else target_frame
            section_widgets = []
            
            for widget in scrollable_frame.winfo_children():
                if hasattr(widget, 'title') and hasattr(widget, '_refresh_listbox'):
                    section_widgets.append(widget)
            
            # Створюємо мапу: назва секції → віджет
            widget_map = {}
            for widget in section_widgets:
                original_title = getattr(widget, '_original_title', widget.title)
                clean_title = original_title.replace(' (Col 2)', '').replace(' (Col 3)', '').replace(' (Col 4)', '')
                widget_map[clean_title] = widget
            
            # Переміщуємо секції у порядку Day Log
            for idx, section_name in enumerate(section_order):
                # Знаходимо відповідний віджет
                matching_widget = None
                for clean_title, widget in widget_map.items():
                    if self._section_matches(clean_title, section_name):
                        matching_widget = widget
                        break
                
                if matching_widget:
                    # Оновлюємо дані секції
                    matching_widget.items = section_tasks[section_name]
                    matching_widget.index = 0
                    matching_widget._refresh_listbox()
                    
                    # Переміщуємо віджет на нову позицію
                    matching_widget.pack_forget()
                    matching_widget.pack(fill='x', pady=4)
                    
                    print(f"  📋 Оновлено і переміщено '{matching_widget.title}' з {len(section_tasks[section_name])} завданнями на позицію {idx + 1}")
                            
        except Exception as e:
            print(f"❌ Помилка оновлення колонки {column_idx}: {str(e)}")
    
    def _section_matches(self, section_title, task_section):
        """Перевіряє чи відповідає секція завданню"""
        try:
            # Прибираємо (Col X) з назви
            clean_title = section_title.replace(' (Col 2)', '').replace(' (Col 3)', '').replace(' (Col 4)', '')
            
            # Маппінг секцій до завдань Day Log
            # Формат: 'назва_секції': ['завдання з Day Log', ...]
            mappings = {
                'Google Alerts': ['🔔 Google Alerts', 'Google Alerts'],
                'google sign': ['🔐 Авторизації через Google', 'Авторизації через Google', 'Авторизації через google', 'авторизації через google'],
                'Email Subscription': ['📧 Підписки на newsletters', 'Підписки на newsletters', '📧 підписки на newsletters', 'підписки на newsletters'],
                'Email for mail': ['✉️ Email активність', 'Email активність', '✉️ Надсилання листів', 'Надсилання листів', 'email активність', 'надсилання листів'],
                'Gmail answer': ['✉️ Email активність', 'Email активність', '📨 Gmail відповіді', 'Gmail відповіді', 'email активність'],  # Підтягується під Email активність
                'Prompts': ['🍪 Нагуляти кукі', 'Нагуляти кукі', '🍪 Нагулювання кукі', 'Нагулювання кукі', '🌐 Відкриття сайтів', 'Відкриття сайтів', 'Prompts', 'нагуляти кукі', 'відкриття сайтів'],
                'youtube comentary': ['🎥 YouTube коментарі', 'YouTube коментарі', 'youtube коментарі'],
                'Weirdo': ['🔮 Дивні запити', 'Дивні запити', 'Weirdo', 'дивні запити']
            }
            
            if clean_title in mappings:
                # Перевіряємо чи є task_section в списку маппінгів
                result = any(mapping.lower() in task_section.lower() or task_section.lower() in mapping.lower() 
                          for mapping in mappings[clean_title])
                if result:
                    print(f"    ✅ Збіг: '{clean_title}' ← '{task_section}'")
                return result
                
            return False
            
        except Exception as e:
            print(f"❌ Помилка зіставлення секцій: {str(e)}")
            return False

    def _create_window_variants(self, items, window_idx):
        """Створює унікальні варіанти завдань для кожного вікна (застаріла функція)"""
        try:
            if not items:
                return []
            
            import random
            
            # Різні стратегії для кожного вікна
            if window_idx == 0:
                # Вікно 1 - оригінальні завдання в прямому порядку
                return items.copy()
                
            elif window_idx == 1:
                # Вікно 2 - пріоритетні завдання з додатковими мітками
                priority_suffixes = [" ⭐ Priority", " 🔥 Important", " ⚡ Urgent", " 📌 First", " 🎯 Focus"]
                modified = []
                for i, item in enumerate(items):
                    suffix = priority_suffixes[i % len(priority_suffixes)]
                    modified.append(f"{item}{suffix}")
                return modified
                
            elif window_idx == 2:
                # Вікно 3 - змінений порядок + розширені завдання
                shuffled_items = items.copy()
                random.shuffle(shuffled_items)  # Змішуємо порядок
                extended_suffixes = [" + Extended", " + Advanced", " + Plus", " + Enhanced", " + Pro"]
                modified = []
                for i, item in enumerate(shuffled_items):
                    suffix = extended_suffixes[i % len(extended_suffixes)]
                    modified.append(f"{item}{suffix}")
                return modified
                
            elif window_idx == 3:
                # Вікно 4 - зворотний порядок + альтернативні завдання
                reversed_items = list(reversed(items))
                alt_suffixes = [" → Alternative", " → Backup", " → Option B", " → Extra", " → Variant"]
                modified = []
                for i, item in enumerate(reversed_items):
                    suffix = alt_suffixes[i % len(alt_suffixes)]
                    modified.append(f"{item}{suffix}")
                return modified
                
            return items
            
        except Exception as e:
            print(f"❌ Помилка створення варіантів: {str(e)}")
            return items

    def _reorder_sections_by_day_log(self, day_entries):
        """Переміщує секції генераторів відповідно до порядку в Day логу"""
        # Запам'ятовуємо стан 4-windows режиму перед перебудовою
        four_windows_active = getattr(self, 'four_windows_var', None) and self.four_windows_var.get()
        print(f"🔄 _reorder_sections: 4-windows режим {'активний' if four_windows_active else 'неактивний'}")
        
        try:
            # Створюємо маппінг назв з логу до назв секцій
            section_mapping = {
                # Основні маппінги
                '🔐 Авторизації через Google': 'google sign',
                'Авторизації через Google': 'google sign', 
                'Авторизація з гугла': 'google sign',
                
                '🔮 Дивні запити': 'Weirdo',
                'Дивні запити': 'Weirdo',
                'Weirdo': 'Weirdo',
                
                '✉️ Надсилання листів': 'Email for mail',
                'Надсилання листів': 'Email for mail',
                
                '🍪 Нагулювання кукі': 'Prompts',
                'Нагулювання кукі': 'Prompts',
                'Prompts': 'Prompts',
                'Промпти пошуку': 'Prompts',
                
                '📧 Підписки на newsletters': 'Email Subscription',
                'Підписки на newsletters': 'Email Subscription',
                'Email Subscription підпишись на розилку': 'Email Subscription',
                
                '🎬 YouTube активність': 'youtube comentary',
                'YouTube активність': 'youtube comentary',
                'Youtube перегляд відео , коменти': 'youtube comentary',
                'Youtube': 'youtube comentary',
                
                # Інші існуючі маппінги
                'Google News': 'Google News', 
                'Google Store': 'Google Store',
                'Google maps збережені місця': 'Google maps',
                'Google drive закинь пару файлів': 'Google drive', 
                'результати для прогріву ADS on google search': 'ADS create',
                'додаємо Google play store books': 'Google play store books',
                
                # Нові секції з Day 1-5
                '📄 Створити Google документи': 'Google drive',
                '📰 Google News підписки': 'Google News',
                '🏢 Google Business Profile': 'Google Store', 
                '📅 Google Calendar': 'Google drive',
                '💾 Google Drive': 'Google drive',
                '💰 Google Ads': 'ADS create',
                '💰 Перевірка Google Ads': 'ADS create',
                '🌍 Google Sites': 'Google drive',
                '📊 Google Analytics': 'Google drive',
                '📝 Google Forms': 'Google drive',
                '📒 Google Keep': 'Google drive',
                '🗺️ Google Maps': 'Google maps',
                'Create ADS Full': 'ADS create'
            }
            
            # Отримуємо порядок секцій з Day логу
            ordered_sections = []
            for title, items in day_entries:
                if title in section_mapping:
                    section_name = section_mapping[title]
                    # Знаходимо відповідну секцію
                    for section in self.sections_objs:
                        if section.title == section_name:
                            ordered_sections.append(section)
                            break
            
            # Додаємо решту секцій які не були в логу (в кінець)
            for section in self.sections_objs:
                if section not in ordered_sections:
                    ordered_sections.append(section)
            
            # Перебудовуємо UI тільки якщо порядок змінився
            if ordered_sections != self.sections_objs:
                self._rebuild_sections_ui(ordered_sections)
                
        except Exception as e:
            print(f"Error reordering sections: {e}")
        
        # Відновлюємо 4-windows режим якщо він був активний (після перебудови)
        if four_windows_active:
            print("🔧 Відновлюємо 4-windows режим після перебудови секцій")
            self.after(100, lambda: self.restore_four_windows_layout())
            
    def _rebuild_sections_ui(self, new_order):
        """Перебудовує UI секцій у новому порядку"""
        try:
            # Зберігаємо новий порядок
            self.sections_objs = new_order
            
            # Очищуємо контейнери секцій
            for widget in self.sections_frame_left.winfo_children():
                widget.destroy()
            for widget in self.sections_frame_right.winfo_children():
                widget.destroy()
                
            # Перебудовуємо секції у новому порядку
            random_stuff = ['subscribe', 'to', 'newsletter', 'updates', 'news', 'emails', 'register', 'sign up']
            
            for i, section in enumerate(self.sections_objs):
                # Визначаємо чи потрібні append_values для цієї секції
                needs_append = section.title == 'Email Subscription' or section.title.lower().startswith('google sign')
                
                if i < 2:  # Перші 2 секції в лівій колонці  
                    section.pack_forget()
                    if needs_append:
                        new_section = SectionFrame(
                            self.sections_frame_left, 
                            section.title, 
                            section.min_items, 
                            section.max_items, 
                            section.column_values, 
                            self.logger, 
                            append_values=random_stuff,
                            font=self.font_default
                        )
                    else:
                        new_section = SectionFrame(
                            self.sections_frame_left, 
                            section.title, 
                            section.min_items, 
                            section.max_items, 
                            section.column_values, 
                            self.logger, 
                            font=self.font_default
                        )
                    new_section.pack(fill='x', pady=3)
                    self.sections_objs[i] = new_section
                else:  # Решта в правій колонці
                    section.pack_forget() 
                    if needs_append:
                        new_section = SectionFrame(
                            self.sections_frame_right,
                            section.title,
                            section.min_items,
                            section.max_items, 
                            section.column_values,
                            self.logger,
                            append_values=random_stuff,
                            font=self.font_default
                        )
                    else:
                        new_section = SectionFrame(
                            self.sections_frame_right,
                            section.title,
                            section.min_items,
                            section.max_items, 
                            section.column_values,
                            self.logger,
                            font=self.font_default
                        )
                    new_section.pack(fill='x', pady=3)
                    self.sections_objs[i] = new_section
            
            # Оновлюємо посилання на секції після перебудови
            self._update_section_references()
                    
            self.logger.log("Секції переміщені відповідно до Day логу")
            
        except Exception as e:
            print(f"Error rebuilding sections UI: {e}")
            
    def _update_section_references(self):
        """Оновлює посилання на секції після перебудови UI"""
        # Скидаємо всі посилання
        self.sec_ga = None
        self.sec_email = None  
        self.sec_email_for_mail = None
        self.sec_gsign = None
        self.sec_gmail = None
        self.sec_prompts = None
        self.sec_yt = None
        self.sec_weirdo = None
        
        # Встановлюємо нові посилання (використовуємо оригінальні назви)
        for s in self.sections_objs:
            # Використовуємо оригінальну назву якщо є, інакше поточну
            original_title = getattr(s, '_original_title', s.title)
            
            if original_title == 'Google Alerts':
                self.sec_ga = s
            elif original_title == 'Email Subscription':
                self.sec_email = s
            elif original_title == 'Email for mail':
                self.sec_email_for_mail = s
            elif original_title.lower().startswith('google sign'):
                self.sec_gsign = s
            elif original_title == 'Gmail answer':
                self.sec_gmail = s
            elif original_title == 'Prompts':
                self.sec_prompts = s
            elif original_title == 'youtube comentary':
                self.sec_yt = s
            elif original_title == 'Weirdo':
                self.sec_weirdo = s

    def _day_next(self):
        # advance internal day index and update highlight
        if not getattr(self, '_last_day_items', None):
            return
        n = len(self._last_day_items)
        self._last_day_index = (self._last_day_index + 1) % n
        idx = self._last_day_index
        val = self._last_day_items[idx]
        # update highlight in scrolled text
        try:
            self.day_log_box.configure(state='normal')
            # remove old highlight
            self.day_log_box.tag_remove('highlight', '1.0', tk.END)
            # compute line number from mapping if available
            line_no = None
            if getattr(self, '_day_line_map', None) and idx < len(self._day_line_map):
                line_no = self._day_line_map[idx]
            if line_no is None:
                line_no = 2 + idx
            self.day_log_box.tag_add('highlight', f'{line_no}.0', f'{line_no}.0 lineend')
            self.day_log_box.see(f'{line_no}.0')
            self.day_log_box.configure(state='disabled')
        except Exception:
            pass
        self.logger.log(f'Day Log: Next -> {val}')

    def _day_copy(self):
        if not getattr(self, '_last_day_items', None):
            return
        val = self._last_day_items[self._last_day_index]
        try:
            self.clipboard_clear()
            self.clipboard_append(val)
            self.logger.log(f'Day Log: Copied -> {val}')
        except Exception:
            self.logger.log('Day Log: Failed to copy')

    def _on_day_log_click(self, event):
        """Click-to-copy для рядків у Day Log (працює навіть з disabled textbox)"""
        try:
            # Тимчасово активуємо textbox для роботи з подіями
            was_disabled = str(self.day_log_box.cget('state')) == 'disabled'
            if was_disabled:
                self.day_log_box.configure(state='normal')
            
            # Отримуємо позицію кліку
            index = self.day_log_box.index("@%s,%s" % (event.x, event.y))
            line_num = int(float(index))
            
            # Отримуємо текст рядка
            line_start = f"{line_num}.0"
            line_end = f"{line_num}.end"
            line_text = self.day_log_box.get(line_start, line_end).strip()
            
            # Копіюємо тільки якщо це рядок з завданням (починається з " - ")
            if line_text.startswith(" - "):
                task_text = line_text[3:]  # Прибираємо " - "
                try:
                    # Використовуємо безпечну операцію clipboard
                    safe_clipboard_operation("set", task_text)
                    self.logger.log(f'📋 Copied from Day Log: {task_text[:50]}{"..." if len(task_text) > 50 else ""}')
                    
                    # Тимчасове підсвічування скопійованого рядка
                    self.day_log_box.tag_configure('copied', background='#90EE90', foreground='#000000')
                    self.day_log_box.tag_add('copied', line_start, line_end)
                    
                    # Прибираємо підсвічування через 800мс і повертаємо стан
                    def restore_state():
                        try:
                            self.day_log_box.tag_remove('copied', '1.0', tk.END)
                            if was_disabled:
                                self.day_log_box.configure(state='disabled')
                        except:
                            pass
                    self.after(800, restore_state)
                    
                except Exception as e:
                    self.logger.log(f'⚠️ Failed to copy from Day Log: {e}')
                    # Відновлюємо стан навіть при помилці
                    if was_disabled:
                        self.day_log_box.configure(state='disabled')
            else:
                # Якщо клік не на завданні, просто відновлюємо стан
                if was_disabled:
                    self.day_log_box.configure(state='disabled')
            
        except Exception as e:
            # При будь-якій помилці відновлюємо стан textbox
            try:
                if str(self.day_log_box.cget('state')) == 'normal':
                    self.day_log_box.configure(state='disabled')
            except:
                pass

    # -------------------- UI Settings (customtkinter fallback) --------------------

    def open_ui_settings(self):
        popup = ctk.CTkToplevel(self)
        popup.title('UI Settings')
        popup.geometry('450x280')
        popup.attributes('-topmost', True)  # Always on top
        
        # Center the popup window
        popup.transient(self)  # Make it a transient window
        popup.grab_set()  # Make it modal
        
        # Center on parent window
        popup.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (450 // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (280 // 2)
        popup.geometry(f'450x280+{x}+{y}')

        # Appearance mode
        mode_lbl = ctk.CTkLabel(popup, text='Appearance Mode:', font=self.font_default)
        mode_lbl.pack(anchor='w', padx=8, pady=(8,0))
        mode_var = tk.StringVar(value='System')
        mode_cb = ctk.CTkOptionMenu(popup, variable=mode_var, values=['System', 'Dark', 'Light'], font=self.font_default)
        mode_cb.pack(fill='x', padx=8, pady=4)

        # Theme files selection
        theme_files_lbl = ctk.CTkLabel(popup, text='Color Theme:', font=self.font_default)
        theme_files_lbl.pack(anchor='w', padx=8, pady=(8,0))
        theme_files_var = tk.StringVar()
        
        # Перевіряємо чи існує папка themes
        themes_dir = os.path.join(WORKDIR, "themes")
        if os.path.exists(themes_dir):
            theme_files = [f for f in os.listdir(themes_dir) if f.endswith('.json')]
        else:
            theme_files = []
            
        # Додаємо базові теми CustomTkinter
        default_themes = ["blue", "green", "dark-blue"]
        all_themes = default_themes + [f"Custom: {f}" for f in theme_files]
        
        theme_files_cb = ctk.CTkOptionMenu(popup, variable=theme_files_var, values=all_themes, font=self.font_default)
        theme_files_cb.pack(fill='x', padx=8, pady=4)
        theme_files_var.set('blue')  # Базова тема

        # Швидкість скролу
        scroll_speed_lbl = ctk.CTkLabel(popup, text='Швидкість скролу (1=повільно, 5=швидко):', font=self.font_default)
        scroll_speed_lbl.pack(anchor='w', padx=8, pady=(8,0))
        
        # Frame для слайдера швидкості
        scroll_frame = ctk.CTkFrame(popup)
        scroll_frame.pack(fill='x', padx=8, pady=4)
        
        # Поточне значення
        current_speed = self.scroll_speed
        speed_value_lbl = ctk.CTkLabel(scroll_frame, text=f"Поточна: {current_speed}x", font=self.font_default)
        speed_value_lbl.pack(side='right', padx=5)
        
        # Слайдер швидкості
        speed_var = tk.DoubleVar(value=current_speed)
        speed_slider = ctk.CTkSlider(scroll_frame, from_=1, to=10, number_of_steps=9, 
                                   variable=speed_var, width=200)
        speed_slider.pack(side='left', padx=5, fill='x', expand=True)
        
        # Оновлення лейбла при зміні слайдера
        def update_speed_label(value):
            speed_value_lbl.configure(text=f"Поточна: {int(value)}x")
        speed_slider.configure(command=update_speed_label)

        def apply_theme_selection():
            selected_theme = theme_files_var.get()
            if selected_theme in default_themes:
                # Базова тема
                ctk.set_default_color_theme(selected_theme)
                ctk.set_appearance_mode(mode_var.get())
                self.save_current_theme({"theme_name": selected_theme, "mode": mode_var.get()})
                self.logger.log(f'Applied built-in theme: {selected_theme}')
                messagebox.showinfo("Theme Applied", f"Theme '{selected_theme}' applied! Restart app for full effect.")
            elif selected_theme.startswith("Custom: "):
                # Кастомна тема
                theme_file = selected_theme[8:]  # Прибираємо "Custom: "
                theme_path = os.path.join(themes_dir, theme_file)
                self.apply_ctk_settings(mode_var.get(), theme_path)
                self.save_current_theme({"theme_path": theme_path, "mode": mode_var.get()})
            
            # Застосовуємо швидкість скролу
            new_speed = int(speed_var.get())
            self.scroll_speed = new_speed
            self.save_scroll_speed_config(new_speed)
            self.logger.log(f'🔄 Швидкість скролу змінено на: {new_speed}x')
            
            popup.destroy()
        
        # Frame для кнопок
        buttons_frame = ctk.CTkFrame(popup, fg_color="transparent")
        buttons_frame.pack(pady=12, fill="x", padx=8)
        
        # Кнопка Features Settings (фіолетова)
        features_btn = ctk.CTkButton(
            buttons_frame, 
            text='⚙️ Features Settings', 
            command=lambda: [popup.destroy(), self.open_features_settings()],
            width=180, 
            height=35, 
            corner_radius=8, 
            font=self.font_default,
            fg_color="#9b59b6",  # Фіолетовий
            hover_color="#8e44ad"
        )
        features_btn.pack(side='left', padx=5)
        
        # Кнопка Apply (зелена)
        apply_btn = ctk.CTkButton(
            buttons_frame, 
            text='Apply', 
            command=apply_theme_selection, 
            width=120, 
            height=35, 
            corner_radius=8, 
            font=self.font_default,
            fg_color="green",
            hover_color="darkgreen"
        )
        apply_btn.pack(side='right', padx=5)

    def apply_ctk_settings(self, mode, theme_path):
        try:
            # Налаштування режиму вигляду
            ctk.set_appearance_mode(mode)
            
            if theme_path and os.path.exists(theme_path):
                # Нормалізація шляху
                full_path = os.path.abspath(theme_path).replace("\\", "/")
                
                # Установка теми CustomTkinter
                ctk.set_default_color_theme(full_path)
                
                # Повідомлення користувача про необхідність перезапуску
                messagebox.showinfo("Theme Applied", "Theme has been applied! Please restart the application to see full changes.")
                
            self.logger.log(f'Applied CTk settings: mode={mode}, theme={os.path.basename(theme_path) if theme_path else "default"}')
            
        except Exception as e:
            messagebox.showerror("Theme Error", f"Failed to load theme: {str(e)}\nUsing default theme.")
            ctk.set_default_color_theme("blue")
            self.logger.log(f'Theme load failed, using default: {str(e)}')

    # -------------------- Features Settings --------------------
    
    def open_features_settings(self):
        """Відкриває діалог налаштування функцій"""
        popup = ctk.CTkToplevel(self)
        popup.title('⚙️ Features Settings')
        popup.geometry('750x700')
        popup.attributes('-topmost', True)
        
        # Центруємо вікно
        popup.transient(self)
        popup.grab_set()
        
        popup.update_idletasks()
        x = self.winfo_x() + (self.winfo_width() // 2) - (750 // 2)
        y = self.winfo_y() + (self.winfo_height() // 2) - (700 // 2)
        popup.geometry(f'750x700+{x}+{y}')
        
        # Заголовок
        title_frame = ctk.CTkFrame(popup, fg_color="transparent")
        title_frame.pack(pady=15, padx=20, fill="x")
        
        title = ctk.CTkLabel(
            title_frame,
            text="⚙️ Features Settings Manager",
            font=ctk.CTkFont(size=22, weight="bold")
        )
        title.pack()
        
        subtitle = ctk.CTkLabel(
            title_frame,
            text="Налаштування функцій New soft 3.0",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        subtitle.pack(pady=(5, 0))
        
        # Scrollable Frame для налаштувань
        scroll_frame = ctk.CTkScrollableFrame(popup, width=700, height=450)
        scroll_frame.pack(pady=10, padx=20, fill="both", expand=True)
        
        checkboxes = {}
        
        # Секція Generators
        self._create_features_section(scroll_frame, checkboxes, 
            "📊 Generators - Плитки генерації", "generators", [
                ("google_alerts", "🚨 Google Alerts"),
                ("google_sign", "📝 Google Sign"),
                ("email_subscription", "📧 Email Subscription"),
                ("email_for_mail", "📬 Email for mail"),
                ("gmail_answer", "✉️ Gmail Answer"),
                ("prompts", "💬 Prompts"),
                ("youtube_commentary", "🎥 YouTube Commentary"),
                ("weirdo", "🤪 Weirdo"),
            ])
        
        # Секція Daily Report
        self._create_features_section(scroll_frame, checkboxes,
            "📈 Daily Report - Елементи звіту", "daily_report", [
                ("enabled", "📊 Відображати Daily Report"),
                ("farmer_name", "👨‍🌾 Поле імені фармера"),
                ("date", "📅 Дата"),
                ("copy_button", "📋 Кнопка копіювання"),
            ])
        
        # Секція Utilities
        self._create_features_section(scroll_frame, checkboxes,
            "🛠️ Utilities - Інструменти", "utilities", [
                ("google_backup_codes", "🔑 Резервні коди Google"),
                ("2fa_generator", "🔐 Генератор 2FA кодів"),
                ("password_generator", "🔒 Генератор паролів"),
            ])
        
        # Секція UI Sections
        self._create_features_section(scroll_frame, checkboxes,
            "🎨 UI Sections - Секції інтерфейсу", "ui_sections", [
                ("action_log", "📝 Лог дій (права колонка)"),
            ])
        
        # Секція Windows Mode
        windows_mode_var = self._create_windows_mode_section(scroll_frame)
        
        # Кнопки управління
        buttons_frame = ctk.CTkFrame(popup, fg_color="transparent")
        buttons_frame.pack(pady=15, padx=20, fill="x")
        
        def save_features():
            """Зберігає налаштування"""
            try:
                # Збираємо generators
                for key in ["google_alerts", "google_sign", "email_subscription", "email_for_mail",
                           "gmail_answer", "prompts", "youtube_commentary", "weirdo"]:
                    var_key = f"generators_{key}"
                    if var_key in checkboxes:
                        _features_config.set_enabled("generators", key, checkboxes[var_key].get())
                
                # Збираємо daily report
                for key in ["enabled", "farmer_name", "date", "copy_button"]:
                    var_key = f"daily_report_{key}"
                    if var_key in checkboxes:
                        _features_config.set_enabled("daily_report", key, checkboxes[var_key].get())
                
                # Збираємо utilities
                for key in ["google_backup_codes", "2fa_generator", "password_generator"]:
                    var_key = f"utilities_{key}"
                    if var_key in checkboxes:
                        _features_config.set_enabled("utilities", key, checkboxes[var_key].get())
                
                # Збираємо ui_sections
                for key in ["action_log"]:
                    var_key = f"ui_sections_{key}"
                    if var_key in checkboxes:
                        _features_config.set_enabled("ui_sections", key, checkboxes[var_key].get())
                
                # Збираємо windows mode
                mode = windows_mode_var.get()
                _features_config.set_enabled("windows_mode", "two_windows", mode == "two")
                _features_config.set_enabled("windows_mode", "four_windows", mode == "four")
                _features_config.set_enabled("windows_mode", "six_windows", mode == "six")
                
                # Зберігаємо в файл
                if _features_config.save_config():
                    messagebox.showinfo(
                        "✅ Успіх",
                        "Налаштування збережено!\n\n"
                        "Для застосування змін перезапустіть програму\n\n"
                        f"Конфіг: {_features_config.config_file}"
                    )
                    self.logger.log('✅ Features Settings збережено')
                    popup.destroy()
                else:
                    messagebox.showerror("❌ Помилка", "Не вдалося зберегти налаштування")
            except Exception as e:
                messagebox.showerror("❌ Помилка", f"Помилка збереження:\n{str(e)}")
        
        def reset_features():
            """Скидає налаштування"""
            if messagebox.askyesno(
                "Підтвердження",
                "Скинути всі налаштування до дефолтних значень?\n\n"
                "Це увімкне всі функції та встановить режим 2 вікна."
            ):
                _features_config.features = DEFAULT_FEATURES.copy()
                if _features_config.save_config():
                    messagebox.showinfo(
                        "✅ Успіх",
                        "Налаштування скинуто!\n\n"
                        "Перезапустіть програму для застосування змін."
                    )
                    self.logger.log('🔄 Features Settings скинуто до дефолтних')
                    popup.destroy()
        
        # Ліва частина - Save та Reset
        left_frame = ctk.CTkFrame(buttons_frame, fg_color="transparent")
        left_frame.pack(side="left")
        
        save_btn = ctk.CTkButton(
            left_frame,
            text="💾 Зберегти",
            command=save_features,
            width=150,
            height=40,
            font=ctk.CTkFont(size=14, weight="bold"),
            fg_color="green",
            hover_color="darkgreen"
        )
        save_btn.pack(side="left", padx=5)
        
        reset_btn = ctk.CTkButton(
            left_frame,
            text="🔄 Скинути",
            command=reset_features,
            width=130,
            height=40,
            font=ctk.CTkFont(size=13),
            fg_color="orange",
            hover_color="darkorange"
        )
        reset_btn.pack(side="left", padx=5)
        
        # Права частина - Close
        close_btn = ctk.CTkButton(
            buttons_frame,
            text="❌ Закрити",
            command=popup.destroy,
            width=120,
            height=40,
            font=ctk.CTkFont(size=13),
            fg_color="gray",
            hover_color="darkgray"
        )
        close_btn.pack(side="right", padx=5)
    
    def _create_features_section(self, parent, checkboxes, title, category, items):
        """Створює секцію з чекбоксами для features"""
        section_frame = ctk.CTkFrame(parent)
        section_frame.pack(fill="x", pady=10, padx=10)
        
        # Заголовок
        header = ctk.CTkFrame(section_frame, fg_color="transparent")
        header.pack(fill="x", padx=10, pady=10)
        
        ctk.CTkLabel(
            header,
            text=title,
            font=ctk.CTkFont(size=14, weight="bold")
        ).pack(side="left")
        
        # Чекбокси
        for key, label in items:
            var = tk.BooleanVar(value=_features_config.is_enabled(category, key))
            checkboxes[f"{category}_{key}"] = var
            
            cb = ctk.CTkCheckBox(
                section_frame,
                text=label,
                variable=var,
                font=ctk.CTkFont(size=12)
            )
            cb.pack(anchor="w", padx=20, pady=4)
    
    def _create_windows_mode_section(self, parent):
        """Створює секцію Windows Mode"""
        section_frame = ctk.CTkFrame(parent)
        section_frame.pack(fill="x", pady=10, padx=10)
        
        # Заголовок
        header = ctk.CTkFrame(section_frame, fg_color="transparent")
        header.pack(fill="x", padx=10, pady=10)
        
        ctk.CTkLabel(
            header,
            text="🪟 Windows Mode - Кількість вікон",
            font=ctk.CTkFont(size=14, weight="bold")
        ).pack(side="left")
        
        # Info
        info = ctk.CTkLabel(
            section_frame,
            text="Оберіть режим відображення генераторів (потребує перезапуску)",
            font=ctk.CTkFont(size=11),
            text_color="gray"
        )
        info.pack(anchor="w", padx=20, pady=(0, 10))
        
        # Radio buttons
        windows_mode_var = tk.StringVar()
        
        if _features_config.is_enabled("windows_mode", "two_windows"):
            windows_mode_var.set("two")
        elif _features_config.is_enabled("windows_mode", "four_windows"):
            windows_mode_var.set("four")
        elif _features_config.is_enabled("windows_mode", "six_windows"):
            windows_mode_var.set("six")
        else:
            windows_mode_var.set("two")
        
        modes = [
            ("two", "🪟 2 вікна (стандартний - ноутбуки)"),
            ("four", "🪟 4 вікна (розширений - існуючий)"),
            ("six", "🪟 6 вікон (максимум - НОВИЙ!)"),
        ]
        
        for value, label in modes:
            rb = ctk.CTkRadioButton(
                section_frame,
                text=label,
                variable=windows_mode_var,
                value=value,
                font=ctk.CTkFont(size=12)
            )
            rb.pack(anchor="w", padx=20, pady=4)
        
        return windows_mode_var

    def load_theme_from_file(self):
        path = filedialog.askopenfilename(filetypes=[('JSON files', '*.json')])
        if not path:
            return
        try:
            with open(path, 'r', encoding='utf-8') as f:
                theme = json.load(f)
            # Check if it's a CustomTkinter theme and convert to app format
            if "CTk" in theme:
                root_bg = theme["CTk"]["fg_color"][0]
                btn_bg = theme["CTkButton"]["fg_color"][0]
                btn_fg = theme["CTkButton"]["text_color"][0]
                label_fg = theme["CTkLabel"]["text_color"][0]
                entry_bg = theme["CTkEntry"]["fg_color"][0]
                entry_fg = theme["CTkEntry"]["text_color"][0]
                text_bg = theme["CTkTextbox"]["fg_color"][0]
                text_fg = theme["CTkTextbox"]["text_color"][0]
                scrollbar_bg = theme["CTkScrollbar"]["button_color"][0]
                frame_top_bg = theme["CTkFrame"]["top_fg_color"][0]
                theme = {
                    "name": "Converted CTk Theme",
                    "root_bg": root_bg,
                    "styles": {
                        "TButton": {"background": btn_bg, "foreground": btn_fg},
                        "TLabel": {"background": root_bg, "foreground": label_fg},
                        "TFrame": {"background": root_bg},
                        "TNotebook": {"background": root_bg, "tabbackground": frame_top_bg},
                        "TCombobox": {"fieldbackground": entry_bg, "background": root_bg, "foreground": entry_fg},
                        "TEntry": {"fieldbackground": entry_bg, "background": root_bg, "foreground": entry_fg},
                        "TText": {"background": text_bg, "foreground": text_fg},
                        "TScrollbar": {"background": scrollbar_bg, "troughcolor": root_bg}
                    }
                }
            self.apply_theme(theme)
            self.current_theme = theme
            self.save_theme()
            self.logger.log(f"Loaded theme from {os.path.basename(path)}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load theme: {e}")

    def get_brightness(self, color):
        if color.startswith('#'):
            hex_color = color.lstrip('#')
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return (r + g + b) / 3
        else:
            # Handle color names like gray92
            if color.startswith('gray'):
                try:
                    num = int(color[4:])
                    return num * 255 / 100
                except ValueError:
                    pass
            # Default to light
            return 200

    def load_theme_from_file_path(self, path):
        if not path or not os.path.exists(path):
            return
        try:
            with open(path, 'r', encoding='utf-8') as f:
                theme = json.load(f)
            # Check if it's a CustomTkinter theme and convert to app format
            if "CTk" in theme:
                # Dynamically select color index based on appearance mode
                index = 1 if ctk.get_appearance_mode() == 'Dark' else 0
                root_bg = theme["CTk"]["fg_color"][index]
                btn_bg = theme["CTkButton"]["fg_color"][index]
                btn_fg = theme["CTkButton"]["text_color"][index]
                label_fg = theme["CTkLabel"]["text_color"][index]
                entry_bg = theme["CTkEntry"]["fg_color"][index]
                entry_fg = theme["CTkEntry"]["text_color"][index]
                text_bg = theme["CTkTextbox"]["fg_color"][index]
                text_fg = theme["CTkTextbox"]["text_color"][index]
                scrollbar_bg = theme["CTkScrollbar"]["button_color"][index]
                frame_top_bg = theme["CTkFrame"]["top_fg_color"][index]
                border_color = theme["CTkButton"]["border_color"][index] if "border_color" in theme["CTkButton"] else btn_bg
                theme = {
                    "name": "Converted CTk Theme",
                    "root_bg": root_bg,
                    "styles": {
                        "TButton": {"background": btn_bg, "foreground": btn_fg, "activebackground": btn_bg, "activeforeground": btn_fg, "highlightbackground": root_bg, "highlightcolor": btn_bg, "lightcolor": btn_bg, "darkcolor": btn_bg, "bordercolor": border_color, "relief": "flat"},
                        "TLabel": {"background": root_bg, "foreground": label_fg},
                        "TFrame": {"background": root_bg},
                        "TNotebook": {"background": root_bg, "tabbackground": frame_top_bg},
                        "TCombobox": {"fieldbackground": entry_bg, "background": root_bg, "foreground": entry_fg},
                        "TEntry": {"fieldbackground": entry_bg, "background": root_bg, "foreground": entry_fg},
                        "TText": {"background": text_bg, "foreground": text_fg},
                        "TScrollbar": {"background": scrollbar_bg, "troughcolor": root_bg}
                    }
                }
            self.apply_theme(theme)
            self.current_theme = theme
            self.logger.log(f"Loaded theme from {os.path.basename(path)}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load theme: {e}")

    def toggle_fullscreen(self):
        self.fullscreen = not self.fullscreen
        self.attributes('-fullscreen', self.fullscreen)
        self.logger.log(f'Full screen: {"ON" if self.fullscreen else "OFF"}')
    
    def on_window_resize(self, event):
        """Обробляє зміну розміру вікна"""
        # Перевіряємо чи це подія зміни розміру самого вікна, а не віджетів
        if event.widget != self:
            return
        
        # Отримуємо поточний розмір вікна
        current_width = self.winfo_width()
        current_height = self.winfo_height()
        
        # Визначаємо мінімальний розмір на основі вмісту
        min_content_width = 750  # Мінімум для відображення всіх табів
        min_content_height = 450  # Мінімум для відображення контенту
        
        # Адаптивний мінімальний розмір залежно від масштабу
        scale_factor = self._get_window_scaling()
        adaptive_min_width = max(min_content_width, int(800 / scale_factor))
        adaptive_min_height = max(min_content_height, int(500 / scale_factor))
        
        # Якщо вікно стало меншим за адаптивний мінімум, коригуємо
        if current_width < adaptive_min_width or current_height < adaptive_min_height:
            new_width = max(current_width, adaptive_min_width)
            new_height = max(current_height, adaptive_min_height)
            # Оновлюємо мінімальний розмір
            self.minsize(new_width, new_height)
        else:
            # Зменшуємо мінімальний розмір якщо вікно достатньо велике
            self.minsize(adaptive_min_width, adaptive_min_height)
    
    def _get_window_scaling(self):
        """Отримує масштаб вікна (DPI scaling)"""
        try:
            # Отримуємо DPI scaling фактор
            return self.tk.call('tk', 'scaling')
        except:
            return 1.0

    def load_csv(self):
        # Визначаємо початкову директорію для діалогу
        initial_dir = self.get_csv_initial_directory()
        
        path = filedialog.askopenfilename(
            filetypes=[('CSV files', '*.csv')],
            initialdir=initial_dir,
            title='Виберіть CSV файл для генерації'
        )
        if not path:
            return
            
        self.csv_path = path
        self.data = load_csv_columns(path)
        
        # Зберігаємо новий шлях в конфіг
        self.save_csv_config(path)
        
        # update sections
        for s in self.sections_objs:
            key = s.title
            if key in self.data:
                s.column_values = self.data[key]
            else:
                s.column_values = []
            s.items = []
            s.index = 0
            s._refresh_listbox()
        # update help label with enhanced info
        csv_info = f'CSV: {os.path.basename(path)} ✅ (завантажено)'
        if hasattr(self, 'help_lbl'):
            self.help_lbl.configure(text=csv_info)
        self.logger.log(f'📁 Завантажено CSV: {os.path.basename(path)}')
    
    def refresh_gen_csv_files(self):
        """Оновлює список CSV файлів для Generators"""
        try:
            self.csv_files = self.scan_csv_files()
            if hasattr(self, 'gen_main_csv_dropdown'):
                self.gen_main_csv_dropdown.configure(values=self.csv_files or ["Немає файлів"])
            print(f"Оновлено список CSV файлів: {len(self.csv_files)} файлів знайдено")
        except Exception as e:
            print(f"Помилка оновлення CSV файлів для Generators: {str(e)}")
    
    def on_gen_main_csv_selected(self, selected_file):
        """Обробляє вибір CSV файлу для Generators"""
        if selected_file and selected_file != "Немає файлів":
            try:
                # Завантажуємо обраний CSV файл
                csv_path = os.path.join(os.getcwd(), selected_file)
                if os.path.exists(csv_path):
                    self.csv_path = csv_path
                    self.data = load_csv_columns(csv_path)
                    
                    # Оновлюємо секції
                    for s in self.sections_objs:
                        key = s.title
                        if key in self.data:
                            s.column_values = self.data[key]
                        else:
                            s.column_values = []
                        s.items = []
                        s.index = 0
                        s._refresh_listbox()
                    
                    # Оновлюємо help label якщо існує
                    if hasattr(self, 'help_lbl'):
                        self.help_lbl.configure(text=f'CSV: {os.path.basename(csv_path)}')
                    
                    self.logger.log(f'Завантажено CSV з dropdown: {os.path.basename(csv_path)}')
                    print(f"CSV файл завантажено: {selected_file}")
                else:
                    print(f"Файл не знайдено: {csv_path}")
            except Exception as e:
                print(f"Помилка завантаження CSV файлу: {str(e)}")
                self.logger.log(f'Помилка завантаження CSV: {str(e)}')
    
    def scan_csv_files(self):
        """Сканує робочу директорію на наявність CSV файлів"""
        try:
            csv_files = []
            current_dir = os.getcwd()
            
            for file in os.listdir(current_dir):
                if file.lower().endswith('.csv'):
                    csv_files.append(file)
            
            # Додаємо Account generation.csv якщо існує
            if not csv_files:
                csv_files = ["Account generation.csv"]
            
            return sorted(csv_files)
        except Exception as e:
            print(f"Помилка сканування CSV файлів: {str(e)}")
            return ["Account generation.csv"]
    
    def show_generators_instruction(self):
        instruction_text = """
🚀 GENERATORS - ІНСТРУКЦІЯ З ВИКОРИСТАННЯ

📊 CSV FILE MANAGER:
• Основний CSV - оберіть файл з даними для генерації
• Кнопка 🔄 - оновлює список доступних CSV файлів
• Load CSV - завантажити обраний файл в систему
• Автоматичне сканування CSV файлів в папці

🎯 DAY GENERATOR:
• Day 1-5 - оберіть день для генерації контенту
• Generate Day - згенерувати контент на обраний день
• Next - перехід до наступного елемента
• Copy - копіювати поточний елемент

📝 СЕКЦІЇ КОНТЕНТУ:
• Google Alerts - алерти для моніторингу
• Google Sign - підписи для Google сервісів  
• Email Subscription - підписки на розсилки
• Email for Mail - адреси для прямих листів
• Gmail Answer - відповіді в Gmail
• Prompts - запити та шаблони
• YouTube Commentary - коментарі для YouTube
• Weirdo - дивні запити та нестандартний контент

⚙️ ГЕНЕРАТОРИ ІНСТРУМЕНТІВ:
• Резервні коди Google - трансформація формату
• 2FA Генератор - створення кодів автентифікації
• Генератор паролів - створення безпечних паролів
• Click-to-copy в лог області

🔧 ОСНОВНІ ФУНКЦІЇ:
• Generate All - генерувати весь контент
• Copy All - скопіювати всі результати
• FullScreen - повноекранний режим
• UI Settings - налаштування інтерфейсу

💡 ПРИНЦИП РОБОТИ: Завантажуємо CSV → Генеруємо контент → Копіюємо результати
        """
        messagebox.showinfo("🚀 Generators - Інструкція", instruction_text)

    




    def paste_codes(self):
        try:
            clipboard_text = safe_clipboard_operation("get").strip()
            safe_text_input(self.input_area, clipboard_text)
        except Exception:
            safe_text_input(self.input_area, "Буфер обміну пустий")

    def process_codes(self):
        self.output_area.configure(state='normal')
        self.output_area.delete('0.0', 'end')
        try:
            codes_text = self.input_area.get('0.0', 'end').strip()
            lines = codes_text.splitlines()
            result = []
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                parts = line.split()
                if len(parts) != 2:
                    continue
                result.append(f"{parts[0].zfill(4)}{parts[1].zfill(4)}")
            output_text = ", ".join(result)
            self.output_area.insert('end', f"Резервні коди:\n{output_text}")
            if output_text:
                safe_clipboard_operation("set", output_text)
        except Exception as e:
            self.output_area.insert('end', f"Помилка: {str(e)}")
        finally:
            self.output_area.configure(state='disabled')

    def paste_secret(self):
        self.secret_key = safe_clipboard_operation("get").strip()
        self.secret_label.configure(text=f"Секрет: {self.secret_key[:8]}...")

    def generate_2fa(self):
        self.output_area.configure(state='normal')
        self.output_area.delete('0.0', 'end')
        secret_key = self.secret_key.strip()
        if not secret_key:
            self.output_area.insert('end', "Введіть секретний ключ!")
            self.output_area.configure(state='disabled')
            return
        try:
            secret_key_clean = re.sub(r'[\s=]+', '', secret_key).upper()
            if not re.match(r'^[A-Z2-7]+$', secret_key_clean):
                raise ValueError("Невірний формат Base32 ключа")
            totp = pyotp.TOTP(secret_key_clean)
            current_code = totp.now()
            self.output_area.insert('end', f"Поточний 2FA код:\n{current_code}")
            safe_clipboard_operation("set", current_code)
        except Exception as e:
            self.output_area.insert('end', f"Невірний ключ:\n{str(e)}")
        finally:
            self.output_area.configure(state='disabled')

    def generate_passwords(self):
        try:
            count = int(self.password_count.get())
            if count < 1 or count > 500:
                messagebox.showwarning("Увага", "Кількість паролів повинна бути від 1 до 500")
                return
            self.password_list.delete('0.0', 'end')
            generated_passwords = []
            for i in range(count):
                length = random.randint(8, 12)
                chars = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789!@#$%^&*()_+-=[]{}|;:,.<>?"
                password = "".join(random.choice(chars) for _ in range(length))
                generated_passwords.append(password)
                self.password_list.insert('end', f"{password}\n")
            
            # Автоматично копіюємо перший пароль якщо тільки 1 пароль
            if count == 1 and generated_passwords:
                safe_clipboard_operation("set", generated_passwords[0])
                self.logger.log(f'🔑 Пароль автоматично скопійовано: {generated_passwords[0]}')
            
            self.logger.log(f'Generated {count} password(s) - клік по паролю = копія')
            
        except ValueError:
            messagebox.showwarning("Увага", "Введіть коректне число для кількості паролів")

    def copy_passwords(self):
        passwords = self.password_list.get('0.0', 'end').strip()
        if passwords:
            safe_clipboard_operation("set", passwords)

    def _on_2fa_click_copy(self, event):
        """Обробка кліку по 2FA області - копіює код"""
        try:
            # Отримуємо внутрішній текстовий віджет
            inner_text = self.output_area._textbox
            
            # Отримуємо позицію кліку
            click_index = inner_text.index(f"@{event.x},{event.y}")
            line_start = inner_text.index(f"{click_index} linestart")
            line_end = inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка
            line_text = inner_text.get(line_start, line_end).strip()
            
            # Шукаємо код (6 цифр)
            import re
            code_match = re.search(r'\b\d{6}\b', line_text)
            if code_match:
                code = code_match.group()
                safe_clipboard_operation("set", code)
                
                # Візуальний фідбек
                inner_text.tag_remove('copied_2fa', '1.0', 'end')
                inner_text.tag_add('copied_2fa', line_start, line_end)
                inner_text.tag_config('copied_2fa', background='#404040', foreground='#90EE90')
                
                # Прибираємо підсвітку через 500мс
                inner_text.after(500, lambda: inner_text.tag_remove('copied_2fa', '1.0', 'end'))
                
                self.logger.log(f'📋 2FA код скопійовано: {code}')
            else:
                # Копіюємо весь рядок
                if line_text:
                    safe_clipboard_operation("set", line_text)
                    self.logger.log(f'📋 2FA текст скопійовано: {line_text}')
                
        except Exception as e:
            print(f"Помилка копіювання 2FA: {e}")
    
    def _on_password_click_copy(self, event):
        """Обробка кліку по списку паролів - копіює пароль"""
        try:
            # Отримуємо внутрішній текстовий віджет
            inner_text = self.password_list._textbox
            
            # Отримуємо позицію кліку
            click_index = inner_text.index(f"@{event.x},{event.y}")
            line_start = inner_text.index(f"{click_index} linestart")
            line_end = inner_text.index(f"{click_index} lineend")
            
            # Отримуємо текст рядка (пароль)
            password = inner_text.get(line_start, line_end).strip()
            
            if password:
                safe_clipboard_operation("set", password)
                
                # Візуальний фідбек
                inner_text.tag_remove('copied_password', '1.0', 'end')
                inner_text.tag_add('copied_password', line_start, line_end)
                inner_text.tag_config('copied_password', background='#404040', foreground='#90EE90')
                
                # Прибираємо підсвітку через 500мс
                inner_text.after(500, lambda: inner_text.tag_remove('copied_password', '1.0', 'end'))
                
                self.logger.log(f'🔑 Пароль скопійовано: {password}')
                
        except Exception as e:
            print(f"Помилка копіювання пароля: {e}")

    def increase_password_count(self):
        try:
            current = int(self.password_count.get())
            if current < 500:
                self.password_count.delete(0, 'end')
                self.password_count.insert(0, str(current + 1))
        except ValueError:
            self.password_count.delete(0, 'end')
            self.password_count.insert(0, "1")

    def decrease_password_count(self):
        try:
            current = int(self.password_count.get())
            if current > 1:
                self.password_count.delete(0, 'end')
                self.password_count.insert(0, str(current - 1))
        except ValueError:
            self.password_count.delete(0, 'end')
            self.password_count.insert(0, "1")

    def on_password_count_mousewheel(self, event):
        try:
            current = int(self.password_count.get())
            # Використовуємо глобальну швидкість скролу
            scroll_speed = get_global_scroll_speed()
            step = max(1, int(scroll_speed))
            
            if event.delta > 0:  # Колесико вгору
                if current < 500:
                    new_value = min(500, current + step)
                    self.password_count.delete(0, 'end')
                    self.password_count.insert(0, str(new_value))
            else:  # Колесико вниз
                if current > 1:
                    new_value = max(1, current - step)
                    self.password_count.delete(0, 'end')
                    self.password_count.insert(0, str(new_value))
        except ValueError:
            self.password_count.delete(0, 'end')
            self.password_count.insert(0, "5")

    # removed advanced contrast helper to keep button colors simple and readable


# ================== Tab Manager ==================
class TabManager:
    """Клас для управління порядком та видимістю вкладок"""
    
    def __init__(self, app):
        self.app = app
        self.config_file = get_config_path("tab_config.json")
        self.default_tabs = [
            {"name": "Generators", "title": "Generators", "visible": True, "class": "GeneratorsTab"},
            {"name": "Gmail Hacks", "title": "Gmail Hacks", "visible": True, "class": "GmailHacksTab"},
            {"name": "Gmail Parser", "title": "Gmail Parser", "visible": True, "class": "GmailParserTab"},
            {"name": "Settings", "title": "⚙️ Settings", "visible": True, "class": "SettingsTab"},
            {"name": "SMS Checker", "title": "DAISYSMS", "visible": True, "class": "SmsCheckerTab"},
            {"name": "Registration", "title": "Mail.TM", "visible": True, "class": "RegistrationTab"},
            {"name": "ChatGPT", "title": "ChatGPT", "visible": True, "class": "ChatGPTTab"},
            {"name": "File Generator", "title": "File Generator", "visible": True, "class": "FileGeneratorTab"}
        ]
        self.load_config()
    
    def load_config(self):
        """Завантажує конфігурацію вкладок"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, "r", encoding="utf-8") as f:
                    saved_config = json.load(f)
                    # Оновлюємо існуючі вкладки та додаємо нові
                    saved_names = {tab["name"] for tab in saved_config}
                    self.tabs = []
                    
                    # Спочатку додаємо збережені вкладки в збереженому порядку
                    for saved_tab in saved_config:
                        if saved_tab["name"] in [t["name"] for t in self.default_tabs]:
                            self.tabs.append(saved_tab)
                    
                    # Потім додаємо нові вкладки, які не були збережені
                    for default_tab in self.default_tabs:
                        if default_tab["name"] not in saved_names:
                            self.tabs.append(default_tab)
            else:
                self.tabs = self.default_tabs.copy()
        except Exception as e:
            print(f"Помилка завантаження конфігурації вкладок: {e}")
            self.tabs = self.default_tabs.copy()
    
    def save_config(self):
        """Зберігає конфігурацію вкладок"""
        try:
            print(f"💾 Збереження конфігурації у файл: {self.config_file}")
            print(f"📋 Стан вкладок перед збереженням:")
            for tab in self.tabs:
                print(f"   - {tab['name']}: visible={tab.get('visible', True)}")
            
            with open(self.config_file, "w", encoding="utf-8") as f:
                json.dump(self.tabs, f, ensure_ascii=False, indent=2)
            
            print(f"✅ Конфігурація успішно збережена!")
        except Exception as e:
            print(f"❌ Помилка збереження конфігурації вкладок: {e}")
    
    def get_visible_tabs(self):
        """Повертає список видимих вкладок у правильному порядку"""
        visible = [tab for tab in self.tabs if tab.get("visible", True)]
        print(f"👁️ Видимі вкладки ({len(visible)}/{len(self.tabs)}):")
        for tab in visible:
            print(f"   ✅ {tab['name']}")
        
        hidden = [tab for tab in self.tabs if not tab.get("visible", True)]
        if hidden:
            print(f"🚫 Приховані вкладки:")
            for tab in hidden:
                print(f"   ❌ {tab['name']}")
        
        return visible
    
    def move_tab_up(self, tab_name):
        """Переміщує вкладку вверх по списку"""
        for i, tab in enumerate(self.tabs):
            if tab["name"] == tab_name and i > 0:
                self.tabs[i], self.tabs[i-1] = self.tabs[i-1], self.tabs[i]
                return True
        return False
    
    def move_tab_down(self, tab_name):
        """Переміщує вкладку вниз по списку"""
        for i, tab in enumerate(self.tabs):
            if tab["name"] == tab_name and i < len(self.tabs) - 1:
                self.tabs[i], self.tabs[i+1] = self.tabs[i+1], self.tabs[i]
                return True
        return False
    
    def toggle_tab_visibility(self, tab_name):
        """Перемикає видимість вкладки"""
        for tab in self.tabs:
            if tab["name"] == tab_name:
                tab["visible"] = not tab.get("visible", True)
                return tab["visible"]
        return False
    
    def toggle_tab_visibility(self, tab_name):
        """Перемикає видимість вкладки"""
        for tab in self.tabs:
            if tab["name"] == tab_name:
                tab["visible"] = not tab.get("visible", True)
                return tab["visible"]
        return False
    
    def show_tab_settings(self):
        """Відкриває діалог налаштування вкладок"""
        TabSettingsDialog(self.app, self)
    
    def add_settings_button_to_tab(self, tab_frame):
        """Додає кнопку налаштувань в правий верхній кут вкладки"""
        try:
            # Створюємо контейнер для кнопки налаштувань (у правому верхньому куті)
            settings_container = ctk.CTkFrame(tab_frame, fg_color="transparent")
            settings_container.place(relx=0.98, rely=0.02, anchor="ne")
            
            # Кнопка налаштувань
            settings_btn = ctk.CTkButton(
                settings_container,
                text="⚙️",
                width=32,
                height=32,
                corner_radius=6,
                command=self.show_tab_settings,
                font=ctk.CTkFont(size=16),
                fg_color=("gray75", "gray25"),
                hover_color=("gray65", "gray35")
            )
            settings_btn.pack()
            
            print(f"✅ Додано кнопку налаштувань до вкладки")
        except Exception as e:
            print(f"⚠️ Помилка додавання кнопки налаштувань: {e}")


class TabSettingsDialog(ctk.CTkToplevel):
    """Діалог для налаштування вкладок"""
    
    def __init__(self, parent, tab_manager):
        super().__init__(parent)
        self.tab_manager = tab_manager
        self.setup_dialog()
        self.create_widgets()
        
    def setup_dialog(self):
        """Налаштування діалогового вікна"""
        self.title("Налаштування вкладок")
        self.geometry("500x650")  # Трохи більше для кнопок
        self.resizable(False, False)
        self.attributes('-topmost', True)
        self.transient(self.master)
        self.grab_set()
        
        # Центруємо вікно
        self.update_idletasks()
        x = (self.winfo_screenwidth() // 2) - (500 // 2)
        y = (self.winfo_screenheight() // 2) - (650 // 2)
        self.geometry(f"500x650+{x}+{y}")
    
    def create_widgets(self):
        """Створює елементи інтерфейсу"""
        # Заголовок
        title_label = ctk.CTkLabel(self, text="Налаштування порядку та видимості вкладок", 
                                  font=ctk.CTkFont(size=16, weight="bold"))
        title_label.pack(pady=(20, 10))
        
        # Інструкція
        info_label = ctk.CTkLabel(self, text="Змініть порядок вкладок або приховайте непотрібні", 
                                 font=ctk.CTkFont(size=12), text_color="gray")
        info_label.pack(pady=(0, 5))
        
        # Підказка про захищені вкладки
        protected_info = ctk.CTkLabel(self, text="🔐 Захищені вкладки (SMS Checker, Gmail Parser) вимагають пароль", 
                                     font=ctk.CTkFont(size=10), text_color="#ff9800")
        protected_info.pack(pady=(0, 15))
        
        # Контейнер для списку вкладок
        self.tabs_frame = ctk.CTkScrollableFrame(self, width=400, height=400)
        self.tabs_frame.pack(pady=10, padx=20, fill="both", expand=True)
        
        self.create_tab_items()
        
        # Кнопки управління
        buttons_frame = ctk.CTkFrame(self)
        buttons_frame.pack(pady=20, padx=20, fill="x")
        
        save_btn = ctk.CTkButton(buttons_frame, text="💾 Зберегти", 
                               command=lambda: self.on_save_click(), width=140, height=40,
                               font=ctk.CTkFont(size=14, weight="bold"),
                               fg_color="#4CAF50", hover_color="#45a049")
        save_btn.pack(side="left", padx=5)
        
        reset_btn = ctk.CTkButton(buttons_frame, text="🔄 Скинути", 
                                command=self.reset_to_default, width=140, height=40,
                                font=ctk.CTkFont(size=13),
                                fg_color="#FF9800", hover_color="#F57C00")
        reset_btn.pack(side="left", padx=5)
        
        cancel_btn = ctk.CTkButton(buttons_frame, text="❌ Скасувати", 
                                 command=self.destroy, width=140, height=40,
                                 font=ctk.CTkFont(size=13),
                                 fg_color="#f44336", hover_color="#d32f2f")
        cancel_btn.pack(side="right", padx=5)
    
    def create_tab_items(self):
        """Створює елементи для кожної вкладки"""
        # Очищуємо попередні елементи
        for widget in self.tabs_frame.winfo_children():
            widget.destroy()
        
        self.tab_vars = {}
        self.tab_frames = {}
        
        for i, tab in enumerate(self.tab_manager.tabs):
            # Контейнер для вкладки
            tab_frame = ctk.CTkFrame(self.tabs_frame)
            tab_frame.pack(fill="x", pady=2, padx=5)
            self.tab_frames[tab["name"]] = tab_frame
            
            # Checkbox для видимості
            var = tk.BooleanVar(value=tab.get("visible", True))
            self.tab_vars[tab["name"]] = var
            
            # Перевіряємо чи треба пароль для цієї вкладки
            protected_tabs = ["SMS Checker", "Gmail Parser"]
            if tab["name"] in protected_tabs:
                # Для захищених вкладок додаємо іконку замка та callback з перевіркою паролю
                checkbox = ctk.CTkCheckBox(tab_frame, text="🔒", variable=var, width=40,
                                          command=lambda t=tab["name"], v=var: self.toggle_protected_tab(t, v))
            else:
                checkbox = ctk.CTkCheckBox(tab_frame, text="", variable=var, width=20)
            
            checkbox.pack(side="left", padx=(10, 5), pady=10)
            
            # Назва вкладки (з позначкою захищеної)
            protected_tabs = ["SMS Checker", "Gmail Parser"]
            if tab["name"] in protected_tabs:
                label_text = f"{tab['title']} 🔐"
                text_color = "#ff9800"  # Помаранчевий для захищених
            else:
                label_text = tab["title"]
                text_color = None
            
            name_label = ctk.CTkLabel(tab_frame, text=label_text, 
                                    font=ctk.CTkFont(size=13, weight="bold"),
                                    anchor="w", width=200,
                                    text_color=text_color)
            name_label.pack(side="left", padx=5, pady=10)
            
            # Кнопки переміщення
            buttons_frame = ctk.CTkFrame(tab_frame)
            buttons_frame.pack(side="right", padx=10, pady=5)
            
            up_btn = ctk.CTkButton(buttons_frame, text="↑", width=30, height=25,
                                 command=lambda name=tab["name"]: self.move_up(name),
                                 state="disabled" if i == 0 else "normal")
            up_btn.pack(side="left", padx=2)
            
            down_btn = ctk.CTkButton(buttons_frame, text="↓", width=30, height=25,
                                   command=lambda name=tab["name"]: self.move_down(name),
                                   state="disabled" if i == len(self.tab_manager.tabs) - 1 else "normal")
            down_btn.pack(side="left", padx=2)
    
    def toggle_protected_tab(self, tab_name, var):
        """Обробляє перемикання захищеної вкладки (з паролем)"""
        # Поточний стан після кліку
        new_state = var.get()
        
        # Знаходимо попередній стан вкладки
        old_state = True  # За замовчуванням активна
        for tab in self.tab_manager.tabs:
            if tab["name"] == tab_name:
                old_state = tab.get("visible", True)
                break
        
        # Якщо змінюється стан (вмикається або вимикається) - запитуємо пароль
        if new_state != old_state:
            password = self.ask_password()
            
            if password == "PunchiteverydayGODMODE":
                # Пароль правильний - дозволяємо зміну та ОДРАЗУ оновлюємо в конфігурації
                print(f"✅ Правильний пароль! Змінюємо видимість {tab_name}: {old_state} -> {new_state}")
                for tab in self.tab_manager.tabs:
                    if tab["name"] == tab_name:
                        tab["visible"] = new_state
                        print(f"✅ Оновлено в конфігурації: {tab_name} visible={new_state}")
                        break
                return
            else:
                # Пароль неправильний - повертаємо попередній стан
                print(f"❌ Неправильний пароль для {tab_name}! Повертаємо стан: {old_state}")
                messagebox.showerror("Помилка", "Неправильний пароль!\n\nПотрібен пароль для зміни видимості\nзахищених вкладок (SMS Checker, Gmail Parser)")
                var.set(old_state)  # Повертаємо в попередній стан
    
    def ask_password(self):
        """Запитує пароль у користувача"""
        password_dialog = ctk.CTkInputDialog(
            text="Введіть пароль для зміни налаштувань захищеної вкладки:",
            title="ПарольRequired"
        )
        password = password_dialog.get_input()
        return password if password else ""
    
    def move_up(self, tab_name):
        """Переміщує вкладку вверх"""
        if self.tab_manager.move_tab_up(tab_name):
            self.create_tab_items()
    
    def move_down(self, tab_name):
        """Переміщує вкладку вниз"""
        if self.tab_manager.move_tab_down(tab_name):
            self.create_tab_items()
    
    def reset_to_default(self):
        """Скидає налаштування до значень за замовчуванням"""
        if messagebox.askyesno("Підтвердження", "Скинути всі налаштування вкладок до значень за замовчуванням?"):
            self.tab_manager.tabs = self.tab_manager.default_tabs.copy()
            self.create_tab_items()
    
    def on_save_click(self):
        """Обробник натискання кнопки Зберегти"""
        print("\n🖱️ НАТИСНУТО КНОПКУ 'ЗБЕРЕГТИ'!")
        self.save_and_close()
    
    def save_and_close(self):
        """Зберігає зміни та закриває діалог"""
        print("\n💾 Збереження змін налаштувань вкладок...")
        
        # Оновлюємо видимість вкладок (враховуємо що для захищених вже оновлено)
        for tab_name, var in self.tab_vars.items():
            for tab in self.tab_manager.tabs:
                if tab["name"] == tab_name:
                    current_visible = tab.get("visible", True)
                    new_visible = var.get()
                    
                    # Для захищених вкладок значення вже оновлено в toggle_protected_tab
                    # Просто підтверджуємо що воно синхронізоване
                    if tab_name in ["SMS Checker", "Gmail Parser"]:
                        print(f"🔒 Захищена вкладка {tab_name}: visible={current_visible} (вже оновлено)")
                    else:
                        tab["visible"] = new_visible
                        print(f"📋 Вкладка {tab_name}: visible={current_visible} -> {new_visible}")
        
        # Зберігаємо швидкість скролу (якщо є такий контрол)
        if hasattr(self, 'scroll_speed_var'):
            scroll_speed_value = self.scroll_speed_var.get()
            set_global_scroll_speed(scroll_speed_value)  # Оновлюємо глобальну швидкість
            self.master.save_scroll_speed_config(scroll_speed_value)  # Зберігаємо в файл
            self.master.update_all_scroll_speeds()  # Синхронізуємо всі компоненти
        
        # Зберігаємо конфігурацію
        print("💾 Збереження конфігурації в файл...")
        self.tab_manager.save_config()
        
        print("🔄 Перестворення вкладок...")
        # Перестворюємо вкладки в основному вікні
        self.master.recreate_tabs()
        
        print("✅ Налаштування збережено та застосовано!\n")
        self.destroy()


# Система збереження позиції та розміру вікна
def save_window_state(geometry_string, is_maximized=False):
    """Зберігає стан вікна (розмір та позицію)"""
    try:
        window_config = {
            "geometry": geometry_string,
            "maximized": is_maximized,
            "saved_at": time.time()
        }
        
        config_path = get_config_path("window_state.json")
        with open(config_path, "w", encoding="utf-8") as f:
            json.dump(window_config, f, ensure_ascii=False, indent=2)
            
        print(f"💾 Window state saved: {geometry_string}")
        
    except Exception as e:
        print(f"Failed to save window state: {e}")

def load_window_state():
    """Завантажує збережений стан вікна"""
    try:
        config_path = get_config_path("window_state.json")
        
        if not os.path.exists(config_path):
            return None
        
        with open(config_path, "r", encoding="utf-8") as f:
            window_config = json.load(f)
        
        # Перевіряємо чи не застарілий конфіг (старіше 30 днів)
        if time.time() - window_config.get("saved_at", 0) > 30 * 24 * 3600:
            print("⚠ Window state config is too old, using defaults")
            return None
        
        print(f"📂 Window state loaded: {window_config.get('geometry', 'N/A')}")
        return window_config
        
    except Exception as e:
        print(f"Failed to load window state: {e}")
        return None

def validate_geometry(geometry_string):
    """Перевіряє чи коректна геометрія вікна для поточного екрану"""
    try:
        import tkinter as tk
        root = tk.Tk()
        root.withdraw()
        
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        root.destroy()
        
        # Парсимо геометрію: "1200x800+100+50" або "1200x800-100-50"
        if not ("+" in geometry_string or "-" in geometry_string.split("x")[1]):
            return False
        
        # Розділяємо на розмір і позицію
        parts = geometry_string.replace("-", "+-").split("+")
        size_part = parts[0]
        
        if "x" not in size_part:
            return False
            
        width, height = map(int, size_part.split("x"))
        
        # Перевіряємо позицію
        if len(parts) >= 3:
            x = int(parts[1]) if parts[1] else 0
            y = int(parts[2]) if parts[2] else 0
        elif len(parts) == 2:
            # Якщо тільки одна координата
            x = int(parts[1]) if parts[1] else 0
            y = 50  # default
        else:
            x, y = 100, 50  # default position
        
        # Перевіряємо розумні межі
        min_width, min_height = 600, 400
        max_width = screen_width + 200  # дозволяємо трохи більше для мультімонітору
        max_height = screen_height + 200
        
        if (width < min_width or height < min_height or 
            width > max_width or height > max_height):
            print(f"⚠ Window size {width}x{height} is outside reasonable bounds")
            return False
        
        # Перевіряємо чи вікно хоча б частково видиме на екрані
        if (x < -width + 100 or y < -height + 100 or 
            x > screen_width or y > screen_height):
            print(f"⚠ Window position ({x}, {y}) is off-screen")
            return False
            
        return True
        
    except Exception as e:
        print(f"Geometry validation error: {e}")
        return False

# Auto-backup система
def create_backup():
    """Створює резервну копію всіх конфігів в папці Документи"""
    try:
        # Папка для бекапів в Документах
        documents_path = os.path.join(os.path.expanduser("~"), "Documents")
        backup_folder = os.path.join(documents_path, "Punch_Gmail_GOD_Backups")
        
        # Створюємо папку якщо не існує
        os.makedirs(backup_folder, exist_ok=True)
        
        # Папка для поточного бекапу з датою
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        current_backup = os.path.join(backup_folder, f"backup_{timestamp}")
        os.makedirs(current_backup, exist_ok=True)
        
        # Список файлів для бекапу
        config_files = [
            "gmail_hacks_config.json",
            "octo_profile_last.json",
            "Account generation.csv",
            "company.txt"
        ]
        
        # Додаємо звукові файли (різні формати)
        sound_extensions = [".mp3", ".wav", ".ogg"]
        for ext in sound_extensions:
            sound_file = f"sms_alert{ext}"
            config_files.append(sound_file)
        
        backup_count = 0
        
        # Копіюємо конфіги
        for config_file in config_files:
            source_path = get_config_path(config_file)
            if os.path.exists(source_path):
                try:
                    destination = os.path.join(current_backup, config_file)
                    shutil.copy2(source_path, destination)
                    backup_count += 1
                except Exception as e:
                    print(f"Backup error for {config_file}: {e}")
        
        # Очищуємо старі бекапи (залишаємо тільки останні 10)
        cleanup_old_backups(backup_folder)
        
        print(f"Backup created: {backup_count} files saved to Documents/Punch_Gmail_GOD_Backups")
        return True
        
    except Exception as e:
        print(f"Backup failed: {e}")
        return False

def cleanup_old_backups(backup_folder, keep_count=10):
    """Видаляє старі бекапи, залишаючи тільки останні keep_count"""
    try:
        # Отримуємо список всіх папок бекапів
        backup_dirs = []
        for item in os.listdir(backup_folder):
            item_path = os.path.join(backup_folder, item)
            if os.path.isdir(item_path) and item.startswith("backup_"):
                backup_dirs.append(item_path)
        
        # Сортуємо по даті створення (найновіші спочатку)
        backup_dirs.sort(key=lambda x: os.path.getctime(x), reverse=True)
        
        # Видаляємо старі бекапи
        for old_backup in backup_dirs[keep_count:]:
            try:
                shutil.rmtree(old_backup)
                print(f"Removed old backup: {os.path.basename(old_backup)}")
            except Exception as e:
                print(f"Failed to remove old backup {old_backup}: {e}")
                
    except Exception as e:
        print(f"Cleanup error: {e}")

def setup_sound_file():
    """Налаштовує звуковий файл для сповіщень"""
    try:
        sound_path = get_config_path("sms_alert.mp3")
        
        # Якщо файл вже є в конфігах, нічого не робимо
        if os.path.exists(sound_path):
            return
        
        # Шукаємо оригінальний файл в Downloads
        original_path = r"C:\Users\alexx\Downloads\Звуки - Внимание.mp3"
        if os.path.exists(original_path):
            shutil.copy2(original_path, sound_path)
            print(f"Sound file setup complete: {sound_path}")
        else:
            print(f"Original sound file not found at: {original_path}")
    except Exception as e:
        print(f"Sound setup error: {e}")

def test_sound():
    """Тестує відтворення звукового сповіщення"""
    try:
        print("🧪 Testing sound playback...")
        play_alert_sound()
        return True
    except Exception as e:
        print(f"Sound test failed: {e}")
        return False

def choose_custom_sound():
    """Дозволяє користувачу вибрати власний звуковий файл"""
    try:
        from tkinter import filedialog
        
        # Відкриваємо діалог вибору файлу
        file_path = filedialog.askopenfilename(
            title="Виберіть звуковий файл для сповіщень",
            filetypes=[
                ("Audio files", "*.mp3 *.wav *.ogg"),
                ("MP3 files", "*.mp3"),
                ("WAV files", "*.wav"),
                ("OGG files", "*.ogg"),
                ("All files", "*.*")
            ]
        )
        
        if not file_path:
            return False  # Користувач скасував вибір
        
        # Отримуємо розширення файлу
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Копіюємо файл в папку конфігів
        config_sound_path = get_config_path(f"sms_alert{file_ext}")
        
        # Якщо це не mp3, то зберігаємо з оригінальним розширенням
        if file_ext != '.mp3':
            # Видаляємо старий mp3 файл якщо є
            old_mp3_path = get_config_path("sms_alert.mp3")
            if os.path.exists(old_mp3_path):
                os.remove(old_mp3_path)
        else:
            config_sound_path = get_config_path("sms_alert.mp3")
        
        # Копіюємо новий файл
        shutil.copy2(file_path, config_sound_path)
        
        print(f"Custom sound installed: {os.path.basename(file_path)}")
        
        # Тестуємо новий звук
        test_sound()
        
        return True
        
    except Exception as e:
        print(f"Failed to install custom sound: {e}")
        return False

def get_current_sound_file():
    """Повертає шлях до поточного звукового файлу"""
    # Шукаємо звуковий файл з різними розширеннями
    extensions = ['.mp3', '.wav', '.ogg']
    for ext in extensions:
        sound_path = get_config_path(f"sms_alert{ext}")
        if os.path.exists(sound_path):
            return sound_path
    return None

def main():
    # Перевірку оновлень тепер робимо через кнопку в Settings
    # try:
    #     has_update, new_version, changelog = check_for_updates()
    #     if has_update:
    #         threading.Thread(target=lambda: show_update_notification(new_version, changelog), daemon=True).start()
    # except:
    #     pass
    pass
    
    # Налаштовуємо звуковий файл
    try:
        setup_sound_file()
    except:
        pass
    
    # Створюємо початковий бекап при запуску
    try:
        create_backup()
    except:
        pass
    
    print("\n" + "="*80)
    print("🚀 PUNCH IT NOW 9.2 - Cross-Platform Generator")
    print("="*80)
    print("⌨️  ГАРЯЧІ КЛАВІШІ:")
    print("   • F11 - Повноекранний режим")
    print("   • Ctrl+, - Налаштування вкладок (показати/приховати вкладки)")
    print("   • ⚙️ на вкладці - Швидкий доступ до налаштувань")
    print("="*80 + "\n")
    
    app = App()
    
    # Налаштовуємо періодичні бекапи кожні 30 хвилин
    def periodic_backup():
        create_backup()
        # Запускаємо наступний бекап через 30 хвилин (1800000 мс)
        app.after(1800000, periodic_backup)
    
    # Запускаємо перший періодичний бекап через 30 хвилин
    app.after(1800000, periodic_backup)
    
    app.mainloop()


if __name__ == '__main__':
    main()


